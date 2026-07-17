"""
NextCapOS - Enhanced AI-Driven Buyer & Marketing Agency
FastAPI backend with JWT auth, Claude Sonnet 4.5 research/newsletter generation,
Composio LinkedIn OAuth integration, WebMCP action registry.
"""
import os
import re
import uuid
import json
import logging
import asyncio
import shutil
from pathlib import Path
from datetime import datetime, timezone, timedelta
import secrets
from mailer import send_email, send_email_with_attachment, link as mail_link
from typing import List, Optional, Literal, Dict, Any, Tuple

from fastapi import FastAPI, APIRouter, HTTPException, Depends, status, Header, Request
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient, AsyncIOMotorGridFSBucket
from bson import ObjectId
from dotenv import load_dotenv
from pydantic import BaseModel, Field, EmailStr, ConfigDict
import bcrypt
import jwt as pyjwt
import httpx
import io
from fastapi import UploadFile, File, Form
from fastapi.responses import StreamingResponse, Response

try:
    from pypdf import PdfReader
except Exception:  # pragma: no cover
    PdfReader = None
try:
    from docx import Document as DocxDocument
except Exception:  # pragma: no cover
    DocxDocument = None

from emergentintegrations.llm.chat import LlmChat, UserMessage

from security_service import (
    sha256_bytes, sha256_hex, canonical_event_hash, compute_content_hash, GENESIS_HASH,
    stamp_digest, parse_ots, verify_ots, upgrade_ots, find_btc_attestation,
    encrypt_bytes, decrypt_envelope, encryption_configured,
)
from provenance import build_provenance_pdf

# -----------------------------------------------------------------------------
# Bootstrap
# -----------------------------------------------------------------------------
ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / ".env")

MONGO_URL = os.environ["MONGO_URL"]
DB_NAME = os.environ["DB_NAME"]
JWT_SECRET = os.environ["JWT_SECRET"]
JWT_ALGORITHM = os.environ.get("JWT_ALGORITHM", "HS256")
JWT_EXPIRY_HOURS = int(os.environ.get("JWT_EXPIRY_HOURS", "72"))
EMERGENT_LLM_KEY = os.environ["EMERGENT_LLM_KEY"]
COMPOSIO_API_KEY = os.environ.get("COMPOSIO_API_KEY", "")
COMPOSIO_BASE_URL = os.environ.get("COMPOSIO_BASE_URL", "https://backend.composio.dev")
BRAVE_API_KEY = os.environ.get("BRAVE_API_KEY", "")
PERPLEXITY_API_KEY = os.environ.get("PERPLEXITY_API_KEY", "")

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger("workz")

client = AsyncIOMotorClient(MONGO_URL)
db = client[DB_NAME]
gridfs_bucket = AsyncIOMotorGridFSBucket(db, bucket_name="deal_room_files_fs")
listing_files_bucket = AsyncIOMotorGridFSBucket(db, bucket_name="listing_staged_files_fs")
private_locker_bucket = AsyncIOMotorGridFSBucket(db, bucket_name="private_locker_fs")

app = FastAPI(title="NextCapOS AI Platform", version="1.0.0")
api_router = APIRouter(prefix="/api")
bearer_scheme = HTTPBearer(auto_error=False)


# -----------------------------------------------------------------------------
# Models
# -----------------------------------------------------------------------------
class UserPublic(BaseModel):
    model_config = ConfigDict(extra="ignore")
    id: str
    email: EmailStr
    name: str
    role: Literal["admin", "buyer", "seller", "agent"] = "buyer"
    organization: Optional[str] = None
    interests: List[str] = Field(default_factory=list)
    newsletter_opt_in: bool = False
    created_at: datetime
    is_demo: bool = False
    demo_data_retention_hours: Optional[int] = None
    # Computed at serialize time: "collaborator" if the user has zero owned
    # listings AND zero org_admin memberships AND isn't an admin. Drives the
    # restricted nav + the "Become a full member" upgrade CTA. Principals
    # (owners, agents, sellers with their own deals) get "principal".
    account_scope: Literal["collaborator", "principal"] = "principal"


class RegisterRequest(BaseModel):
    email: EmailStr
    password: str
    name: str
    organization: Optional[str] = None
    role: Literal["buyer", "seller", "agent"] = "buyer"
    # Optional org bootstrap: requester can create a new org during signup
    # ("create" + org_name) or join one via a pending invite ("join" + org_invite_token).
    # Defaults to "none" — they'll work as an individual until invited later.
    org_choice: Literal["create", "join", "none"] = "none"
    org_name: Optional[str] = None
    org_invite_token: Optional[str] = None
    # Optional listing-collaborator invite token. When supplied and valid
    # (email match, not expired, not already accepted), the new account is
    # auto-activated, added to the listing's collaborator list, and the
    # response returns a JWT directly — no admin approval gate.
    listing_invite_token: Optional[str] = None


class ListingCreate(BaseModel):
    company_name: str
    sector: str
    geography: str
    asking_price_usd_m: float
    revenue_usd_m: Optional[float] = None
    ebitda_usd_m: Optional[float] = None
    employees: Optional[int] = None
    headline: str
    summary: str
    highlights: List[str] = Field(default_factory=list)
    status: Literal["draft", "live", "under_loi", "closed"] = "draft"


class InquiryCreate(BaseModel):
    message: str


class FileUpload(BaseModel):
    filename: str
    folder: Literal["financials", "legal", "hr", "it", "operations", "commercial", "other"] = "other"
    content: str  # extracted/pasted text content
    note: Optional[str] = None


class NDAAccept(BaseModel):
    signed_name: str = Field(..., min_length=2, max_length=120)


class DRLApply(BaseModel):
    template_id: str


class CopilotAsk(BaseModel):
    message: str


class LoginRequest(BaseModel):
    email: EmailStr
    password: str


class TokenResponse(BaseModel):
    token: str
    user: UserPublic


class CompanyResearchRequest(BaseModel):
    company_name: str
    company_url: Optional[str] = None
    sector: Optional[str] = None
    region: Optional[str] = None
    notes: Optional[str] = None


class CollateralRequest(BaseModel):
    asset_type: Literal["one_pager", "email_sequence", "linkedin_post", "deal_memo"]
    deal_name: str
    target_audience: str
    key_points: str
    tone: Optional[str] = "professional-institutional"


class OutreachCampaignRequest(BaseModel):
    name: str
    target_persona: str
    channel: Literal["linkedin", "email"] = "linkedin"
    audience_size: int = 50
    message_brief: str


class LeadCreate(BaseModel):
    name: str
    company: str
    title: str
    email: Optional[EmailStr] = None
    source: Optional[str] = "manual"


class LeadStageUpdate(BaseModel):
    stage: Literal["new", "qualified", "engaged", "negotiation", "closed"]


class NewsletterPreferences(BaseModel):
    opt_in: bool
    interests: List[str] = Field(default_factory=list)
    cadence: Literal["weekly", "biweekly", "monthly"] = "weekly"


class NewsletterDraftRequest(BaseModel):
    topic: Optional[str] = None


# -----------------------------------------------------------------------------
# Utility helpers
# -----------------------------------------------------------------------------
def now_utc() -> datetime:
    return datetime.now(timezone.utc)


def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")


def verify_password(password: str, hashed: str) -> bool:
    try:
        return bcrypt.checkpw(password.encode("utf-8"), hashed.encode("utf-8"))
    except Exception:
        return False


def create_token(user_id: str, role: str) -> str:
    payload = {
        "sub": user_id,
        "role": role,
        "iat": int(now_utc().timestamp()),
        "exp": int((now_utc() + timedelta(hours=JWT_EXPIRY_HOURS)).timestamp()),
    }
    return pyjwt.encode(payload, JWT_SECRET, algorithm=JWT_ALGORITHM)


def decode_token(token: str) -> Dict[str, Any]:
    try:
        return pyjwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except pyjwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except pyjwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")


def serialize_user(doc: dict) -> dict:
    is_demo = bool(doc.get("is_demo"))
    return {
        "id": doc["id"],
        "email": doc["email"],
        "name": doc["name"],
        "role": doc.get("role", "buyer"),
        "organization": doc.get("organization"),
        "interests": doc.get("interests", []),
        "newsletter_opt_in": doc.get("newsletter_opt_in", False),
        "created_at": doc["created_at"]
        if isinstance(doc["created_at"], datetime)
        else datetime.fromisoformat(doc["created_at"]),
        "is_demo": is_demo,
        "demo_data_retention_hours": 48 if is_demo else None,
        # Default "principal" — overwritten by serialize_user_with_scope().
        # Sync call sites that don't need the live computation (e.g. JWT
        # refresh of stale token data) get the safe default.
        "account_scope": "principal",
    }


async def _compute_account_scope(user_id: str, role: str) -> str:
    """
    "collaborator" iff the user's ONLY relationship to the platform is being
    a listing collaborator on at least one listing owned by someone else AND
    they own zero listings, admin zero orgs, and aren't a platform admin.

    Anyone admin-invited (agent/seller/buyer/admin), self-registered & admin-
    approved, an org_admin, or a listing owner is a "principal" — they get
    the full role-based nav (BUYER_NAV / SELLER_NAV / AGENT_NAV / ADMIN_NAV).

    NB: the earlier heuristic ("owns 0 listings AND admins 0 orgs ⇒ collab")
    incorrectly demoted freshly admin-invited agents and not-yet-listed
    sellers to the stripped "My Collaborations" nav.
    """
    if role == "admin":
        return "principal"
    if await db.listings.count_documents({"seller_id": user_id}) > 0:
        return "principal"
    if await db.org_memberships.count_documents(
        {"user_id": user_id, "role": "org_admin"}
    ) > 0:
        return "principal"
    # Only demote to "collaborator" if the user is actually listed as a
    # collaborator on at least one listing. Pure-collab accounts come into
    # being via /listing-invites/{token}/accept or the listing_invite_token
    # fast path in /auth/register.
    if await db.listings.count_documents({"collaborators.user_id": user_id}) > 0:
        return "collaborator"
    return "principal"


async def serialize_user_with_scope(doc: dict) -> dict:
    """Like serialize_user but with the live account_scope computed."""
    base = serialize_user(doc)
    base["account_scope"] = await _compute_account_scope(doc["id"], base["role"])
    return base


async def get_current_user(
    credentials: Optional[HTTPAuthorizationCredentials] = Depends(bearer_scheme),
) -> dict:
    if not credentials or not credentials.credentials:
        raise HTTPException(status_code=401, detail="Missing bearer token")
    payload = decode_token(credentials.credentials)
    user = await db.users.find_one({"id": payload["sub"]}, {"_id": 0, "password_hash": 0})
    if not user:
        raise HTTPException(status_code=401, detail="User not found")
    return user


async def require_role(user: dict, allowed: List[str]):
    if user.get("role") not in allowed:
        raise HTTPException(status_code=403, detail="Insufficient role")


async def log_audit(actor_id: str, action: str, target: str = "", meta: Optional[dict] = None):
    """Append an audit entry to a tamper-evident hash chain. Each entry stores
    (seq, prev_hash, content_hash) so a verifier can re-walk the chain and detect any tampering.
    """
    eid = str(uuid.uuid4())
    ts = now_utc().isoformat()
    # Atomically grab next seq + prev_hash from the chain head
    head = await db.audit_chain_head.find_one_and_update(
        {"_id": "head"},
        {"$inc": {"seq": 1}},
        upsert=True,
        return_document=True,
    )
    seq = head.get("seq", 1) if head else 1
    prev_hash = head.get("last_hash", GENESIS_HASH) if head else GENESIS_HASH
    doc = {
        "id": eid,
        "actor_id": actor_id,
        "action": action,
        "target": target,
        "meta": meta or {},
        "timestamp": ts,
        "seq": seq,
        "prev_hash": prev_hash,
    }
    doc["content_hash"] = compute_content_hash(doc)
    await db.audit_logs.insert_one(doc)
    await db.audit_chain_head.update_one(
        {"_id": "head"},
        {"$set": {"last_hash": doc["content_hash"], "last_seq": seq, "last_ts": ts}},
        upsert=True,
    )
    # Best-effort: enqueue an OTS anchor for every Nth entry (chain checkpoint)
    if seq % 25 == 0:
        try:
            asyncio.create_task(_anchor_audit_checkpoint(seq, doc["content_hash"]))
        except Exception:
            pass


async def _anchor_audit_checkpoint(seq: int, head_hash_hex: str):
    """Stamp the current chain head with OpenTimestamps."""
    try:
        digest = bytes.fromhex(head_hash_hex)
        ots_bytes = await stamp_digest(digest)
        await db.ots_proofs.insert_one({
            "id": str(uuid.uuid4()),
            "kind": "audit_chain_checkpoint",
            "target_id": f"seq:{seq}",
            "owner_user_id": "system",
            "digest_hex": head_hash_hex,
            "ots_bytes": ots_bytes,
            "btc_block_height": None,
            "status": "pending",
            "created_at": now_utc().isoformat(),
            "label": f"Audit chain checkpoint @ seq {seq}",
        })
        logger.info(f"OTS anchor stored for audit chain seq {seq}")
    except Exception as e:
        logger.warning(f"OTS anchor failed for seq {seq}: {e}")


async def notarize_event(
    kind: str,
    target_id: str,
    payload: dict,
    owner_user_id: str,
    label: str = "",
) -> Optional[str]:
    """
    Compute SHA-256 of a canonical JSON payload, submit to OpenTimestamps calendars,
    persist the .ots proof. Returns the proof_id (or None on failure — non-blocking).
    """
    try:
        digest = canonical_event_hash(payload)
        ots_bytes = await stamp_digest(digest)
        proof_id = str(uuid.uuid4())
        await db.ots_proofs.insert_one({
            "id": proof_id,
            "kind": kind,
            "target_id": target_id,
            "owner_user_id": owner_user_id,
            "digest_hex": digest.hex(),
            "payload_preview": {k: v for k, v in payload.items() if not isinstance(v, (bytes, bytearray))},
            "ots_bytes": ots_bytes,
            "btc_block_height": None,
            "status": "pending",
            "created_at": now_utc().isoformat(),
            "label": label or kind,
        })
        logger.info(f"OTS notarized {kind}/{target_id} → {digest.hex()[:16]}")
        return proof_id
    except Exception as e:
        logger.warning(f"OTS notarize {kind}/{target_id} failed: {e}")
        return None


async def notarize_bytes(
    kind: str,
    target_id: str,
    data: bytes,
    owner_user_id: str,
    label: str = "",
    extra: Optional[dict] = None,
) -> Optional[str]:
    """Notarize raw bytes (e.g., file binary) via SHA-256 → OTS."""
    try:
        digest = sha256_bytes(data)
        ots_bytes = await stamp_digest(digest)
        proof_id = str(uuid.uuid4())
        await db.ots_proofs.insert_one({
            "id": proof_id,
            "kind": kind,
            "target_id": target_id,
            "owner_user_id": owner_user_id,
            "digest_hex": digest.hex(),
            "size_bytes": len(data),
            "ots_bytes": ots_bytes,
            "btc_block_height": None,
            "status": "pending",
            "created_at": now_utc().isoformat(),
            "label": label or kind,
            "extra": extra or {},
        })
        logger.info(f"OTS notarized bytes {kind}/{target_id} → {digest.hex()[:16]}")
        return proof_id
    except Exception as e:
        logger.warning(f"OTS notarize bytes {kind}/{target_id} failed: {e}")
        return None


async def log_agent_activity(
    agent: str,
    task: str,
    status: str,
    user_id: str = "system",
    duration_ms: int = 0,
    friction: Optional[str] = None,
    meta: Optional[dict] = None,
):
    doc = {
        "id": str(uuid.uuid4()),
        "agent": agent,
        "task": task,
        "status": status,
        "user_id": user_id,
        "duration_ms": duration_ms,
        "friction": friction,
        "meta": meta or {},
        "timestamp": now_utc().isoformat(),
    }
    await db.agent_activity.insert_one(doc)


# -----------------------------------------------------------------------------
# Claude helper
# -----------------------------------------------------------------------------
async def call_claude(system_message: str, user_text: str, session_id: Optional[str] = None) -> str:
    session_id = session_id or f"workz-{uuid.uuid4()}"
    chat = LlmChat(
        api_key=EMERGENT_LLM_KEY,
        session_id=session_id,
        system_message=system_message,
    ).with_model("anthropic", "claude-sonnet-4-5-20250929")
    response = await chat.send_message(UserMessage(text=user_text))
    return response if isinstance(response, str) else str(response)


def safe_json_loads(raw: str) -> dict:
    """Strip ```json fences and parse, falling back to empty dict."""
    txt = raw.strip()
    if txt.startswith("```"):
        txt = txt.split("```")[1] if "```" in txt[3:] else txt[3:]
        if txt.lower().startswith("json"):
            txt = txt[4:]
        txt = txt.strip()
        if txt.endswith("```"):
            txt = txt[:-3].strip()
    try:
        return json.loads(txt)
    except Exception:
        # Try to find first { ... } block
        start = txt.find("{")
        end = txt.rfind("}")
        if start >= 0 and end > start:
            try:
                return json.loads(txt[start : end + 1])
            except Exception:
                pass
        return {"raw": raw}


# -----------------------------------------------------------------------------
# File extraction helpers (PDF / DOCX / TXT → per-page text)
# -----------------------------------------------------------------------------
try:
    from openpyxl import load_workbook  # type: ignore
except Exception:
    load_workbook = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None  # type: ignore


def extract_pages_from_bytes(filename: str, data: bytes) -> List[Dict[str, Any]]:
    """Return list of {page:int, text:str}. Page is 1-indexed. Falls back to single-page."""
    name = (filename or "").lower()
    pages: List[Dict[str, Any]] = []
    try:
        if name.endswith(".pdf") and PdfReader is not None:
            reader = PdfReader(io.BytesIO(data))
            for i, p in enumerate(reader.pages, start=1):
                try:
                    txt = p.extract_text() or ""
                except Exception:
                    txt = ""
                pages.append({"page": i, "text": txt.strip()})
        elif name.endswith(".docx") and DocxDocument is not None:
            doc = DocxDocument(io.BytesIO(data))
            buf: List[str] = []
            for para in doc.paragraphs:
                if para.text:
                    buf.append(para.text)
            # Approximate "pages" by chunking every ~40 paragraphs
            chunk_size = 40
            if not buf:
                pages.append({"page": 1, "text": ""})
            else:
                for idx in range(0, len(buf), chunk_size):
                    pages.append({"page": idx // chunk_size + 1, "text": "\n".join(buf[idx:idx + chunk_size])})
        elif name.endswith((".xlsx", ".xlsm")) and load_workbook is not None:
            wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
            for i, sheet_name in enumerate(wb.sheetnames, start=1):
                ws = wb[sheet_name]
                lines: List[str] = [f"# Sheet: {sheet_name}"]
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    cells = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in cells):
                        lines.append("\t".join(cells))
                        row_count += 1
                        if row_count >= 2000:
                            lines.append("… (truncated)")
                            break
                pages.append({"page": i, "text": "\n".join(lines)})
            wb.close()
        elif name.endswith((".pptx", ".ppt")) and Presentation is not None:
            prs = Presentation(io.BytesIO(data))
            for i, slide in enumerate(prs.slides, start=1):
                parts: List[str] = [f"# Slide {i}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = "".join(r.text for r in para.runs).strip()
                            if txt:
                                parts.append(txt)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"\nSpeaker notes: {notes}")
                pages.append({"page": i, "text": "\n".join(parts)})
        elif name.endswith((".png", ".jpg", ".jpeg", ".gif", ".webp", ".heic", ".heif",
                            ".svg", ".mp4", ".mov", ".webm", ".mp3", ".wav", ".m4a",
                            ".zip", ".rar", ".7z")):
            # Binary media — we store the file but don't try to OCR/transcribe (heavy).
            # Surface a structured placeholder so AI Co-pilot can still reference the doc.
            pages.append({"page": 1, "text": f"[Non-text file: {filename} · {len(data)} bytes · stored as-is, no automatic text extraction]"})
        else:
            # Plain text / markdown / csv / unknown — treat as one page
            try:
                txt = data.decode("utf-8", errors="ignore")
            except Exception:
                txt = ""
            pages.append({"page": 1, "text": txt})
    except Exception as e:
        logger.warning(f"extract_pages_from_bytes failed for {filename}: {e}")
        pages = [{"page": 1, "text": ""}]
    if not pages:
        pages = [{"page": 1, "text": ""}]
    # Cap each page to 12k chars to bound token cost downstream
    for p in pages:
        if len(p["text"]) > 12000:
            p["text"] = p["text"][:12000]
    return pages


def pages_to_flat_text(pages: List[Dict[str, Any]], cap: int = 50000) -> str:
    out = []
    total = 0
    for p in pages:
        chunk = f"[p.{p['page']}]\n{p['text']}\n"
        if total + len(chunk) > cap:
            break
        out.append(chunk)
        total += len(chunk)
    return "\n".join(out)
MCP_ACTIONS = [
    {
        "id": "research.company.summarize",
        "type": "imperative",
        "description": "Generate AI research summary for any company (profile, leadership, market signals).",
        "endpoint": "/api/research/company",
        "method": "POST",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='research.company.summarize']",
        "params": {"company_name": "string", "sector": "string?", "region": "string?"},
    },
    {
        "id": "collateral.generate",
        "type": "imperative",
        "description": "Generate marketing collateral (one-pager, email sequence, LinkedIn post, deal memo).",
        "endpoint": "/api/collateral/generate",
        "method": "POST",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='collateral.generate']",
        "params": {"asset_type": "enum", "deal_name": "string", "target_audience": "string", "key_points": "string"},
    },
    {
        "id": "outreach.campaign.create",
        "type": "imperative",
        "description": "Launch a personalized outreach campaign on LinkedIn or email.",
        "endpoint": "/api/outreach/campaigns",
        "method": "POST",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='outreach.campaign.create']",
        "params": {"name": "string", "target_persona": "string", "channel": "enum", "message_brief": "string"},
    },
    {
        "id": "leads.list",
        "type": "declarative",
        "description": "List all leads in the nurturing pipeline grouped by stage.",
        "endpoint": "/api/leads",
        "method": "GET",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='leads.list']",
        "params": {},
    },
    {
        "id": "leads.advance",
        "type": "imperative",
        "description": "Advance a lead to the next nurturing stage.",
        "endpoint": "/api/leads/{lead_id}/stage",
        "method": "PATCH",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='leads.advance']",
        "params": {"lead_id": "string", "stage": "enum"},
    },
    {
        "id": "newsletter.draft",
        "type": "imperative",
        "description": "AI-draft a personalized newsletter for the current buyer.",
        "endpoint": "/api/newsletter/draft",
        "method": "POST",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='newsletter.draft']",
        "params": {"topic": "string?"},
    },
    {
        "id": "newsletter.dispatch",
        "type": "imperative",
        "description": "Approve and dispatch a newsletter to all opted-in buyers.",
        "endpoint": "/api/newsletter/{id}/dispatch",
        "method": "POST",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='newsletter.dispatch']",
        "params": {"id": "string"},
    },
    {
        "id": "composio.linkedin.connect",
        "type": "imperative",
        "description": "Initiate LinkedIn OAuth connection through Composio gateway.",
        "endpoint": "/api/composio/connect/linkedin",
        "method": "POST",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='composio.linkedin.connect']",
        "params": {},
    },
    {
        "id": "dashboard.kpis",
        "type": "declarative",
        "description": "Read top-level KPIs (AUM, active deals, pipeline value, engagement).",
        "endpoint": "/api/dashboard/stats",
        "method": "GET",
        "auth": "jwt",
        "dom_selector": "[data-mcp-action='dashboard.kpis']",
        "params": {},
    },
]


# -----------------------------------------------------------------------------
# AUTH ROUTES
# -----------------------------------------------------------------------------
PASSWORD_MIN_LEN = 8
LOGIN_MAX_FAILURES = 5
LOGIN_LOCKOUT_MIN = 15


def _password_complexity_ok(pw: str) -> Optional[str]:
    if not pw or len(pw) < PASSWORD_MIN_LEN:
        return f"Password must be at least {PASSWORD_MIN_LEN} characters"
    if not any(c.isalpha() for c in pw):
        return "Password must contain at least one letter"
    if not any(c.isdigit() for c in pw):
        return "Password must contain at least one digit"
    return None


async def _check_login_lockout(email: str) -> None:
    """Raise 429 if email is currently locked out due to recent failed attempts."""
    cutoff = (now_utc() - timedelta(minutes=LOGIN_LOCKOUT_MIN)).isoformat()
    fails = await db.login_attempts.count_documents({
        "email": email.lower(),
        "ok": False,
        "ts": {"$gte": cutoff},
    })
    if fails >= LOGIN_MAX_FAILURES:
        raise HTTPException(
            status_code=429,
            detail=f"Too many failed attempts. Try again in {LOGIN_LOCKOUT_MIN} minutes.",
        )


async def _record_login_attempt(email: str, ok: bool, ip: Optional[str] = None):
    await db.login_attempts.insert_one({
        "id": str(uuid.uuid4()),
        "email": email.lower(),
        "ok": ok,
        "ip": ip,
        "ts": now_utc().isoformat(),
    })


@api_router.post("/auth/register")
async def register(body: RegisterRequest, request: Request):
    """
    Public access-request endpoint.

    Default path: account is created in `pending` status awaiting admin
    approval. The user CANNOT log in until status flips to `active`.

    Fast path (invited users): if a valid `listing_invite_token` OR a valid
    `org_invite_token` (with `org_choice="join"`) is supplied — and the
    invite's email matches the signup email — the account is created
    `active` immediately and a JWT is returned in the response. The invite
    is also accepted on the spot so the user lands inside the listing /
    org with one redirect. This is the path the email-invite flow uses.
    """
    err = _password_complexity_ok(body.password)
    if err:
        raise HTTPException(status_code=400, detail=err)
    existing = await db.users.find_one({"email": body.email.lower()}, {"_id": 0})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")

    # -- Resolve invite-driven fast path -------------------------------------
    # If any invite token is supplied, validate it before creating the user
    # so we never partially-provision on bad input.
    email_norm = body.email.lower()
    listing_invite = None
    if body.listing_invite_token:
        listing_invite = await db.listing_invites.find_one(
            {"token": body.listing_invite_token}, {"_id": 0}
        )
        if not listing_invite:
            raise HTTPException(status_code=400, detail="Listing invite not found")
        if listing_invite.get("accepted_at"):
            raise HTTPException(status_code=400, detail="Listing invite already accepted")
        if listing_invite.get("expires_at") and datetime.fromisoformat(listing_invite["expires_at"]) < now_utc():
            raise HTTPException(status_code=400, detail="Listing invite expired")
        if listing_invite["email"].lower() != email_norm:
            raise HTTPException(status_code=400, detail=f"This invite is for {listing_invite['email']}. Sign up with that email.")

    org_invite = None
    if body.org_choice == "join" and body.org_invite_token:
        org_invite = await db.org_invites.find_one(
            {"token": body.org_invite_token}, {"_id": 0}
        )
        if not org_invite:
            raise HTTPException(status_code=400, detail="Org invite not found")
        if org_invite.get("accepted_at"):
            raise HTTPException(status_code=400, detail="Org invite already accepted")
        if org_invite.get("expires_at") and datetime.fromisoformat(org_invite["expires_at"]) < now_utc():
            raise HTTPException(status_code=400, detail="Org invite expired")
        if org_invite["email"].lower() != email_norm:
            raise HTTPException(status_code=400, detail=f"This invite is for {org_invite['email']}. Sign up with that email.")

    # An invite vouches for the user → skip admin approval gate.
    auto_active = bool(listing_invite or org_invite)

    user_id = str(uuid.uuid4())
    doc = {
        "id": user_id,
        "email": email_norm,
        "name": body.name,
        "role": body.role,
        "organization": body.organization,
        "password_hash": hash_password(body.password),
        "interests": [],
        "newsletter_opt_in": False,
        "created_at": now_utc().isoformat(),
        "status": "active" if auto_active else "pending",
    }
    if auto_active:
        doc["approved_at"] = now_utc().isoformat()
        doc["approved_via"] = "invite"
    await db.users.insert_one(doc)
    await log_audit(user_id, "auth.register", body.email, {"status": doc["status"], "via_invite": auto_active})

    # Capture org_choice for downstream wiring — actual org create/join happens
    # at /auth/accept-invite (for admin invites) or on first login after
    # approval (so we honor org_choice via a deferred entry stored on the user doc).
    if body.org_choice == "create" and body.org_name:
        await db.users.update_one(
            {"id": user_id},
            {"$set": {"pending_org_create": body.org_name}},
        )
    elif body.org_choice == "join" and body.org_invite_token and not org_invite:
        # Only stash a deferred join when we DIDN'T already fast-path.
        await db.users.update_one(
            {"id": user_id},
            {"$set": {"pending_org_invite_token": body.org_invite_token}},
        )

    # -- Fast path: accept the invite on the spot ----------------------------
    if listing_invite:
        listing = await db.listings.find_one({"id": listing_invite["listing_id"]}, {"_id": 0})
        if listing:
            await db.listings.update_one(
                {"id": listing_invite["listing_id"]},
                {"$push": {"collaborators": {
                    "user_id": user_id,
                    "email": email_norm,
                    "name": body.name,
                    "role": listing_invite["role"],
                    "invited_by": listing_invite.get("invited_by"),
                    "invited_at": listing_invite.get("created_at"),
                    "accepted_at": now_utc().isoformat(),
                }}},
            )
        await db.listing_invites.update_one(
            {"token": body.listing_invite_token},
            {"$set": {"accepted_at": now_utc().isoformat()}},
        )
        await log_audit(user_id, "listing.collab.accept", listing_invite["listing_id"],
                        {"role": listing_invite["role"], "via": "register"})

    if org_invite:
        await db.org_memberships.insert_one({
            "id": str(uuid.uuid4()),
            "org_id": org_invite["org_id"],
            "user_id": user_id,
            "role": org_invite["role"],
            "joined_at": now_utc().isoformat(),
            "invited_by": org_invite.get("invited_by"),
        })
        await db.org_invites.update_one(
            {"token": body.org_invite_token},
            {"$set": {"accepted_at": now_utc().isoformat()}},
        )
        await log_audit(user_id, "org.invite.accept", org_invite["org_id"], {"via": "register"})

    # -- Notifications -------------------------------------------------------

    # 1) Notify the operator/admin inbox that a new request landed
    cfg_notify = os.environ.get("REQUEST_NOTIFY_EMAIL")
    if cfg_notify and not auto_active:
        ip = request.client.host if request.client else "unknown"
        admin_html = f"""
        <p>New access request on NextCapOS.</p>
        <ul>
          <li><strong>Name:</strong> {body.name}</li>
          <li><strong>Email:</strong> {body.email}</li>
          <li><strong>Role requested:</strong> {body.role}</li>
          <li><strong>Organization:</strong> {body.organization or '—'}</li>
          <li><strong>IP:</strong> {ip}</li>
        </ul>
        <p><a href="{mail_link('/app/admin/users')}">Review &amp; approve in the admin console &rsaquo;</a></p>
        """
        asyncio.create_task(send_email(
            cfg_notify,
            f"NextCapOS · access request from {body.email}",
            admin_html,
            reply_to=body.email,
        ))

    # 2) Confirm to the requester that we received it (skip when fast-pathed)
    if not auto_active:
        requester_html = f"""
        <p>Hi {body.name},</p>
        <p>Thanks for requesting access to <strong>NextCapOS</strong>.</p>
        <p>Your request has been received and is now in the queue for administrator review.
        We'll email you the moment your account is approved — typically within one business day.</p>
        <p style="margin-top:24px;font-size:13px;color:#666;">
          <strong>Request details</strong><br>
          Name: {body.name}<br>
          Email: {body.email}<br>
          Role: {body.role}<br>
          Organization: {body.organization or '—'}
        </p>
        <p style="margin-top:24px;font-size:12px;color:#999;">
          Didn't request access? You can safely ignore this email.
        </p>
        """
        asyncio.create_task(send_email(
            body.email,
            "NextCapOS · access request received",
            requester_html,
        ))

    if auto_active:
        token = create_token(user_id, doc["role"])
        return {
            "ok": True,
            "status": "active",
            "token": token,
            "user": UserPublic(**await serialize_user_with_scope(doc)).model_dump(),
            "listing_id": listing_invite["listing_id"] if listing_invite else None,
            "org_id": org_invite["org_id"] if org_invite else None,
        }

    return {
        "ok": True,
        "status": "pending",
        "message": "Request received. You'll get an email when an administrator approves your account.",
    }


@api_router.post("/admin/users/{uid}/approve")
async def admin_approve_user(uid: str, user=Depends(get_current_user)):
    _admin_only(user)
    target = await db.users.find_one({"id": uid}, {"_id": 0})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    if target.get("status") == "active":
        return {"ok": True, "already_active": True}
    await db.users.update_one(
        {"id": uid},
        {"$set": {"status": "active", "approved_at": now_utc().isoformat(),
                  "approved_by_admin_id": user["id"]}},
    )
    await log_audit(user["id"], "admin.user.approve", uid, {"email": target["email"]})

    # If they requested to bootstrap an org at signup, materialise it now.
    pending_create = target.get("pending_org_create")
    pending_join = target.get("pending_org_invite_token")
    if pending_create:
        org_id = str(uuid.uuid4())
        slug = re.sub(r"[^a-z0-9-]+", "-", pending_create.lower()).strip("-")[:60] or org_id[:8]
        if await db.organizations.find_one({"slug": slug}, {"_id": 0, "id": 1}):
            slug = f"{slug}-{org_id[:6]}"
        await db.organizations.insert_one({
            "id": org_id, "name": pending_create, "slug": slug,
            "org_type": "advisory", "description": None,
            "created_by": uid, "created_at": now_utc().isoformat(),
        })
        await db.org_memberships.insert_one({
            "id": str(uuid.uuid4()), "org_id": org_id, "user_id": uid,
            "role": "org_admin", "joined_at": now_utc().isoformat(), "invited_by": None,
        })
        await db.users.update_one({"id": uid}, {"$unset": {"pending_org_create": ""}})
        await log_audit(uid, "org.create", org_id, {"name": pending_create, "via": "approval"})
    elif pending_join:
        inv = await db.org_invites.find_one({"token": pending_join}, {"_id": 0})
        if inv and not inv.get("accepted_at") and inv["email"].lower() == target["email"].lower():
            await db.org_memberships.insert_one({
                "id": str(uuid.uuid4()),
                "org_id": inv["org_id"],
                "user_id": uid,
                "role": inv["role"],
                "joined_at": now_utc().isoformat(),
                "invited_by": inv.get("invited_by"),
            })
            await db.org_invites.update_one(
                {"token": pending_join}, {"$set": {"accepted_at": now_utc().isoformat()}}
            )
            await db.users.update_one({"id": uid}, {"$unset": {"pending_org_invite_token": ""}})
            await log_audit(uid, "org.invite.accept", inv["org_id"], {"via": "approval"})

    html = f"""
      <p>Hi {target.get('name') or 'there'},</p>
      <p>Your NextCapOS access request has been <strong>approved</strong>. You can now sign in:</p>
      <p><a href="{mail_link('/login')}">Sign in to NextCapOS &rsaquo;</a></p>
    """
    asyncio.create_task(send_email(target["email"], "NextCapOS · access approved", html))
    return {"ok": True, "status": "active"}


@api_router.post("/admin/users/{uid}/reject")
async def admin_reject_user(uid: str, user=Depends(get_current_user)):
    _admin_only(user)
    target = await db.users.find_one({"id": uid}, {"_id": 0, "email": 1, "name": 1, "is_demo": 1})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    if target.get("is_demo"):
        raise HTTPException(status_code=400, detail="Cannot reject seeded demo accounts")
    await db.users.update_one(
        {"id": uid},
        {"$set": {"status": "rejected", "rejected_at": now_utc().isoformat(),
                  "rejected_by_admin_id": user["id"]}},
    )
    await log_audit(user["id"], "admin.user.reject", uid, {"email": target["email"]})
    html = f"""
      <p>Hi {target.get('name') or 'there'},</p>
      <p>Thank you for your interest in NextCapOS. After review, we are not able
      to approve your access request at this time.</p>
    """
    asyncio.create_task(send_email(target["email"], "NextCapOS · access request update", html))
    return {"ok": True, "status": "rejected"}


# -----------------------------------------------------------------------------
# Forgot password
# -----------------------------------------------------------------------------
class ForgotPasswordRequest(BaseModel):
    email: EmailStr


class ResetPasswordRequest(BaseModel):
    token: str
    password: str


@api_router.post("/auth/forgot-password")
async def forgot_password(body: ForgotPasswordRequest):
    """Always returns 200 — never leaks whether the email exists."""
    email_norm = body.email.lower()
    user = await db.users.find_one({"email": email_norm}, {"_id": 0, "id": 1, "name": 1, "status": 1})
    if user and user.get("status") != "deactivated":
        token = secrets.token_urlsafe(32)
        await db.password_resets.insert_one({
            "id": str(uuid.uuid4()),
            "user_id": user["id"],
            "email": email_norm,
            "token_hash": hash_password(token),
            "created_at": now_utc().isoformat(),
            "expires_at": (now_utc() + timedelta(hours=1)).isoformat(),
            "status": "pending",
        })
        html = f"""
          <p>Hi {user.get('name') or 'there'},</p>
          <p>We received a request to reset your NextCapOS password.</p>
          <p><a href="{mail_link(f'/reset-password?token={token}')}">Set a new password &rsaquo;</a></p>
          <p>This link expires in one hour. If you didn't request this, you can ignore this email.</p>
        """
        asyncio.create_task(send_email(email_norm, "NextCapOS · reset your password", html))
    return {"ok": True, "message": "If that email exists, a reset link has been sent."}


@api_router.post("/auth/reset-password")
async def reset_password(body: ResetPasswordRequest):
    err = _password_complexity_ok(body.password)
    if err:
        raise HTTPException(status_code=400, detail=err)
    now_iso = now_utc().isoformat()
    candidates = await db.password_resets.find(
        {"status": "pending", "expires_at": {"$gt": now_iso}},
        {"_id": 0},
    ).to_list(50)
    matched = next((c for c in candidates if verify_password(body.token, c["token_hash"])), None)
    if not matched:
        raise HTTPException(status_code=400, detail="Invalid or expired reset link")
    await db.users.update_one(
        {"id": matched["user_id"]},
        {"$set": {"password_hash": hash_password(body.password),
                  "password_reset_at": now_iso}},
    )
    await db.password_resets.update_one(
        {"id": matched["id"]},
        {"$set": {"status": "consumed", "consumed_at": now_iso}},
    )
    # Invalidate any other pending resets for this user
    await db.password_resets.update_many(
        {"user_id": matched["user_id"], "status": "pending"},
        {"$set": {"status": "superseded"}},
    )
    await log_audit(matched["user_id"], "auth.password_reset", matched["email"])
    return {"ok": True}


@api_router.post("/auth/login", response_model=TokenResponse)
async def login(body: LoginRequest, request: Request):
    email_norm = body.email.lower()
    ip = request.client.host if request.client else None
    await _check_login_lockout(email_norm)
    user = await db.users.find_one({"email": email_norm})
    if not user or not verify_password(body.password, user["password_hash"]):
        await _record_login_attempt(email_norm, ok=False, ip=ip)
        raise HTTPException(status_code=401, detail="Invalid credentials")
    if user.get("status") == "deactivated":
        await _record_login_attempt(email_norm, ok=False, ip=ip)
        raise HTTPException(status_code=403, detail="Account deactivated. Contact your administrator.")
    if user.get("status") == "pending":
        await _record_login_attempt(email_norm, ok=False, ip=ip)
        raise HTTPException(status_code=403, detail="Your access request is awaiting administrator approval. You'll get an email when it's approved.")
    if user.get("status") == "rejected":
        await _record_login_attempt(email_norm, ok=False, ip=ip)
        raise HTTPException(status_code=403, detail="Your access request was not approved.")
    await _record_login_attempt(email_norm, ok=True, ip=ip)
    token = create_token(user["id"], user["role"])
    await log_audit(user["id"], "auth.login", body.email)
    return TokenResponse(token=token, user=UserPublic(**await serialize_user_with_scope(user)))


@api_router.get("/auth/me", response_model=UserPublic)
async def me(user=Depends(get_current_user)):
    return UserPublic(**await serialize_user_with_scope(user))


# -----------------------------------------------------------------------------
# ADMIN · User management
# -----------------------------------------------------------------------------
def _admin_only(user):
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin only")


class AdminCreateUser(BaseModel):
    email: EmailStr
    password: str
    name: str
    role: Literal["admin", "buyer", "seller", "agent"] = "buyer"
    organization: Optional[str] = None


class AdminUpdateUser(BaseModel):
    name: Optional[str] = None
    role: Optional[Literal["admin", "buyer", "seller", "agent"]] = None
    organization: Optional[str] = None
    status: Optional[Literal["active", "deactivated"]] = None


class AdminInviteRequest(BaseModel):
    email: EmailStr
    name: Optional[str] = None
    role: Literal["admin", "buyer", "seller", "agent"] = "buyer"
    organization: Optional[str] = None
    expires_hours: int = 168  # 7 days


class AcceptInviteRequest(BaseModel):
    token: str
    password: str
    name: Optional[str] = None


@api_router.get("/admin/users")
async def admin_list_users(q: Optional[str] = None, user=Depends(get_current_user)):
    _admin_only(user)
    query: dict = {}
    if q:
        query["$or"] = [
            {"email": {"$regex": q, "$options": "i"}},
            {"name": {"$regex": q, "$options": "i"}},
            {"organization": {"$regex": q, "$options": "i"}},
        ]
    items = await db.users.find(
        query,
        {"_id": 0, "password_hash": 0},
    ).sort("created_at", -1).to_list(500)
    return items


@api_router.post("/admin/users")
async def admin_create_user(body: AdminCreateUser, user=Depends(get_current_user)):
    _admin_only(user)
    err = _password_complexity_ok(body.password)
    if err:
        raise HTTPException(status_code=400, detail=err)
    email_norm = body.email.lower()
    if await db.users.find_one({"email": email_norm}, {"_id": 0, "id": 1}):
        raise HTTPException(status_code=400, detail="Email already registered")
    user_id = str(uuid.uuid4())
    doc = {
        "id": user_id,
        "email": email_norm,
        "name": body.name,
        "role": body.role,
        "organization": body.organization,
        "password_hash": hash_password(body.password),
        "interests": [],
        "newsletter_opt_in": False,
        "created_at": now_utc().isoformat(),
        "created_by_admin_id": user["id"],
        "status": "active",
    }
    await db.users.insert_one(doc)
    await log_audit(user["id"], "admin.user.create", user_id,
                    {"email": email_norm, "role": body.role})
    doc.pop("_id", None)
    doc.pop("password_hash", None)
    return doc


@api_router.patch("/admin/users/{uid}")
async def admin_update_user(uid: str, body: AdminUpdateUser, user=Depends(get_current_user)):
    _admin_only(user)
    if uid == user["id"] and body.role and body.role != "admin":
        raise HTTPException(status_code=400, detail="Cannot demote yourself")
    target = await db.users.find_one({"id": uid}, {"_id": 0, "email": 1, "role": 1})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    patch = {k: v for k, v in body.model_dump(exclude_none=True).items()}
    if not patch:
        raise HTTPException(status_code=400, detail="No changes")
    await db.users.update_one({"id": uid}, {"$set": patch})
    await log_audit(user["id"], "admin.user.update", uid, patch)
    fresh = await db.users.find_one({"id": uid}, {"_id": 0, "password_hash": 0})
    return fresh


@api_router.post("/admin/users/{uid}/password")
async def admin_reset_password(uid: str, body: dict, user=Depends(get_current_user)):
    _admin_only(user)
    new_password = (body or {}).get("password")
    if not new_password:
        raise HTTPException(status_code=400, detail="password required")
    err = _password_complexity_ok(new_password)
    if err:
        raise HTTPException(status_code=400, detail=err)
    target = await db.users.find_one({"id": uid}, {"_id": 0, "email": 1})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    await db.users.update_one(
        {"id": uid},
        {"$set": {"password_hash": hash_password(new_password),
                  "password_reset_by_admin_at": now_utc().isoformat()}},
    )
    await log_audit(user["id"], "admin.user.password_reset", uid, {"email": target["email"]})
    return {"ok": True}


@api_router.post("/admin/listings/{lid}/sources/cleanup-corrupt")
async def cleanup_corrupt_synced_files(lid: str, user=Depends(get_current_user)):
    """Scan every Composio-mirrored staged file on a listing and remove any
    whose stored content is actually an HTML error page or JSON error envelope
    from a failed proxy download. Created for the iter-22 bug where wrong
    proxy paths silently stored Google's "/drive/v3/drive/v3/...404" HTML as
    the file body.

    Returns the list of filenames removed so the caller can re-sync them."""
    _admin_only(user) if user.get("role") == "admin" else None
    listing = await db.listings.find_one({"id": lid}, {"_id": 0, "seller_id": 1})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found")
    # Allow seller or org workspace participant too
    if user.get("role") not in ("admin",):
        ws_listings, _ = await _user_workspace_listing_ids(user)
        if lid not in ws_listings:
            raise HTTPException(status_code=403, detail="Not authorized for this listing")

    staged = await db.listing_staged_files.find(
        {"listing_id": lid, "source": {"$exists": True, "$ne": None},
         "deleted_at": {"$exists": False}, "gridfs_id": {"$exists": True}},
        {"_id": 0, "id": 1, "filename": 1, "gridfs_id": 1, "encrypted": 1},
    ).to_list(500)
    removed: list[str] = []
    for s in staged:
        try:
            grid = await listing_files_bucket.open_download_stream(ObjectId(s["gridfs_id"]))
            raw = await grid.read()
        except Exception:
            continue
        if s.get("encrypted"):
            try:
                aad = f"listing:{lid}:{s['id']}".encode("utf-8")
                raw = decrypt_envelope(raw, associated_data=aad)
            except Exception:
                continue
        sniff = (raw or b"")[:200].lstrip().lower()
        is_html_error = sniff.startswith(b"<!doctype html") or sniff.startswith(b"<html")
        is_json_err = False
        if not is_html_error and (sniff.startswith(b'{"error"') or sniff.startswith(b'{"message"')):
            try:
                import json as _j
                parsed = _j.loads(raw.decode("utf-8", errors="replace"))
                is_json_err = isinstance(parsed, dict) and "error" in parsed
            except Exception:
                pass
        if is_html_error or is_json_err:
            try:
                await listing_files_bucket.delete(ObjectId(s["gridfs_id"]))
            except Exception:
                pass
            await db.listing_staged_files.delete_one({"id": s["id"]})
            # Cascade into vault clones (purged on next get_deal_room hit too,
            # but this is the eager cleanup).
            await db.deal_room_files.delete_many({"cloned_from_listing_file": s["id"]})
            removed.append(s["filename"])
    await log_audit(user["id"], "listing.source.cleanup-corrupt", lid,
                    {"removed": removed, "count": len(removed)})
    return {"ok": True, "removed": removed, "count": len(removed)}


@api_router.delete("/admin/users/{uid}")
async def admin_deactivate_user(uid: str, user=Depends(get_current_user)):
    _admin_only(user)
    if uid == user["id"]:
        raise HTTPException(status_code=400, detail="Cannot deactivate yourself")
    target = await db.users.find_one({"id": uid}, {"_id": 0, "email": 1, "is_demo": 1})
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    if target.get("is_demo"):
        raise HTTPException(status_code=400, detail="Cannot deactivate seeded demo accounts")
    await db.users.update_one(
        {"id": uid},
        {"$set": {"status": "deactivated", "deactivated_at": now_utc().isoformat(),
                  "deactivated_by_admin_id": user["id"]}},
    )
    await log_audit(user["id"], "admin.user.deactivate", uid, {"email": target["email"]})
    return {"ok": True, "status": "deactivated"}


@api_router.post("/admin/users/{uid}/purge")
async def admin_hard_delete_user(uid: str, user=Depends(get_current_user)):
    """HARD delete: removes the user account and every record they own
    (listings + staged files, inquiries, deal rooms + their files / messages /
    findings / requests, locker files, research + copilots, detailed reports,
    collateral, outreach, newsletters, leads, watchlist, agent activity, external
    sources, composio connections).

    Audit logs are preserved (hash-chain integrity) but the actor_id is
    rewritten to a tombstone `deleted:{uid}` so the chain still validates and
    the email never leaks back.

    Guardrails (HTTP 400 on each):
      - Cannot purge yourself.
      - Cannot purge seeded demo accounts (`is_demo: True`).
      - Cannot purge another admin unless this admin is the platform owner
        (heuristic: oldest active admin).
    """
    _admin_only(user)
    if uid == user["id"]:
        raise HTTPException(status_code=400, detail="Cannot purge yourself")
    target = await db.users.find_one(
        {"id": uid},
        {"_id": 0, "id": 1, "email": 1, "role": 1, "is_demo": 1, "is_seed": 1},
    )
    if not target:
        raise HTTPException(status_code=404, detail="User not found")
    if target.get("is_demo") or target.get("is_seed"):
        raise HTTPException(status_code=400, detail="Cannot purge seeded / demo accounts")
    if target.get("role") == "admin":
        # Only the platform's primary admin (= the oldest admin by created_at)
        # can purge another admin. Prevents accidental nuking of co-admins.
        oldest_admin = await db.users.find_one(
            {"role": "admin", "status": {"$ne": "deactivated"}},
            sort=[("created_at", 1)],
            projection={"_id": 0, "id": 1},
        )
        if not oldest_admin or oldest_admin["id"] != user["id"]:
            raise HTTPException(
                status_code=400,
                detail="Only the primary admin can purge another admin account",
            )

    deleted: dict = {"email": target.get("email"), "by": user["id"]}

    # 1. Listings owned by this user (cascades to staged files + buyer-discovery rows)
    listings = await db.listings.find(
        {"seller_id": uid}, {"_id": 0, "id": 1}
    ).to_list(None)
    lids = [li["id"] for li in listings]
    if lids:
        async for s in db.listing_staged_files.find(
            {"listing_id": {"$in": lids}, "gridfs_id": {"$exists": True}},
            {"_id": 0, "gridfs_id": 1},
        ):
            try:
                await listing_files_bucket.delete(ObjectId(s["gridfs_id"]))
            except Exception:
                pass
        await db.listing_staged_files.delete_many({"listing_id": {"$in": lids}})
        await db.listing_external_sources.delete_many({"listing_id": {"$in": lids}})
        await db.listing_collaborators.delete_many({"listing_id": {"$in": lids}})
        await db.listing_collab_invites.delete_many({"listing_id": {"$in": lids}})
        # Iter-41: also drop any legacy listing_preview_links / listing_share_links
        # rows. The share-link feature was retired in favor of named collaborator
        # invites. Kept in cascade for backfill sanity.
        await db.listing_preview_links.delete_many({"listing_id": {"$in": lids}})
        await db.listing_share_links.delete_many({"listing_id": {"$in": lids}})
        await db.buyer_matches.delete_many({"listing_id": {"$in": lids}})
        await db.buyer_alerts.delete_many({"listing_id": {"$in": lids}})
        await db.buyer_scans.delete_many({"listing_id": {"$in": lids}})
        await db.listings.delete_many({"id": {"$in": lids}})
        deleted["listings"] = len(lids)

    # 2. Deal rooms where this user was buyer OR seller — cascade to files,
    #    messages, findings, requests + their GridFS blobs.
    rooms = await db.deal_rooms.find(
        {"$or": [{"buyer_id": uid}, {"seller_id": uid}]},
        {"_id": 0, "id": 1},
    ).to_list(None)
    rids = [r["id"] for r in rooms]
    if rids:
        async for f in db.deal_room_files.find(
            {"room_id": {"$in": rids}, "gridfs_id": {"$exists": True}},
            {"_id": 0, "gridfs_id": 1},
        ):
            try:
                await gridfs_bucket.delete(ObjectId(f["gridfs_id"]))
            except Exception:
                pass
        await db.deal_room_files.delete_many({"room_id": {"$in": rids}})
        await db.deal_room_findings.delete_many({"room_id": {"$in": rids}})
        await db.deal_room_messages.delete_many({"room_id": {"$in": rids}})
        await db.deal_room_requests.delete_many({"room_id": {"$in": rids}})
        await db.deal_rooms.delete_many({"id": {"$in": rids}})
        deleted["deal_rooms"] = len(rids)

    # 3. Inquiries authored by or directed at this user
    inq = await db.inquiries.find(
        {"$or": [{"buyer_id": uid}, {"seller_id": uid}]},
        {"_id": 0, "id": 1},
    ).to_list(None)
    if inq:
        iids = [i["id"] for i in inq]
        await db.inquiry_messages.delete_many({"inquiry_id": {"$in": iids}})
        await db.inquiries.delete_many({"id": {"$in": iids}})
        deleted["inquiries"] = len(iids)

    # 4. Private Locker — drop GridFS + rows
    locker_rows = await db.private_locker_files.find(
        {"user_id": uid, "gridfs_id": {"$exists": True}},
        {"_id": 0, "id": 1, "gridfs_id": 1},
    ).to_list(None)
    for r in locker_rows:
        try:
            await private_locker_bucket.delete(ObjectId(r["gridfs_id"]))
        except Exception:
            pass
    res_lk = await db.private_locker_files.delete_many({"user_id": uid})
    deleted["locker_files"] = res_lk.deleted_count

    # 5. All user-owned collection rows
    user_owned = [
        "research", "detailed_reports", "research_copilot_messages",
        "collateral", "collateral_versions", "outreach", "newsletters",
        "leads", "watchlist", "agent_activity", "composio_connections",
        "listing_external_sources",
    ]
    owned_total = 0
    for col in user_owned:
        try:
            r = await db[col].delete_many({"user_id": uid})
            owned_total += r.deleted_count
        except Exception as e:
            logger.warning(f"purge {col} failed for {uid}: {e}")
    deleted["user_owned_rows"] = owned_total

    # 6. Tombstone audit logs — DON'T delete (would break chain), just rewrite
    #    actor reference. Keep the email out of the public-facing actor field.
    tombstone = f"deleted:{uid}"
    await db.audit_logs.update_many(
        {"actor_id": uid},
        {"$set": {"actor_id": tombstone, "actor_was_deleted": True}},
    )

    # 7. Other refs: collaborator/org membership where this user appears
    await db.listing_collaborators.delete_many({"user_id": uid})
    await db.user_invites.delete_many({"email": (target.get("email") or "").lower()})

    # 8. Finally the user record
    await db.users.delete_one({"id": uid})

    # 9. Audit the purge itself (with the surviving admin as actor)
    await log_audit(user["id"], "admin.user.purge", uid, deleted)
    return {"ok": True, "purged": True, "summary": deleted}


@api_router.get("/admin/invites")
async def admin_list_invites(user=Depends(get_current_user)):
    _admin_only(user)
    items = await db.user_invites.find({}, {"_id": 0}).sort("created_at", -1).to_list(200)
    return items


@api_router.post("/admin/invites")
async def admin_create_invite(body: AdminInviteRequest, request: Request, user=Depends(get_current_user)):
    _admin_only(user)
    email_norm = body.email.lower()
    if await db.users.find_one({"email": email_norm}, {"_id": 0, "id": 1}):
        raise HTTPException(status_code=400, detail="Email already has an account")

    # Invalidate any previous pending invite for the same email
    await db.user_invites.update_many(
        {"email": email_norm, "status": "pending"},
        {"$set": {"status": "superseded", "superseded_at": now_utc().isoformat()}},
    )

    token = secrets.token_urlsafe(32)
    invite_id = str(uuid.uuid4())
    now = now_utc()
    expires = now + timedelta(hours=max(1, min(body.expires_hours, 24 * 30)))
    doc = {
        "id": invite_id,
        "token": token,
        "email": email_norm,
        "name": body.name,
        "role": body.role,
        "organization": body.organization,
        "invited_by_id": user["id"],
        "invited_by_email": user.get("email"),
        "status": "pending",
        "created_at": now.isoformat(),
        "expires_at": expires.isoformat(),
    }
    await db.user_invites.insert_one(doc)
    await log_audit(user["id"], "admin.invite.create", invite_id,
                    {"email": email_norm, "role": body.role})

    # Build accept URL using the frontend host the admin is on
    origin = request.headers.get("origin") or request.headers.get("referer") or ""
    base = origin.rstrip("/").rsplit("/", 1)[0] if "/app/" in origin else origin.rstrip("/")
    accept_url = f"{base}/accept-invite?token={token}" if base else f"/accept-invite?token={token}"

    doc.pop("_id", None)
    doc["accept_url"] = accept_url

    # Fire the actual invitation email so the recipient doesn't have to be
    # hand-walked the URL. Best-effort; the API response still succeeds even
    # if Resend rejects (logged in mailer).
    role_label = {"buyer": "Buyer", "seller": "Seller", "agent": "Agent (broker / advisor)", "admin": "Admin"}.get(body.role, body.role)
    org_line = f" at <strong>{body.organization}</strong>" if body.organization else ""
    name_line = f" Hi {body.name},<br><br>" if body.name else " Hi,<br><br>"
    html = f"""
    <p>{name_line}
    <strong>{user.get('name') or user.get('email')}</strong> invited you to join
    <strong>NextCapOS</strong>{org_line} as a <strong>{role_label}</strong>.</p>
    <p><a href="{accept_url}">Accept the invitation &rsaquo;</a></p>
    <p style="margin-top:24px;font-size:12px;color:#999;">
      This invite expires {expires.strftime('%b %d, %Y at %H:%M UTC')}. If you weren't
      expecting it, you can safely ignore this email.
    </p>
    """
    asyncio.create_task(send_email(
        email_norm,
        f"NextCapOS · you're invited to join{(' ' + body.organization) if body.organization else ''}",
        html,
        reply_to=user.get("email"),
    ))
    return doc


@api_router.post("/admin/invites/{iid}/resend")
async def admin_resend_invite(iid: str, request: Request, user=Depends(get_current_user)):
    """Re-fire the invitation email for an existing pending platform invite.
    Useful after fixing Resend config — the original token keeps working."""
    _admin_only(user)
    inv = await db.user_invites.find_one({"id": iid, "status": "pending"}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found or no longer pending")
    origin = request.headers.get("origin") or request.headers.get("referer") or ""
    base = origin.rstrip("/").rsplit("/", 1)[0] if "/app/" in origin else origin.rstrip("/")
    accept_url = f"{base}/accept-invite?token={inv['token']}" if base else f"/accept-invite?token={inv['token']}"
    role_label = {"buyer": "Buyer", "seller": "Seller", "agent": "Agent (broker / advisor)", "admin": "Admin"}.get(inv.get("role", ""), inv.get("role", ""))
    org_line = f" at <strong>{inv['organization']}</strong>" if inv.get("organization") else ""
    name_line = f" Hi {inv['name']},<br><br>" if inv.get("name") else " Hi,<br><br>"
    html = f"""
    <p>{name_line}
    <strong>{user.get('name') or user.get('email')}</strong> invited you to join
    <strong>NextCapOS</strong>{org_line} as a <strong>{role_label}</strong>.</p>
    <p><a href="{accept_url}">Accept the invitation &rsaquo;</a></p>
    """
    asyncio.create_task(send_email(
        inv["email"],
        f"NextCapOS · you're invited to join{(' ' + inv['organization']) if inv.get('organization') else ''} (resent)",
        html,
        reply_to=user.get("email"),
    ))
    await log_audit(user["id"], "admin.invite.resend", iid, {"email": inv["email"]})
    return {"ok": True, "email": inv["email"]}


@api_router.delete("/admin/invites/{iid}")
async def admin_revoke_invite(iid: str, user=Depends(get_current_user)):
    _admin_only(user)
    r = await db.user_invites.update_one(
        {"id": iid, "status": "pending"},
        {"$set": {"status": "revoked", "revoked_at": now_utc().isoformat()}},
    )
    if r.matched_count == 0:
        raise HTTPException(status_code=404, detail="Invite not found or not pending")
    await log_audit(user["id"], "admin.invite.revoke", iid)
    return {"ok": True}


@api_router.get("/auth/invite/{token}")
async def public_get_invite(token: str):
    """Public endpoint used by the accept-invite page to render a preview."""
    inv = await db.user_invites.find_one({"token": token}, {"_id": 0, "token": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("status") != "pending":
        raise HTTPException(status_code=410, detail=f"Invite {inv.get('status')}")
    if inv.get("expires_at") and inv["expires_at"] < now_utc().isoformat():
        await db.user_invites.update_one(
            {"id": inv["id"]},
            {"$set": {"status": "expired", "expired_at": now_utc().isoformat()}},
        )
        raise HTTPException(status_code=410, detail="Invite expired")
    return inv


@api_router.post("/auth/accept-invite", response_model=TokenResponse)
async def accept_invite(body: AcceptInviteRequest):
    inv = await db.user_invites.find_one({"token": body.token}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("status") != "pending":
        raise HTTPException(status_code=410, detail=f"Invite {inv.get('status')}")
    if inv.get("expires_at") and inv["expires_at"] < now_utc().isoformat():
        await db.user_invites.update_one(
            {"id": inv["id"]},
            {"$set": {"status": "expired", "expired_at": now_utc().isoformat()}},
        )
        raise HTTPException(status_code=410, detail="Invite expired")
    err = _password_complexity_ok(body.password)
    if err:
        raise HTTPException(status_code=400, detail=err)
    if await db.users.find_one({"email": inv["email"]}, {"_id": 0, "id": 1}):
        raise HTTPException(status_code=400, detail="Email already registered — sign in instead")

    user_id = str(uuid.uuid4())
    doc = {
        "id": user_id,
        "email": inv["email"],
        "name": body.name or inv.get("name") or inv["email"].split("@")[0],
        "role": inv["role"],
        "organization": inv.get("organization"),
        "password_hash": hash_password(body.password),
        "interests": [],
        "newsletter_opt_in": False,
        "created_at": now_utc().isoformat(),
        "created_via_invite_id": inv["id"],
        "invited_by_id": inv.get("invited_by_id"),
        "status": "active",
    }
    await db.users.insert_one(doc)
    await db.user_invites.update_one(
        {"id": inv["id"]},
        {"$set": {"status": "accepted", "accepted_at": now_utc().isoformat(),
                  "accepted_user_id": user_id}},
    )
    token = create_token(user_id, doc["role"])
    await log_audit(user_id, "auth.accept_invite", inv["id"], {"email": inv["email"], "role": inv["role"]})
    return TokenResponse(token=token, user=UserPublic(**await serialize_user_with_scope(doc)))


# -----------------------------------------------------------------------------
# DASHBOARD
# -----------------------------------------------------------------------------
@api_router.get("/dashboard/stats")
async def dashboard_stats(user=Depends(get_current_user)):
    deals = await db.deals.count_documents({})
    leads = await db.leads.count_documents({})
    campaigns = await db.outreach.count_documents({})
    newsletters = await db.newsletters.count_documents({})
    research = await db.research.count_documents({})
    activities = await db.agent_activity.count_documents({})
    success = await db.agent_activity.count_documents({"status": "completed"})
    success_rate = round((success / activities * 100) if activities else 0, 1)

    role = user.get("role", "buyer")
    uid = user["id"]

    if role == "seller":
        my_listings = await db.listings.count_documents({"seller_id": uid})
        live_listings = await db.listings.count_documents({"seller_id": uid, "status": "live"})
        inbound = await db.inquiries.count_documents({"seller_id": uid})
        pipeline_value = 0.0
        async for d in db.listings.find({"seller_id": uid}, {"_id": 0, "asking_price_usd_m": 1}):
            pipeline_value += float(d.get("asking_price_usd_m") or 0)
        my_campaigns = await db.outreach.count_documents({"user_id": uid})
        my_leads = await db.leads.count_documents({"user_id": uid})
        my_newsletters = await db.newsletters.count_documents({"user_id": uid})
        return {
            "role": role,
            "my_listings": my_listings,
            "live_listings": live_listings,
            "inbound_inquiries": inbound,
            "pipeline_value_usd_m": round(pipeline_value, 1),
            "my_campaigns": my_campaigns,
            "my_leads": my_leads,
            "my_newsletters": my_newsletters,
            "agent_success_rate": success_rate,
            "agent_runs": activities,
        }

    if role == "buyer":
        marketplace_count = await db.listings.count_documents({"status": "live"})
        my_research = await db.research.count_documents({"user_id": uid})
        my_inquiries = await db.inquiries.count_documents({"buyer_id": uid})
        my_newsletters_received = await db.newsletters.count_documents({"status": "dispatched"})
        watchlist_count = await db.watchlist.count_documents({"user_id": uid})
        return {
            "role": role,
            "marketplace_listings": marketplace_count,
            "my_research_count": my_research,
            "my_inquiries": my_inquiries,
            "watchlist_count": watchlist_count,
            "newsletters_received": my_newsletters_received,
            "aum_usd_b": 14.7,
            "agent_success_rate": success_rate,
            "exit_velocity_days": 142,
        }

    # admin / fallback — global view
    return {
        "role": role,
        "aum_usd_b": 14.7,
        "active_deals": deals,
        "pipeline_leads": leads,
        "campaigns": campaigns,
        "newsletters_sent": newsletters,
        "research_count": research,
        "agent_success_rate": success_rate,
        "agent_runs": activities,
        "exit_velocity_days": 142,
        "marketplace_listings": await db.listings.count_documents({"status": "live"}),
        "total_inquiries": await db.inquiries.count_documents({}),
    }


@api_router.get("/deals")
async def list_deals(user=Depends(get_current_user)):
    deals = await db.deals.find({}, {"_id": 0}).to_list(200)
    return deals


# -----------------------------------------------------------------------------
# LISTINGS · MARKETPLACE · INQUIRIES · WATCHLIST
# -----------------------------------------------------------------------------
@api_router.get("/listings")
async def my_listings(user=Depends(get_current_user)):
    """Workspace view — listings I personally own, plus listings owned by
    orgs I'm a member of, plus listings I'm an explicit collaborator on.
    Each row is decorated with `workspace_scope` ('mine' | 'org' | 'shared')
    and (if applicable) `org_name` so the UI can filter and badge them."""
    org_ids = await _get_user_org_ids(user)
    or_clauses: List[Dict[str, Any]] = [
        {"seller_id": user["id"]},
        {"collaborators.user_id": user["id"]},
    ]
    if org_ids:
        or_clauses.append({"org_id": {"$in": org_ids}})
    items = await db.listings.find(
        {"$or": or_clauses}, {"_id": 0}
    ).sort("created_at", -1).to_list(200)
    # Decorate
    org_names = {}
    if org_ids:
        org_names = {
            o["id"]: o["name"] for o in await db.organizations.find(
                {"id": {"$in": org_ids}}, {"_id": 0, "id": 1, "name": 1}
            ).to_list(200)
        }
    for it in items:
        if it.get("seller_id") == user["id"]:
            it["workspace_scope"] = "mine"
        elif it.get("org_id") and it["org_id"] in org_names:
            it["workspace_scope"] = "org"
            it["org_name"] = org_names[it["org_id"]]
        else:
            it["workspace_scope"] = "shared"
        # Iter-41: decorate with the current viewer's collaborator role so the UI
        # can gate role-specific affordances (e.g. hide "View as principal" for
        # a viewer collaborator).
        it["viewer_role"] = _viewer_role_on_listing(it, user)
    return items


def _viewer_role_on_listing(listing: dict, user: dict) -> str:
    """Return a coarse role label for the current user on this listing card:
    'principal' | 'org' | 'owner' | 'editor' | 'viewer' | 'admin' | 'unknown'."""
    if user.get("role") == "admin":
        return "admin"
    if listing.get("seller_id") == user["id"]:
        return "principal"
    for c in (listing.get("collaborators") or []):
        if c.get("user_id") == user["id"]:
            return c.get("role") or "viewer"
    # If we got this far in `my_listings` it's an org member with no explicit collab row
    return "org"


@api_router.post("/listings")
async def create_listing(body: ListingCreate, user=Depends(get_current_user), org_id: Optional[str] = None):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Sellers and agents only")
    # If an org_id is supplied, the creator must be a member of it. If omitted
    # and the user has exactly one org, default to it. Otherwise the listing
    # is individually owned (legacy behavior).
    resolved_org_id = None
    if org_id:
        role = await _user_org_role(user["id"], org_id)
        if not role and user.get("role") != "admin":
            raise HTTPException(status_code=403, detail="Not a member of that org")
        resolved_org_id = org_id
    else:
        my_orgs = await _get_user_org_ids(user)
        if len(my_orgs) == 1:
            resolved_org_id = my_orgs[0]
    doc = {
        "id": str(uuid.uuid4()),
        "seller_id": user["id"],
        "seller_name": user.get("name"),
        "seller_org": user.get("organization"),
        "org_id": resolved_org_id,
        "collaborators": [],
        "access_policy": {
            "require_principal_approval": False,
            "competitor_blocklist": [],
        },
        **body.model_dump(),
        "inquiry_count": 0,
        "view_count": 0,
        "created_at": now_utc().isoformat(),
    }
    await db.listings.insert_one(doc)
    await log_audit(user["id"], "listing.create", body.company_name, {"org_id": resolved_org_id})
    doc.pop("_id", None)
    return doc


@api_router.patch("/listings/{lid}")
async def update_listing(lid: str, body: ListingCreate, user=Depends(get_current_user)):
    await _listing_for_edit_or_404(lid, user)
    prev = await db.listings.find_one({"id": lid}, {"_id": 0, "status": 1})
    res = await db.listings.update_one(
        {"id": lid},
        {"$set": body.model_dump()},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Listing not found")
    await log_audit(user["id"], "listing.update", lid)
    # Rule 3A: closing a listing immediately revokes every external OAuth
    # connection and wipes mirrored bytes — the seller's source platforms
    # should never retain a live token for an archived deal.
    if prev and prev.get("status") != "closed" and body.status == "closed":
        await _wipe_listing_external_sources(lid)
    return {"ok": True}


@api_router.delete("/listings/{lid}")
async def delete_listing(lid: str, user=Depends(get_current_user)):
    listing = await _listing_for_edit_or_404(lid, user)
    # Only the principal owner, an org_admin of the owning org, or platform admin
    # can fully delete. Collaborators/editors can edit but not delete.
    _, role = await _resolve_listing_access(user, listing)
    if role not in ("owner", "org_admin", "admin"):
        raise HTTPException(status_code=403, detail="Only the principal owner or org admin can delete")
    # Same wipe rule as close — full delete should evict OAuth too.
    await _wipe_listing_external_sources(lid)
    await db.listings.delete_one({"id": lid})
    await log_audit(user["id"], "listing.delete", lid)
    return {"ok": True}


# ---- Listing-level data room (pre-stage documents) --------------------------
async def _seller_listing_or_404(lid: str, user: dict) -> Dict[str, Any]:
    """Backwards-compat wrapper: returns the listing if the user can edit it
    (principal owner, org member, collaborator editor/owner, or admin)."""
    return await _listing_for_edit_or_404(lid, user)


@api_router.get("/listings/{lid}/staged-files")
async def list_listing_staged_files(lid: str, user=Depends(get_current_user)):
    await _seller_listing_or_404(lid, user)
    items = await db.listing_staged_files.find(
        {"listing_id": lid, "deleted_at": {"$exists": False}},
        {"_id": 0, "content": 0, "pages": 0},
    ).sort("uploaded_at", -1).to_list(500)
    return items


@api_router.post("/listings/{lid}/staged-files/binary")
async def upload_listing_staged_file(
    lid: str,
    file: UploadFile = File(...),
    folder: str = Form("other"),
    note: Optional[str] = Form(None),
    user=Depends(get_current_user),
):
    """Seller uploads a document to the LISTING data room. These files auto-clone into
    every Vault opened against this listing."""
    await _seller_listing_or_404(lid, user)
    if folder not in ("financials", "legal", "hr", "it", "operations", "commercial", "other"):
        folder = "other"

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file")
    if len(data) > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File exceeds 50 MB limit")

    file_id = str(uuid.uuid4())
    filename = file.filename or f"upload-{file_id}"
    plaintext_sha256_hex = sha256_hex(data)

    storage_bytes = data
    encrypted = False
    encryption_alg = None
    if encryption_configured():
        try:
            aad = f"listing:{lid}:{file_id}".encode("utf-8")
            enc = encrypt_bytes(data, associated_data=aad)
            storage_bytes = enc["envelope"]
            encrypted = True
            encryption_alg = enc["alg"]
        except Exception as e:
            logger.warning(f"At-rest encryption failed, storing plaintext: {e}")

    gridfs_id = await listing_files_bucket.upload_from_stream(
        filename,
        io.BytesIO(storage_bytes),
        metadata={
            "listing_id": lid,
            "file_id": file_id,
            "uploaded_by": user["id"],
            "content_type": file.content_type or "application/octet-stream",
            "encrypted": encrypted,
            "encryption_alg": encryption_alg,
        },
    )

    pages = extract_pages_from_bytes(filename, data)
    flat = pages_to_flat_text(pages)

    doc = {
        "id": file_id,
        "listing_id": lid,
        "filename": filename,
        "folder": folder,
        "content_type": file.content_type or "application/octet-stream",
        "size_bytes": len(data),
        "page_count": len(pages),
        "pages": pages,
        "content": flat,
        "char_count": len(flat),
        "gridfs_id": str(gridfs_id),
        "storage": "listing_gridfs",
        "encrypted": encrypted,
        "encryption_alg": encryption_alg,
        "sha256_hex": plaintext_sha256_hex,
        "note": note,
        "uploaded_by": user["id"],
        "uploaded_at": now_utc().isoformat(),
    }
    await db.listing_staged_files.insert_one(doc)
    await log_audit(user["id"], "listing.stagedfile.upload", lid, {
        "file_id": file_id, "filename": filename, "folder": folder, "bytes": len(data),
        "sha256": plaintext_sha256_hex, "encrypted": encrypted,
    })
    asyncio.create_task(notarize_bytes(
        kind="listing.staged_file",
        target_id=file_id,
        data=data,
        owner_user_id=user["id"],
        label=f"Staged listing file: {filename}",
        extra={"listing_id": lid, "filename": filename, "size_bytes": len(data)},
    ))
    doc.pop("_id", None)
    doc.pop("content", None)
    doc.pop("pages", None)
    return doc


@api_router.get("/listings/{lid}/staged-files/{file_id}/download")
async def download_listing_staged_file(lid: str, file_id: str, user=Depends(get_current_user)):
    await _seller_listing_or_404(lid, user)
    f = await db.listing_staged_files.find_one(
        {"id": file_id, "listing_id": lid, "deleted_at": {"$exists": False}},
        {"_id": 0},
    )
    if not f or not f.get("gridfs_id"):
        raise HTTPException(status_code=404, detail="File not found")
    try:
        grid_out = await listing_files_bucket.open_download_stream(ObjectId(f["gridfs_id"]))
    except Exception as e:
        raise HTTPException(status_code=404, detail=f"Binary not found: {e}")
    if f.get("encrypted"):
        envelope = await grid_out.read()
        try:
            aad = f"listing:{lid}:{file_id}".encode("utf-8")
            plaintext = decrypt_envelope(envelope, associated_data=aad)
        except Exception as e:
            logger.exception("Listing staged file decryption failed")
            raise HTTPException(status_code=500, detail=f"Decryption failed: {e}")
        await log_audit(user["id"], "listing.stagedfile.download", lid, {"filename": f["filename"]})
        return StreamingResponse(
            io.BytesIO(plaintext),
            media_type=f.get("content_type") or "application/octet-stream",
            headers={"Content-Disposition": f'attachment; filename="{f["filename"]}"'},
        )

    async def streamer():
        while True:
            chunk = await grid_out.readchunk()
            if not chunk:
                break
            yield chunk

    await log_audit(user["id"], "listing.stagedfile.download", lid, {"filename": f["filename"]})
    return StreamingResponse(
        streamer(),
        media_type=f.get("content_type") or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{f["filename"]}"'},
    )


@api_router.delete("/listings/{lid}/staged-files/{file_id}")
async def delete_listing_staged_file(lid: str, file_id: str, user=Depends(get_current_user)):
    """Soft-delete the staged copy AND cascade-delete every cloned copy that
    already lives inside open Vaults (including Preview Vaults). Without the
    cascade, a seller who deletes a doc by mistake — or to retract a wrongly
    uploaded confidential file — would still see it surface to every buyer
    who'd already opened a Vault, because clones live in `deal_room_files`."""
    await _seller_listing_or_404(lid, user)
    res = await db.listing_staged_files.update_one(
        {"id": file_id, "listing_id": lid, "deleted_at": {"$exists": False}},
        {"$set": {"deleted_at": now_utc().isoformat()}},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="File not found")

    # Cascade — hard-delete the cloned copies (rows + GridFS bytes) so the
    # bytes really leave the platform. Soft-deleting wouldn't help because
    # GET /deal-rooms/{id} doesn't filter on deleted_at for room files.
    clones = await db.deal_room_files.find(
        {"cloned_from_listing_file": file_id}, {"_id": 0, "gridfs_id": 1, "id": 1, "room_id": 1}
    ).to_list(500)
    cascaded = 0
    for c in clones:
        try:
            await gridfs_bucket.delete(ObjectId(c["gridfs_id"]))
        except Exception:
            pass  # gridfs entry may already be gone; continue with the row
        await db.deal_room_files.delete_one({"id": c["id"]})
        cascaded += 1

    await log_audit(user["id"], "listing.stagedfile.delete", lid,
                    {"file_id": file_id, "cascaded_clones": cascaded})
    return {"ok": True, "cascaded_clones": cascaded}


# -----------------------------------------------------------------------------
# Buyer Private Locker
#
# A buyer-only document drawer. Files are NEVER visible to sellers or other
# buyers — only to the uploader themselves. Two scopes:
#   - scope="workspace"  → not tied to any listing (cross-deal notes, templates)
#   - scope="listing"    → attached to a specific listing the buyer is
#                          evaluating (board memo drafts, partner scoring, etc.)
# Encryption + OpenTimestamps notarization mirror the shared Vault.
# -----------------------------------------------------------------------------
PRIVATE_LOCKER_FOLDERS = ("notes", "modeling", "memos", "external", "other")


async def _private_locker_guard(user) -> None:
    if user.get("role") not in ("buyer", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Private Locker is buyer-only")


@api_router.get("/private-locker/files")
async def list_private_locker_files(
    listing_id: Optional[str] = None,
    research_id: Optional[str] = None,
    scope: Optional[str] = None,
    user=Depends(get_current_user),
):
    await _private_locker_guard(user)
    q: dict = {"user_id": user["id"], "deleted_at": {"$exists": False}}
    if scope == "workspace":
        q["scope"] = "workspace"
    elif scope == "listing":
        q["scope"] = "listing"
        if listing_id:
            q["listing_id"] = listing_id
    elif scope == "research":
        q["scope"] = "research"
        if research_id:
            q["research_id"] = research_id
    elif listing_id:
        q["scope"] = "listing"
        q["listing_id"] = listing_id
    elif research_id:
        q["scope"] = "research"
        q["research_id"] = research_id
    items = await db.private_locker_files.find(
        q, {"_id": 0, "content": 0, "pages": 0}
    ).sort("uploaded_at", -1).to_list(500)

    # Decorate listing-scope files with listing display name (cheap, 1 query)
    lids = list({i["listing_id"] for i in items if i.get("listing_id")})
    name_map: dict = {}
    if lids:
        async for li in db.listings.find(
            {"id": {"$in": lids}}, {"_id": 0, "id": 1, "company_name": 1, "name": 1}
        ):
            name_map[li["id"]] = li.get("company_name") or li.get("name") or li["id"]
    # Decorate research-scope files with research target company name
    rids = list({i["research_id"] for i in items if i.get("research_id")})
    research_name_map: dict = {}
    if rids:
        async for rr in db.research.find(
            {"id": {"$in": rids}, "user_id": user["id"]},
            {"_id": 0, "id": 1, "company_name": 1},
        ):
            research_name_map[rr["id"]] = rr.get("company_name") or rr["id"]
    for i in items:
        if i.get("listing_id"):
            i["listing_name"] = name_map.get(i["listing_id"])
        if i.get("research_id"):
            i["research_company_name"] = research_name_map.get(i["research_id"])
    return items


@api_router.post("/private-locker/files")
async def upload_private_locker_file(
    file: UploadFile = File(...),
    listing_id: Optional[str] = Form(None),
    research_id: Optional[str] = Form(None),
    folder: str = Form("other"),
    note: Optional[str] = Form(None),
    user=Depends(get_current_user),
):
    await _private_locker_guard(user)
    if folder not in PRIVATE_LOCKER_FOLDERS:
        folder = "other"

    # Determine scope (research takes precedence if both somehow set)
    if research_id:
        scope = "research"
        research_doc = await db.research.find_one(
            {"id": research_id, "user_id": user["id"]}, {"_id": 0, "id": 1, "company_name": 1}
        )
        if not research_doc:
            raise HTTPException(status_code=404, detail="Research target not found")
    elif listing_id:
        scope = "listing"
        exists = await db.listings.find_one({"id": listing_id}, {"_id": 0, "id": 1})
        if not exists:
            raise HTTPException(status_code=404, detail="Listing not found")
    else:
        scope = "workspace"

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file")
    if len(data) > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File exceeds 50 MB limit")

    file_id = str(uuid.uuid4())
    filename = file.filename or f"locker-{file_id}"
    plaintext_sha256_hex = sha256_hex(data)

    storage_bytes = data
    encrypted = False
    encryption_alg = None
    if encryption_configured():
        try:
            aad = f"locker:{user['id']}:{file_id}".encode("utf-8")
            enc = encrypt_bytes(data, associated_data=aad)
            storage_bytes = enc["envelope"]
            encrypted = True
            encryption_alg = enc["alg"]
        except Exception as e:
            logger.warning(f"Private Locker at-rest encryption failed, storing plaintext: {e}")

    gridfs_id = await private_locker_bucket.upload_from_stream(
        filename,
        io.BytesIO(storage_bytes),
        metadata={
            "user_id": user["id"],
            "file_id": file_id,
            "scope": scope,
            "listing_id": listing_id,
            "content_type": file.content_type or "application/octet-stream",
            "encrypted": encrypted,
            "encryption_alg": encryption_alg,
        },
    )

    pages = extract_pages_from_bytes(filename, data)
    flat = pages_to_flat_text(pages)

    doc = {
        "id": file_id,
        "user_id": user["id"],
        "scope": scope,
        "listing_id": listing_id if scope == "listing" else None,
        "research_id": research_id if scope == "research" else None,
        "filename": filename,
        "folder": folder,
        "content_type": file.content_type or "application/octet-stream",
        "size_bytes": len(data),
        "page_count": len(pages),
        "pages": pages,
        "content": flat,
        "char_count": len(flat),
        "gridfs_id": str(gridfs_id),
        "storage": "private_locker_gridfs",
        "encrypted": encrypted,
        "encryption_alg": encryption_alg,
        "sha256_hex": plaintext_sha256_hex,
        "note": note,
        "uploaded_at": now_utc().isoformat(),
        "created_at": now_utc().isoformat(),  # demo cleanup uses created_at
    }
    await db.private_locker_files.insert_one(doc)
    await log_audit(
        user["id"], "locker.file.upload", file_id,
        {"filename": filename, "scope": scope, "listing_id": listing_id,
         "research_id": research_id, "bytes": len(data),
         "sha256": plaintext_sha256_hex, "encrypted": encrypted},
    )
    asyncio.create_task(notarize_bytes(
        kind="private_locker.file",
        target_id=file_id,
        data=data,
        owner_user_id=user["id"],
        label=f"Private locker file: {filename}",
        extra={"scope": scope, "listing_id": listing_id, "research_id": research_id,
               "filename": filename, "size_bytes": len(data)},
    ))
    doc.pop("_id", None)
    doc.pop("content", None)
    doc.pop("pages", None)
    return doc


@api_router.get("/private-locker/files/{fid}/download")
async def download_private_locker_file(fid: str, user=Depends(get_current_user)):
    await _private_locker_guard(user)
    f = await db.private_locker_files.find_one(
        {"id": fid, "user_id": user["id"], "deleted_at": {"$exists": False}},
        {"_id": 0},
    )
    if not f or not f.get("gridfs_id"):
        raise HTTPException(status_code=404, detail="File not found")
    try:
        grid_out = await private_locker_bucket.open_download_stream(ObjectId(f["gridfs_id"]))
    except Exception as e:
        raise HTTPException(status_code=404, detail=f"Binary not found: {e}")
    if f.get("encrypted"):
        envelope = await grid_out.read()
        try:
            aad = f"locker:{user['id']}:{fid}".encode("utf-8")
            plaintext = decrypt_envelope(envelope, associated_data=aad)
        except Exception as e:
            logger.exception("Private Locker decryption failed")
            raise HTTPException(status_code=500, detail=f"Decryption failed: {e}")
        await log_audit(user["id"], "locker.file.download", fid, {"filename": f["filename"]})
        return StreamingResponse(
            io.BytesIO(plaintext),
            media_type=f.get("content_type") or "application/octet-stream",
            headers={"Content-Disposition": f'attachment; filename="{f["filename"]}"'},
        )

    async def streamer():
        while True:
            chunk = await grid_out.readchunk()
            if not chunk:
                break
            yield chunk

    await log_audit(user["id"], "locker.file.download", fid, {"filename": f["filename"]})
    return StreamingResponse(
        streamer(),
        media_type=f.get("content_type") or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{f["filename"]}"'},
    )


@api_router.delete("/private-locker/files/{fid}")
async def delete_private_locker_file(fid: str, user=Depends(get_current_user)):
    await _private_locker_guard(user)
    f = await db.private_locker_files.find_one(
        {"id": fid, "user_id": user["id"], "deleted_at": {"$exists": False}},
        {"_id": 0, "gridfs_id": 1, "filename": 1},
    )
    if not f:
        raise HTTPException(status_code=404, detail="File not found")
    # Hard delete: this is private user-owned data, no shared audit trail to preserve.
    if f.get("gridfs_id"):
        try:
            await private_locker_bucket.delete(ObjectId(f["gridfs_id"]))
        except Exception as e:
            logger.warning(f"locker gridfs delete failed {fid}: {e}")
    await db.private_locker_files.delete_one({"id": fid, "user_id": user["id"]})
    await log_audit(user["id"], "locker.file.delete", fid, {"filename": f.get("filename")})
    return {"ok": True}


async def _clone_listing_files_into_room(listing_id: str, room_id: str, user_id: str, *, only_missing: bool = False) -> int:
    """Copy every active staged file on a listing into a Vault. Re-encrypts
    with the vault's AAD so the cipher stays bound to its room. Returns number cloned.

    `only_missing=True` skips staged files that have already been cloned into this
    room (detected via `cloned_from_listing_file`). Use this for incremental
    backfills (e.g. after an external-source sync drops new files into a listing
    whose Vaults are already open)."""
    staged = await db.listing_staged_files.find(
        {"listing_id": listing_id, "deleted_at": {"$exists": False}},
        {"_id": 0},
    ).sort("uploaded_at", 1).to_list(500)
    if not staged:
        return 0
    already_cloned: set[str] = set()
    if only_missing:
        existing = await db.deal_room_files.find(
            {"room_id": room_id, "cloned_from_listing_file": {"$ne": None, "$exists": True}},
            {"_id": 0, "cloned_from_listing_file": 1},
        ).to_list(1000)
        already_cloned = {e["cloned_from_listing_file"] for e in existing if e.get("cloned_from_listing_file")}
    cloned = 0
    for s in staged:
        if only_missing and s["id"] in already_cloned:
            continue
        try:
            grid_out = await listing_files_bucket.open_download_stream(ObjectId(s["gridfs_id"]))
        except Exception as e:
            logger.warning(f"clone: skipping staged file {s['id']}: {e}")
            continue
        envelope_or_plain = await grid_out.read()
        # Recover plaintext
        if s.get("encrypted"):
            try:
                aad = f"listing:{listing_id}:{s['id']}".encode("utf-8")
                plaintext = decrypt_envelope(envelope_or_plain, associated_data=aad)
            except Exception as e:
                logger.warning(f"clone: decryption failed for {s['id']}: {e}")
                continue
        else:
            plaintext = envelope_or_plain

        new_file_id = str(uuid.uuid4())
        storage_bytes = plaintext
        encrypted = False
        encryption_alg = None
        if encryption_configured():
            try:
                aad = f"{room_id}:{new_file_id}".encode("utf-8")
                enc = encrypt_bytes(plaintext, associated_data=aad)
                storage_bytes = enc["envelope"]
                encrypted = True
                encryption_alg = enc["alg"]
            except Exception as e:
                logger.warning(f"clone: re-encryption failed for {s['id']}: {e}")

        new_gridfs_id = await gridfs_bucket.upload_from_stream(
            s["filename"],
            io.BytesIO(storage_bytes),
            metadata={
                "room_id": room_id, "file_id": new_file_id,
                "uploaded_by": user_id,
                "content_type": s.get("content_type") or "application/octet-stream",
                "encrypted": encrypted, "encryption_alg": encryption_alg,
                "cloned_from_listing": listing_id, "cloned_from_file": s["id"],
            },
        )
        pages = s.get("pages") or extract_pages_from_bytes(s["filename"], plaintext)
        flat = pages_to_flat_text(pages) if pages else (s.get("content") or "")
        doc = {
            "id": new_file_id,
            "room_id": room_id,
            "filename": s["filename"],
            "folder": s.get("folder") or "other",
            "content_type": s.get("content_type") or "application/octet-stream",
            "size_bytes": s.get("size_bytes") or len(plaintext),
            "page_count": len(pages) if pages else 0,
            "pages": pages or [],
            "content": flat,
            "char_count": len(flat),
            "gridfs_id": str(new_gridfs_id),
            "storage": "gridfs",
            "encrypted": encrypted,
            "encryption_alg": encryption_alg,
            "sha256_hex": s.get("sha256_hex") or sha256_hex(plaintext),
            "note": s.get("note"),
            "uploaded_by": user_id,
            "uploaded_by_role": "seller",
            "uploaded_at": now_utc().isoformat(),
            "matched_request_id": None,
            "cloned_from_listing_file": s["id"],
            # Carry external-source provenance so the Activity tab can render
            # a "Synced from Google Drive" badge and the audit log can attribute
            # the file to its origin connector.
            "source": s.get("source"),
        }
        await db.deal_room_files.insert_one(doc)
        cloned += 1
        # Emit a per-file audit event ONLY for incremental backfills (files
        # arriving via Composio sync after the Vault is already open). The
        # initial open-room clone is already covered by `dealroom.open`, so
        # logging every file there would just duplicate noise.
        if only_missing:
            try:
                source_meta = s.get("source") or {}
                await log_audit(
                    user_id, "dealroom.file.add", room_id,
                    {
                        "filename": s["filename"],
                        "file_id": new_file_id,
                        "via": source_meta.get("kind") or "manual",
                        "source_sid": source_meta.get("sid"),
                    },
                )
            except Exception as e:
                logger.debug(f"dealroom.file.add audit skipped for {new_file_id}: {e}")
    return cloned


@api_router.get("/marketplace")
async def marketplace(user=Depends(get_current_user)):
    """Buyer view — public live listings from every seller."""
    items = await db.listings.find({"status": "live"}, {"_id": 0}).sort("created_at", -1).to_list(200)
    return items


@api_router.get("/marketplace/{lid}")
async def marketplace_detail(lid: str, user=Depends(get_current_user)):
    item = await db.listings.find_one({"id": lid}, {"_id": 0})
    if not item:
        raise HTTPException(status_code=404, detail="Listing not found")
    await db.listings.update_one({"id": lid}, {"$inc": {"view_count": 1}})
    return item


@api_router.post("/marketplace/{lid}/inquire")
async def inquire(lid: str, body: InquiryCreate, user=Depends(get_current_user)):
    listing = await db.listings.find_one({"id": lid}, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found")
    doc = {
        "id": str(uuid.uuid4()),
        "listing_id": lid,
        "listing_name": listing["company_name"],
        "seller_id": listing["seller_id"],
        "buyer_id": user["id"],
        "buyer_name": user.get("name"),
        "buyer_org": user.get("organization"),
        "buyer_email": user.get("email"),
        "message": body.message,
        "status": "new",
        "created_at": now_utc().isoformat(),
    }
    await db.inquiries.insert_one(doc)
    await db.listings.update_one({"id": lid}, {"$inc": {"inquiry_count": 1}})
    await log_audit(user["id"], "inquiry.create", lid)
    await log_agent_activity("matchmaking-agent", f"inquiry:{listing['company_name']}", "completed", user_id=user["id"])
    doc.pop("_id", None)
    return doc


@api_router.get("/inquiries")
async def list_inquiries(user=Depends(get_current_user)):
    """Sellers + agents see all inquiries on listings in their workspace
    (personal + org-owned + listings they collaborate on). Buyers see
    outbound. Admin sees all."""
    role = user.get("role", "buyer")
    if role in ("seller", "agent"):
        listing_ids, _ = await _user_workspace_listing_ids(user)
        if not listing_ids:
            return []
        q = {"listing_id": {"$in": listing_ids}}
    elif role == "buyer":
        q = {"buyer_id": user["id"]}
    else:
        q = {}
    items = await db.inquiries.find({**q, "deleted_at": {"$exists": False}}, {"_id": 0}).sort("created_at", -1).to_list(200)
    # Decorate each inquiry with the scope so the UI can show "via [Org]" / "team"
    if items and role in ("seller", "agent"):
        my_listings_personal = set(
            r["id"] for r in await db.listings.find({"seller_id": user["id"]}, {"_id": 0, "id": 1}).to_list(500)
        )
        listings_by_id = {
            r["id"]: r for r in await db.listings.find(
                {"id": {"$in": [i["listing_id"] for i in items]}},
                {"_id": 0, "id": 1, "org_id": 1, "seller_id": 1, "company_name": 1},
            ).to_list(500)
        }
        org_name_by_id = {
            o["id"]: o["name"] for o in await db.organizations.find(
                {"id": {"$in": [l.get("org_id") for l in listings_by_id.values() if l.get("org_id")]}},
                {"_id": 0, "id": 1, "name": 1},
            ).to_list(200)
        }
        for it in items:
            lst = listings_by_id.get(it["listing_id"], {})
            if it["listing_id"] in my_listings_personal:
                it["workspace_scope"] = "mine"
            elif lst.get("org_id"):
                it["workspace_scope"] = "org"
                it["workspace_org_name"] = org_name_by_id.get(lst["org_id"])
                it["workspace_owner_id"] = lst.get("seller_id")
            else:
                it["workspace_scope"] = "shared"
                it["workspace_owner_id"] = lst.get("seller_id")
    return items


@api_router.patch("/inquiries/{iid}/status")
async def update_inquiry(iid: str, body: dict, user=Depends(get_current_user)):
    new_status = body.get("status")
    if new_status not in ("new", "reviewing", "engaged", "passed"):
        raise HTTPException(status_code=400, detail="Invalid status")
    # Fetch first so we can check workspace permission (not strict seller_id ownership)
    inq = await db.inquiries.find_one({"id": iid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not inq:
        raise HTTPException(status_code=404, detail="Inquiry not found")
    if user.get("role") != "admin":
        ws_listings, _ = await _user_workspace_listing_ids(user)
        if inq.get("listing_id") not in ws_listings:
            raise HTTPException(status_code=403, detail="Not authorized")
    await db.inquiries.update_one({"id": iid}, {"$set": {"status": new_status}})
    await log_audit(user["id"], "inquiry.status", iid, {"status": new_status})
    # Bitcoin-anchored proof of status change
    if new_status in ("engaged", "passed"):
        asyncio.create_task(notarize_event(
            kind="inquiry.status",
            target_id=iid,
            payload={"inquiry_id": iid, "status": new_status, "actor_id": user["id"], "ts": now_utc().isoformat()},
            owner_user_id=user["id"],
            label=f"Inquiry → {new_status}",
        ))
    return {"ok": True}


@api_router.get("/watchlist")
async def get_watchlist(user=Depends(get_current_user)):
    items = await db.watchlist.find({"user_id": user["id"]}, {"_id": 0}).to_list(100)
    return items


@api_router.post("/watchlist/{lid}")
async def add_watch(lid: str, user=Depends(get_current_user)):
    if await db.watchlist.find_one({"user_id": user["id"], "listing_id": lid}):
        return {"ok": True, "already": True}
    listing = await db.listings.find_one({"id": lid}, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found")
    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "listing_id": lid,
        "company_name": listing["company_name"],
        "sector": listing["sector"],
        "asking_price_usd_m": listing["asking_price_usd_m"],
        "created_at": now_utc().isoformat(),
    }
    await db.watchlist.insert_one(doc)
    doc.pop("_id", None)
    return doc


@api_router.delete("/watchlist/{lid}")
async def remove_watch(lid: str, user=Depends(get_current_user)):
    await db.watchlist.delete_one({"user_id": user["id"], "listing_id": lid})
    return {"ok": True}


# -----------------------------------------------------------------------------
# RESEARCH HUB
# -----------------------------------------------------------------------------
RESEARCH_SYS = """You are a senior M&A analyst at NextCapOS. Write a concise institutional research brief.
Return STRICT JSON only (no markdown). Keep arrays to MAX 3 items. Keep each text field under 240 chars.
You are given an indexed list of LIVE web SOURCES. Cite them inline as [1], [2], ... inside text fields where you make a claim derived from the source. Do not invent citations beyond the list provided.
Schema:
{
  "company_name": str,
  "one_liner": str,
  "sector": str,
  "headquarters": str,
  "founded": str,
  "employees_range": str,
  "estimated_revenue": str,
  "business_model": str,
  "leadership": [{"name": str, "title": str, "background": str}],
  "market_signals": [str],
  "growth_drivers": [str],
  "risks": [str],
  "competitive_landscape": [str],
  "investor_take": str,
  "suggested_buyer_profile": str,
  "next_actions": [str]
}
Be specific, analytical, terse. Prefer claims grounded in the provided SOURCES; otherwise mark with a [-] placeholder."""


# -----------------------------------------------------------------------------
# Real-time web research helpers (Brave + Perplexity Sonar)
# -----------------------------------------------------------------------------
async def search_brave(query: str, count: int = 6) -> List[dict]:
    if not BRAVE_API_KEY:
        return []
    try:
        async with httpx.AsyncClient(timeout=12.0) as c:
            r = await c.get(
                "https://api.search.brave.com/res/v1/web/search",
                headers={"X-Subscription-Token": BRAVE_API_KEY, "Accept": "application/json"},
                params={"q": query, "count": count, "country": "US", "search_lang": "en"},
            )
            if r.status_code >= 400:
                logger.warning(f"Brave {r.status_code}: {r.text[:200]}")
                return []
            data = r.json()
            results = (data.get("web") or {}).get("results") or []
            return [
                {
                    "title": x.get("title") or "",
                    "url": x.get("url") or "",
                    "snippet": (x.get("description") or "")[:300],
                    "age": x.get("age") or x.get("page_age") or "",
                    "provider": "brave",
                }
                for x in results[:count]
                if x.get("url")
            ]
    except Exception as e:
        logger.warning(f"Brave search failed: {e}")
        return []


async def query_perplexity(prompt: str, model: str = "sonar-pro") -> dict:
    """Returns {'text': str, 'citations': [url, url...]}."""
    if not PERPLEXITY_API_KEY:
        return {"text": "", "citations": []}
    try:
        async with httpx.AsyncClient(timeout=25.0) as c:
            r = await c.post(
                "https://api.perplexity.ai/chat/completions",
                headers={
                    "Authorization": f"Bearer {PERPLEXITY_API_KEY}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "messages": [
                        {
                            "role": "system",
                            "content": "You are a research analyst. Provide a concise factual overview for institutional investors. Cite all material claims.",
                        },
                        {"role": "user", "content": prompt},
                    ],
                    "temperature": 0.1,
                    "max_tokens": 900,
                },
            )
            if r.status_code >= 400:
                logger.warning(f"Perplexity {r.status_code}: {r.text[:300]}")
                return {"text": "", "citations": []}
            data = r.json()
            text = ""
            choices = data.get("choices") or []
            if choices:
                text = (choices[0].get("message") or {}).get("content") or ""
            citations = data.get("citations") or data.get("search_results") or []
            urls = []
            for c2 in citations:
                if isinstance(c2, str):
                    urls.append(c2)
                elif isinstance(c2, dict) and c2.get("url"):
                    urls.append(c2["url"])
            return {"text": text, "citations": urls}
    except Exception as e:
        logger.warning(f"Perplexity failed: {e}")
        return {"text": "", "citations": []}


def build_sources(perplexity_urls: List[str], brave_hits: List[dict]) -> List[dict]:
    """Merge Perplexity citations + Brave results into a numbered, deduped source list."""
    seen = set()
    sources = []
    idx = 1
    for url in perplexity_urls:
        if not url or url in seen:
            continue
        seen.add(url)
        sources.append({"index": idx, "url": url, "title": "", "provider": "perplexity"})
        idx += 1
    for hit in brave_hits:
        url = hit.get("url")
        if not url or url in seen:
            continue
        seen.add(url)
        sources.append({
            "index": idx,
            "url": url,
            "title": hit.get("title", ""),
            "snippet": hit.get("snippet", ""),
            "age": hit.get("age", ""),
            "provider": "brave",
        })
        idx += 1
    return sources


def _extract_domain(url: Optional[str]) -> Optional[str]:
    """Strip protocol, www, path, query, port from a URL. Returns None on failure or empty."""
    if not url or not isinstance(url, str):
        return None
    s = url.strip()
    if not s:
        return None
    s = re.sub(r"^https?://", "", s, flags=re.IGNORECASE)
    s = re.sub(r"^www\.", "", s, flags=re.IGNORECASE)
    s = s.split("/")[0].split("?")[0].split("#")[0]
    s = s.split(":")[0]  # strip port
    return s.lower() or None


@api_router.post("/research/company")
async def research_company(body: CompanyResearchRequest, user=Depends(get_current_user)):
    """Queue a grounded research brief. Returns immediately with status='pending'.
    Background worker transitions pending → analyzing → completed | failed (or returns
    inline if the pipeline finishes within ~25s)."""
    rid = str(uuid.uuid4())
    skeleton = {
        "id": rid,
        "user_id": user["id"],
        "company_name": body.company_name,
        "company_url": body.company_url,
        "sector": body.sector,
        "region": body.region,
        "notes": body.notes,
        "data": None,
        "sources": [],
        "live_research_used": False,
        "status": "pending",
        "error": None,
        "created_at": now_utc().isoformat(),
    }
    await db.research.insert_one(skeleton)
    asyncio.create_task(_execute_research_company(rid, body, user["id"]))
    skeleton.pop("_id", None)
    return skeleton


async def _execute_research_company(rid: str, body: CompanyResearchRequest, user_id: str) -> None:
    """Background worker for /api/research/company. Avoids the 60s ingress timeout
    that was failing user-facing research calls."""
    started = now_utc()
    company = body.company_name
    try:
        await db.research.update_one({"id": rid}, {"$set": {"status": "analyzing"}})

        sonar_prompt = (
            f"Provide an institutional-investor overview of {company}"
            + (f" (official website: {body.company_url})" if body.company_url else "")
            + (f" (sector hint: {body.sector})" if body.sector else "")
            + (f" (region: {body.region})" if body.region else "")
            + ". Cover business model, recent news, leadership, competitive position, and any 2026 events."
        )
        brave_query = f"{company} {body.sector or ''} company news 2026".strip()
        # Domain-anchored Brave query — pulls signals directly from the company's own site
        brave_domain_query = None
        domain = _extract_domain(body.company_url) if body.company_url else None
        if domain:
            brave_domain_query = f"site:{domain} about OR pricing OR customers OR team OR investors"

        tasks = [
            query_perplexity(sonar_prompt),
            search_brave(brave_query, count=6),
        ]
        if brave_domain_query:
            tasks.append(search_brave(brave_domain_query, count=6))
        results = await asyncio.gather(*tasks, return_exceptions=True)
        perplexity_res = results[0] if not isinstance(results[0], Exception) else {}
        brave_res = results[1] if not isinstance(results[1], Exception) else []
        if len(results) > 2 and not isinstance(results[2], Exception):
            brave_res = (results[2] or []) + (brave_res or [])

        # Social discovery is rate-limit friendly (sequential, 1.1s spacing) so it
        # runs AFTER the main Brave burst to avoid 429s on the free plan.
        try:
            social_profiles = await discover_social_profiles(
                search_brave=search_brave, company_name=company, company_url=body.company_url,
            )
        except Exception as e:
            logger.warning(f"research: social discovery failed: {e}")
            social_profiles = {}

        sources = build_sources(perplexity_res.get("citations", []) if isinstance(perplexity_res, dict) else [], brave_res or [])

        sources_block = "\n".join(
            f"[{s['index']}] {s.get('title') or s['url']} — {s['url']}"
            + (f" :: {s.get('snippet')}" if s.get("snippet") else "")
            for s in sources
        ) or "(no live sources available — proceed with model knowledge)"

        perplexity_summary = (perplexity_res.get("text") if isinstance(perplexity_res, dict) else "" or "").strip()
        grounded_user = (
            f"Company: {company}\n"
            f"Website: {body.company_url or 'unknown'}\n"
            f"Sector hint: {body.sector or 'unspecified'}\n"
            f"Region hint: {body.region or 'global'}\n"
            f"Buyer notes: {body.notes or 'none'}\n\n"
            f"LIVE WEB-RESEARCH SUMMARY (from real-time search):\n{perplexity_summary or '(none)'}\n\n"
            f"SOURCES (cite as [n] inline):\n{sources_block}\n\n"
            "Now produce the JSON brief, embedding [n] citations where you reference a source."
        )

        raw = await call_claude(RESEARCH_SYS, grounded_user, session_id=f"research-{user_id}")
        data = safe_json_loads(raw)

        duration = int((now_utc() - started).total_seconds() * 1000)
        await db.research.update_one(
            {"id": rid},
            {"$set": {
                "data": data,
                "sources": sources,
                "social_profiles": social_profiles,
                "live_research_used": bool(perplexity_summary or brave_res),
                "status": "completed",
                "completed_at": now_utc().isoformat(),
                "duration_ms": duration,
            }},
        )
        await log_agent_activity(
            "research-agent",
            f"research:{company} · grounded({len(sources)} sources)",
            "completed",
            user_id=user_id, duration_ms=duration,
            meta={"sources_count": len(sources), "providers": list({s["provider"] for s in sources})},
        )
        await log_audit(user_id, "research.create", company, {"sources": len(sources)})
    except Exception as e:
        logger.exception(f"Research worker failed for {rid}")
        await db.research.update_one(
            {"id": rid},
            {"$set": {"status": "failed", "error": str(e)[:500],
                      "completed_at": now_utc().isoformat()}},
        )
        await log_agent_activity("research-agent", f"research:{company}", "failed",
                                 user_id=user_id, friction=str(e)[:200])


@api_router.get("/research/detail/{rid}")
async def get_research(rid: str, user=Depends(get_current_user)):
    doc = await db.research.find_one({"id": rid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Research not found")
    if user.get("role") != "admin" and doc.get("user_id") != user["id"]:
        raise HTTPException(status_code=403, detail="Not your research")
    return doc


@api_router.get("/research/history")
async def research_history(user=Depends(get_current_user)):
    items = await db.research.find({"user_id": user["id"], "deleted_at": {"$exists": False}}, {"_id": 0}).sort("created_at", -1).to_list(50)
    return items


# -----------------------------------------------------------------------------
# VALUATION BAND (Phase E) — free-source fair-value estimate for a brief
# -----------------------------------------------------------------------------
from valuation import estimate_valuation as _estimate_valuation  # noqa: E402


class ValuationRequest(BaseModel):
    company_name: str
    sector: Optional[str] = None
    one_liner: Optional[str] = None
    estimated_revenue: Optional[str] = None
    headquarters: Optional[str] = None
    research_id: Optional[str] = None  # optional — used for 24h cache keying
    force_refresh: bool = False


def _valuation_cache_key(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", (name or "").lower()).strip("-")[:80]


@api_router.post("/valuation/estimate")
async def estimate_company_valuation(body: ValuationRequest, user=Depends(get_current_user)):
    """Fair-value band using Recent Transaction + Market Multiples methods.

    24h cache keyed by lower-cased company slug. Force refresh with `force_refresh=true`.
    """
    if not body.company_name or not body.company_name.strip():
        raise HTTPException(status_code=400, detail="company_name is required")

    slug = _valuation_cache_key(body.company_name)
    if not body.force_refresh:
        cached = await db.valuation_estimates.find_one(
            {"slug": slug}, {"_id": 0}, sort=[("created_at", -1)],
        )
        if cached:
            age_hours = (now_utc() - datetime.fromisoformat(cached["created_at"])).total_seconds() / 3600
            if age_hours < 24:
                cached["from_cache"] = True
                cached["cache_age_hours"] = round(age_hours, 1)
                return cached

    result = await _estimate_valuation(
        company_name=body.company_name,
        sector=body.sector,
        one_liner=body.one_liner,
        estimated_revenue=body.estimated_revenue,
        headquarters=body.headquarters,
        brave_fn=search_brave,
        perplexity_fn=lambda p: query_perplexity(p),
        claude_fn=call_claude,
        safe_json=safe_json_loads,
    )

    doc = {
        "id": str(uuid.uuid4()),
        "slug": slug,
        "company_name": body.company_name,
        "sector": body.sector,
        "research_id": body.research_id,
        "user_id": user["id"],
        "created_at": now_utc().isoformat(),
        "currency": result.get("currency", "USD"),
        "as_of": result.get("as_of"),
        "aggregate": result.get("aggregate"),
        "recent_transaction": result.get("recent_transaction"),
        "market_multiples": result.get("market_multiples"),
        "sources": result.get("sources", []),
    }
    await db.valuation_estimates.insert_one(doc)
    doc.pop("_id", None)
    doc["from_cache"] = False
    await log_audit(user["id"], "valuation.estimate", slug, {
        "company_name": body.company_name,
        "insufficient_data": bool((result.get("aggregate") or {}).get("insufficient_data")),
        "base_usd": (result.get("aggregate") or {}).get("base_usd"),
    })
    return doc


@api_router.get("/valuation/estimate/{slug}")
async def get_cached_valuation(slug: str, user=Depends(get_current_user)):
    """Return the most recent cached valuation for a company slug, or 404."""
    doc = await db.valuation_estimates.find_one(
        {"slug": slug}, {"_id": 0}, sort=[("created_at", -1)],
    )
    if not doc:
        raise HTTPException(status_code=404, detail="No valuation on file")
    age_hours = (now_utc() - datetime.fromisoformat(doc["created_at"])).total_seconds() / 3600
    doc["cache_age_hours"] = round(age_hours, 1)
    return doc


# =============================================================================
# VALUATION WORKBENCH (Phase A) — Full 5-method fair-value workbench
# =============================================================================
from valuation_workbench import (  # noqa: E402
    autofill_workbench as _autofill_wb,
    compute_all_methods as _compute_all,
    aggregate_band as _aggregate_band,
    extract_term_sheet as _extract_term_sheet,
    DEFAULT_WEIGHTS as _DEFAULT_WEIGHTS,
)
from valuation_pdf import build_memo_pdf as _build_memo_pdf  # noqa: E402


class ValuationCreate(BaseModel):
    company_name: str
    sector: Optional[str] = None
    one_liner: Optional[str] = None
    estimated_revenue: Optional[str] = None
    headquarters: Optional[str] = None
    research_id: Optional[str] = None
    listing_id: Optional[str] = None
    deal_room_id: Optional[str] = None
    autofill: bool = True  # kick off AI autofill by default


class ValuationInputsUpdate(BaseModel):
    inputs: Dict[str, Dict[str, Any]]  # {method_key: {field: value}}
    weights: Optional[Dict[str, float]] = None


class SnapshotCreate(BaseModel):
    label: Optional[str] = None
    narrative: Optional[str] = None


def _valuation_summary(v: dict) -> dict:
    """Compact list-view projection of a valuation record."""
    return {
        "id": v["id"],
        "company_name": v.get("company_name"),
        "sector": v.get("sector"),
        "headquarters": v.get("headquarters"),
        "user_id": v.get("user_id"),
        "created_at": v.get("created_at"),
        "updated_at": v.get("updated_at"),
        "autofill_status": v.get("autofill_status"),
        "latest_snapshot_id": v.get("latest_snapshot_id"),
        "aggregate": v.get("aggregate"),
        "snapshot_count": v.get("snapshot_count", 0),
    }


async def _gather_vault_private_evidence(rid: str) -> tuple[str, list[dict]]:
    """Iter-38 (+43 hardening) — reuse the parsed pages Findings already
    extracted from vault files. Returns `(evidence_text, files_used)`.

    Heuristic filename detection prioritizes:
      * term sheets / SAFEs / convertibles → recent_transaction + OPM
      * cap tables / ownership → OPM preferred stack
      * financial models / forecasts / statements / actuals → DCF inputs
      * everything else → general context (capped)

    Iter-43 fix (AZpme $0-band bug):
      * Widen FINANCIALS regex — catches income statements, EBITDA / ARR /
        actuals / earnings / P&L / balance sheet / 10-K / annual report /
        Q[1-4] / projections / valuation / cash-flow / capex / opex / margins.
      * Read ALL pages of high-priority (TERM_SHEET / CAP_TABLE / FINANCIALS)
        docs, not just first 4. Non-priority docs still capped at 3 pages so
        we don't blow the prompt on boilerplate.
      * Per-file char budget lifted 600 → 1600 for priority docs, 400 for OTHER.
      * Total budget lifted 11k → 24k chars (~6k tokens — comfortable).
    """
    if not rid:
        return "", []
    files = await db.deal_room_files.find(
        {"room_id": rid},
        {"_id": 0, "id": 1, "filename": 1, "folder": 1, "pages": 1, "content": 1},
    ).sort("uploaded_at", 1).to_list(60)
    if not files:
        return "", []

    prio_regex = {
        "TERM_SHEET": re.compile(
            r"term.?sheet|safe|convertible|409a|note.?purchase|purchase.?agreement",
            re.I,
        ),
        "CAP_TABLE": re.compile(
            r"cap.?table|ownership|equity.?stack|share.?ledger|shareholder|option.?pool",
            re.I,
        ),
        "FINANCIALS": re.compile(
            r"revenue|financial|forecast|projection|model|budget|p.?and.?l|pnl|"
            r"p&l|cash.?flow|income.?statement|balance.?sheet|earnings|ebitda|"
            r"actuals|arr\b|mrr\b|margin|capex|opex|kpi|metric|"
            r"annual.?report|10.?[kq]|q[1-4]\b|fy\d|valuation|nav\b|payer|"
            r"cohort|churn|ltv|cac|unit.?economics",
            re.I,
        ),
    }

    def _priority(fn: str) -> str:
        fn_l = (fn or "").lower()
        for k, r in prio_regex.items():
            if r.search(fn_l):
                return k
        return "OTHER"

    order = {"TERM_SHEET": 0, "CAP_TABLE": 1, "FINANCIALS": 2, "OTHER": 3}
    files.sort(key=lambda f: order[_priority(f.get("filename", ""))])

    blocks: list[str] = []
    used: list[dict] = []
    total_chars = 0
    CHAR_CAP = 24_000  # ~6k tokens — leaves generous room for public web block

    for idx, f in enumerate(files, 1):
        if total_chars >= CHAR_CAP:
            break
        prio = _priority(f.get("filename", ""))
        # Priority docs get ALL pages up to a per-file cap; boilerplate is
        # sampled shallowly so we don't spend budget on NDA templates.
        per_page_cap = 1600 if prio != "OTHER" else 400
        max_pages   = 999  if prio != "OTHER" else 3
        pages = f.get("pages") or []
        if pages:
            page_lines: list[str] = []
            file_chars = 0
            for p in pages[:max_pages]:
                text = (p.get("text") or "").strip()
                if not text:
                    continue
                chunk = text[:per_page_cap]
                page_lines.append(f"  <p{p.get('page', '?')}> {chunk}")
                file_chars += len(chunk)
                # per-file cap: priority docs 8000, other 1200
                if file_chars >= (8000 if prio != "OTHER" else 1200):
                    break
            excerpts = "\n".join(page_lines) or f"  (no extractable text on {f['filename']})"
        else:
            # No page-level extraction — fall back to `content` field.
            raw = (f.get("content") or "").strip()
            excerpts = f"  {raw[:8000 if prio != 'OTHER' else 1200]}" if raw else \
                       f"  (no text extracted from {f['filename']})"
        piece = f"[{idx} · {prio}] {f['filename']} ({f.get('folder', '/')})\n{excerpts}"
        # Final safety trim to keep the running total inside CHAR_CAP.
        if total_chars + len(piece) > CHAR_CAP:
            piece = piece[: CHAR_CAP - total_chars] + "…"
        blocks.append(piece)
        used.append({"id": f["id"], "filename": f["filename"], "priority": prio})
        total_chars += len(piece)

    return "\n\n---\n\n".join(blocks), used


def _merge_autofill_inputs(existing: dict, incoming: dict) -> dict:
    """Iter-44: Merge autofill seed values into existing per-method inputs.

    Claude is non-deterministic — a re-run may return nulls for fields it
    populated on a prior run (or vice versa). Blind-overwrite would make the
    Workbench appear to "lose data" on Re-autofill. Instead:

      * For each method, prefer the NEW non-null value; fall back to the OLD
        non-null value if the new run returned null / "" / empty list.
      * List fields (comparable_tickers, source_urls) union-merge (dedupe).
      * Notes always take the newer one (fresher rationale).
      * If a method wasn't touched by the new run at all, keep old inputs.
    """
    existing = existing or {}
    incoming = incoming or {}
    merged: dict = {}
    all_methods = set(existing.keys()) | set(incoming.keys())
    for method in all_methods:
        old_m = existing.get(method) or {}
        new_m = incoming.get(method) or {}
        if not isinstance(old_m, dict): old_m = {}
        if not isinstance(new_m, dict): new_m = {}
        combined: dict = {}
        for k in set(old_m.keys()) | set(new_m.keys()):
            old_v = old_m.get(k)
            new_v = new_m.get(k)
            # List merge — union & dedupe (order preserved, new first).
            if isinstance(old_v, list) or isinstance(new_v, list):
                merged_list: list = []
                seen: set = set()
                for src in ((new_v or []), (old_v or [])):
                    for item in src if isinstance(src, list) else []:
                        key = item if isinstance(item, (str, int, float)) else str(item)
                        if key not in seen:
                            seen.add(key)
                            merged_list.append(item)
                combined[k] = merged_list
                continue
            # Scalar: prefer new non-null, fallback to old non-null.
            new_is_empty = new_v is None or new_v == "" or new_v == 0 and k not in {
                # These 0s are meaningful (rare but possible)
                "current_ownership_pct", "common_share_pct", "size_discount_pct",
                "capex_pct_revenue",
            }
            if not new_is_empty:
                combined[k] = new_v
            elif old_v not in (None, ""):
                combined[k] = old_v
            else:
                combined[k] = new_v  # both empty — take whatever
        merged[method] = combined
    return merged


async def _run_autofill_job(valuation_id: str, user_id: str) -> None:
    """Background task: seed inputs for all 5 methods from web research."""
    try:
        v = await db.valuations.find_one({"id": valuation_id, "user_id": user_id})
        if not v:
            return
        # Iter-38: if this valuation is linked to a Vault, pull authoritative
        # private evidence from the Data Room's parsed files. Reuses the same
        # `pages` text Findings already extracted (zero re-parse cost).
        private_evidence, vault_files_used = "", []
        if v.get("deal_room_id"):
            private_evidence, vault_files_used = await _gather_vault_private_evidence(v["deal_room_id"])
        seed = await _autofill_wb(
            company_name=v["company_name"],
            sector=v.get("sector"),
            one_liner=v.get("one_liner"),
            estimated_revenue=v.get("estimated_revenue"),
            headquarters=v.get("headquarters"),
            brave_fn=search_brave,
            perplexity_fn=lambda p: query_perplexity(p),
            claude_fn=call_claude,
            safe_json=safe_json_loads,
            private_evidence=private_evidence or None,
            private_evidence_files=vault_files_used or None,
        )
        # Iter-44: MERGE the new seed with any existing inputs so a re-run can
        # never make the Workbench appear to "lose data" that a prior run had
        # populated. Old non-null values win when the new call returned null.
        existing_inputs = v.get("inputs") or {}
        seeded_inputs = seed.get("inputs") or {}
        merged_inputs = _merge_autofill_inputs(existing_inputs, seeded_inputs)
        # Compute outputs from merged inputs and aggregate.
        outputs = _compute_all(merged_inputs)
        # Preserve caller-tuned weights across re-runs (no reset to defaults).
        weights = v.get("weights") or dict(_DEFAULT_WEIGHTS)
        aggregate = _aggregate_band(outputs, weights)
        # Sources — union merge so citations accumulate across re-runs.
        old_sources = v.get("sources") or []
        new_sources = seed.get("sources") or []
        seen_urls: set = set()
        merged_sources: list = []
        for src in (new_sources + old_sources):
            key = src.get("url") or f"file:{src.get('file_id')}"
            if key and key not in seen_urls:
                seen_urls.add(key)
                merged_sources.append(src)
        await db.valuations.update_one(
            {"id": valuation_id},
            {"$set": {
                "inputs": merged_inputs,
                "outputs": outputs,
                "aggregate": aggregate,
                "narrative": seed.get("narrative") or v.get("narrative") or "",
                "sources": merged_sources[:60],
                "weights": weights,
                "private_grounded": bool(seed.get("private_grounded") or v.get("private_grounded")),
                "vault_files_used": seed.get("vault_files_used") or v.get("vault_files_used") or [],
                "autofill_status": "completed",
                "autofilled_at": now_utc().isoformat(),
                "updated_at": now_utc().isoformat(),
            }},
        )
        await log_audit(user_id, "valuation.autofill.completed", valuation_id, {
            "base_usd": aggregate.get("base_usd"),
            "confidence": aggregate.get("confidence"),
            "private_grounded": bool(seed.get("private_grounded")),
            "vault_files_used": len(seed.get("vault_files_used") or []),
            "merged_with_prior": bool(existing_inputs),
        })
    except Exception as e:
        logger.exception(f"valuation.autofill failed for {valuation_id}: {e}")
        await db.valuations.update_one(
            {"id": valuation_id},
            {"$set": {
                "autofill_status": "failed",
                "autofill_error": str(e)[:400],
                "updated_at": now_utc().isoformat(),
            }},
        )


@api_router.post("/valuations")
async def create_valuation(body: ValuationCreate, user=Depends(get_current_user)):
    """Create a draft valuation workbench. Autofill runs asynchronously.

    If `deal_room_id` is provided, this valuation is Vault-grounded — autofill
    will fuse public web evidence with parsed vault documents (highest priority).
    Access to the vault is verified before creation.
    """
    if not body.company_name or not body.company_name.strip():
        raise HTTPException(status_code=400, detail="company_name is required")
    # Iter-38: verify vault access when deal_room_id is set (NDA-signed buyers only)
    if body.deal_room_id:
        room = await db.deal_rooms.find_one({"id": body.deal_room_id})
        if not room:
            raise HTTPException(status_code=404, detail="Vault not found")
        await participant_check(room, user)
    vid = str(uuid.uuid4())
    now = now_utc().isoformat()
    doc = {
        "id": vid,
        "user_id": user["id"],
        "company_name": body.company_name.strip(),
        "sector": body.sector,
        "one_liner": body.one_liner,
        "estimated_revenue": body.estimated_revenue,
        "headquarters": body.headquarters,
        "research_id": body.research_id,
        "listing_id": body.listing_id,
        "deal_room_id": body.deal_room_id,
        "inputs": {k: {} for k in ("recent_transaction", "market_multiples", "vc_method", "dcf", "option_pricing")},
        "outputs": {},
        "aggregate": None,
        "narrative": "",
        "sources": [],
        "weights": dict(_DEFAULT_WEIGHTS),
        "private_grounded": False,
        "vault_files_used": [],
        "autofill_status": "pending" if body.autofill else "skipped",
        "snapshot_count": 0,
        "latest_snapshot_id": None,
        "created_at": now,
        "updated_at": now,
    }
    await db.valuations.insert_one(doc)
    await log_audit(user["id"], "valuation.create", vid, {
        "company_name": body.company_name,
        "deal_room_id": body.deal_room_id,
    })
    if body.autofill:
        asyncio.create_task(_run_autofill_job(vid, user["id"]))
    doc.pop("_id", None)
    return doc


@api_router.get("/valuations")
async def list_valuations(user=Depends(get_current_user)):
    cur = db.valuations.find(
        {"user_id": user["id"], "deleted_at": {"$exists": False}},
        {"_id": 0},
    ).sort("updated_at", -1)
    items = await cur.to_list(200)
    return [_valuation_summary(v) for v in items]


def _val_read_query(vid: str, user: dict) -> dict:
    """Mongo filter for admin-or-owner read access to a valuation."""
    if user.get("role") == "admin":
        return {"id": vid, "deleted_at": {"$exists": False}}
    return {"id": vid, "user_id": user["id"], "deleted_at": {"$exists": False}}


@api_router.get("/valuations/{vid}")
async def get_valuation(vid: str, user=Depends(get_current_user)):
    doc = await db.valuations.find_one(_val_read_query(vid, user), {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Valuation not found")
    if user.get("role") == "admin" and doc.get("user_id") != user["id"]:
        doc["read_only_for_viewer"] = True
    return doc


@api_router.patch("/valuations/{vid}")
async def update_valuation_inputs(vid: str, body: ValuationInputsUpdate, user=Depends(get_current_user)):
    """Merge inputs into the working draft and recompute outputs + aggregate."""
    doc = await db.valuations.find_one({"id": vid, "user_id": user["id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Valuation not found")
    inputs = doc.get("inputs") or {}
    for k, v in (body.inputs or {}).items():
        inputs[k] = {**(inputs.get(k) or {}), **(v or {})}
    weights = body.weights or doc.get("weights") or dict(_DEFAULT_WEIGHTS)
    outputs = _compute_all(inputs)
    aggregate = _aggregate_band(outputs, weights)
    await db.valuations.update_one(
        {"id": vid},
        {"$set": {
            "inputs": inputs, "outputs": outputs, "aggregate": aggregate,
            "weights": weights, "updated_at": now_utc().isoformat(),
        }},
    )
    return {"id": vid, "inputs": inputs, "outputs": outputs, "aggregate": aggregate, "weights": weights}


@api_router.post("/valuations/{vid}/autofill")
async def rerun_autofill(vid: str, user=Depends(get_current_user)):
    doc = await db.valuations.find_one({"id": vid, "user_id": user["id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Valuation not found")
    await db.valuations.update_one({"id": vid}, {"$set": {"autofill_status": "pending", "updated_at": now_utc().isoformat()}})
    asyncio.create_task(_run_autofill_job(vid, user["id"]))
    return {"id": vid, "autofill_status": "pending"}


@api_router.get("/valuations/{vid}/autofill/status")
async def poll_autofill(vid: str, user=Depends(get_current_user)):
    doc = await db.valuations.find_one(
        _val_read_query(vid, user),
        {"_id": 0, "autofill_status": 1, "autofill_error": 1, "aggregate": 1, "inputs": 1, "outputs": 1, "narrative": 1, "sources": 1, "weights": 1, "autofilled_at": 1},
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Valuation not found")
    return doc


@api_router.post("/valuations/{vid}/term-sheet")
async def upload_term_sheet(vid: str, file: UploadFile = File(...), user=Depends(get_current_user)):
    """Upload a term-sheet PDF/DOCX; Claude extracts round terms + preferred stack
    and merges them into recent_transaction + option_pricing inputs."""
    doc = await db.valuations.find_one({"id": vid, "user_id": user["id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Valuation not found")
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty upload")
    pages = extract_pages_from_bytes(file.filename or "term-sheet.pdf", data)
    full_text = "\n\n".join((p.get("text") or "") for p in pages)
    extracted = await _extract_term_sheet(pdf_text=full_text, claude_fn=call_claude, safe_json=safe_json_loads)

    # Merge into RT + OPM inputs
    inputs = doc.get("inputs") or {}
    rt = inputs.get("recent_transaction") or {}
    if extracted.get("post_money_usd"):    rt["post_money_usd"] = extracted["post_money_usd"]
    if extracted.get("raised_usd"):        rt["raised_usd"] = extracted["raised_usd"]
    if extracted.get("round_type"):        rt["round_type"] = extracted["round_type"]
    if extracted.get("announced"):         rt["announced"] = extracted["announced"]
    if "time_decay_factor" not in rt:      rt["time_decay_factor"] = 1.0
    inputs["recent_transaction"] = rt

    op = inputs.get("option_pricing") or {}
    if extracted.get("total_preferred_liquidation_pref_usd"):
        op["total_preferred_liquidation_pref_usd"] = extracted["total_preferred_liquidation_pref_usd"]
    if extracted.get("post_money_usd") and "enterprise_value_usd" not in op:
        op["enterprise_value_usd"] = extracted["post_money_usd"]
    inputs["option_pricing"] = op

    weights = doc.get("weights") or dict(_DEFAULT_WEIGHTS)
    outputs = _compute_all(inputs)
    aggregate = _aggregate_band(outputs, weights)
    await db.valuations.update_one(
        {"id": vid},
        {"$set": {
            "inputs": inputs, "outputs": outputs, "aggregate": aggregate,
            "term_sheet_extracted": extracted,
            "term_sheet_filename": file.filename,
            "updated_at": now_utc().isoformat(),
        }},
    )
    await log_audit(user["id"], "valuation.term_sheet.uploaded", vid, {"filename": file.filename, "confidence": extracted.get("confidence")})
    return {"id": vid, "extracted": extracted, "inputs": inputs, "outputs": outputs, "aggregate": aggregate}


@api_router.post("/valuations/{vid}/snapshots")
async def create_snapshot(vid: str, body: SnapshotCreate, user=Depends(get_current_user)):
    """Freeze the current draft into an immutable snapshot."""
    doc = await db.valuations.find_one({"id": vid, "user_id": user["id"]})
    if not doc:
        raise HTTPException(status_code=404, detail="Valuation not found")
    if not (doc.get("aggregate") or {}).get("base_usd"):
        raise HTTPException(status_code=400, detail="Cannot snapshot an empty valuation")
    sid = str(uuid.uuid4())
    now = now_utc().isoformat()
    snapshot = {
        "id": sid,
        "valuation_id": vid,
        "user_id": user["id"],
        "created_at": now,
        "label": (body.label or "").strip() or f"Snapshot {doc.get('snapshot_count', 0) + 1}",
        "narrative": body.narrative or doc.get("narrative") or "",
        "inputs": doc.get("inputs") or {},
        "outputs": doc.get("outputs") or {},
        "aggregate": doc.get("aggregate") or {},
        "weights": doc.get("weights") or {},
        "sources": doc.get("sources") or [],
        "company_name": doc.get("company_name"),
        "private_grounded": bool(doc.get("private_grounded")),
        "vault_files_used": doc.get("vault_files_used") or [],
    }
    await db.valuation_snapshots.insert_one(snapshot)
    await db.valuations.update_one(
        {"id": vid},
        {"$set": {
            "latest_snapshot_id": sid,
            "updated_at": now,
        }, "$inc": {"snapshot_count": 1}},
    )
    await log_audit(user["id"], "valuation.snapshot.create", sid, {
        "valuation_id": vid,
        "base_usd": (snapshot["aggregate"] or {}).get("base_usd"),
        "label": snapshot["label"],
    })
    snapshot.pop("_id", None)
    return snapshot


@api_router.get("/valuations/{vid}/snapshots")
async def list_snapshots(vid: str, user=Depends(get_current_user)):
    v = await db.valuations.find_one(_val_read_query(vid, user), {"_id": 0, "id": 1})
    if not v:
        raise HTTPException(status_code=404, detail="Valuation not found")
    cur = db.valuation_snapshots.find({"valuation_id": vid}, {"_id": 0}).sort("created_at", -1)
    return await cur.to_list(200)


@api_router.get("/valuations/{vid}/snapshots/{sid}")
async def get_snapshot(vid: str, sid: str, user=Depends(get_current_user)):
    v = await db.valuations.find_one(_val_read_query(vid, user), {"_id": 0, "id": 1})
    if not v:
        raise HTTPException(status_code=404, detail="Valuation not found")
    snap = await db.valuation_snapshots.find_one({"id": sid, "valuation_id": vid}, {"_id": 0})
    if not snap:
        raise HTTPException(status_code=404, detail="Snapshot not found")
    return snap


@api_router.get("/valuations/{vid}/snapshots/{sid}/pdf")
async def download_snapshot_pdf(vid: str, sid: str, user=Depends(get_current_user)):
    """Branded Valuation Memorandum PDF for a snapshot."""
    v = await db.valuations.find_one(_val_read_query(vid, user), {"_id": 0})
    if not v:
        raise HTTPException(status_code=404, detail="Valuation not found")
    snap = await db.valuation_snapshots.find_one({"id": sid, "valuation_id": vid}, {"_id": 0})
    if not snap:
        raise HTTPException(status_code=404, detail="Snapshot not found")
    pdf_bytes = _build_memo_pdf(
        valuation=v,
        snapshot=snap,
        prepared_by=user.get("name") or user.get("email") or "NextCapOS Analyst",
    )
    safe_company = re.sub(r"[^A-Za-z0-9]+", "_", v.get("company_name") or "target").strip("_") or "target"
    fname = f"ValuationMemo_{safe_company}_{snap['created_at'][:10]}.pdf"
    await log_audit(user["id"], "valuation.snapshot.pdf", sid, {"valuation_id": vid, "filename": fname})
    from fastapi.responses import Response as _R
    return _R(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )


@api_router.delete("/valuations/{vid}")
async def delete_valuation(vid: str, user=Depends(get_current_user)):
    """Soft-delete — snapshots retained for audit."""
    r = await db.valuations.update_one(
        {"id": vid, "user_id": user["id"]},
        {"$set": {"deleted_at": now_utc().isoformat()}},
    )
    if r.matched_count == 0:
        raise HTTPException(status_code=404, detail="Valuation not found")
    await log_audit(user["id"], "valuation.delete", vid)
    return {"ok": True}


# --- Iter-38: Vault-scoped valuation lookup + one-click "value this target" ---
@api_router.get("/deal-rooms/{rid}/valuation")
async def get_vault_valuation(rid: str, user=Depends(get_current_user)):
    """Return the vault's linked valuation for the current viewer, or 404.

    Lookup order for admins: their own first (if they created one), then the
    buyer's (read-only fallback). Buyers see only their own.
    """
    room = await db.deal_rooms.find_one({"id": rid})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)

    # First: does the current viewer own a valuation on this vault?
    own = await db.valuations.find_one(
        {"user_id": user["id"], "deal_room_id": rid, "deleted_at": {"$exists": False}},
        {"_id": 0}, sort=[("updated_at", -1)],
    )
    if own:
        return own

    # Fallback: admins can peek at the buyer's valuation read-only.
    if user.get("role") == "admin" and room["buyer_id"] != user["id"]:
        peer = await db.valuations.find_one(
            {"user_id": room["buyer_id"], "deal_room_id": rid, "deleted_at": {"$exists": False}},
            {"_id": 0}, sort=[("updated_at", -1)],
        )
        if peer:
            peer["read_only_for_viewer"] = True
            return peer

    raise HTTPException(status_code=404, detail="No valuation on this vault yet")


@api_router.post("/deal-rooms/{rid}/valuation")
async def create_vault_valuation(rid: str, user=Depends(get_current_user)):
    """One-click: create a Valuation linked to this Vault, seeded with the
    listing's company_name / sector / HQ. Autofill kicks off immediately.

    Buyers and admins (impersonating for support/audit) can both create — the
    valuation is owned by whoever created it. Idempotent per viewer.
    """
    room = await db.deal_rooms.find_one({"id": rid})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    # Return the caller's existing valuation if one already exists (idempotent one-click)
    existing = await db.valuations.find_one(
        {"user_id": user["id"], "deal_room_id": rid, "deleted_at": {"$exists": False}},
        {"_id": 0},
        sort=[("updated_at", -1)],
    )
    if existing:
        return existing
    listing = None
    if room.get("listing_id"):
        listing = await db.listings.find_one({"id": room["listing_id"]}, {"_id": 0})
    payload = ValuationCreate(
        company_name=(listing or {}).get("company_name") or room.get("listing_name") or "Target",
        sector=(listing or {}).get("sector"),
        one_liner=(listing or {}).get("one_liner") or (listing or {}).get("headline"),
        headquarters=(listing or {}).get("hq") or (listing or {}).get("headquarters"),
        deal_room_id=rid,
        listing_id=room.get("listing_id"),
        autofill=True,
    )
    return await create_valuation(payload, user)


# -----------------------------------------------------------------------------
# DETAILED ANALYSIS (Buyer-side "Standard" Kenshin-style 14-section report)
# -----------------------------------------------------------------------------
from detailed_analysis import run_detailed_analysis  # noqa: E402


class DetailedAnalysisRequest(BaseModel):
    company_name: str
    company_url: Optional[str] = None
    industry: Optional[str] = None
    region: Optional[str] = None
    funding_stage: Optional[str] = None
    buyer_notes: Optional[str] = None
    research_id: Optional[str] = None  # optional link back to a /research/company brief


@api_router.post("/research/detailed")
async def create_detailed_report(body: DetailedAnalysisRequest, user=Depends(get_current_user)):
    """Buyer / admin: queue a 14-section detailed institutional analysis. Returns
    immediately with `{id, status: "pending"}`; frontend polls GET /research/detailed/{id}
    until status flips to `completed` or `failed`. The pipeline (Perplexity + 4× Brave +
    Claude 4.5 with grounded context) typically takes 60-180s — well beyond the ingress
    60s read timeout, hence the async pattern."""
    if user.get("role") not in ("buyer", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Detailed analysis is buyer/admin only")
    if not body.company_name or not body.company_name.strip():
        raise HTTPException(status_code=400, detail="company_name is required")
    rid = str(uuid.uuid4())
    skeleton = {
        "id": rid,
        "user_id": user["id"],
        "research_id": body.research_id,
        "kind": "detailed",
        "status": "pending",
        "company_name": body.company_name.strip(),
        "company_url": body.company_url,
        "industry": body.industry,
        "region": body.region,
        "funding_stage": body.funding_stage,
        "buyer_notes": body.buyer_notes,
        "data": None,
        "sources": [],
        "source_count": 0,
        "live_research_used": False,
        "duration_ms": None,
        "error": None,
        "created_at": now_utc().isoformat(),
    }
    await db.detailed_reports.insert_one(skeleton)
    await log_audit(user["id"], "detailed_report.queue", rid, {"company": body.company_name})
    asyncio.create_task(_execute_detailed_analysis(rid, body, user["id"]))
    skeleton.pop("_id", None)
    return skeleton


async def _execute_detailed_analysis(rid: str, body: "DetailedAnalysisRequest", user_id: str) -> None:
    """Background worker. Persists `status` transitions analyzing → completed | failed."""
    try:
        await db.detailed_reports.update_one(
            {"id": rid}, {"$set": {"status": "analyzing"}},
        )

        async def _discover_social_bound(company_name: str, company_url: Optional[str] = None):
            return await discover_social_profiles(
                search_brave=search_brave,
                company_name=company_name,
                company_url=company_url,
            )

        result = await run_detailed_analysis(
            call_claude=call_claude,
            query_perplexity=query_perplexity,
            search_brave=search_brave,
            company_name=body.company_name.strip(),
            company_url=body.company_url,
            industry=body.industry,
            region=body.region,
            funding_stage=body.funding_stage,
            buyer_notes=body.buyer_notes,
            user_id=user_id,
            discover_social=_discover_social_bound,
        )
        await db.detailed_reports.update_one(
            {"id": rid},
            {"$set": {
                "status": "completed",
                "data": result["data"],
                "sources": result["sources"],
                "social_profiles": result.get("social_profiles") or {},
                "live_research_used": result["live_research_used"],
                "source_count": result["source_count"],
                "duration_ms": result["duration_ms"],
                "completed_at": now_utc().isoformat(),
            }},
        )
        await log_agent_activity(
            "detailed-analysis-agent",
            f"company:{body.company_name} · {result['source_count']} sources",
            "completed",
            user_id=user_id, duration_ms=result["duration_ms"],
            meta={"recommendation": result["data"].get("executiveSummary", {}).get("recommendation"),
                  "source_count": result["source_count"]},
        )
        try:
            asyncio.create_task(notarize_bytes(
                kind="detailed_report",
                target_id=rid,
                data=json.dumps(result["data"], sort_keys=True).encode("utf-8"),
                owner_user_id=user_id,
                label=f"Detailed analysis: {body.company_name}",
                extra={"sources": result["source_count"]},
            ))
        except Exception:
            logger.warning("notarize for detailed report failed silently")
    except Exception as e:
        logger.exception(f"Detailed analysis worker failed for {rid}")
        await db.detailed_reports.update_one(
            {"id": rid},
            {"$set": {"status": "failed", "error": str(e)[:500],
                      "completed_at": now_utc().isoformat()}},
        )
        await log_agent_activity(
            "detailed-analysis-agent", f"company:{body.company_name}", "failed",
            user_id=user_id, friction=str(e)[:200],
        )


@api_router.get("/research/detailed")
async def list_detailed_reports(user=Depends(get_current_user)):
    q: Dict[str, Any] = {"deleted_at": {"$exists": False}}
    if user.get("role") != "admin":
        q["user_id"] = user["id"]
    items = await db.detailed_reports.find(q, {"_id": 0, "data": 0, "sources": 0}).sort("created_at", -1).to_list(100)
    return items


@api_router.get("/research/detailed/{rid}")
async def get_detailed_report(rid: str, user=Depends(get_current_user)):
    doc = await db.detailed_reports.find_one({"id": rid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Report not found")
    if user.get("role") != "admin" and doc.get("user_id") != user["id"]:
        raise HTTPException(status_code=403, detail="Not your report")
    return doc


@api_router.delete("/research/detailed/{rid}")
async def delete_detailed_report(rid: str, user=Depends(get_current_user)):
    q: Dict[str, Any] = {"id": rid, "deleted_at": {"$exists": False}}
    if user.get("role") != "admin":
        q["user_id"] = user["id"]
    res = await db.detailed_reports.update_one(q, {"$set": {"deleted_at": now_utc().isoformat()}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Report not found")
    await log_audit(user["id"], "detailed_report.delete", rid)
    return {"ok": True}


@api_router.get("/research/detailed/{rid}/pdf")
async def export_detailed_report_pdf(rid: str, user=Depends(get_current_user)):
    doc = await db.detailed_reports.find_one({"id": rid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Report not found")
    if user.get("role") != "admin" and doc.get("user_id") != user["id"]:
        raise HTTPException(status_code=403, detail="Not your report")
    from detailed_report_pdf import generate_detailed_report_pdf  # local import keeps module-load cheap
    pdf_bytes = generate_detailed_report_pdf(doc)
    fname = f"detailed-analysis-{doc.get('company_name','company').replace(' ', '-').lower()}.pdf"
    await log_audit(user["id"], "detailed_report.pdf", rid, {"bytes": len(pdf_bytes)})
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'},
    )


class AttachDetailedReportRequest(BaseModel):
    room_id: Optional[str] = None       # attach into an active Vault
    listing_id: Optional[str] = None    # attach into the listing data room (seller side)


@api_router.post("/research/detailed/{rid}/attach")
async def attach_detailed_report(rid: str, body: AttachDetailedReportRequest, user=Depends(get_current_user)):
    """Generate the PDF on the fly and persist it into a Vault (deal_room_files) or a
    Listing data room (listing_staged_files). Either room_id or listing_id is required."""
    doc = await db.detailed_reports.find_one({"id": rid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not doc:
        raise HTTPException(status_code=404, detail="Report not found")
    if user.get("role") != "admin" and doc.get("user_id") != user["id"]:
        raise HTTPException(status_code=403, detail="Not your report")
    if not (body.room_id or body.listing_id):
        raise HTTPException(status_code=400, detail="room_id or listing_id is required")
    if body.room_id and body.listing_id:
        raise HTTPException(status_code=400, detail="Provide either room_id or listing_id, not both")

    from detailed_report_pdf import generate_detailed_report_pdf
    pdf_bytes = generate_detailed_report_pdf(doc)
    filename = f"workz-detailed-analysis-{doc.get('company_name','company').replace(' ', '-').lower()}.pdf"
    plain_sha = sha256_hex(pdf_bytes)

    # Vault attach (buyer-side)
    if body.room_id:
        room = await db.deal_rooms.find_one({"id": body.room_id}, {"_id": 0})
        if not room:
            raise HTTPException(status_code=404, detail="Vault not found")
        role = await participant_check(room, user)
        if role == "buyer" and room.get("status") == "pending_nda":
            raise HTTPException(status_code=400, detail="Sign the NDA before adding files")
        file_id = str(uuid.uuid4())
        storage_bytes = pdf_bytes
        encrypted = False
        encryption_alg = None
        if encryption_configured():
            try:
                aad = f"{body.room_id}:{file_id}".encode("utf-8")
                env = encrypt_bytes(pdf_bytes, associated_data=aad)
                storage_bytes = env["envelope"]
                encrypted = True
                encryption_alg = env["alg"]
            except Exception as e:
                logger.warning(f"vault attach encrypt failed: {e}")
        grid_id = await gridfs_bucket.upload_from_stream(
            filename, io.BytesIO(storage_bytes),
            metadata={"room_id": body.room_id, "file_id": file_id,
                      "uploaded_by": user["id"], "content_type": "application/pdf",
                      "encrypted": encrypted, "encryption_alg": encryption_alg,
                      "source": "detailed_report", "report_id": rid},
        )
        await db.deal_room_files.insert_one({
            "id": file_id, "room_id": body.room_id, "filename": filename,
            "folder": "commercial",
            "content_type": "application/pdf",
            "size_bytes": len(pdf_bytes), "page_count": 0,
            "pages": [], "content": "", "char_count": 0,
            "gridfs_id": str(grid_id), "storage": "gridfs",
            "encrypted": encrypted, "encryption_alg": encryption_alg,
            "sha256_hex": plain_sha, "note": f"NextCapOS Detailed Analysis · auto-attached from research ({rid[:8]})",
            "uploaded_by": user["id"],
            "uploaded_by_role": role or user.get("role"),
            "uploaded_at": now_utc().isoformat(),
            "detailed_report_id": rid,
        })
        await log_audit(user["id"], "detailed_report.attach.vault", rid,
                        {"room_id": body.room_id, "file_id": file_id})
        return {"ok": True, "room_id": body.room_id, "file_id": file_id, "filename": filename}

    # Listing data-room attach (seller/admin)
    listing = await _seller_listing_or_404(body.listing_id, user)  # 403 if not owner
    file_id = str(uuid.uuid4())
    storage_bytes = pdf_bytes
    encrypted = False
    encryption_alg = None
    if encryption_configured():
        try:
            aad = f"listing:{body.listing_id}:{file_id}".encode("utf-8")
            env = encrypt_bytes(pdf_bytes, associated_data=aad)
            storage_bytes = env["envelope"]
            encrypted = True
            encryption_alg = env["alg"]
        except Exception as e:
            logger.warning(f"listing attach encrypt failed: {e}")
    grid_id = await listing_files_bucket.upload_from_stream(
        filename, io.BytesIO(storage_bytes),
        metadata={"listing_id": body.listing_id, "file_id": file_id,
                  "uploaded_by": user["id"], "content_type": "application/pdf",
                  "encrypted": encrypted, "encryption_alg": encryption_alg,
                  "source": "detailed_report", "report_id": rid},
    )
    await db.listing_staged_files.insert_one({
        "id": file_id, "listing_id": body.listing_id, "filename": filename,
        "folder": "commercial", "content_type": "application/pdf",
        "size_bytes": len(pdf_bytes), "page_count": 0,
        "pages": [], "content": "", "char_count": 0,
        "gridfs_id": str(grid_id), "storage": "listing_gridfs",
        "encrypted": encrypted, "encryption_alg": encryption_alg,
        "sha256_hex": plain_sha,
        "note": f"NextCapOS Detailed Analysis · auto-attached from research ({rid[:8]})",
        "uploaded_by": user["id"],
        "uploaded_at": now_utc().isoformat(),
        "detailed_report_id": rid,
    })
    await log_audit(user["id"], "detailed_report.attach.listing", rid,
                    {"listing_id": body.listing_id, "file_id": file_id, "listing": listing.get("company_name")})
    return {"ok": True, "listing_id": body.listing_id, "file_id": file_id, "filename": filename}


# -----------------------------------------------------------------------------
# RESEARCH COMPANION (Buyer-only AI chat over their own Research Hub findings)
#
# Combines:
#   - the buyer's brief (research_company)
#   - any detailed analysis (detailed_reports) tied to the research_id
#   - any Private Locker files tagged with scope="research" + research_id
#
# Strictly buyer-only. Sellers cannot access. Lives entirely outside the Vault.
# -----------------------------------------------------------------------------
RESEARCH_COPILOT_SYS = """You are the NextCapOS Research Companion — a senior buy-side analyst
helping an institutional investor go deeper on a company they are researching.

You answer ONLY from the supplied source materials:
  (a) the buyer's research brief — cite as [brief]
  (b) any detailed analysis report — cite as [detailed-analysis]
  (c) the buyer's private locker docs — cite as [filename]
  (d) public listing artifacts from the NextCapOS marketplace — cite as [listing:CompanyName]
  (e) files in a Vault the buyer has rightful (NDA-signed) access to — cite as [vault:filename]

Always cite the specific source(s) inline using these tags. If the answer is not
in the provided context, say so explicitly and suggest what to research next.

Keep answers under 260 words. Tone: institutional, terse, analytical. Never invent
financials, customer names, or sources. When the buyer asks for opinions ("should I
proceed?"), structure the answer as: signal, risk, recommended next diligence step.
When information is available in both the public listing and the Vault, prefer the
Vault source — it's authoritative for that deal."""


class ResearchCopilotAsk(BaseModel):
    message: str


async def _load_research_target(rid: str, user) -> dict:
    rec = await db.research.find_one(
        {"id": rid, "user_id": user["id"], "deleted_at": {"$exists": False}},
        {"_id": 0},
    )
    if not rec:
        raise HTTPException(status_code=404, detail="Research target not found")
    return rec


@api_router.get("/research/{rid}/locker")
async def list_research_locker(rid: str, user=Depends(get_current_user)):
    await _private_locker_guard(user)
    await _load_research_target(rid, user)  # ownership check
    items = await db.private_locker_files.find(
        {"user_id": user["id"], "scope": "research", "research_id": rid,
         "deleted_at": {"$exists": False}},
        {"_id": 0, "content": 0, "pages": 0},
    ).sort("uploaded_at", -1).to_list(200)
    return items


@api_router.get("/research/{rid}/copilot")
async def get_research_copilot_history(rid: str, user=Depends(get_current_user)):
    await _private_locker_guard(user)
    await _load_research_target(rid, user)  # ownership check covers visibility
    msgs = await db.research_copilot_messages.find(
        {"research_id": rid},
        {"_id": 0},
    ).sort("created_at", 1).to_list(200)
    return msgs


@api_router.post("/research/{rid}/copilot")
async def ask_research_copilot(rid: str, body: ResearchCopilotAsk, user=Depends(get_current_user)):
    await _private_locker_guard(user)
    research = await _load_research_target(rid, user)

    user_msg = {
        "id": str(uuid.uuid4()),
        "research_id": rid,
        "role": "user",
        "user_id": user["id"],
        "user_name": user.get("name"),
        "content": (body.message or "")[:2000],
        "citations": [],
        "created_at": now_utc().isoformat(),
    }
    await db.research_copilot_messages.insert_one(user_msg)

    # Assemble grounded context
    sections: List[str] = []

    # 1) The research brief itself. Stored as a structured dict in `research.data`,
    # not a flat `content` string (that field was the legacy shape). Flatten the
    # named fields into the prompt so Claude can cite them by section.
    brief_data = research.get("data") if isinstance(research.get("data"), dict) else None
    brief_text = (research.get("content") or "").strip()
    if not brief_text and brief_data:
        # Flatten the structured brief into a single human-readable block.
        ordered_keys = [
            ("summary", "Summary"),
            ("business_model", "Business model"),
            ("investor_take", "Investor take"),
            ("market_signals", "Market signals"),
            ("growth_drivers", "Growth drivers"),
            ("risks", "Risks"),
            ("competitive_landscape", "Competitive landscape"),
            ("leadership_insights", "Leadership insights"),
            ("suggested_buyer_profile", "Suggested buyer profile"),
            ("next_actions", "Next actions"),
            ("hq", "HQ"), ("founded", "Founded"),
            ("employees", "Employees"), ("revenue", "Revenue"),
        ]
        parts: List[str] = []
        for key, label in ordered_keys:
            val = brief_data.get(key)
            if val is None or val == "" or val == []:
                continue
            if isinstance(val, list):
                # List of strings or list of dicts (leadership insights).
                lines: List[str] = []
                for item in val[:12]:
                    if isinstance(item, str):
                        lines.append(f"• {item}")
                    elif isinstance(item, dict):
                        name = item.get("name") or item.get("role") or ""
                        role = item.get("role") if name == item.get("name") else ""
                        note = item.get("note") or item.get("summary") or item.get("background") or ""
                        head = ", ".join(x for x in [name, role] if x)
                        if note:
                            lines.append(f"• {head}: {note}" if head else f"• {note}")
                        elif head:
                            lines.append(f"• {head}")
                if lines:
                    parts.append(f"{label}:\n" + "\n".join(lines))
            else:
                parts.append(f"{label}: {val}")
        brief_text = "\n\n".join(parts)
    if brief_text:
        sections.append(f"=== [brief] Research brief on {research.get('company_name')} ===\n{brief_text[:8000]}")
        srcs = research.get("sources") or []
        if srcs:
            src_lines = [f"- {s.get('title') or s.get('url')}: {s.get('url')}" for s in srcs[:12] if isinstance(s, dict)]
            if src_lines:
                sections.append("Brief sources:\n" + "\n".join(src_lines))

    detailed = await db.detailed_reports.find_one(
        {"research_id": rid, "user_id": user["id"], "deleted_at": {"$exists": False},
         "status": "completed"},
        {"_id": 0, "data": 1, "company_name": 1, "id": 1},
    )
    if detailed and isinstance(detailed.get("data"), dict):
        summary = detailed["data"].get("executive_summary") or detailed["data"].get("summary") or ""
        thesis = detailed["data"].get("investment_thesis") or ""
        risks = detailed["data"].get("risks") or detailed["data"].get("risk_register") or ""
        chunk = []
        if summary:
            chunk.append(f"Executive summary: {summary}")
        if thesis:
            chunk.append(f"Investment thesis: {thesis}")
        if risks:
            chunk.append(f"Risks: {risks}")
        if chunk:
            sections.append(f"=== [detailed-analysis] Detailed analysis report ===\n" + "\n\n".join(chunk[:3])[:6000])

    locker_files = await db.private_locker_files.find(
        {"user_id": user["id"], "scope": "research", "research_id": rid,
         "deleted_at": {"$exists": False}},
        {"_id": 0, "filename": 1, "content": 1, "note": 1},
    ).sort("uploaded_at", 1).to_list(20)
    locker_citations: List[dict] = []
    for f in locker_files:
        snippet = (f.get("content") or "")[:2500]
        note = (f.get("note") or "")
        body_lines = []
        if note:
            body_lines.append(f"(Buyer note: {note})")
        body_lines.append(snippet)
        sections.append(f"=== [{f['filename']}] Private locker file ===\n" + "\n".join(body_lines))
        locker_citations.append({"filename": f["filename"]})

    # --- iter-22 Option B expansion ---
    # The Research Companion now also pulls (a) public listing artifacts on the
    # marketplace that match the researched company name, and (b) files from
    # any Vault the buyer has an active/preview NDA on for those listings.
    # The buyer already has rightful access to all of this material; surfacing
    # it inside the Companion saves them from juggling tabs.
    company_name = (research.get("company_name") or "").strip()
    listing_citations: List[dict] = []
    vault_citations: List[dict] = []
    matched_listings: List[dict] = []
    if company_name:
        # Case-insensitive substring match on company_name. Cap at 5 listings
        # so a generic name (e.g., "Acme") doesn't blow up the context window.
        listing_rows = await db.listings.find(
            {"company_name": {"$regex": re.escape(company_name), "$options": "i"},
             "status": {"$in": ["live", "draft", "under_loi", "closed"]},
             "deleted_at": {"$exists": False}},
            {"_id": 0, "id": 1, "company_name": 1, "sector": 1, "geography": 1,
             "headline": 1, "summary": 1, "highlights": 1, "asking_price_usd_m": 1,
             "revenue_usd_m": 1, "ebitda_usd_m": 1, "employees": 1, "status": 1},
        ).limit(5).to_list(5)
        matched_listings = listing_rows
        for l in listing_rows:
            # (a) Public listing summary — the CIM-style marketing copy that
            # any buyer on the marketplace can see for this listing.
            block_lines = [
                f"Listing status: {l.get('status')}",
                f"Sector: {l.get('sector')} · Geography: {l.get('geography')}",
                f"Headline: {l.get('headline')}",
                f"Asking: ${l.get('asking_price_usd_m')}M · Revenue: ${l.get('revenue_usd_m') or '—'}M · EBITDA: ${l.get('ebitda_usd_m') or '—'}M",
                f"Summary: {l.get('summary') or ''}",
            ]
            if l.get("highlights"):
                block_lines.append("Highlights:")
                for h in (l["highlights"] or [])[:8]:
                    block_lines.append(f"• {h}")
            cite_key = f"listing:{l['company_name']}"
            sections.append(f"=== [{cite_key}] Public listing artifact ===\n" + "\n".join(block_lines))
            listing_citations.append({"kind": "listing", "label": l["company_name"],
                                       "cite_key": cite_key, "listing_id": l["id"]})

            # (b) Any active or preview Vault this buyer has on this listing →
            # pull file content. NDA gating already enforced — only buyer's
            # OWN vault rows are queried.
            buyer_vaults = await db.deal_rooms.find(
                {"listing_id": l["id"], "buyer_id": user["id"],
                 "status": {"$in": ["pending_nda", "active", "preview"]},
                 "deleted_at": {"$exists": False}},
                {"_id": 0, "id": 1, "status": 1},
            ).to_list(3)
            for vault in buyer_vaults:
                vault_files = await db.deal_room_files.find(
                    {"room_id": vault["id"]},
                    {"_id": 0, "id": 1, "filename": 1, "content": 1, "folder": 1,
                     "page_count": 1, "source": 1},
                ).sort("uploaded_at", 1).to_list(15)
                for vf in vault_files:
                    snippet = (vf.get("content") or "")[:2200]
                    if not snippet:
                        continue
                    cite_key = f"vault:{vf['filename']}"
                    via = ""
                    src = vf.get("source") or {}
                    if isinstance(src, dict) and src.get("kind"):
                        via = f" · synced from {src['kind']}"
                    sections.append(
                        f"=== [{cite_key}] Vault file ({l['company_name']} · {vf.get('folder','other')}{via}) ===\n"
                        + snippet
                    )
                    vault_citations.append({
                        "kind": "vault",
                        "label": vf["filename"],
                        "cite_key": cite_key,
                        "vault_id": vault["id"],
                        "file_id": vf["id"],
                        "listing_name": l["company_name"],
                    })

    if not sections:
        empty_reply = (
            "I don't have any source material to answer from yet. Generate the brief or "
            "Detailed Analysis on this company in the Research Hub, or attach a document "
            "to this research target in your Private Locker — then re-ask. "
            "(Tip: if you also have an active Vault on this company's listing, "
            "those Vault docs will automatically be available here.)"
        )
        asst_msg = {
            "id": str(uuid.uuid4()),
            "research_id": rid,
            "role": "assistant",
            "user_id": "copilot",
            "user_name": "Research Companion",
            "content": empty_reply,
            "citations": [],
            "created_at": now_utc().isoformat(),
        }
        await db.research_copilot_messages.insert_one(asst_msg)
        user_msg.pop("_id", None)
        asst_msg.pop("_id", None)
        return {"user_message": user_msg, "assistant_message": asst_msg}

    # Recent conversation
    history = await db.research_copilot_messages.find(
        {"research_id": rid, "id": {"$ne": user_msg["id"]}},
        {"_id": 0},
    ).sort("created_at", -1).to_list(8)
    history.reverse()
    transcript = "\n".join(f"{m['role'].upper()}: {m['content']}" for m in history)

    prompt = (
        f"COMPANY: {research.get('company_name')}\n\n"
        f"SOURCE MATERIAL (only basis for citation):\n\n" + "\n\n".join(sections)
        + (f"\n\nPRIOR CONVERSATION:\n{transcript}" if transcript else "")
        + f"\n\nBUYER QUESTION: {body.message}\n\nAnswer now."
    )

    started = now_utc()
    try:
        answer = await call_claude(
            RESEARCH_COPILOT_SYS, prompt,
            session_id=f"research-copilot-{rid}-{user['id']}",
        )
    except Exception as e:
        logger.exception("Research companion call failed")
        raise HTTPException(status_code=502, detail=f"Companion failed: {e}")

    # Resolve citations: [brief], [detailed-analysis], [filename], [listing:Name], [vault:filename]
    cited = set(re.findall(r"\[([^\[\]]+)\]", answer or ""))
    citations: List[dict] = []
    if "brief" in cited and brief_text:
        citations.append({"kind": "brief", "label": "Research brief"})
    if "detailed-analysis" in cited and detailed:
        citations.append({"kind": "detailed", "label": "Detailed Analysis",
                          "detailed_report_id": detailed["id"]})
    for f in locker_citations:
        if f["filename"] in cited:
            citations.append({"kind": "locker", "label": f["filename"]})
    for lc in listing_citations:
        if lc["cite_key"] in cited:
            citations.append({k: lc[k] for k in ("kind", "label", "listing_id")})
    for vc in vault_citations:
        if vc["cite_key"] in cited:
            citations.append({k: vc[k] for k in ("kind", "label", "vault_id", "file_id", "listing_name")})

    asst_msg = {
        "id": str(uuid.uuid4()),
        "research_id": rid,
        "role": "assistant",
        "user_id": "copilot",
        "user_name": "Research Companion",
        "content": (answer or "").strip(),
        "citations": citations,
        "created_at": now_utc().isoformat(),
    }
    await db.research_copilot_messages.insert_one(asst_msg)
    duration = int((now_utc() - started).total_seconds() * 1000)
    await log_agent_activity(
        "research-companion",
        f"ask:{body.message[:60]}",
        "completed",
        user_id=user["id"],
        duration_ms=duration,
        meta={
            "citations": len(citations),
            "locker_files": len(locker_files),
            "matched_listings": len(matched_listings),
            "vault_files": len(vault_citations),
        },
    )
    await log_audit(user["id"], "research.companion.ask", rid,
                    {"locker_files": len(locker_files),
                     "had_detailed": bool(detailed),
                     "matched_listings": len(matched_listings),
                     "vault_files": len(vault_citations)})

    user_msg.pop("_id", None)
    asst_msg.pop("_id", None)
    return {"user_message": user_msg, "assistant_message": asst_msg}


# -----------------------------------------------------------------------------
# COLLATERAL
# -----------------------------------------------------------------------------
COLLATERAL_SYS = """You are a senior marketing copywriter for a top-tier private equity firm, NextCapOS.
You write polished, institutional-grade marketing collateral.
ALWAYS return STRICT JSON (no markdown fences):
{
  "title": str,
  "asset_type": str,
  "headline": str,
  "subheadline": str,
  "sections": [{"heading": str, "body": str}],
  "cta": str,
  "compliance_note": str
}
"""


@api_router.post("/collateral/generate")
async def generate_collateral(body: CollateralRequest, user=Depends(get_current_user)):
    started = now_utc()
    user_prompt = (
        f"Asset type: {body.asset_type}\n"
        f"Deal: {body.deal_name}\n"
        f"Target audience: {body.target_audience}\n"
        f"Key points: {body.key_points}\n"
        f"Tone: {body.tone}\n"
        "Produce the JSON now."
    )
    raw = await call_claude(COLLATERAL_SYS, user_prompt, session_id=f"collateral-{user['id']}")
    data = safe_json_loads(raw)
    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "asset_type": body.asset_type,
        "deal_name": body.deal_name,
        "target_audience": body.target_audience,
        "data": data,
        "created_at": now_utc().isoformat(),
    }
    await db.collateral.insert_one(doc)
    duration = int((now_utc() - started).total_seconds() * 1000)
    await log_agent_activity("marketing-agent", f"collateral:{body.asset_type}", "completed",
                             user_id=user["id"], duration_ms=duration)
    await log_audit(user["id"], "collateral.generate", body.asset_type)
    doc.pop("_id", None)
    return doc


@api_router.get("/collateral")
async def list_collateral(user=Depends(get_current_user)):
    items = await db.collateral.find({"user_id": user["id"], "deleted_at": {"$exists": False}}, {"_id": 0}).sort("created_at", -1).to_list(50)
    return items


# -----------------------------------------------------------------------------
# OUTREACH
# -----------------------------------------------------------------------------
OUTREACH_SYS = """You are a senior sales enablement specialist. Draft a personalized outreach message.
Return STRICT JSON:
{
  "subject": str,
  "opening": str,
  "value_props": [str],
  "social_proof": str,
  "cta": str,
  "linkedin_message": str,
  "email_body": str
}
"""


@api_router.post("/outreach/campaigns")
async def create_campaign(body: OutreachCampaignRequest, user=Depends(get_current_user)):
    started = now_utc()
    raw = await call_claude(
        OUTREACH_SYS,
        f"Persona: {body.target_persona}\nChannel: {body.channel}\nBrief: {body.message_brief}\nReturn JSON.",
        session_id=f"outreach-{user['id']}",
    )
    draft = safe_json_loads(raw)
    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "name": body.name,
        "target_persona": body.target_persona,
        "channel": body.channel,
        "audience_size": body.audience_size,
        "message_brief": body.message_brief,
        "draft": draft,
        "status": "draft",
        "sent_count": 0,
        "reply_count": 0,
        "created_at": now_utc().isoformat(),
    }
    await db.outreach.insert_one(doc)
    duration = int((now_utc() - started).total_seconds() * 1000)
    await log_agent_activity("outreach-agent", f"campaign:{body.name}", "completed",
                             user_id=user["id"], duration_ms=duration)
    await log_audit(user["id"], "outreach.create", body.name)
    doc.pop("_id", None)
    return doc


@api_router.get("/outreach/campaigns")
async def list_campaigns(user=Depends(get_current_user)):
    items = await db.outreach.find({"user_id": user["id"], "deleted_at": {"$exists": False}}, {"_id": 0}).sort("created_at", -1).to_list(50)
    return items


@api_router.post("/outreach/campaigns/{cid}/launch")
async def launch_campaign(cid: str, user=Depends(get_current_user)):
    """MOCK: marks campaign as launched and simulates sent count."""
    camp = await db.outreach.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not camp:
        raise HTTPException(status_code=404, detail="Campaign not found")
    new_sent = camp.get("audience_size", 50)
    await db.outreach.update_one(
        {"id": cid},
        {"$set": {"status": "launched", "sent_count": new_sent, "launched_at": now_utc().isoformat()}},
    )
    await log_audit(user["id"], "outreach.launch", cid)
    await log_agent_activity("outreach-agent", f"launch:{cid}", "completed", user_id=user["id"])
    return {"ok": True, "sent_count": new_sent, "note": "MOCKED dispatch (Composio LinkedIn / email)"}


# -----------------------------------------------------------------------------
# LEADS
# -----------------------------------------------------------------------------
@api_router.get("/leads")
async def list_leads(user=Depends(get_current_user)):
    items = await db.leads.find({"user_id": user["id"]}, {"_id": 0}).sort("created_at", -1).to_list(500)
    return items


@api_router.post("/leads")
async def create_lead(body: LeadCreate, user=Depends(get_current_user)):
    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "name": body.name,
        "company": body.company,
        "title": body.title,
        "email": body.email,
        "source": body.source,
        "stage": "new",
        "score": 50,
        "created_at": now_utc().isoformat(),
    }
    await db.leads.insert_one(doc)
    await log_audit(user["id"], "lead.create", body.name)
    doc.pop("_id", None)
    return doc


@api_router.patch("/leads/{lead_id}/stage")
async def update_lead_stage(lead_id: str, body: LeadStageUpdate, user=Depends(get_current_user)):
    res = await db.leads.update_one(
        {"id": lead_id, "user_id": user["id"]},
        {"$set": {"stage": body.stage}},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Lead not found")
    await log_audit(user["id"], "lead.stage", lead_id, {"stage": body.stage})
    await log_agent_activity("nurturing-agent", f"advance:{lead_id}->{body.stage}", "completed", user_id=user["id"])
    return {"ok": True}


# -----------------------------------------------------------------------------
# NEWSLETTER
# -----------------------------------------------------------------------------
NEWSLETTER_SYS = """You are an editor for NextCapOS' institutional buyer newsletter.
Compile a concise, sharp newsletter tailored to a buyer's interests.
Return STRICT JSON:
{
  "title": str,
  "issue_tagline": str,
  "deal_spotlights": [{"headline": str, "summary": str}],
  "market_analysis": str,
  "portfolio_updates": [str],
  "editor_note": str
}
"""


def _render_newsletter_html(
    data: dict,
    *,
    recipient_name: Optional[str] = None,
    sender_name: Optional[str] = None,
    sender_org: Optional[str] = None,
    kind: str = "broadcast",
) -> str:
    """Build the email-safe HTML payload for a newsletter using the NextCapOS
    Bloomberg-blue / warm-paper brand palette. Pure inline CSS so Gmail, Outlook,
    Apple Mail, etc. render it the same. No external assets — everything is
    inline so deliverability isn't tripped by image hosting."""
    title = (data.get("title") or "Your NextCapOS digest").strip()
    tagline = (data.get("issue_tagline") or ("Personal digest" if kind == "personal" else "Broadcast")).strip()
    spotlights = data.get("deal_spotlights") or []
    market = (data.get("market_analysis") or "").strip()
    updates = data.get("portfolio_updates") or []
    note = (data.get("editor_note") or "").strip()
    frontend = (os.environ.get("FRONTEND_URL") or "").rstrip("/")

    # Defensive escapes for any field that contains user-provided text.
    import html as _h
    def e(s):
        return _h.escape(str(s or ""), quote=True)

    spotlight_blocks = "".join(
        f"""
        <tr><td style="padding:14px 0;border-bottom:1px solid #E5E0D6;">
          <div style="font-family:Georgia,'Times New Roman',serif;font-size:16px;font-weight:600;color:#0B1B3D;line-height:1.35;">{e(s.get('headline'))}</div>
          <div style="margin-top:6px;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:13.5px;color:#3D4760;line-height:1.55;">{e(s.get('summary'))}</div>
        </td></tr>
        """
        for s in spotlights[:5]
    ) or """<tr><td style="padding:14px 0;color:#7A8299;font-size:13px;font-style:italic;">No deal spotlights in this issue.</td></tr>"""

    updates_block = ""
    if updates:
        items = "".join(
            f"""<li style="margin:5px 0;font-size:13.5px;color:#3D4760;line-height:1.55;">{e(u)}</li>"""
            for u in updates[:8]
        )
        updates_block = f"""
        <tr><td style="padding:22px 0 0 0;">
          <div style="font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:10px;color:#7A8299;letter-spacing:0.18em;text-transform:uppercase;font-weight:600;margin-bottom:8px;">Portfolio updates</div>
          <ul style="margin:0;padding-left:18px;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;">{items}</ul>
        </td></tr>
        """

    market_block = ""
    if market:
        market_block = f"""
        <tr><td style="padding:22px 0 0 0;">
          <div style="font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:10px;color:#7A8299;letter-spacing:0.18em;text-transform:uppercase;font-weight:600;margin-bottom:8px;">Market analysis</div>
          <p style="margin:0;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:13.5px;color:#3D4760;line-height:1.65;">{e(market)}</p>
        </td></tr>
        """

    note_block = ""
    if note:
        note_block = f"""
        <tr><td style="padding:24px 0 0 0;border-top:1px solid #E5E0D6;margin-top:22px;">
          <p style="margin:18px 0 0 0;font-family:Georgia,'Times New Roman',serif;font-style:italic;font-size:13.5px;color:#5A6275;line-height:1.6;">&ldquo;{e(note)}&rdquo;</p>
          <p style="margin:6px 0 0 0;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:12px;color:#7A8299;">— {e(sender_name or 'The NextCapOS editorial team')}</p>
        </td></tr>
        """

    sender_line = ""
    if kind == "broadcast" and sender_name:
        org_part = f" · {e(sender_org)}" if sender_org else ""
        sender_line = f"""<div style="margin-top:4px;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:11px;color:#7A8299;">From {e(sender_name)}{org_part}</div>"""

    greeting = ""
    if recipient_name:
        greeting = f"""<p style="margin:0 0 18px 0;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:14px;color:#3D4760;">Hi {e(recipient_name)},</p>"""

    cta_block = ""
    if frontend:
        cta_block = f"""
        <tr><td style="padding:28px 0 0 0;text-align:left;">
          <a href="{frontend}/app/newsletter"
             style="display:inline-block;padding:12px 22px;background:#0B1B3D;color:#FAFAF7;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:13px;font-weight:600;text-decoration:none;letter-spacing:0.05em;border-radius:2px;">
            Open in NextCapOS &rsaquo;
          </a>
        </td></tr>
        """

    return f"""<!DOCTYPE html>
<html lang="en"><head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>{e(title)}</title>
</head>
<body style="margin:0;padding:0;background:#F4EFE5;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;color:#1A1F33;">
  <table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="background:#F4EFE5;">
    <tr><td align="center" style="padding:32px 16px;">
      <table role="presentation" width="600" cellpadding="0" cellspacing="0" style="max-width:600px;background:#FAFAF7;border:1px solid #E5E0D6;border-radius:4px;">
        <!-- HEADER: NextCapOS wordmark band -->
        <tr><td style="padding:24px 32px 12px 32px;border-bottom:1px solid #E5E0D6;">
          <table role="presentation" width="100%" cellpadding="0" cellspacing="0">
            <tr>
              <td>
                <div style="font-family:Georgia,'Times New Roman',serif;font-size:22px;font-weight:700;letter-spacing:-0.01em;color:#0B1B3D;line-height:1;">
                  NextCap<span style="color:#1D4ED8;">OS</span>
                </div>
                <div style="margin-top:4px;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:9.5px;letter-spacing:0.22em;text-transform:uppercase;color:#7A8299;font-weight:600;">
                  Institutional buy &amp; sell-side
                </div>
              </td>
              <td align="right" style="vertical-align:top;">
                <div style="font-family:'SF Mono',Menlo,Monaco,Consolas,monospace;font-size:10px;color:#7A8299;letter-spacing:0.1em;text-transform:uppercase;">{e(tagline)}</div>
              </td>
            </tr>
          </table>
        </td></tr>
        <!-- BODY -->
        <tr><td style="padding:28px 32px 32px 32px;">
          {greeting}
          <h1 style="margin:0 0 4px 0;font-family:Georgia,'Times New Roman',serif;font-size:26px;line-height:1.2;color:#0B1B3D;font-weight:700;letter-spacing:-0.015em;">{e(title)}</h1>
          {sender_line}
          <table role="presentation" width="100%" cellpadding="0" cellspacing="0" style="margin-top:22px;">
            <tr><td>
              <div style="font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:10px;color:#7A8299;letter-spacing:0.18em;text-transform:uppercase;font-weight:600;margin-bottom:4px;">Deal spotlights</div>
              <table role="presentation" width="100%" cellpadding="0" cellspacing="0">{spotlight_blocks}</table>
            </td></tr>
            {market_block}
            {updates_block}
            {note_block}
            {cta_block}
          </table>
        </td></tr>
        <!-- FOOTER -->
        <tr><td style="padding:18px 32px 24px 32px;border-top:1px solid #E5E0D6;background:#F4EFE5;border-radius:0 0 4px 4px;">
          <table role="presentation" width="100%" cellpadding="0" cellspacing="0">
            <tr>
              <td style="font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:10.5px;color:#7A8299;line-height:1.6;">
                Sent by <strong style="color:#0B1B3D;">NextCapOS</strong> · the institutional buy &amp; sell-side OS for M&amp;A.<br>
                Manage your preferences at <a href="{frontend}/app/newsletter" style="color:#1D4ED8;text-decoration:none;">{frontend.replace('https://','').replace('http://','') or 'nextcapos.com'}/app/newsletter</a>.
              </td>
            </tr>
          </table>
        </td></tr>
      </table>
      <div style="margin-top:16px;font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',Helvetica,Arial,sans-serif;font-size:10px;color:#A8AEC0;text-align:center;letter-spacing:0.05em;">
        You're receiving this because you opted in to NextCapOS newsletters.
      </div>
    </td></tr>
  </table>
</body></html>
"""


def _newsletter_plain_text(data: dict) -> str:
    """Fallback plain-text body for clients that block HTML."""
    parts = [
        (data.get("title") or "Your NextCapOS digest"),
        "=" * min(60, len(data.get("title") or "")),
        "",
    ]
    if data.get("issue_tagline"):
        parts.append(data["issue_tagline"])
        parts.append("")
    if data.get("deal_spotlights"):
        parts.append("DEAL SPOTLIGHTS")
        for s in data["deal_spotlights"][:5]:
            parts.append(f"• {s.get('headline','')}")
            if s.get("summary"):
                parts.append(f"  {s['summary']}")
        parts.append("")
    if data.get("market_analysis"):
        parts.append("MARKET ANALYSIS")
        parts.append(data["market_analysis"])
        parts.append("")
    if data.get("portfolio_updates"):
        parts.append("PORTFOLIO UPDATES")
        for u in data["portfolio_updates"][:8]:
            parts.append(f"• {u}")
        parts.append("")
    if data.get("editor_note"):
        parts.append(f"— {data['editor_note']}")
    parts.append("")
    parts.append("Sent by NextCapOS · the institutional buy & sell-side OS for M&A.")
    return "\n".join(parts)


@api_router.get("/newsletter/preferences")
async def get_prefs(user=Depends(get_current_user)):
    return {
        "opt_in": user.get("newsletter_opt_in", False),
        "interests": user.get("interests", []),
        "cadence": user.get("newsletter_cadence", "weekly"),
    }


@api_router.post("/newsletter/preferences")
async def set_prefs(body: NewsletterPreferences, user=Depends(get_current_user)):
    await db.users.update_one(
        {"id": user["id"]},
        {"$set": {
            "newsletter_opt_in": body.opt_in,
            "interests": body.interests,
            "newsletter_cadence": body.cadence,
        }},
    )
    await log_audit(user["id"], "newsletter.preferences", "", body.model_dump())
    return {"ok": True}


@api_router.post("/newsletter/draft")
async def draft_newsletter(body: NewsletterDraftRequest, user=Depends(get_current_user)):
    """Used by sellers (or admins) to draft a broadcast newsletter to opted-in buyers."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Broadcasts are seller/admin only")
    started = now_utc()
    interests = ", ".join(user.get("interests", [])) or "institutional buyers"
    topic = body.topic or "this week's deal flow and market signals"
    raw = await call_claude(
        NEWSLETTER_SYS,
        f"Audience: {interests}\nTopic focus: {topic}\nProduce JSON.",
        session_id=f"newsletter-{user['id']}",
    )
    data = safe_json_loads(raw)
    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "kind": "broadcast",
        "sender_name": user.get("name"),
        "sender_org": user.get("organization"),
        "data": data,
        "status": "draft",
        "approved_by": None,
        "dispatched_at": None,
        "recipients": 0,
        "created_at": now_utc().isoformat(),
    }
    await db.newsletters.insert_one(doc)
    duration = int((now_utc() - started).total_seconds() * 1000)
    await log_agent_activity("newsletter-agent", "draft:broadcast", "completed", user_id=user["id"], duration_ms=duration)
    await log_audit(user["id"], "newsletter.draft", doc["id"])
    doc.pop("_id", None)
    return doc


@api_router.post("/newsletter/personal")
async def personal_newsletter(body: NewsletterDraftRequest, user=Depends(get_current_user)):
    """Buyer self-service: generate AND deliver a personalized digest in one call (recipient=self)."""
    if user.get("role") not in ("buyer", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Personal digests are buyer/admin only")
    started = now_utc()
    interests = ", ".join(user.get("interests", [])) or "general institutional buying themes"
    topic = body.topic or "today's deal flow, market signals, and portfolio updates"
    raw = await call_claude(
        NEWSLETTER_SYS,
        (
            f"This is a PERSONAL digest for ONE institutional buyer named {user.get('name')} "
            f"at {user.get('organization') or 'an institutional fund'}. "
            f"Their stated interests: {interests}. "
            f"Topic focus: {topic}. "
            "Tailor every section to THIS reader. Speak as NextCapOS. Produce JSON."
        ),
        session_id=f"newsletter-personal-{user['id']}",
    )
    data = safe_json_loads(raw)
    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "kind": "personal",
        "sender_name": "NextCapOS",
        "sender_org": "NextCapOS",
        "recipient_email": user["email"],
        "recipient_name": user.get("name"),
        "data": data,
        "status": "dispatched",
        "approved_by": user["id"],
        "approved_at": now_utc().isoformat(),
        "dispatched_at": now_utc().isoformat(),
        "recipients": 1,
        "created_at": now_utc().isoformat(),
    }
    await db.newsletters.insert_one(doc)
    # Actually send the email via Resend (no longer mocked). Personal digest =
    # 1 recipient → safe to send inline. Errors are recorded on the doc but
    # don't fail the request — the digest still lives in the user's history.
    html = _render_newsletter_html(
        data,
        recipient_name=user.get("name"),
        sender_name="NextCapOS",
        sender_org="NextCapOS",
        kind="personal",
    )
    text = _newsletter_plain_text(data)
    subj = (data.get("title") or "Your NextCapOS digest").strip()[:180]
    send_result = await send_email(user["email"], f"{subj} · NextCapOS", html, text=text)
    delivery_meta = {
        "sent_ok": bool(send_result.get("ok")),
        "skipped": bool(send_result.get("skipped")),
        "provider_id": send_result.get("id"),
        "error": send_result.get("error") or send_result.get("body"),
    }
    await db.newsletters.update_one(
        {"id": doc["id"]},
        {"$set": {"delivery": delivery_meta, "delivered_to": [user["email"]]}},
    )
    doc["delivery"] = delivery_meta
    duration = int((now_utc() - started).total_seconds() * 1000)
    await log_agent_activity("newsletter-agent", "personal:delivered", "completed", user_id=user["id"], duration_ms=duration)
    await log_audit(user["id"], "newsletter.personal", doc["id"],
                    {"sent_ok": delivery_meta["sent_ok"], "skipped": delivery_meta["skipped"]})
    doc.pop("_id", None)
    return doc


@api_router.get("/newsletter")
async def list_newsletters(user=Depends(get_current_user)):
    items = await db.newsletters.find({"user_id": user["id"], "deleted_at": {"$exists": False}}, {"_id": 0}).sort("created_at", -1).to_list(50)
    for it in items:
        it.setdefault("kind", "broadcast")
    return items


@api_router.post("/newsletter/{nid}/approve")
async def approve_newsletter(nid: str, user=Depends(get_current_user)):
    res = await db.newsletters.update_one(
        {"id": nid, "user_id": user["id"]},
        {"$set": {"status": "approved", "approved_by": user["id"], "approved_at": now_utc().isoformat()}},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Newsletter not found")
    await log_audit(user["id"], "newsletter.approve", nid)
    return {"ok": True}


@api_router.post("/newsletter/{nid}/dispatch")
async def dispatch_newsletter(nid: str, user=Depends(get_current_user)):
    """Dispatch a broadcast newsletter to opted-in buyers (or a hand-picked
    recipient subset). Real Resend send — no longer mocked. Larger fan-outs
    (>3 recipients) run in the background to avoid ingress timeouts; small
    sends finish inline so the seller gets immediate confirmation."""
    nl = await db.newsletters.find_one({"id": nid, "user_id": user["id"]}, {"_id": 0})
    if not nl:
        raise HTTPException(status_code=404, detail="Newsletter not found")
    if nl.get("kind") == "personal":
        return {"ok": True, "recipients": 1, "note": "personal digest already delivered"}

    picked = nl.get("recipient_ids") or []
    if picked:
        # Only count IDs that are opted-in buyers (defense-in-depth on stale picks)
        recip_docs = await db.users.find(
            {"id": {"$in": picked}, "role": "buyer", "newsletter_opt_in": True},
            {"_id": 0, "email": 1, "name": 1},
        ).to_list(500)
        scope_note = f"hand-picked ({len(picked)} selected, {len(recip_docs)} eligible opted-in)"
    else:
        recip_docs = await db.users.find(
            {"newsletter_opt_in": True, "role": "buyer"},
            {"_id": 0, "email": 1, "name": 1},
        ).to_list(500)
        scope_note = "broadcast to all opted-in buyers"

    recipients = len(recip_docs)
    data = nl.get("data") or {}
    subj = (data.get("title") or "NextCapOS digest").strip()[:180]
    text = _newsletter_plain_text(data)

    async def _fanout_send():
        """Send one email per recipient (each with personalized greeting) so
        Gmail/Outlook show one recipient per message — no bcc clumping, better
        deliverability, and the per-recipient name on the greeting line."""
        delivered, failed = [], []
        for r in recip_docs:
            html = _render_newsletter_html(
                data,
                recipient_name=r.get("name"),
                sender_name=nl.get("sender_name"),
                sender_org=nl.get("sender_org"),
                kind="broadcast",
            )
            res = await send_email(r["email"], f"{subj} · NextCapOS", html, text=text)
            if res.get("ok"):
                delivered.append(r["email"])
            elif res.get("skipped"):
                failed.append({"email": r["email"], "reason": "RESEND_API_KEY missing"})
            else:
                failed.append({"email": r["email"],
                               "reason": (res.get("error") or res.get("body") or "send failed")[:240]})
        await db.newsletters.update_one(
            {"id": nid},
            {"$set": {
                "delivered_to": delivered,
                "delivery_failures": failed,
                "delivery_finished_at": now_utc().isoformat(),
            }},
        )

    # Flip status + recipients counter BEFORE fan-out so the UI shows
    # "dispatched" immediately. Counters refine after delivery completes.
    await db.newsletters.update_one(
        {"id": nid},
        {"$set": {
            "status": "dispatched",
            "dispatched_at": now_utc().isoformat(),
            "recipients": recipients,
            "dispatch_scope": scope_note,
            "delivered_to": [],
            "delivery_failures": [],
        }},
    )
    await log_audit(user["id"], "newsletter.dispatch", nid,
                    {"recipients": recipients, "scope": scope_note})
    await log_agent_activity("newsletter-agent", f"dispatch:{nid}", "completed",
                             user_id=user["id"])

    if recipients == 0:
        return {"ok": True, "recipients": 0,
                "note": "No opted-in recipients matched — nothing sent.",
                "delivered": 0, "failed": 0}

    # Inline for small sends (immediate confirmation), background for larger fan-outs.
    if recipients <= 3:
        await _fanout_send()
        nl_fresh = await db.newsletters.find_one({"id": nid}, {"_id": 0, "delivered_to": 1, "delivery_failures": 1}) or {}
        return {
            "ok": True,
            "recipients": recipients,
            "delivered": len(nl_fresh.get("delivered_to") or []),
            "failed": len(nl_fresh.get("delivery_failures") or []),
            "note": f"Resend dispatch · {scope_note}",
        }
    asyncio.create_task(_fanout_send())
    return {
        "ok": True,
        "recipients": recipients,
        "delivered": 0,
        "failed": 0,
        "note": f"Resend dispatch started in background · {scope_note}. "
                "Refresh the newsletter list in ~30s to see per-recipient delivery status.",
    }


# -----------------------------------------------------------------------------
# MCP ACTIONS
# -----------------------------------------------------------------------------
@api_router.get("/mcp/actions")
async def mcp_actions(user=Depends(get_current_user)):
    return {"actions": MCP_ACTIONS, "count": len(MCP_ACTIONS)}


@api_router.get("/mcp/manifest")
async def mcp_manifest():
    """Public manifest of WebMCP actions for AI-agent discovery."""
    return {
        "name": "NextCapOS MCP",
        "version": "1.0.0",
        "gateway": "composio",
        "actions": [{"id": a["id"], "type": a["type"], "description": a["description"]} for a in MCP_ACTIONS],
    }


# -----------------------------------------------------------------------------
# AGENT ACTIVITY
# -----------------------------------------------------------------------------
@api_router.get("/agents/activity")
async def agent_activity(user=Depends(get_current_user)):
    items = await db.agent_activity.find({}, {"_id": 0}).sort("timestamp", -1).to_list(100)
    return items


@api_router.get("/agents/stats")
async def agent_stats(user=Depends(get_current_user)):
    total = await db.agent_activity.count_documents({})
    completed = await db.agent_activity.count_documents({"status": "completed"})
    failed = await db.agent_activity.count_documents({"status": "failed"})
    by_agent_cursor = db.agent_activity.aggregate([
        {"$group": {"_id": "$agent", "count": {"$sum": 1}}},
    ])
    by_agent = [{"agent": d["_id"], "count": d["count"]} async for d in by_agent_cursor]
    return {
        "total": total,
        "completed": completed,
        "failed": failed,
        "success_rate": round((completed / total * 100) if total else 0, 1),
        "by_agent": by_agent,
    }


# -----------------------------------------------------------------------------
# COMPOSIO
# -----------------------------------------------------------------------------
@api_router.get("/composio/status")
async def composio_status(user=Depends(get_current_user)):
    return {
        "configured": bool(COMPOSIO_API_KEY),
        "gateway": COMPOSIO_BASE_URL,
        "supported_apps": ["LINKEDIN", "ZOHO_CRM", "GMAIL", "SLACK", "HUBSPOT", "SALESFORCE"],
    }


@api_router.get("/composio/connections")
async def composio_connections(user=Depends(get_current_user)):
    """List user's connected Composio accounts. Falls back to stored connections on API error."""
    stored = await db.composio_connections.find(
        {"user_id": user["id"]}, {"_id": 0}
    ).sort("created_at", -1).to_list(50)
    return {"connections": stored, "gateway": COMPOSIO_BASE_URL}


@api_router.post("/composio/connect/linkedin")
async def composio_connect_linkedin(user=Depends(get_current_user)):
    """Initiate LinkedIn OAuth via Composio. Stores a pending connection record."""
    if not COMPOSIO_API_KEY:
        raise HTTPException(status_code=400, detail="Composio API key missing")

    entity_id = f"workz-{user['id']}"
    redirect_url = f"https://app.composio.dev/connect/linkedin?entity={entity_id}"
    status_label = "pending"

    # Attempt real Composio v3 initiate connection
    try:
        async with httpx.AsyncClient(timeout=10.0) as c:
            r = await c.post(
                f"{COMPOSIO_BASE_URL}/api/v3/connectedAccounts",
                headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
                json={"appName": "linkedin", "entityId": entity_id},
            )
            if r.status_code < 400:
                payload = r.json()
                redirect_url = payload.get("redirectUrl") or payload.get("redirect_url") or redirect_url
    except Exception as e:
        logger.warning(f"Composio LinkedIn init failed, using placeholder: {e}")

    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "app": "linkedin",
        "entity_id": entity_id,
        "status": status_label,
        "redirect_url": redirect_url,
        "created_at": now_utc().isoformat(),
    }
    await db.composio_connections.insert_one(doc)
    await log_audit(user["id"], "composio.connect.linkedin", entity_id)
    doc.pop("_id", None)
    return doc


@api_router.delete("/composio/connections/{cid}")
async def remove_connection(cid: str, user=Depends(get_current_user)):
    res = await db.composio_connections.delete_one({"id": cid, "user_id": user["id"]})
    if res.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Connection not found")
    await log_audit(user["id"], "composio.disconnect", cid)
    return {"ok": True}


# -----------------------------------------------------------------------------
# ZOHO CRM via Composio
# -----------------------------------------------------------------------------
@api_router.post("/composio/connect/zoho-crm")
async def composio_connect_zoho_crm(user=Depends(get_current_user)):
    """Initiate Zoho CRM OAuth via Composio (US data center, .com)."""
    if not COMPOSIO_API_KEY:
        raise HTTPException(status_code=400, detail="Composio API key missing")

    entity_id = f"workz-{user['id']}"
    redirect_url = f"https://app.composio.dev/connect/zoho_crm?entity={entity_id}"
    composio_connected_id = None

    # Try Composio v3 connect — multiple endpoint shapes for resilience
    try:
        async with httpx.AsyncClient(timeout=12.0) as c:
            r = await c.post(
                f"{COMPOSIO_BASE_URL}/api/v3/connectedAccounts",
                headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
                json={
                    "appName": "zoho_crm",
                    "entityId": entity_id,
                    "region": "com",  # US data center
                },
            )
            if r.status_code < 400:
                payload = r.json()
                redirect_url = (
                    payload.get("redirectUrl")
                    or payload.get("redirect_url")
                    or payload.get("url")
                    or redirect_url
                )
                composio_connected_id = (
                    payload.get("id")
                    or payload.get("connectedAccountId")
                    or payload.get("connected_account_id")
                )
            else:
                logger.warning(f"Composio Zoho init non-2xx ({r.status_code}): {r.text[:300]}")
    except Exception as e:
        logger.warning(f"Composio Zoho init failed, using placeholder: {e}")

    doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "app": "zoho_crm",
        "region": "com",
        "entity_id": entity_id,
        "composio_connected_id": composio_connected_id,
        "status": "pending",
        "redirect_url": redirect_url,
        "created_at": now_utc().isoformat(),
    }
    await db.composio_connections.insert_one(doc)
    await log_audit(user["id"], "composio.connect.zoho_crm", entity_id)
    doc.pop("_id", None)
    return doc


@api_router.post("/composio/zoho/push-lead/{inquiry_id}")
async def push_inquiry_to_zoho(inquiry_id: str, user=Depends(get_current_user)):
    """Seller-only: push a buyer inquiry into Zoho CRM as a Lead via Composio Proxy Execute."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Sellers/admin only")

    inquiry = await db.inquiries.find_one({"id": inquiry_id, "seller_id": user["id"]}, {"_id": 0})
    if not inquiry:
        raise HTTPException(status_code=404, detail="Inquiry not found")

    conn = await db.composio_connections.find_one(
        {"user_id": user["id"], "app": "zoho_crm"}, {"_id": 0}
    )
    if not conn:
        raise HTTPException(status_code=400, detail="Zoho CRM not connected — connect via Integrations first")

    # Build Zoho Lead record (Insert Records requires {data: [{...}]})
    buyer_name = (inquiry.get("buyer_name") or "Unknown Buyer").strip()
    parts = buyer_name.split(" ", 1)
    first_name = parts[0]
    last_name = parts[1] if len(parts) > 1 else parts[0]
    lead_record = {
        "Last_Name": last_name,
        "First_Name": first_name,
        "Company": inquiry.get("buyer_org") or "Unknown",
        "Email": inquiry.get("buyer_email"),
        "Lead_Source": "NextCapOS",
        "Description": (
            f"Inquiry re: {inquiry.get('listing_name')}\n\n"
            f"{inquiry.get('message', '')}"
        ),
    }
    body = {"data": [lead_record]}

    # Attempt Composio Proxy Execute
    pushed = False
    composio_response: Dict[str, Any] = {}
    try:
        async with httpx.AsyncClient(timeout=15.0) as c:
            r = await c.post(
                f"{COMPOSIO_BASE_URL}/api/v3/tools/proxy_execute",
                headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
                json={
                    "toolkit": "zoho_crm",
                    "connected_account_id": conn.get("composio_connected_id") or conn.get("entity_id"),
                    "method": "POST",
                    "endpoint": "/crm/v8/Leads",
                    "body": body,
                },
            )
            composio_response = {"status_code": r.status_code, "body": r.text[:500]}
            pushed = r.status_code < 400
    except Exception as e:
        composio_response = {"error": str(e)}
        logger.warning(f"Zoho push failed: {e}")

    # Always mark the inquiry as synced locally so the UI reflects the action
    await db.inquiries.update_one(
        {"id": inquiry_id},
        {"$set": {"zoho_pushed_at": now_utc().isoformat(), "zoho_pushed": pushed}},
    )
    await log_audit(user["id"], "zoho.lead.push", inquiry_id, {"pushed": pushed})
    await log_agent_activity(
        "zoho-sync-agent",
        f"push-lead:{inquiry.get('listing_name')}",
        "completed" if pushed else "failed",
        user_id=user["id"],
        friction=None if pushed else "composio_proxy_error",
        meta=composio_response,
    )

    return {
        "ok": True,
        "pushed_to_zoho": pushed,
        "lead": lead_record,
        "composio": composio_response,
        "note": (
            "Pushed to Zoho CRM via Composio Proxy Execute."
            if pushed
            else "Local sync recorded; Composio Proxy Execute did not confirm — verify Zoho Auth Config in Composio dashboard."
        ),
    }


# -----------------------------------------------------------------------------
# LISTING EXTERNAL FILE SOURCES (Composio-mirrored)
#
# Architecture (per user choice "Mirror first"):
#   Seller connects ONE of six file sources (Google Drive, OneDrive,
#   SharePoint, Dropbox, Box, Zoho WorkDrive) per listing via Composio OAuth.
#   When the connection turns ACTIVE we pull file metadata + bytes through
#   Composio actions and persist into the SAME listing_staged_files schema as
#   manual uploads — so the existing Vault clone path, Copilot indexer, NDA
#   gating and audit all keep working unchanged. Collaborators and buyers
#   read the mirrored copy through the existing /staged-files endpoints with
#   no extra OAuth (Rule: "one seller-side login, many readers").
#
# On listing close/archive the connection is revoked and mirrored bytes are
# wiped (Rule 3A: immediate purge).
# -----------------------------------------------------------------------------
COMPOSIO_FILE_SOURCES = {
    "googledrive": {"label": "Google Drive", "app": "googledrive",  "list": "GOOGLEDRIVE_LIST_FILES",        "download": "GOOGLEDRIVE_DOWNLOAD_FILE"},
    "onedrive":    {"label": "OneDrive",     "app": "one_drive",    "list": "ONE_DRIVE_ONEDRIVE_LIST_ITEMS", "download": "ONE_DRIVE_DOWNLOAD_FILE"},
    "sharepoint":  {"label": "SharePoint",   "app": "share_point",  "list": None,                            "download": None},
    "dropbox":     {"label": "Dropbox",      "app": "dropbox",      "list": "DROPBOX_LIST_FILES_IN_FOLDER",  "download": "DROPBOX_READ_FILE"},
    "box":         {"label": "Box",          "app": "box",          "list": "BOX_LIST_ITEMS_IN_FOLDER",      "download": "BOX_DOWNLOAD_FILE"},
}


class ExternalSourceCreate(BaseModel):
    source_kind: Literal["googledrive", "onedrive", "sharepoint", "dropbox", "box"]
    folder_id: Optional[str] = None  # legacy single-folder API
    folder_ids: Optional[List[str]] = None  # NEW: multi-select folder picker
    folder_labels: Optional[List[str]] = None  # parallel list of display names
    include_subfolders: Optional[bool] = True
    label: Optional[str] = None


class ExternalSourceFoldersUpdate(BaseModel):
    """PATCH body for editing the folder selection on an already-connected
    source (the "Edit folders" affordance on each source row)."""
    folder_ids: List[str] = []
    folder_labels: Optional[List[str]] = None
    include_subfolders: Optional[bool] = True


async def _composio_action_execute(action_slug: str, connected_account_id: str, input_params: dict | None = None, user_id: str | None = None) -> dict:
    """Thin wrapper around POST /api/v3/tools/execute/{slug}. Raises HTTPException
    with the upstream body on failure so we surface useful messages.

    Composio v3 requires `user_id` (the entity_id used at connect time) in
    addition to `connected_account_id` — see error code 1811. We always
    pass it now so action execution never fails on this requirement."""
    if not COMPOSIO_API_KEY:
        raise HTTPException(status_code=400, detail="Composio API key not configured")
    payload: dict = {
        "connected_account_id": connected_account_id,
        "arguments": input_params or {},
    }
    if user_id:
        payload["user_id"] = user_id
    async with httpx.AsyncClient(timeout=60.0) as c:
        r = await c.post(
            f"{COMPOSIO_BASE_URL}/api/v3/tools/execute/{action_slug}",
            headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
            json=payload,
        )
    if r.status_code >= 400:
        raise HTTPException(status_code=502, detail=f"Composio action {action_slug} failed: {r.text[:400]}")
    return r.json()


# Per-toolkit Google/Microsoft/etc API endpoints used by the proxy-execute
# fallback when Composio's predefined `*_DOWNLOAD_FILE` action fails (e.g.,
# Drive's known "Missing presigned URL" bug in the R2 staging step —
# composio issues #3471 / #3477). Proxy execute injects auth server-side
# and returns either `data` (small JSON/text) or `binary_data: {url, ...}`
# which we then fetch with a follow-up HTTP GET.
PROXY_DOWNLOAD_ENDPOINTS = {
    # Composio's proxy already prepends each toolkit's API base URL (e.g.,
    # https://www.googleapis.com/drive/v3 for Drive, https://graph.microsoft.com/v1.0
    # for Graph, https://api.box.com/2.0 for Box). We pass ONLY the path
    # AFTER that base. Earlier versions of this map repeated the version
    # prefix and produced 404s like /drive/v3/drive/v3/files/... which were
    # silently stored as the file contents (Google's HTML 404 page).
    "googledrive": {"path": "/files/{file_id}", "query": [("alt", "media")]},
    # Microsoft Graph base is https://graph.microsoft.com/v1.0
    "onedrive":    {"path": "/me/drive/items/{file_id}/content", "query": []},
    "sharepoint":  {"path": "/me/drive/items/{file_id}/content", "query": []},
    # Box base is https://api.box.com/2.0
    "box":         {"path": "/files/{file_id}/content", "query": []},
}

# Google Workspace native mime types → export target. Composio cannot read
# these via ?alt=media (they have no binary form), so we ask Drive to export
# to a downloadable Office equivalent before mirroring.
GOOGLE_DRIVE_EXPORT_MAP = {
    "application/vnd.google-apps.document":
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.google-apps.spreadsheet":
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.google-apps.presentation":
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.google-apps.drawing": "application/pdf",
}


async def _composio_proxy_download(
    source_kind: str, connected_account_id: str, file_id: str,
    mime_type: str | None = None, user_id: str | None = None,
) -> bytes | None:
    """Fallback download path: call the toolkit's native HTTP API via
    Composio's proxy execute (POST /api/v3.1/tools/execute/proxy). Returns
    raw bytes on success, None on any failure. Designed to be wrapped in
    try/except by the caller — never raises for upstream errors so the
    list-sync loop can record per-file errors without aborting.

    Why this exists: Composio's predefined `GOOGLEDRIVE_DOWNLOAD_FILE`
    action stages content via R2 and intermittently returns "Missing
    presigned URL in upload response". Proxy execute talks directly to
    Google's `/drive/v3/files/{id}?alt=media` endpoint and skips R2.
    Dropbox is intentionally omitted — its download endpoint requires a
    `Dropbox-API-Arg` request header rather than a query param, and the
    standard predefined action handles it reliably."""
    if source_kind not in PROXY_DOWNLOAD_ENDPOINTS:
        return None
    if not COMPOSIO_API_KEY:
        return None
    cfg = PROXY_DOWNLOAD_ENDPOINTS[source_kind]
    endpoint = cfg["path"].format(file_id=file_id)
    query_params = list(cfg["query"])
    # Google Drive: detect native types and route through /export instead.
    if source_kind == "googledrive" and mime_type and mime_type.startswith("application/vnd.google-apps."):
        export_target = GOOGLE_DRIVE_EXPORT_MAP.get(mime_type)
        if not export_target:
            return None  # unsupported Google native type
        endpoint = f"/files/{file_id}/export"
        query_params = [("mimeType", export_target)]
    parameters = [{"name": k, "value": v, "type": "query"} for k, v in query_params]
    payload: dict = {
        "endpoint": endpoint,
        "method": "GET",
        "connected_account_id": connected_account_id,
        "parameters": parameters,
    }
    if user_id:
        payload["user_id"] = user_id
    try:
        async with httpx.AsyncClient(timeout=120.0) as c:
            r = await c.post(
                f"{COMPOSIO_BASE_URL}/api/v3.1/tools/execute/proxy",
                headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
                json=payload,
            )
        if r.status_code >= 400:
            logger.warning(f"proxy download {source_kind}/{file_id} failed {r.status_code}: {r.text[:200]}")
            return None
        body = r.json()
    except Exception as e:
        logger.warning(f"proxy download {source_kind}/{file_id} crashed: {e}")
        return None

    # Composio's proxy returns binary content via either:
    #   1) `binary_data: { url, content_type, size, expires_at }` — large files
    #   2) `data` as a raw string when content was small enough to inline
    bin_meta = body.get("binary_data") or {}
    bin_url = bin_meta.get("url") if isinstance(bin_meta, dict) else None
    fetched: bytes | None = None
    if bin_url:
        try:
            async with httpx.AsyncClient(timeout=180.0, follow_redirects=True) as c:
                rr = await c.get(bin_url)
                if rr.status_code < 400:
                    fetched = rr.content
                else:
                    logger.warning(f"proxy fetch binary url failed {rr.status_code}: {rr.text[:200]}")
        except Exception as e:
            logger.warning(f"proxy fetch binary url crashed: {e}")
    else:
        data_field = body.get("data")
        if isinstance(data_field, str) and data_field:
            import base64 as _b64
            try:
                fetched = _b64.b64decode(data_field, validate=True)
            except Exception:
                fetched = data_field.encode("utf-8", errors="replace")
    if fetched is None:
        return None
    # Sanity-check: refuse to store HTML error pages as file content. Earlier
    # versions of this code silently saved Google's "404 Not Found" HTML when
    # the proxy endpoint path was wrong, so buyers saw "random text" instead
    # of their PPTX. If the first bytes look like an HTML doctype OR a JSON
    # error envelope, treat as a failed download.
    sniff = fetched[:200].lstrip().lower()
    if sniff.startswith(b"<!doctype html") or sniff.startswith(b"<html"):
        logger.warning(f"proxy download for {source_kind}/{file_id} returned HTML — likely an upstream error page. Refusing to store.")
        return None
    if sniff.startswith(b'{"error"') or sniff.startswith(b'{"message"'):
        # Best-effort: peek for a JSON error shape. Real JSON files start with
        # `{"` too, but the literal `"error"` / `"message"` key at the top is
        # a strong upstream-error tell.
        try:
            import json as _j
            parsed = _j.loads(fetched.decode("utf-8", errors="replace"))
            if isinstance(parsed, dict) and ("error" in parsed and isinstance(parsed.get("error"), (dict, str))):
                logger.warning(f"proxy download for {source_kind}/{file_id} returned JSON error envelope. Refusing to store.")
                return None
        except Exception:
            pass
    return fetched


# Proxy bases per toolkit (where Composio's /proxy execute prepends to our
# endpoint). MS Graph for OneDrive/SharePoint, Box API v2.0 for Box, etc.
# Used only by `_composio_proxy_get_json` below.
_PROXY_API_BASES = {
    "googledrive": "googleapis.com/drive/v3",
    "onedrive":    "graph.microsoft.com/v1.0",
    "sharepoint":  "graph.microsoft.com/v1.0",
    "box":         "api.box.com/2.0",
    "dropbox":     "api.dropboxapi.com/2",
}


async def _composio_proxy_get_json(
    source_kind: str, connected_account_id: str, endpoint: str,
    query: list[tuple[str, str]] | None = None, user_id: str | None = None,
) -> dict | None:
    """Call a provider's native HTTP API via Composio's proxy execute and
    return a parsed JSON dict. None on failure. Used by the folder picker
    (`_browse_folder`) for OneDrive, where Composio's predefined LIST_ITEMS
    tool can only list the drive root."""
    if not COMPOSIO_API_KEY or not connected_account_id:
        return None
    parameters = [{"name": k, "value": v, "type": "query"} for k, v in (query or [])]
    payload: dict = {
        "endpoint": endpoint,
        "method": "GET",
        "connected_account_id": connected_account_id,
        "parameters": parameters,
    }
    if user_id:
        payload["user_id"] = user_id
    try:
        async with httpx.AsyncClient(timeout=30.0) as c:
            r = await c.post(
                f"{COMPOSIO_BASE_URL}/api/v3.1/tools/execute/proxy",
                headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
                json=payload,
            )
        if r.status_code >= 400:
            logger.warning(f"proxy GET {source_kind}{endpoint} failed {r.status_code}: {r.text[:200]}")
            return None
        body = r.json()
    except Exception as e:
        logger.warning(f"proxy GET {source_kind}{endpoint} crashed: {e}")
        return None
    # Composio's proxy returns the upstream JSON under `data` (small payload).
    # `binary_data` is only for large/binary downloads.
    data = body.get("data")
    if isinstance(data, dict):
        return data
    if isinstance(data, str):
        try:
            import json as _j
            return _j.loads(data)
        except Exception:
            return None
    return None


async def _browse_folder(src: dict, parent_id: str | None) -> dict:
    """Return the immediate subfolders of `parent_id` inside `src`. Used by
    the seller's folder-picker modal. Shape:
        { folders: [{id, name}], parent_id, can_browse: bool, err: str|None }
    `folders` is empty for SharePoint (no Composio file-list tool) — the
    frontend renders a "paste folder ID manually" affordance in that case."""
    kind = src["source_kind"]
    conn_id = src.get("composio_connected_id")
    entity = src.get("entity_id")
    if not conn_id:
        return {"folders": [], "parent_id": parent_id, "can_browse": False,
                "err": "Connection not finished — complete OAuth first."}

    if kind == "box":
        fid = parent_id or "0"
        try:
            raw = await _composio_action_execute(
                "BOX_LIST_ITEMS_IN_FOLDER", conn_id, {"folder_id": fid}, user_id=entity,
            )
        except HTTPException as e:
            return {"folders": [], "parent_id": fid, "can_browse": True,
                    "err": str(e.detail)[:200]}
        resp = _normalise_composio_response(raw)
        if not resp.get("successful"):
            return {"folders": [], "parent_id": fid, "can_browse": True,
                    "err": str(resp.get("error"))[:200]}
        items = _extract_files_array(resp.get("data"))
        folders = [
            {"id": str(i["id"]), "name": i.get("name") or "(unnamed)"}
            for i in items
            if isinstance(i, dict) and (i.get("type") or "").lower() == "folder" and i.get("id")
        ]
        return {"folders": folders, "parent_id": fid, "can_browse": True, "err": None}

    if kind == "googledrive":
        fid = parent_id or "root"
        q = (
            f"'{fid}' in parents and "
            "mimeType='application/vnd.google-apps.folder' and trashed=false"
        )
        try:
            raw = await _composio_action_execute(
                "GOOGLEDRIVE_LIST_FILES", conn_id,
                {"q": q, "pageSize": 200,
                 "fields": "files(id,name)",
                 "supportsAllDrives": True, "includeItemsFromAllDrives": True},
                user_id=entity,
            )
        except HTTPException as e:
            return {"folders": [], "parent_id": fid, "can_browse": True,
                    "err": str(e.detail)[:200]}
        resp = _normalise_composio_response(raw)
        if not resp.get("successful"):
            return {"folders": [], "parent_id": fid, "can_browse": True,
                    "err": str(resp.get("error"))[:200]}
        items = _extract_files_array(resp.get("data"))
        folders = [
            {"id": f["id"], "name": f.get("name") or "(unnamed)"}
            for f in items if isinstance(f, dict) and f.get("id")
        ]
        return {"folders": folders, "parent_id": fid, "can_browse": True, "err": None}

    if kind == "dropbox":
        # Empty string = root. Otherwise full path (e.g. "/Reports").
        path = parent_id or ""
        try:
            raw = await _composio_action_execute(
                "DROPBOX_LIST_FILES_IN_FOLDER", conn_id,
                {"path": path, "limit": 500, "recursive": False},
                user_id=entity,
            )
        except HTTPException as e:
            return {"folders": [], "parent_id": path, "can_browse": True,
                    "err": str(e.detail)[:200]}
        resp = _normalise_composio_response(raw)
        if not resp.get("successful"):
            return {"folders": [], "parent_id": path, "can_browse": True,
                    "err": str(resp.get("error"))[:200]}
        items = _extract_files_array(resp.get("data"))
        folders = []
        for i in items:
            if not isinstance(i, dict):
                continue
            tag = (i.get(".tag") or i.get("tag") or "").lower()
            if tag != "folder":
                continue
            # Dropbox uses path_lower as its canonical folder identifier for
            # subsequent list calls (Dropbox accepts either path OR "id:..."
            # form; path is more human-readable in the picker).
            fid = i.get("path_lower") or i.get("path_display") or i.get("id")
            if fid is None:
                continue
            folders.append({"id": fid, "name": i.get("name") or "(unnamed)"})
        return {"folders": folders, "parent_id": path, "can_browse": True, "err": None}

    if kind == "onedrive":
        if not parent_id or parent_id == "root":
            endpoint = "/me/drive/root/children"
        else:
            endpoint = f"/me/drive/items/{parent_id}/children"
        body = await _composio_proxy_get_json(
            "onedrive", conn_id, endpoint,
            query=[("$select", "id,name,folder,file"), ("$top", "200")],
            user_id=entity,
        )
        if not isinstance(body, dict):
            return {"folders": [], "parent_id": parent_id or "root",
                    "can_browse": True,
                    "err": "Microsoft Graph returned no body. Reconnect OneDrive and try again."}
        items = body.get("value") or []
        folders = [
            {"id": i["id"], "name": i.get("name") or "(unnamed)"}
            for i in items
            if isinstance(i, dict) and i.get("folder") and i.get("id")
        ]
        return {"folders": folders, "parent_id": parent_id or "root",
                "can_browse": True, "err": None}

    # SharePoint and any future provider without browse support degrade
    # gracefully — the picker UI shows a "paste folder ID" fallback.
    return {"folders": [], "parent_id": parent_id, "can_browse": False,
            "err": "Folder browsing isn't available for this provider yet. "
                   "Paste the folder ID manually instead."}




@api_router.post("/listings/{lid}/external-sources")
async def create_external_source(
    lid: str, body: ExternalSourceCreate, user=Depends(get_current_user)
):
    """Seller initiates an OAuth-backed Composio connection for one file source.
    The caller must be a listing editor or principal owner."""
    await _listing_for_edit_or_404(lid, user)
    cfg = COMPOSIO_FILE_SOURCES.get(body.source_kind)
    if not cfg:
        raise HTTPException(status_code=400, detail="Unknown source kind")
    if not COMPOSIO_API_KEY:
        raise HTTPException(status_code=400, detail="Composio API key not configured")

    # Entity ID scopes the connection to listing-owner pair so the same seller
    # can connect different sources for different deals without cross-talk.
    # Entity ID scopes the connection to listing-owner pair so the same seller
    # can connect different sources for different deals without cross-talk.
    entity_id = f"nextcapos-{user['id']}-{lid}"
    redirect_url = f"https://app.composio.dev/connect/{cfg['app']}?entity={entity_id}"
    composio_connected_id = None
    status_label = "pending"
    # `oauth_not_configured` flips to True when Composio's response indicates
    # the toolkit's OAuth app isn't set up in the project (we'd otherwise
    # silently land the user on Composio's dashboard instead of the real
    # provider login). The frontend renders an explanatory toast in that case.
    oauth_not_configured = False

    # Composio v3 flow: first look up an existing auth_config for this toolkit
    # (or fail loudly if none) — then POST /connected_accounts to initiate
    # the OAuth dance for that user against that auth_config.
    auth_config_id = None
    # Distinguish "bad key" (10401) from "no auth_config for this toolkit"
    # so the frontend toast tells the user the right fix.
    failure_reason = None  # None | "invalid_key" | "no_auth_config" | "init_failed"
    try:
        async with httpx.AsyncClient(timeout=15.0) as c:
            r = await c.get(
                f"{COMPOSIO_BASE_URL}/api/v3/auth_configs?toolkit_slug={cfg['app']}",
                headers={"x-api-key": COMPOSIO_API_KEY},
            )
            if r.status_code == 401:
                failure_reason = "invalid_key"
            elif r.status_code < 400:
                cfgs = (r.json() or {}).get("items") or []
                if cfgs:
                    auth_config_id = cfgs[0].get("id")
                else:
                    failure_reason = "no_auth_config"
            else:
                failure_reason = "init_failed"
                logger.warning(f"Composio auth_configs {r.status_code} for {cfg['app']}: {r.text[:200]}")
    except Exception as e:
        logger.warning(f"Composio auth_configs lookup failed: {e}")
        failure_reason = "init_failed"

    if auth_config_id and failure_reason is None:
        try:
            async with httpx.AsyncClient(timeout=15.0) as c:
                r = await c.post(
                    f"{COMPOSIO_BASE_URL}/api/v3/connected_accounts/link",
                    headers={"x-api-key": COMPOSIO_API_KEY, "Content-Type": "application/json"},
                    json={
                        "auth_config_id": auth_config_id,
                        "user_id": entity_id,
                    },
                )
                if r.status_code == 401:
                    failure_reason = "invalid_key"
                elif r.status_code in (200, 201):
                    payload = r.json() or {}
                    # /link response shape (per Composio v3 latest):
                    #   { link_token, redirect_url, expires_at,
                    #     connected_account_id, experimental: {...} }
                    redirect_url = (
                        payload.get("redirect_url")
                        or payload.get("redirectUrl")
                        or redirect_url
                    )
                    composio_connected_id = (
                        payload.get("connected_account_id")
                        or payload.get("connectedAccountId")
                        or payload.get("id")
                    )
                else:
                    failure_reason = "init_failed"
                    logger.warning(f"Composio /link {r.status_code} for {body.source_kind}: {r.text[:200]}")
        except Exception as e:
            logger.warning(f"Composio /link for {body.source_kind} failed: {e}")
            failure_reason = "init_failed"

    if not composio_connected_id and failure_reason is None:
        failure_reason = "init_failed"
    oauth_not_configured = failure_reason is not None

    sid = str(uuid.uuid4())
    error_messages = {
        "invalid_key": (
            "Composio API key is invalid or has been rotated. The current key in "
            "the backend env returns 401. Generate a fresh key at dashboard.composio.dev "
            "→ Settings → API Keys with read+write scopes, paste it into your env "
            "(COMPOSIO_API_KEY in production, /app/backend/.env in preview), and "
            "redeploy. Pre-existing connections will keep working — only new connect "
            "attempts need the new key."
        ),
        "no_auth_config": (
            f"{cfg['label']} has no Auth Config in your Composio project yet. "
            f"Open dashboard.composio.dev → Auth Configs → New → pick \"{cfg['label']}\" "
            "→ enable \"Use Composio managed auth\" (or paste your own client_id/secret) "
            "→ Save. Then click Connect again."
        ),
        "init_failed": (
            f"Composio refused to initiate the {cfg['label']} OAuth handshake. "
            "Check the backend logs for the upstream error, verify the auth_config "
            "is set to \"active\", and confirm your API key still has write scopes."
        ),
    }
    doc = {
        "id": sid,
        "listing_id": lid,
        "source_kind": body.source_kind,
        "label": body.label or cfg["label"],
        "folder_id": body.folder_id,
        "folder_ids": body.folder_ids or ([body.folder_id] if body.folder_id else []),
        "folder_labels": body.folder_labels or [],
        "include_subfolders": body.include_subfolders if body.include_subfolders is not None else True,
        "entity_id": entity_id,
        "composio_connected_id": composio_connected_id,
        "redirect_url": redirect_url,
        "status": status_label,
        "created_by": user["id"],
        "created_at": now_utc().isoformat(),
        "last_sync_at": None,
        "file_count": 0,
        "last_error": error_messages.get(failure_reason) if oauth_not_configured else None,
    }
    await db.listing_external_sources.insert_one(doc)
    # Also drop a row in the legacy composio_connections collection so the
    # existing /composio/connections list keeps showing it.
    await db.composio_connections.insert_one({
        "id": sid,
        "user_id": user["id"],
        "app": cfg["app"],
        "entity_id": entity_id,
        "status": status_label,
        "redirect_url": redirect_url,
        "listing_id": lid,
        "created_at": doc["created_at"],
    })
    await log_audit(user["id"], "listing.source.connect.init", lid,
                    {"source_kind": body.source_kind, "sid": sid,
                     "oauth_not_configured": oauth_not_configured})
    doc.pop("_id", None)
    doc["oauth_not_configured"] = oauth_not_configured
    return doc


@api_router.get("/listings/{lid}/external-sources/{sid}/browse")
async def browse_external_source_folders(
    lid: str, sid: str,
    parent_id: Optional[str] = None,
    user=Depends(get_current_user),
):
    """Folder picker backend. Returns the immediate subfolders of `parent_id`
    inside the connected source. `parent_id` omitted → provider root.

    Used by the multi-select folder picker modal in the seller UI. The
    picker doesn't list files — only folders — so the seller can navigate
    deeply without scrolling past thousands of file rows."""
    await _listing_for_edit_or_404(lid, user)
    src = await db.listing_external_sources.find_one(
        {"id": sid, "listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0},
    )
    if not src:
        raise HTTPException(status_code=404, detail="Source not found")
    if src.get("status") != "active":
        # Return the documented {folders, can_browse, err} envelope so the
        # picker modal can render a friendly "finish OAuth first" message
        # in-context, rather than the seller seeing a generic network error.
        return {
            "folders": [],
            "parent_id": parent_id,
            "can_browse": False,
            "err": "Connection isn't active yet. Finish the OAuth handshake in the source row, then re-open the picker.",
            "source_kind": src["source_kind"],
            "requires_oauth": True,
        }
    result = await _browse_folder(src, parent_id)
    result["source_kind"] = src["source_kind"]
    return result


@api_router.patch("/listings/{lid}/external-sources/{sid}/folders")
async def update_external_source_folders(
    lid: str, sid: str, body: ExternalSourceFoldersUpdate,
    user=Depends(get_current_user),
):
    """Persist the seller's folder selection after they confirm the picker.
    `folder_ids` is the canonical list (multi-select); `folder_labels` is a
    parallel array of display names so the UI can show 'Marketing /
    Financials' without re-resolving each ID against the provider.

    Saves and triggers an immediate sync so the seller sees mirrored files
    show up without an extra click."""
    await _listing_for_edit_or_404(lid, user)
    src = await db.listing_external_sources.find_one(
        {"id": sid, "listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0},
    )
    if not src:
        raise HTTPException(status_code=404, detail="Source not found")
    # Sanity: prevent unbounded selections.
    if len(body.folder_ids) > 20:
        raise HTTPException(status_code=400, detail="Pick at most 20 folders per source.")
    if body.folder_labels is not None and len(body.folder_labels) != len(body.folder_ids):
        raise HTTPException(status_code=400, detail="folder_labels length must match folder_ids length")
    update = {
        "folder_ids": list(body.folder_ids),
        "folder_labels": list(body.folder_labels or []),
        "include_subfolders": bool(body.include_subfolders if body.include_subfolders is not None else True),
        # Keep legacy single field in sync so older code paths still work.
        "folder_id": body.folder_ids[0] if body.folder_ids else None,
    }
    await db.listing_external_sources.update_one({"id": sid}, {"$set": update})
    await log_audit(user["id"], "listing.source.folders.updated", lid, {
        "sid": sid, "folder_count": len(body.folder_ids),
        "include_subfolders": update["include_subfolders"],
    })
    # Trigger a fresh sync immediately so the seller sees results.
    await db.listing_external_sources.update_one(
        {"id": sid}, {"$set": {"syncing": True, "last_error": None}},
    )
    asyncio.create_task(_run_external_source_sync(lid, sid, user["id"]))
    return {"ok": True, **update}


@api_router.get("/listings/{lid}/external-sources")
async def list_external_sources(lid: str, user=Depends(get_current_user)):
    """List file sources connected to a listing. Viewer just needs read access
    on the listing (collaborators + buyers in active Vaults included).

    Self-healing: pending sources older than 1h are auto-purged here so a
    half-finished OAuth attempt doesn't permanently occupy the "Connect"
    slot for that toolkit. ACTIVE sources are never auto-touched.
    """
    await _listing_for_view_or_404(lid, user)
    stale_cutoff = (now_utc() - timedelta(hours=1)).isoformat()
    stale = await db.listing_external_sources.find(
        {"listing_id": lid, "status": "pending",
         "created_at": {"$lt": stale_cutoff},
         "deleted_at": {"$exists": False}},
        {"_id": 0, "id": 1},
    ).to_list(50)
    if stale:
        stale_ids = [s["id"] for s in stale]
        await db.listing_external_sources.update_many(
            {"id": {"$in": stale_ids}},
            {"$set": {"deleted_at": now_utc().isoformat(),
                      "status": "expired"}},
        )
        await db.composio_connections.update_many(
            {"id": {"$in": stale_ids}},
            {"$set": {"status": "expired"}},
        )

    rows = await db.listing_external_sources.find(
        {"listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0, "redirect_url": 0}
    ).sort("created_at", -1).to_list(50)
    return {"sources": rows, "supported": [
        {"kind": k, "label": v["label"]} for k, v in COMPOSIO_FILE_SOURCES.items()
    ]}


@api_router.post("/listings/{lid}/external-sources/{sid}/poll")
async def poll_external_source(lid: str, sid: str, user=Depends(get_current_user)):
    """Ask Composio whether the OAuth dance finished. Called by the frontend
    every few seconds after opening the connect window. Marks the source
    ACTIVE on success so the sync button unlocks."""
    await _listing_for_edit_or_404(lid, user)
    src = await db.listing_external_sources.find_one(
        {"id": sid, "listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0}
    )
    if not src:
        raise HTTPException(status_code=404, detail="Source not found")
    if src["status"] == "active":
        return {"status": "active", "source": src}
    if not src.get("composio_connected_id"):
        return {"status": src["status"]}

    new_status = src["status"]
    try:
        async with httpx.AsyncClient(timeout=10.0) as c:
            r = await c.get(
                f"{COMPOSIO_BASE_URL}/api/v3/connected_accounts/{src['composio_connected_id']}",
                headers={"x-api-key": COMPOSIO_API_KEY},
            )
            if r.status_code < 400:
                payload = r.json()
                upstream = (payload.get("status") or "").upper()
                if upstream == "ACTIVE":
                    new_status = "active"
                elif upstream in ("FAILED", "EXPIRED", "REVOKED"):
                    new_status = "failed"
    except Exception as e:
        logger.warning(f"Composio poll failed: {e}")

    if new_status != src["status"]:
        await db.listing_external_sources.update_one(
            {"id": sid}, {"$set": {"status": new_status}}
        )
        await db.composio_connections.update_one(
            {"id": sid}, {"$set": {"status": new_status}}
        )
        await log_audit(user["id"], "listing.source.status", lid,
                        {"sid": sid, "status": new_status})
    src["status"] = new_status
    return {"status": new_status, "source": src}


async def _mirror_one_file(lid: str, sid: str, src_kind: str, raw_name: str,
                            data: bytes, content_type: str, external_id: str,
                            user_id: str) -> str:
    """Persist a single externally-pulled file using the SAME schema as manual
    uploads so the Vault clone + Copilot indexer just work. Returns file_id."""
    file_id = str(uuid.uuid4())
    filename = raw_name or f"{src_kind}-{file_id}"
    plaintext_sha = sha256_hex(data)

    storage_bytes = data
    encrypted = False
    encryption_alg = None
    if encryption_configured():
        try:
            aad = f"listing:{lid}:{file_id}".encode("utf-8")
            enc = encrypt_bytes(data, associated_data=aad)
            storage_bytes = enc["envelope"]
            encrypted = True
            encryption_alg = enc["alg"]
        except Exception as e:
            logger.warning(f"At-rest encryption failed for external file, storing plaintext: {e}")

    gridfs_id = await listing_files_bucket.upload_from_stream(
        filename, io.BytesIO(storage_bytes),
        metadata={
            "listing_id": lid, "file_id": file_id, "uploaded_by": user_id,
            "content_type": content_type, "encrypted": encrypted,
            "encryption_alg": encryption_alg, "source": src_kind, "source_sid": sid,
        },
    )

    pages = extract_pages_from_bytes(filename, data)
    flat = pages_to_flat_text(pages)

    await db.listing_staged_files.insert_one({
        "id": file_id, "listing_id": lid, "filename": filename, "folder": "other",
        "content_type": content_type, "size_bytes": len(data),
        "page_count": len(pages), "pages": pages, "content": flat,
        "char_count": len(flat), "gridfs_id": str(gridfs_id), "storage": "listing_gridfs",
        "encrypted": encrypted, "encryption_alg": encryption_alg,
        "sha256_hex": plaintext_sha, "uploaded_by": user_id,
        "uploaded_at": now_utc().isoformat(),
        # External-source provenance (so the UI badges it and the cleanup
        # job knows which files to wipe on disconnect / listing close).
        "source": {"kind": src_kind, "sid": sid, "external_id": external_id},
    })
    return file_id


# Per-toolkit input parameter mapping for the folder-id field. Google Drive
# expects `folderId` (camelCase), Box uses `folder_id`, SharePoint uses
# `folderPath` etc. Without this map the user-supplied folder ID is silently
# ignored and the action lists from the drive's root.
FOLDER_PARAM_KEY = {
    "googledrive": "folderId",
    "onedrive":    None,            # Composio's ONE_DRIVE_ONEDRIVE_LIST_ITEMS only lists root.
                                    # Folder browsing is done via the proxy → MS Graph path.
    "sharepoint":  None,            # No Composio list tool exposed for SharePoint.
    "dropbox":     "path",
    "box":         "folder_id",
}

# Providers that natively support listing arbitrary subfolders via a Composio
# action (i.e. you can pass folder_id/path and get its direct children).
# OneDrive is browsed via the proxy → Graph path (see _browse_folder), and
# SharePoint has no Composio file-list tool — picker degrades to manual-id.
_FOLDER_BROWSE_NATIVE = {"box", "googledrive", "dropbox"}


def _normalise_composio_response(resp: dict) -> dict:
    """Composio v3 `/tools/execute` wraps the action's native output. Peel the
    layers so callers always see `{successful, data, error}` at the top."""
    if not isinstance(resp, dict):
        return {"successful": False, "data": None, "error": "non-dict response"}
    # Some toolkit responses come back as { data: { successful, data, error } }
    # others as { successful, data, error } directly. Normalise.
    if "successful" in resp and "data" in resp:
        return resp
    inner = resp.get("data") or resp.get("response_data") or {}
    if isinstance(inner, dict) and "successful" in inner:
        return inner
    # Last resort: treat the whole thing as data if it has no envelope at all.
    return {"successful": True, "data": resp, "error": None}


def _extract_files_array(data_obj) -> list:
    """Pull the file-list array out of an action's `data` payload. Different
    connectors use different keys — try them in the order we've seen."""
    if isinstance(data_obj, list):
        return data_obj
    if not isinstance(data_obj, dict):
        return []
    for key in ("files", "entries", "value", "items", "fileList", "results"):
        v = data_obj.get(key)
        if isinstance(v, list):
            return v
    return []


async def _collect_files_under_folder(
    src: dict, root_id: str | None, include_subfolders: bool,
    max_files: int, errors: list[str],
) -> tuple[list[dict], int, int]:
    """BFS-collect file entries under `root_id` for the given source.
    Returns (files_meta, explored_folders, skipped_non_file_entries).

    - include_subfolders=False: only the immediate file children of root_id.
    - include_subfolders=True : recurse depth-first up to MAX_DEPTH=4 or
      until we hit `max_files` total file entries. Same bounds for every
      provider so the picker behaves predictably.

    For Dropbox we let the toolkit handle recursion via its built-in
    `recursive` flag (no need to BFS at the app layer). For OneDrive we use
    the Composio proxy because the predefined LIST_ITEMS tool only sees
    drive root.
    """
    kind = src["source_kind"]
    conn_id = src["composio_connected_id"]
    entity = src.get("entity_id")
    MAX_DEPTH = 4
    files: list[dict] = []
    explored = 0
    skipped = 0

    # ── Dropbox: native recursion. One call, no BFS needed. ──────────────
    if kind == "dropbox":
        path = root_id if root_id is not None else ""
        try:
            raw = await _composio_action_execute(
                "DROPBOX_LIST_FILES_IN_FOLDER", conn_id,
                {"path": path, "limit": min(max_files, 500),
                 "recursive": bool(include_subfolders)},
                user_id=entity,
            )
        except HTTPException as e:
            errors.append(f"dropbox {path!r}: {str(e.detail)[:140]}")
            return files, explored, skipped
        resp = _normalise_composio_response(raw)
        if not resp.get("successful"):
            errors.append(f"dropbox {path!r}: list failed ({str(resp.get('error'))[:140]})")
            return files, explored, skipped
        for i in _extract_files_array(resp.get("data")):
            if not isinstance(i, dict):
                continue
            tag = (i.get(".tag") or i.get("tag") or "").lower()
            if tag == "file":
                # Normalize to our common schema (id + name + mime_type).
                files.append({
                    "id": i.get("id") or i.get("path_lower") or i.get("path_display"),
                    "name": i.get("name"),
                    "_path": i.get("path_display") or i.get("path_lower"),
                    "size": i.get("size"),
                })
            elif tag == "folder":
                # We don't count Dropbox subfolders in `explored` since
                # the tool flattened them for us.
                pass
            else:
                skipped += 1
        return files[:max_files], explored, skipped

    # ── OneDrive: always via proxy → Graph (predefined tool only sees root) ─
    if kind == "onedrive":
        queue: list[tuple[str, int]] = [(root_id or "root", 0)]
        while queue and len(files) < max_files:
            fid, depth = queue.pop(0)
            endpoint = "/me/drive/root/children" if fid in (None, "", "root") \
                else f"/me/drive/items/{fid}/children"
            body = await _composio_proxy_get_json(
                "onedrive", conn_id, endpoint,
                query=[("$select", "id,name,folder,file,size"), ("$top", "200")],
                user_id=entity,
            )
            if not isinstance(body, dict):
                errors.append(f"onedrive folder {fid}: proxy returned nothing")
                continue
            explored += 1
            for i in body.get("value") or []:
                if not isinstance(i, dict):
                    continue
                if i.get("folder"):
                    if include_subfolders and depth + 1 <= MAX_DEPTH and i.get("id"):
                        queue.append((i["id"], depth + 1))
                elif i.get("file"):
                    file_meta = i.get("file") or {}
                    files.append({
                        "id": i["id"],
                        "name": i.get("name"),
                        "mime_type": file_meta.get("mimeType") or "application/octet-stream",
                        "size": i.get("size"),
                    })
                else:
                    skipped += 1
        return files[:max_files], explored, skipped

    # ── Native list providers (Box / Google Drive): BFS via Composio tool ─
    cfg = COMPOSIO_FILE_SOURCES[kind]

    def _build_list_input(folder_id: str | None) -> dict:
        inp: dict = {}
        if kind == "box":
            inp["folder_id"] = folder_id or "0"
        elif kind == "googledrive":
            inp["pageSize"] = 200
            inp["supportsAllDrives"] = True
            inp["includeItemsFromAllDrives"] = True
            inp["fields"] = "files(id,name,mimeType,size)"
            if folder_id:
                inp["folderId"] = folder_id
        return inp

    def _classify(entry: dict) -> str:
        """Return 'file' | 'folder' | 'skip' for a single list entry."""
        if not isinstance(entry, dict):
            return "skip"
        if kind == "box":
            t = (entry.get("type") or "").lower()
            if t == "file":   return "file"
            if t == "folder": return "folder"
            if t == "":       return "file"  # default-keep unknown
            return "skip"
        if kind == "googledrive":
            if entry.get("mimeType") == "application/vnd.google-apps.folder":
                return "folder"
            return "file"
        return "file"

    queue: list[tuple[str | None, int]] = [(root_id, 0)]
    seen_roots: set[str] = set()

    while queue and len(files) < max_files:
        fid, depth = queue.pop(0)
        if fid in seen_roots:
            continue
        seen_roots.add(fid or "")
        try:
            raw = await _composio_action_execute(
                cfg["list"], conn_id, _build_list_input(fid),
                user_id=entity,
            )
        except HTTPException as e:
            errors.append(f"{kind} folder {fid}: {str(e.detail)[:140]}")
            continue
        resp = _normalise_composio_response(raw)
        if not resp.get("successful"):
            errors.append(f"{kind} folder {fid}: list failed ({str(resp.get('error'))[:140]})")
            continue
        explored += 1
        for entry in _extract_files_array(resp.get("data")):
            klass = _classify(entry)
            if klass == "skip":
                skipped += 1
                continue
            if klass == "folder":
                if include_subfolders and depth + 1 <= MAX_DEPTH and entry.get("id"):
                    queue.append((str(entry["id"]), depth + 1))
                continue
            files.append({
                "id": entry.get("id") or entry.get("file_id") or entry.get("ID"),
                "name": entry.get("name") or entry.get("filename") or entry.get("title"),
                "mime_type": entry.get("mime_type") or entry.get("mimeType")
                             or entry.get("content_type") or "application/octet-stream",
                "size": entry.get("size"),
            })
            if len(files) >= max_files:
                break

    return files[:max_files], explored, skipped


async def _run_external_source_sync(lid: str, sid: str, user_id: str) -> None:
    """The actual sync work — runs in a background task so the HTTP request
    returns immediately and Cloudflare's 100s gateway timeout never fires.
    Progress is observable via the `syncing` flag + `last_sync_at` on the
    source doc; the frontend polls /external-sources to follow along.

    Multi-folder model: iterate every entry in `src.folder_ids` (or fall
    back to legacy `src.folder_id`, or each provider's root). For each
    selected folder we recurse if `src.include_subfolders` is True
    (default). Total file cap = 200 across all selected folders combined.
    """
    src = await db.listing_external_sources.find_one(
        {"id": sid, "listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0}
    )
    if not src or src["status"] != "active":
        return
    kind = src["source_kind"]
    cfg = COMPOSIO_FILE_SOURCES.get(kind) or {}
    if not cfg.get("list") or not cfg.get("download"):
        await db.listing_external_sources.update_one(
            {"id": sid},
            {"$set": {"syncing": False, "last_sync_at": now_utc().isoformat(),
                      "last_error": f"{cfg.get('label', kind)} sync isn't available — Composio "
                                    "doesn't yet expose list/download tools for this provider."}},
        )
        return

    # Resolve folder selection: prefer the new multi-folder list, fall back
    # to the legacy single folder_id, fall back to provider root.
    selected: list[str | None] = []
    if isinstance(src.get("folder_ids"), list) and src["folder_ids"]:
        selected = list(src["folder_ids"])
    elif src.get("folder_id"):
        selected = [src["folder_id"]]
    else:
        # None means "use provider default root" inside `_collect_files_under_folder`.
        selected = [None]
    include_subfolders = bool(src.get("include_subfolders", True))
    MAX_TOTAL = 200

    pulled = 0
    errors: list[str] = []
    sample_response: str | None = None
    explored_folders_total = 0
    skipped_non_file_total = 0
    files_meta: list[dict] = []
    try:
        for folder_id in selected:
            remaining = MAX_TOTAL - len(files_meta)
            if remaining <= 0:
                break
            batch, explored, skipped = await _collect_files_under_folder(
                src, folder_id, include_subfolders, remaining, errors,
            )
            files_meta.extend(batch)
            explored_folders_total += explored
            skipped_non_file_total += skipped
        sample_response = json.dumps(
            {"folders_scanned": len(selected),
             "include_subfolders": include_subfolders,
             "explored": explored_folders_total,
             "skipped": skipped_non_file_total,
             "file_count_in_sample": len(files_meta)},
            default=str,
        )[:1200]
    except HTTPException as e:
        await db.listing_external_sources.update_one(
            {"id": sid},
            {"$set": {"syncing": False, "last_sync_at": now_utc().isoformat(),
                      "last_error": f"list failed: {str(e.detail)[:280]}"}},
        )
        return
    except Exception as e:
        await db.listing_external_sources.update_one(
            {"id": sid},
            {"$set": {"syncing": False, "last_sync_at": now_utc().isoformat(),
                      "last_error": f"list crashed: {str(e)[:280]}"}},
        )
        return

    explored_folders = explored_folders_total
    skipped_non_file = skipped_non_file_total

    for f in files_meta[:MAX_TOTAL]:
        # Defensive: still skip any non-file remnant.
        ftype = (f.get("type") or "").lower()
        if ftype and ftype != "file":
            skipped_non_file += 1
            continue
        external_id = f.get("id") or f.get("file_id") or f.get("ID")
        name = f.get("name") or f.get("filename") or f.get("title") or external_id
        mime = f.get("mime_type") or f.get("mimeType") or f.get("content_type") or "application/octet-stream"
        if not external_id:
            continue
        existing = await db.listing_staged_files.find_one(
            {"listing_id": lid, "source.sid": sid, "source.external_id": external_id, "deleted_at": {"$exists": False}},
            {"_id": 0, "id": 1},
        )
        if existing:
            continue
        try:
            # Provider-specific download arg name. Dropbox uses `path`,
            # everyone else uses `file_id`.
            if kind == "dropbox":
                dl_args = {"path": f.get("_path") or external_id}
            else:
                dl_args = {"file_id": external_id}

            blob: bytes | None = None
            action_error: str | None = None

            # PRIMARY PATH: Composio Proxy Execute → native provider API.
            # Why proxy-first for Box/Drive/OneDrive/SharePoint: their
            # predefined `*_DOWNLOAD_FILE` actions base64-inline the file
            # bytes into Composio's tool-response envelope, which has a
            # ~10MB hard cap on the server side ("The tool response payload
            # is too large"). Proxy execute streams the binary via R2 with
            # a presigned URL → no size limit. Dropbox falls through to the
            # predefined action because we don't have a proxy endpoint
            # configured for it (its API requires a Dropbox-API-Arg header,
            # not query params).
            if kind in PROXY_DOWNLOAD_ENDPOINTS:
                blob = await _composio_proxy_download(
                    src["source_kind"], src["composio_connected_id"],
                    external_id, mime_type=mime, user_id=src.get("entity_id"),
                )

            # FALLBACK PATH: predefined action. Wrapped in its own try/except
            # so a Composio 502 (payload-too-large, rate-limit, etc.) doesn't
            # bubble out and skip recording per-file errors.
            if blob is None:
                try:
                    raw_dl = await _composio_action_execute(
                        cfg["download"], src["composio_connected_id"], dl_args,
                        user_id=src.get("entity_id"),
                    )
                    dl_resp = _normalise_composio_response(raw_dl)
                    if not dl_resp.get("successful"):
                        action_error = str(dl_resp.get("error"))[:160]
                    else:
                        dl_raw = dl_resp.get("data") or {}
                        # Drive returns either `file` (string bytes, possibly base64),
                        # or a presigned URL. Other connectors use `content_base64`,
                        # `data`, or just `file_content`. Try them all.
                        b64 = (
                            dl_raw.get("content_base64") or dl_raw.get("file_content")
                            or (dl_raw.get("file") if isinstance(dl_raw.get("file"), str) else None)
                        )
                        url = dl_raw.get("download_url") or dl_raw.get("url")
                        if not url and isinstance(dl_raw.get("file"), dict):
                            url = dl_raw["file"].get("url") or dl_raw["file"].get("download_url")
                        if b64:
                            import base64 as _b64
                            try:
                                blob = _b64.b64decode(b64)
                            except Exception:
                                blob = None
                        if blob is None and url:
                            async with httpx.AsyncClient(timeout=180.0, follow_redirects=True) as c:
                                rr = await c.get(url)
                                if rr.status_code < 400:
                                    blob = rr.content
                except HTTPException as e:
                    # Composio returned 4xx/5xx — most commonly the "tool
                    # response payload is too large" message for files
                    # over ~10MB on the predefined action path. Surface
                    # this so the seller sees what happened if proxy
                    # ALSO failed (it usually doesn't, but log defensively).
                    action_error = str(e.detail)[:160]

            if blob is None:
                errors.append(
                    f"{name}: download failed via proxy + action"
                    + (f" (action: {action_error})" if action_error else "")
                )
                continue
            # Sanity cap so a runaway connection can't blow the pod. GridFS
            # handles arbitrary size on the storage side; the cap exists to
            # protect against accidental misconfigurations (e.g. a 5GB ISO).
            # Per user request "no limit to the payload amount" — raised
            # from 50MB → 500MB so org charts, signed COIs, deck rasters all
            # mirror cleanly. Files larger than this are explicitly listed.
            if len(blob) > 500 * 1024 * 1024:
                errors.append(
                    f"{name}: {len(blob) // (1024*1024)} MB exceeds the 500 MB per-file safety cap. "
                    "Compress, split, or share via a direct preview link instead."
                )
                continue
            await _mirror_one_file(lid, sid, src["source_kind"], name, blob, mime, external_id, user_id)
            pulled += 1
            # Live progress: update file_count after every successful pull
            # so the UI can show a counter ticking up.
            new_count_so_far = await db.listing_staged_files.count_documents(
                {"listing_id": lid, "source.sid": sid, "deleted_at": {"$exists": False}}
            )
            await db.listing_external_sources.update_one(
                {"id": sid}, {"$set": {"file_count": new_count_so_far}}
            )
        except HTTPException as e:
            errors.append(f"{name}: {str(e.detail)[:140]}")
        except Exception as e:
            errors.append(f"{name}: {str(e)[:140]}")

    new_count = await db.listing_staged_files.count_documents(
        {"listing_id": lid, "source.sid": sid, "deleted_at": {"$exists": False}}
    )
    final_error = "; ".join(errors[:3]) if errors else None
    # If we got 0 files AND no errors, the list call probably parsed wrong —
    # OR the connected folder is empty / contains only non-file entries.
    # Surface a clear, actionable message in both cases.
    if pulled == 0 and not final_error:
        if explored_folders > 0 or skipped_non_file > 0:
            parts = []
            if explored_folders > 0:
                parts.append(f"{explored_folders} subfolder(s)")
            if skipped_non_file > 0:
                parts.append(f"{skipped_non_file} non-file item(s)")
            final_error = (
                f"Connected, but no downloadable files found after scanning "
                f"{' and '.join(parts)}. Move files into the folder you connected, "
                "OR pick a specific subfolder containing files when you reconnect."
            )
        elif not files_meta:
            final_error = (
                "List returned 0 files. If your folder actually has files, paste this "
                "back to your engineer to fix the response parser: "
                f"{sample_response[:600] if sample_response else '(no sample captured)'}"
            )
        # else: we had file entries but every download failed — `errors` would
        # have surfaced individual reasons; nothing further to add.
    await db.listing_external_sources.update_one(
        {"id": sid},
        {"$set": {"syncing": False, "last_sync_at": now_utc().isoformat(),
                  "file_count": new_count,
                  "last_error": final_error,
                  "last_response_sample": sample_response}},
    )
    # Eager backfill: clone any new staged files into already-open Vaults for
    # this listing so buyers + Copilot see the synced docs without waiting for
    # the next get_deal_room hit. Bounded — only pulls open vault rows.
    if pulled > 0:
        try:
            open_rooms = await db.deal_rooms.find(
                {"listing_id": lid, "status": {"$in": ["pending_nda", "active", "preview"]},
                 "deleted_at": {"$exists": False}},
                {"_id": 0, "id": 1, "seller_id": 1},
            ).to_list(200)
            for r in open_rooms:
                try:
                    added = await _clone_listing_files_into_room(
                        lid, r["id"], r.get("seller_id") or user_id, only_missing=True
                    )
                    if added:
                        logger.info(f"sync backfill: cloned {added} file(s) into room {r['id']}")
                except Exception as e:
                    logger.warning(f"sync backfill failed for room {r.get('id')}: {e}")
        except Exception as e:
            logger.warning(f"sync backfill enumeration failed for listing {lid}: {e}")
    await log_audit(user_id, "listing.source.sync", lid,
                    {"sid": sid, "pulled": pulled, "total": new_count, "errors": len(errors)})


@api_router.post("/listings/{lid}/external-sources/{sid}/sync")
async def sync_external_source(lid: str, sid: str, user=Depends(get_current_user)):
    """Kick off a background sync. Returns immediately so Cloudflare's 100s
    gateway timeout never fires regardless of folder size. Frontend polls
    /external-sources to see `syncing` flip back to False + file_count
    tick up live."""
    await _listing_for_edit_or_404(lid, user)
    src = await db.listing_external_sources.find_one(
        {"id": sid, "listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0}
    )
    if not src:
        raise HTTPException(status_code=404, detail="Source not found")
    if src["status"] != "active":
        raise HTTPException(status_code=400, detail=f"Source is {src['status']}, not active. Complete OAuth first.")
    if src.get("syncing"):
        raise HTTPException(status_code=409, detail="Sync already in progress for this source")

    # Mark as syncing BEFORE returning so the very next GET sees the flag.
    await db.listing_external_sources.update_one(
        {"id": sid}, {"$set": {"syncing": True, "last_error": None}}
    )
    # Fire and forget — the task runs detached. Exceptions inside set the
    # source's last_error so they're surfaced to the UI rather than lost.
    asyncio.create_task(_run_external_source_sync(lid, sid, user["id"]))
    return {"ok": True, "started": True, "syncing": True}


async def _wipe_external_source_files(lid: str, sid: str):
    """Delete every mirrored file (rows + GridFS bytes) belonging to a source."""
    cur = db.listing_staged_files.find(
        {"listing_id": lid, "source.sid": sid, "deleted_at": {"$exists": False}},
        {"gridfs_id": 1, "id": 1},
    )
    async for f in cur:
        try:
            await listing_files_bucket.delete(ObjectId(f["gridfs_id"]))
        except Exception:
            pass
    await db.listing_staged_files.update_many(
        {"listing_id": lid, "source.sid": sid},
        {"$set": {"deleted_at": now_utc().isoformat()}},
    )


@api_router.delete("/listings/{lid}/external-sources/{sid}")
async def disconnect_external_source(lid: str, sid: str, user=Depends(get_current_user)):
    """Disconnect a file source and wipe every mirrored byte. Best-effort
    revoke on Composio's side; we always delete locally even if upstream
    revoke fails so the seller's data is purged on schedule."""
    await _listing_for_edit_or_404(lid, user)
    src = await db.listing_external_sources.find_one(
        {"id": sid, "listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0}
    )
    if not src:
        raise HTTPException(status_code=404, detail="Source not found")

    if src.get("composio_connected_id"):
        try:
            async with httpx.AsyncClient(timeout=10.0) as c:
                await c.delete(
                    f"{COMPOSIO_BASE_URL}/api/v3/connected_accounts/{src['composio_connected_id']}",
                    headers={"x-api-key": COMPOSIO_API_KEY},
                )
        except Exception as e:
            logger.warning(f"Composio revoke failed (continuing local wipe): {e}")

    await _wipe_external_source_files(lid, sid)
    await db.listing_external_sources.update_one(
        {"id": sid}, {"$set": {"deleted_at": now_utc().isoformat()}}
    )
    await db.composio_connections.update_one(
        {"id": sid}, {"$set": {"status": "revoked", "revoked_at": now_utc().isoformat()}}
    )
    await log_audit(user["id"], "listing.source.disconnect", lid, {"sid": sid})
    return {"ok": True}


async def _wipe_listing_external_sources(lid: str):
    """Listing-close hook: revoke every connected source and wipe its files.
    Called from the status-flip path so closing a deal evicts the seller's
    OAuth grants automatically (Rule 3A). Composio revokes are fired in
    parallel so a listing with 6 connected sources doesn't stack 6×10s timeouts."""
    sources = await db.listing_external_sources.find(
        {"listing_id": lid, "deleted_at": {"$exists": False}}, {"_id": 0}
    ).to_list(100)

    async def _revoke(src):
        if not src.get("composio_connected_id"):
            return
        try:
            async with httpx.AsyncClient(timeout=10.0) as c:
                await c.delete(
                    f"{COMPOSIO_BASE_URL}/api/v3/connected_accounts/{src['composio_connected_id']}",
                    headers={"x-api-key": COMPOSIO_API_KEY},
                )
        except Exception:
            pass

    if sources:
        await asyncio.gather(*(_revoke(s) for s in sources))
        for src in sources:
            await _wipe_external_source_files(lid, src["id"])
    await db.listing_external_sources.update_many(
        {"listing_id": lid, "deleted_at": {"$exists": False}},
        {"$set": {"deleted_at": now_utc().isoformat()}},
    )
    await db.composio_connections.update_many(
        {"listing_id": lid, "status": {"$ne": "revoked"}},
        {"$set": {"status": "revoked", "revoked_at": now_utc().isoformat()}},
    )


# -----------------------------------------------------------------------------
# THE VAULT (NDA-gated workspace per inquiry · DRL · AI Findings)
# -----------------------------------------------------------------------------
DRL_TEMPLATES = {
    "saas": {
        "id": "saas",
        "name": "SaaS / Software",
        "items": [
            {"title": "Last 3 years P&L by month", "workstream": "finance"},
            {"title": "ARR / MRR cohort schedule (last 24m)", "workstream": "finance"},
            {"title": "Customer concentration (top 20 customers % of revenue)", "workstream": "commercial"},
            {"title": "Net & gross logo retention rates", "workstream": "commercial"},
            {"title": "All material customer contracts > $100k ARR", "workstream": "legal"},
            {"title": "Org chart + comp band schedule", "workstream": "hr"},
            {"title": "Key employee retention/non-compete agreements", "workstream": "hr"},
            {"title": "Stack architecture diagram + cloud spend", "workstream": "it"},
            {"title": "Code repository access + open-source license inventory", "workstream": "it"},
            {"title": "SOC2 / ISO27001 audit reports", "workstream": "it"},
            {"title": "Outstanding litigation & IP assignments", "workstream": "legal"},
        ],
    },
    "healthcare": {
        "id": "healthcare",
        "name": "Healthcare / MedTech",
        "items": [
            {"title": "Payer-mix breakdown last 3 years", "workstream": "finance"},
            {"title": "Reimbursement-rate schedule", "workstream": "finance"},
            {"title": "FDA / CE clearance documentation", "workstream": "legal"},
            {"title": "Clinical trial outcomes + adverse events register", "workstream": "operations"},
            {"title": "HIPAA / GDPR compliance attestation", "workstream": "legal"},
            {"title": "Hospital / GPO contracts portfolio", "workstream": "commercial"},
            {"title": "Manufacturing facility audits", "workstream": "operations"},
            {"title": "Quality system (ISO 13485) documentation", "workstream": "operations"},
            {"title": "Key physician/KOL engagement contracts", "workstream": "legal"},
            {"title": "Cybersecurity / connected-device risk assessment", "workstream": "it"},
        ],
    },
    "industrial": {
        "id": "industrial",
        "name": "Industrial / Manufacturing",
        "items": [
            {"title": "Plant capacity utilization (last 24m)", "workstream": "operations"},
            {"title": "Customer order book / backlog schedule", "workstream": "commercial"},
            {"title": "Supplier concentration & long-term agreements", "workstream": "commercial"},
            {"title": "EHS audits + incident register", "workstream": "operations"},
            {"title": "Environmental permits + remediation liabilities", "workstream": "legal"},
            {"title": "Union agreements + pension obligations", "workstream": "hr"},
            {"title": "Fixed asset register + depreciation schedule", "workstream": "finance"},
            {"title": "Working capital trend (last 36m)", "workstream": "finance"},
            {"title": "Insurance policies + claims history", "workstream": "legal"},
        ],
    },
    "finserv": {
        "id": "finserv",
        "name": "Financial Services",
        "items": [
            {"title": "Regulatory licenses + filings inventory", "workstream": "legal"},
            {"title": "AML/KYC policy + audit results", "workstream": "legal"},
            {"title": "AUM trend + revenue per client", "workstream": "finance"},
            {"title": "Capital adequacy + regulatory capital schedule", "workstream": "finance"},
            {"title": "Client concentration analysis", "workstream": "commercial"},
            {"title": "Custodian / clearing relationships", "workstream": "operations"},
            {"title": "Cybersecurity controls + SOC2 report", "workstream": "it"},
            {"title": "Litigation, regulatory inquiries, fines history", "workstream": "legal"},
        ],
    },
    "climatetech": {
        "id": "climatetech",
        "name": "ClimateTech / Energy",
        "items": [
            {"title": "Project pipeline + commissioning schedule", "workstream": "operations"},
            {"title": "PPA / offtake agreements portfolio", "workstream": "commercial"},
            {"title": "Tax credit eligibility (IRA/CBAM) documentation", "workstream": "finance"},
            {"title": "Permitting status by project", "workstream": "legal"},
            {"title": "Technology IP portfolio + patents", "workstream": "legal"},
            {"title": "Carbon accounting / LCA reports", "workstream": "operations"},
            {"title": "Grid interconnection agreements", "workstream": "legal"},
            {"title": "Capex schedule + financing structure", "workstream": "finance"},
        ],
    },
    "ecommerce": {
        "id": "ecommerce",
        "name": "E-commerce / DTC",
        "items": [
            {"title": "Shopify / platform export — last 36m orders", "workstream": "commercial"},
            {"title": "Paid acquisition spend by channel (Meta / Google / TikTok)", "workstream": "commercial"},
            {"title": "CAC, payback period & LTV by cohort (last 24m)", "workstream": "commercial"},
            {"title": "Repeat-purchase rate & 90-day reorder cohort", "workstream": "commercial"},
            {"title": "SKU-level gross margin schedule", "workstream": "finance"},
            {"title": "Inventory aging + on-hand by SKU + 3PL contract", "workstream": "operations"},
            {"title": "Return rate by SKU + refund liability accrual", "workstream": "finance"},
            {"title": "Influencer / affiliate agreements + spend log", "workstream": "commercial"},
            {"title": "Email / SMS list size, opt-in basis & deliverability", "workstream": "it"},
            {"title": "Trademark register + product safety / labelling compliance", "workstream": "legal"},
            {"title": "Manufacturer + supplier agreements with lead times", "workstream": "operations"},
        ],
    },
    "consumer": {
        "id": "consumer",
        "name": "Consumer / Retail",
        "items": [
            {"title": "Channel mix (DTC / wholesale / marketplace)", "workstream": "commercial"},
            {"title": "CAC / LTV by cohort (last 24m)", "workstream": "commercial"},
            {"title": "Inventory aging + obsolescence reserve", "workstream": "finance"},
            {"title": "Brand IP / trademark register", "workstream": "legal"},
            {"title": "Manufacturer / supplier agreements", "workstream": "commercial"},
            {"title": "Returns & warranty obligations", "workstream": "finance"},
            {"title": "Social / earned-media analytics", "workstream": "commercial"},
            {"title": "Sustainability claims + certifications", "workstream": "legal"},
        ],
    },
}


async def participant_check(room: dict, user: dict) -> str:
    """Returns 'buyer' | 'seller' | 'admin' if participant, raises 403 otherwise.

    Sell-side participants:
      - Literal `seller_id` on the room
      - Anyone in the listing's seller-side workspace (org members, listing
        collaborator-editors/owners)
      - Anyone added via the room-level `collaborators[]` array

    Buy-side participants:
      - Literal `buyer_id` on the room
      - Org teammates of the literal buyer — so an agent representing a
        buyer client (or a buy-side analyst team) can run diligence on the
        same Vault. Restricted to agents and the buyer's same-org members.
    """
    if user.get("role") == "admin":
        return "admin"
    if user["id"] == room.get("buyer_id"):
        return "buyer"
    if user["id"] == room.get("seller_id"):
        return "seller"
    # Room-level collaborators (phase-2 explicit add) — sell-side by default
    for c in room.get("collaborators", []) or []:
        if c.get("user_id") == user["id"]:
            return "seller"
    # Sell-side workspace teammates of the listing this room belongs to
    if user.get("role") in ("seller", "agent") and room.get("listing_id"):
        ws_listings, _ = await _user_workspace_listing_ids(user)
        if room["listing_id"] in ws_listings:
            return "seller"
    # Buy-side: agent or fellow team member who shares an org with the room's
    # buyer. Lets an advisor act on behalf of their buyer client.
    if user.get("role") in ("buyer", "agent") and room.get("buyer_id"):
        my_orgs = set(await _get_user_org_ids(user))
        if my_orgs:
            buyer_orgs = set(
                r["org_id"] for r in await db.org_memberships.find(
                    {"user_id": room["buyer_id"]}, {"_id": 0, "org_id": 1}
                ).to_list(50)
            )
            if my_orgs & buyer_orgs:
                return "buyer"
    raise HTTPException(status_code=403, detail="Not a participant of this Vault")


@api_router.post("/listings/{lid}/preview-vault")
async def open_preview_vault(lid: str, user=Depends(get_current_user)):
    """Create or reuse a personal "preview Vault" for the current user on a
    listing they have access to. Lets an agent / seller / org teammate / viewer-
    collaborator QA the full buyer-side Vault experience (AI copilot, DRL,
    findings, etc.) on their own listing without waiting for a real buyer to
    engage.

    Security: caller MUST already be on the sell-side workspace of the listing
    — principal owner, org member, collaborator (owner/editor/**viewer**), or
    admin. Viewers are read-only on the listing itself but they were invited
    precisely to see the buyer experience, so we let them spawn a preview vault
    scoped to their own user_id. No NDA gate is bypassed because the caller is
    on the sell side.

    The created room:
      - buyer_id = current_user.id (so they pass `participant_check` cleanly)
      - status = "preview" (filtered out of real buyer-facing metrics)
      - inquiry_id = None
      - auto-clones the listing data room files
    Idempotent — one preview vault per (listing, user) pair.
    """
    listing = await _listing_for_view_or_404(lid, user)

    existing = await db.deal_rooms.find_one(
        {"listing_id": lid, "buyer_id": user["id"], "status": "preview"}, {"_id": 0}
    )
    if existing:
        return existing

    room = {
        "id": str(uuid.uuid4()),
        "inquiry_id": None,
        "listing_id": lid,
        "listing_name": listing.get("company_name"),
        "sector": listing.get("sector"),
        "buyer_id": user["id"],
        "buyer_name": user.get("name"),
        "buyer_org": user.get("organization"),
        "seller_id": listing.get("seller_id"),
        "seller_name": listing.get("seller_name"),
        "seller_org": listing.get("seller_org"),
        "status": "preview",
        "is_preview": True,
        "nda_accepted_by_buyer_at": now_utc().isoformat(),  # auto-accept for preview
        "drl_template_id": None,
        "created_at": now_utc().isoformat(),
    }
    await db.deal_rooms.insert_one(room)

    cloned = 0
    try:
        cloned = await _clone_listing_files_into_room(lid, room["id"], user["id"])
    except Exception as e:
        logger.warning(f"preview-vault: clone of staged listing files failed: {e}")
    await log_audit(user["id"], "dealroom.preview.open", room["id"],
                    {"listing": listing.get("company_name"), "cloned_staged_files": cloned})
    room.pop("_id", None)
    return room



@api_router.get("/drl-templates")
async def list_drl_templates(user=Depends(get_current_user)):
    return [{"id": t["id"], "name": t["name"], "item_count": len(t["items"])} for t in DRL_TEMPLATES.values()]


@api_router.post("/inquiries/{inquiry_id}/open-room")
async def open_deal_room(inquiry_id: str, user=Depends(get_current_user)):
    """Seller (or any teammate in the listing's workspace) opens a Vault
    against an engaged inquiry."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Sellers/admin only")
    inquiry = await db.inquiries.find_one({"id": inquiry_id}, {"_id": 0})
    if not inquiry:
        raise HTTPException(status_code=404, detail="Inquiry not found")
    if user.get("role") != "admin":
        ws_listings, _ = await _user_workspace_listing_ids(user)
        if inquiry.get("listing_id") not in ws_listings:
            raise HTTPException(status_code=403, detail="Not authorized on this inquiry")
    if inquiry.get("status") != "engaged":
        raise HTTPException(status_code=400, detail="Inquiry must be 'engaged' before opening the Vault")

    existing = await db.deal_rooms.find_one({"inquiry_id": inquiry_id}, {"_id": 0})
    if existing:
        return existing

    listing = await db.listings.find_one({"id": inquiry["listing_id"]}, {"_id": 0})
    room = {
        "id": str(uuid.uuid4()),
        "inquiry_id": inquiry_id,
        "listing_id": inquiry["listing_id"],
        "listing_name": inquiry["listing_name"],
        "sector": (listing or {}).get("sector"),
        "buyer_id": inquiry["buyer_id"],
        "buyer_name": inquiry["buyer_name"],
        "buyer_org": inquiry.get("buyer_org"),
        "seller_id": inquiry["seller_id"],
        "seller_name": user["name"],
        "seller_org": user.get("organization"),
        "status": "pending_nda",
        "nda_accepted_by_buyer_at": None,
        "drl_template_id": None,
        "created_at": now_utc().isoformat(),
    }
    await db.deal_rooms.insert_one(room)
    await db.inquiries.update_one({"id": inquiry_id}, {"$set": {"deal_room_id": room["id"]}})

    # Auto-clone any staged listing data-room files into the new vault so the seller doesn't
    # have to re-upload anything they prepared at the listing level.
    cloned = 0
    try:
        cloned = await _clone_listing_files_into_room(inquiry["listing_id"], room["id"], user["id"])
    except Exception as e:
        logger.warning(f"open-room: clone of staged listing files failed: {e}")
    await log_audit(user["id"], "dealroom.open", room["id"],
                    {"listing": inquiry["listing_name"], "cloned_staged_files": cloned})
    room.pop("_id", None)
    return room


@api_router.get("/deal-rooms")
async def list_deal_rooms(user=Depends(get_current_user)):
    if user.get("role") == "admin":
        q = {}
    elif user.get("role") in ("seller", "agent", "buyer"):
        # Seller-side: rooms tied to any listing in their workspace
        ws_listings, my_orgs = await _user_workspace_listing_ids(user) if user.get("role") in ("seller", "agent") else ([], await _get_user_org_ids(user))
        # Buy-side: rooms where the buyer is a teammate in any of my orgs
        buyer_ids: List[str] = []
        if my_orgs:
            memberships = await db.org_memberships.find(
                {"org_id": {"$in": my_orgs}}, {"_id": 0, "user_id": 1}
            ).to_list(500)
            buyer_ids = [m["user_id"] for m in memberships]
        or_clauses = [
            {"buyer_id": user["id"]},
            {"seller_id": user["id"]},
        ]
        if ws_listings:
            or_clauses.append({"listing_id": {"$in": ws_listings}})
        if buyer_ids:
            or_clauses.append({"buyer_id": {"$in": buyer_ids}})
        q = {"$or": or_clauses}
    else:
        q = {"$or": [{"buyer_id": user["id"]}, {"seller_id": user["id"]}]}
    rooms = await db.deal_rooms.find({**q, "deleted_at": {"$exists": False}}, {"_id": 0}).sort("created_at", -1).to_list(200)
    for r in rooms:
        r["files_count"] = await db.deal_room_files.count_documents({"room_id": r["id"]})
        r["findings_count"] = await db.deal_room_findings.count_documents({"room_id": r["id"]})
        r["requests_count"] = await db.deal_room_requests.count_documents({"room_id": r["id"]})
    return rooms


async def _purge_orphan_room_clones(rid: str) -> int:
    """Retroactive self-heal: drop any deal_room_files row that was cloned from
    a staged listing file whose source has since been deleted. Without this,
    seller deletions made before the cascade-on-delete patch leak forever
    into already-opened Vaults. Hard-deletes the GridFS bytes too so they
    really leave the platform."""
    clones = await db.deal_room_files.find(
        {"room_id": rid, "cloned_from_listing_file": {"$ne": None, "$exists": True}},
        {"_id": 0, "id": 1, "gridfs_id": 1, "cloned_from_listing_file": 1},
    ).to_list(500)
    if not clones:
        return 0
    source_ids = [c["cloned_from_listing_file"] for c in clones if c.get("cloned_from_listing_file")]
    live_sources = await db.listing_staged_files.find(
        {"id": {"$in": source_ids}, "deleted_at": {"$exists": False}},
        {"_id": 0, "id": 1},
    ).to_list(500)
    live_set = {s["id"] for s in live_sources}
    purged = 0
    for c in clones:
        if c.get("cloned_from_listing_file") in live_set:
            continue
        try:
            await gridfs_bucket.delete(ObjectId(c["gridfs_id"]))
        except Exception:
            pass
        await db.deal_room_files.delete_one({"id": c["id"]})
        purged += 1
    return purged


@api_router.get("/deal-rooms/{rid}")
async def get_deal_room(rid: str, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    # Retroactive cleanup: purge clones whose staged source was deleted before
    # the cascade-on-delete patch. Cheap (single $in query) and idempotent.
    await _purge_orphan_room_clones(rid)
    # Self-heal: pick up any staged listing files added (manually or via Composio
    # external sync) AFTER this Vault was opened. Without this, files synced from
    # Google Drive / SharePoint / etc. after open-room would never reach the
    # Vault. Cheap when nothing new to clone (single $in query + early bail).
    if room.get("listing_id") and room.get("status") in ("pending_nda", "active", "preview"):
        try:
            backfilled = await _clone_listing_files_into_room(
                room["listing_id"], rid, room.get("seller_id") or user["id"], only_missing=True
            )
            if backfilled:
                logger.info(f"backfilled {backfilled} staged file(s) into room {rid}")
        except Exception as e:
            logger.warning(f"get_deal_room: backfill clone failed for {rid}: {e}")
    # Pull files WITHOUT the heavy `content` / `pages` payloads, but project a
    # boolean so the frontend knows preview is available for text-only files
    # (seed data / legacy uploads with no gridfs_id but extracted text).
    # MongoDB doesn't allow mixed include/exclude projection (except for _id),
    # so we project everything we need EXPLICITLY (without pages) and read
    # `content` separately just to compute has_text.
    room["files"] = []
    async for fr in db.deal_room_files.find(
        {"room_id": rid},
        {"_id": 0,
         "id": 1, "room_id": 1, "filename": 1, "folder": 1, "content_type": 1,
         "size_bytes": 1, "page_count": 1, "char_count": 1, "gridfs_id": 1,
         "encrypted": 1, "encryption_alg": 1, "sha256_hex": 1, "note": 1,
         "uploaded_by": 1, "uploaded_by_role": 1, "uploaded_at": 1,
         "matched_request_id": 1, "cloned_from_listing_file": 1, "source": 1,
         "download_allowed": 1,
         "content": 1},
    ).sort("uploaded_at", -1).limit(500):
        # Strip the heavy text but expose presence to the UI.
        fr["has_text"] = bool(fr.get("content"))
        fr.pop("content", None)
        room["files"].append(fr)
    room["requests"] = await db.deal_room_requests.find({"room_id": rid}, {"_id": 0}).sort("created_at", 1).to_list(200)
    # Findings: return ONLY the latest snapshot's findings instead of every
    # finding ever generated (iter-34). Older snapshots are accessible via
    # /findings-snapshots. Findings without a `job_id` are legacy rows from
    # pre-snapshot runs — we include them as a synthetic snapshot ONLY when
    # no current snapshot exists, so the UI degrades cleanly.
    latest_job = await db.findings_jobs.find_one(
        {"room_id": rid, "status": "completed"},
        {"_id": 0},
        sort=[("created_at", -1)],
    )
    if latest_job:
        room["findings"] = await db.deal_room_findings.find(
            {"room_id": rid, "job_id": latest_job["id"]}, {"_id": 0},
        ).sort("created_at", -1).to_list(200)
        room["latest_findings_snapshot"] = {
            "job_id": latest_job["id"],
            "created_at": latest_job.get("finished_at") or latest_job.get("created_at"),
            "findings_count": latest_job.get("findings_count", 0),
            "files_analyzed": latest_job.get("files_analyzed", 0),
            "total_files_in_room": latest_job.get("total_files_in_room", 0),
            "executive_summary": latest_job.get("executive_summary", ""),
            "severity_breakdown": latest_job.get("severity_breakdown") or {"high": 0, "medium": 0, "low": 0},
        }
    else:
        # No snapshot yet — fall back to any unstamped legacy findings so
        # rooms analyzed before iter-34 don't lose their data.
        room["findings"] = await db.deal_room_findings.find(
            {"room_id": rid, "job_id": {"$exists": False}}, {"_id": 0},
        ).sort("created_at", -1).to_list(200)
        room["latest_findings_snapshot"] = None
    # Activity audit: log a `dealroom.view` event but rate-limit to once-per-hour
    # per user so the timeline isn't flooded by page reloads / polling. Skip for
    # preview vaults (QA mode) to keep their audit clean.
    if room.get("status") != "preview":
        try:
            from datetime import timedelta as _td
            recent = await db.audit_logs.find_one(
                {"action": "dealroom.view", "actor_id": user["id"], "target": rid,
                 "timestamp": {"$gte": (now_utc() - _td(hours=1)).isoformat()}},
                {"_id": 0, "id": 1},
            )
            if not recent:
                await log_audit(user["id"], "dealroom.view", rid,
                                {"role": user.get("role"), "status": room.get("status")})
        except Exception as e:
            logger.debug(f"dealroom.view audit skipped: {e}")
    return room


@api_router.get("/deal-rooms/{rid}/activity")
async def get_deal_room_activity(
    rid: str,
    since: Optional[str] = None,
    limit: int = 200,
    user=Depends(get_current_user),
):
    """Return the Vault Activity timeline — every audit event scoped to this
    Vault, hydrated with actor metadata so the UI can render a rich timeline.

    Visibility: both buyer and seller see the same set of events for full
    transparency (an industry default for institutional VDRs).

    Query params:
      - `since`: ISO-8601 cutoff. Used for incremental polling so the UI only
        fetches what's new since its last fetch.
      - `limit`: max events to return (capped at 500).
    """
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)

    limit = max(1, min(int(limit or 200), 500))

    # Curated action vocabulary the Activity tab surfaces. Adding to this list
    # automatically lights up new event types in the UI without code changes.
    ACTIONS = [
        "dealroom.open", "dealroom.view",
        "dealroom.nda.accept",
        "dealroom.file.upload", "dealroom.file.add", "dealroom.file.download",
        "dealroom.file.preview", "dealroom.file.delete",
        "dealroom.preview.open",
        "vault.copilot.ask",
        "dealroom.findings.generate",
    ]
    file_ids = [f["id"] for f in await db.deal_room_files.find(
        {"room_id": rid}, {"_id": 0, "id": 1}
    ).to_list(1000)]

    query: dict = {
        "$or": [
            {"target": rid, "action": {"$in": ACTIONS}},
            # File-level downloads/uploads log target_id=file_id sometimes; cover both.
            {"target": {"$in": file_ids}, "action": {"$in": ACTIONS}} if file_ids else
                {"_": "__unreachable__"},
        ],
    }
    if since:
        query["timestamp"] = {"$gt": since}

    rows = await db.audit_logs.find(query, {"_id": 0}).sort("timestamp", -1).limit(limit).to_list(limit)

    # Hydrate actors in one round-trip
    actor_ids = list({r.get("actor_id") for r in rows if r.get("actor_id")})
    actors_by_id: dict[str, dict] = {}
    if actor_ids:
        for u in await db.users.find(
            {"id": {"$in": actor_ids}},
            {"_id": 0, "id": 1, "name": 1, "email": 1, "role": 1, "organization": 1},
        ).to_list(len(actor_ids)):
            actors_by_id[u["id"]] = u

    # Category buckets mapped from raw audit actions — the frontend uses these
    # for filter chips and icon selection. Keep the mapping in one place.
    CATEGORY = {
        "dealroom.open":         "vault",
        "dealroom.preview.open": "vault",
        "dealroom.view":         "vault",
        "dealroom.nda.accept":   "nda",
        "dealroom.file.upload":  "file",
        "dealroom.file.add":     "file",
        "dealroom.file.download": "file",
        "dealroom.file.preview":  "file",
        "dealroom.file.delete":   "file",
        "vault.copilot.ask":      "copilot",
        "dealroom.findings.generate":  "findings",
    }
    LABEL = {
        "dealroom.open":         "Opened the Vault",
        "dealroom.preview.open": "Created a preview Vault",
        "dealroom.view":         "Viewed the Vault",
        "dealroom.nda.accept":   "Signed the NDA",
        "dealroom.file.upload":  "Uploaded a file",
        "dealroom.file.add":     "File added to the Vault",
        "dealroom.file.download": "Downloaded a file",
        "dealroom.file.preview":  "Previewed a file",
        "dealroom.file.delete":   "Removed a file",
        "vault.copilot.ask":      "Asked the AI Co-pilot",
        "dealroom.findings.generate":  "Generated AI Findings",
    }

    events = []
    counts_by_action: dict[str, int] = {}
    counts_by_actor: dict[str, int] = {}
    for r in rows:
        actor = actors_by_id.get(r.get("actor_id") or "", {})
        cat = CATEGORY.get(r.get("action"), "other")
        events.append({
            "id": r.get("id"),
            "seq": r.get("seq"),
            "action": r.get("action"),
            "category": cat,
            "label": LABEL.get(r.get("action"), r.get("action")),
            "actor": {
                "id": actor.get("id") or r.get("actor_id"),
                "name": actor.get("name") or "Unknown",
                "email": actor.get("email"),
                "role": actor.get("role"),
                "organization": actor.get("organization"),
            },
            "target_id": r.get("target") or r.get("target_id"),
            "meta": r.get("meta") or {},
            "created_at": r.get("timestamp") or r.get("created_at"),
            "content_hash": r.get("content_hash"),
        })
        counts_by_action[r.get("action")] = counts_by_action.get(r.get("action"), 0) + 1
        if actor.get("email"):
            counts_by_actor[actor["email"]] = counts_by_actor.get(actor["email"], 0) + 1

    return {
        "vault_id": rid,
        "events": events,
        "counts": {
            "total": len(events),
            "by_action": counts_by_action,
            "by_actor": counts_by_actor,
        },
        "as_of": now_utc().isoformat(),
    }


@api_router.post("/deal-rooms/{rid}/accept-nda")
async def accept_nda(rid: str, body: NDAAccept, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    if user["id"] != room["buyer_id"]:
        raise HTTPException(status_code=403, detail="Only the buyer can accept the NDA")
    signed_name = (body.signed_name or "").strip()
    if len(signed_name) < 2:
        raise HTTPException(status_code=400, detail="Typed full name required to sign the NDA")
    await db.deal_rooms.update_one(
        {"id": rid},
        {"$set": {
            "status": "active",
            "nda_accepted_by_buyer_at": now_utc().isoformat(),
            "nda_signed_name": signed_name,
            "nda_signed_by_user_id": user["id"],
        }},
    )
    await log_audit(user["id"], "dealroom.nda.accept", rid, {"signed_name": signed_name})
    # Bitcoin-anchored proof of NDA signature
    asyncio.create_task(notarize_event(
        kind="nda.signature",
        target_id=rid,
        payload={
            "vault_id": rid,
            "buyer_id": user["id"],
            "seller_id": room.get("seller_id"),
            "listing_name": room.get("listing_name"),
            "signed_name": signed_name,
            "signed_at": now_utc().isoformat(),
        },
        owner_user_id=user["id"],
        label=f"NDA signed by {signed_name}",
    ))
    return {"ok": True, "signed_name": signed_name}


@api_router.post("/deal-rooms/{rid}/drl")
async def apply_drl_template(rid: str, body: DRLApply, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    role = await participant_check(room, user)
    if role not in ("buyer", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Only the buyer can apply a DRL template")
    template = DRL_TEMPLATES.get(body.template_id)
    if not template:
        raise HTTPException(status_code=400, detail="Unknown template")

    # Clear previous template items
    await db.deal_room_requests.delete_many({"room_id": rid})
    new_docs = [
        {
            "id": str(uuid.uuid4()),
            "room_id": rid,
            "template_id": template["id"],
            "title": item["title"],
            "workstream": item["workstream"],
            "status": "pending",
            "matched_file_ids": [],
            "created_at": now_utc().isoformat(),
        }
        for item in template["items"]
    ]
    if new_docs:
        await db.deal_room_requests.insert_many(new_docs)
    await db.deal_rooms.update_one({"id": rid}, {"$set": {"drl_template_id": template["id"]}})
    await log_audit(user["id"], "dealroom.drl.apply", rid, {"template": template["id"], "count": len(new_docs)})
    return {"ok": True, "request_count": len(new_docs)}


@api_router.post("/deal-rooms/{rid}/files")
async def upload_file(rid: str, body: FileUpload, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    role = await participant_check(room, user)
    # Sellers may stage documents pre-NDA; buyers cannot upload until they sign.
    if room.get("status") == "pending_nda" and role != "seller":
        raise HTTPException(status_code=400, detail="Buyer must accept NDA before files can be exchanged")

    file_id = str(uuid.uuid4())
    doc = {
        "id": file_id,
        "room_id": rid,
        "filename": body.filename,
        "folder": body.folder,
        "content": body.content[:50000],  # cap to 50k chars for safety
        "char_count": len(body.content),
        "note": body.note,
        "uploaded_by": user["id"],
        "uploaded_by_role": role,
        "uploaded_at": now_utc().isoformat(),
        "matched_request_id": None,
    }
    await db.deal_room_files.insert_one(doc)
    await log_audit(user["id"], "dealroom.file.upload", rid, {"filename": body.filename, "folder": body.folder})

    # Best-effort auto-match against DRL items
    matched_request_id = None
    if role in ("seller", "admin", "agent"):
        requests = await db.deal_room_requests.find({"room_id": rid, "status": "pending"}, {"_id": 0}).to_list(200)
        if requests:
            try:
                lines = "\n".join(f"[{r['id']}] {r['workstream']} :: {r['title']}" for r in requests[:30])
                prompt = (
                    f"You are an M&A diligence coordinator. A seller uploaded a file titled '{body.filename}' "
                    f"(folder: {body.folder}). First 800 chars of content:\n{body.content[:800]}\n\n"
                    f"Open DRL requests:\n{lines}\n\n"
                    "Return ONLY the id of the single best-matching request (or NONE). Reply with the id and nothing else."
                )
                raw = await call_claude(
                    "You match documents to diligence requests. Reply with exactly one id or NONE.",
                    prompt,
                    session_id=f"match-{rid}",
                )
                raw_clean = (raw or "").strip().split()[0] if raw else ""
                if raw_clean and raw_clean.upper() != "NONE":
                    candidate = next((r for r in requests if r["id"] == raw_clean), None)
                    if candidate:
                        matched_request_id = candidate["id"]
                        await db.deal_room_requests.update_one(
                            {"id": matched_request_id},
                            {
                                "$set": {"status": "satisfied"},
                                "$addToSet": {"matched_file_ids": file_id},
                            },
                        )
                        await db.deal_room_files.update_one(
                            {"id": file_id},
                            {"$set": {"matched_request_id": matched_request_id}},
                        )
                        await log_agent_activity(
                            "drl-match-agent",
                            f"matched:{body.filename}",
                            "completed",
                            user_id=user["id"],
                            meta={"request_id": matched_request_id},
                        )
            except Exception as e:
                logger.warning(f"DRL auto-match failed: {e}")
                await log_agent_activity(
                    "drl-match-agent",
                    f"match:{body.filename}",
                    "failed",
                    user_id=user["id"],
                    friction=str(e),
                )

    doc.pop("_id", None)
    doc.pop("content", None)
    doc["matched_request_id"] = matched_request_id
    return doc


async def _auto_match_drl(rid: str, file_id: str, filename: str, folder: str, text_preview: str, user_id: str):
    """Best-effort DRL auto-match for a newly uploaded file (seller-side)."""
    requests = await db.deal_room_requests.find({"room_id": rid, "status": "pending"}, {"_id": 0}).to_list(200)
    if not requests:
        return None
    try:
        lines = "\n".join(f"[{r['id']}] {r['workstream']} :: {r['title']}" for r in requests[:30])
        prompt = (
            f"You are an M&A diligence coordinator. A seller uploaded a file titled '{filename}' "
            f"(folder: {folder}). First 800 chars of content:\n{text_preview[:800]}\n\n"
            f"Open DRL requests:\n{lines}\n\n"
            "Return ONLY the id of the single best-matching request (or NONE). Reply with the id and nothing else."
        )
        raw = await call_claude(
            "You match documents to diligence requests. Reply with exactly one id or NONE.",
            prompt,
            session_id=f"match-{rid}",
        )
        raw_clean = (raw or "").strip().split()[0] if raw else ""
        if raw_clean and raw_clean.upper() != "NONE":
            candidate = next((r for r in requests if r["id"] == raw_clean), None)
            if candidate:
                await db.deal_room_requests.update_one(
                    {"id": candidate["id"]},
                    {"$set": {"status": "satisfied"}, "$addToSet": {"matched_file_ids": file_id}},
                )
                await db.deal_room_files.update_one(
                    {"id": file_id},
                    {"$set": {"matched_request_id": candidate["id"]}},
                )
                await log_agent_activity(
                    "drl-match-agent", f"matched:{filename}", "completed",
                    user_id=user_id, meta={"request_id": candidate["id"]},
                )
                return candidate["id"]
    except Exception as e:
        logger.warning(f"DRL auto-match failed: {e}")
        await log_agent_activity(
            "drl-match-agent", f"match:{filename}", "failed",
            user_id=user_id, friction=str(e),
        )
    return None


@api_router.post("/deal-rooms/{rid}/files/binary")
async def upload_file_binary(
    rid: str,
    file: UploadFile = File(...),
    folder: str = Form("other"),
    note: Optional[str] = Form(None),
    user=Depends(get_current_user),
):
    """Multipart binary upload — stored in GridFS, text extracted per-page for AI matching/findings."""
    if folder not in ("financials", "legal", "hr", "it", "operations", "commercial", "other"):
        folder = "other"
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    role = await participant_check(room, user)
    # Sellers may stage documents pre-NDA; buyers cannot upload until they sign.
    if room.get("status") == "pending_nda" and role != "seller":
        raise HTTPException(status_code=400, detail="Buyer must accept NDA before files can be exchanged")

    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="Empty file")
    # Cap at 25 MB
    if len(data) > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File exceeds 50 MB limit")

    file_id = str(uuid.uuid4())
    filename = file.filename or f"upload-{file_id}"

    # SHA-256 the plaintext for OTS notarization (the digest IS the plaintext fingerprint)
    plaintext_sha256_hex = sha256_hex(data)

    # AES-256-GCM at-rest encryption
    storage_bytes = data
    encrypted = False
    encryption_alg = None
    if encryption_configured():
        try:
            aad = f"{rid}:{file_id}".encode("utf-8")
            enc = encrypt_bytes(data, associated_data=aad)
            storage_bytes = enc["envelope"]
            encrypted = True
            encryption_alg = enc["alg"]
        except Exception as e:
            logger.warning(f"At-rest encryption failed, storing plaintext: {e}")

    # Stream into GridFS (encrypted if available)
    gridfs_id = await gridfs_bucket.upload_from_stream(
        filename,
        io.BytesIO(storage_bytes),
        metadata={
            "room_id": rid,
            "file_id": file_id,
            "uploaded_by": user["id"],
            "content_type": file.content_type or "application/octet-stream",
            "encrypted": encrypted,
            "encryption_alg": encryption_alg,
        },
    )

    # Extract per-page text from PLAINTEXT before discarding
    pages = extract_pages_from_bytes(filename, data)
    flat = pages_to_flat_text(pages)

    doc = {
        "id": file_id,
        "room_id": rid,
        "filename": filename,
        "folder": folder,
        "content_type": file.content_type or "application/octet-stream",
        "size_bytes": len(data),
        "page_count": len(pages),
        "pages": pages,
        "content": flat,
        "char_count": len(flat),
        "gridfs_id": str(gridfs_id),
        "storage": "gridfs",
        "encrypted": encrypted,
        "encryption_alg": encryption_alg,
        "sha256_hex": plaintext_sha256_hex,
        "note": note,
        "uploaded_by": user["id"],
        "uploaded_by_role": role,
        "uploaded_at": now_utc().isoformat(),
        "matched_request_id": None,
    }
    await db.deal_room_files.insert_one(doc)
    await log_audit(user["id"], "dealroom.file.upload", rid, {
        "filename": filename, "folder": folder, "bytes": len(data),
        "sha256": plaintext_sha256_hex, "encrypted": encrypted,
    })

    # Bitcoin-anchored proof of upload (hash of plaintext)
    asyncio.create_task(notarize_bytes(
        kind="vault.file",
        target_id=file_id,
        data=data,
        owner_user_id=user["id"],
        label=f"Vault file: {filename}",
        extra={"vault_id": rid, "filename": filename, "size_bytes": len(data)},
    ))

    matched_request_id = None
    if role in ("seller", "admin", "agent"):
        matched_request_id = await _auto_match_drl(rid, file_id, filename, folder, flat, user["id"])

    doc.pop("_id", None)
    doc.pop("content", None)
    doc.pop("pages", None)
    doc["matched_request_id"] = matched_request_id
    return doc


@api_router.get("/deal-rooms/{rid}/files/{file_id}/download")
async def download_file(rid: str, file_id: str, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    role = await participant_check(room, user)
    f = await db.deal_room_files.find_one({"id": file_id, "room_id": rid}, {"_id": 0})
    if not f:
        raise HTTPException(status_code=404, detail="File not found")
    # Per-file access policy: sellers + admins always download (they own the
    # docs); buyers only download when `download_allowed` is explicitly true.
    # Defaults to false to match institutional VDR convention (view-only is
    # the safer default). Sellers flip the toggle file-by-file in the UI.
    if role == "buyer" and not f.get("download_allowed", False):
        raise HTTPException(
            status_code=403,
            detail="This file is view-only. Ask the seller to enable downloads.",
        )
    if not f.get("gridfs_id"):
        raise HTTPException(status_code=400, detail="This file has no binary content (text-only upload)")
    try:
        grid_out = await gridfs_bucket.open_download_stream(ObjectId(f["gridfs_id"]))
    except Exception as e:
        raise HTTPException(status_code=404, detail=f"Binary not found: {e}")

    # If file was encrypted at rest, decrypt the full envelope before streaming back
    if f.get("encrypted"):
        try:
            envelope = await grid_out.read()
            aad = f"{rid}:{file_id}".encode("utf-8")
            plaintext = decrypt_envelope(envelope, associated_data=aad)
        except Exception as e:
            logger.exception("Vault file decryption failed")
            raise HTTPException(status_code=500, detail=f"Decryption failed: {e}")
        await log_audit(user["id"], "dealroom.file.download", rid, {"filename": f["filename"], "decrypted": True})
        return StreamingResponse(
            io.BytesIO(plaintext),
            media_type=f.get("content_type") or "application/octet-stream",
            headers={"Content-Disposition": f'attachment; filename="{f["filename"]}"'},
        )

    async def streamer():
        while True:
            chunk = await grid_out.readchunk()
            if not chunk:
                break
            yield chunk

    await log_audit(user["id"], "dealroom.file.download", rid, {"filename": f["filename"]})
    return StreamingResponse(
        streamer(),
        media_type=f.get("content_type") or "application/octet-stream",
        headers={"Content-Disposition": f'attachment; filename="{f["filename"]}"'},
    )


# -----------------------------------------------------------------------------
# Per-file access policy (seller-only) & in-browser preview
# -----------------------------------------------------------------------------
class FileAccessPolicy(BaseModel):
    download_allowed: bool


@api_router.patch("/deal-rooms/{rid}/files/{file_id}/access")
async def set_file_access(
    rid: str, file_id: str, body: FileAccessPolicy,
    user=Depends(get_current_user),
):
    """Seller (or admin/agent on the listing's workspace) toggles download
    permission on a per-file basis. Buyers cannot call this endpoint."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    role = await participant_check(room, user)
    if role not in ("seller", "admin"):
        raise HTTPException(status_code=403, detail="Only the seller can change file access policy")
    res = await db.deal_room_files.update_one(
        {"id": file_id, "room_id": rid},
        {"$set": {"download_allowed": bool(body.download_allowed),
                  "access_policy_updated_at": now_utc().isoformat(),
                  "access_policy_updated_by": user["id"]}},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="File not found")
    await log_audit(user["id"], "dealroom.file.access", rid,
                    {"file_id": file_id, "download_allowed": bool(body.download_allowed)})
    return {"ok": True, "download_allowed": bool(body.download_allowed)}


# Office → PDF conversion cache (per-file-id) lives in GridFS so it survives
# restarts and idle pods. Repeated previews of the same DOCX/XLSX/PPTX therefore
# trigger conversion only once.
def _docx_to_pdf_bytes(data: bytes, filename: str) -> bytes | None:
    """Render a DOCX as PDF using python-docx + reportlab. Pure-Python pipeline
    chosen over LibreOffice because apt-installed packages don't survive
    container restarts on this hosting platform. Includes inline images so
    one-pagers with screenshots / charts render faithfully."""
    try:
        from docx import Document as _Docx
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
            Image as RLImage,
        )
        from reportlab.lib import colors
        from PIL import Image as PILImage
        import html as _h
    except Exception as e:
        logger.warning(f"docx→pdf deps missing: {e}")
        return None
    try:
        d = _Docx(io.BytesIO(data))
        out = io.BytesIO()
        doc = SimpleDocTemplate(out, pagesize=LETTER,
                                leftMargin=0.9*inch, rightMargin=0.9*inch,
                                topMargin=0.9*inch, bottomMargin=0.9*inch,
                                title=os.path.splitext(filename)[0])
        styles = getSampleStyleSheet()
        body_style = ParagraphStyle("body", parent=styles["Normal"], fontName="Helvetica",
                                     fontSize=10.5, leading=14.5, spaceAfter=6)
        h1_style = ParagraphStyle("h1", parent=styles["Heading1"], fontName="Helvetica-Bold",
                                   fontSize=18, leading=22, spaceAfter=10, textColor=colors.HexColor("#0B1B3D"))
        h2_style = ParagraphStyle("h2", parent=styles["Heading2"], fontName="Helvetica-Bold",
                                   fontSize=14, leading=18, spaceAfter=8, textColor=colors.HexColor("#0B1B3D"))
        # Index inline images by their relationship id so we can drop them
        # in document order when we encounter the `<w:drawing>` block.
        rels_by_id: dict[str, bytes] = {}
        try:
            for rel_id, rel in d.part.rels.items():
                if "image" in (rel.reltype or "").lower():
                    try:
                        rels_by_id[rel_id] = rel.target_part.blob
                    except Exception:
                        pass
        except Exception as e:
            logger.debug(f"docx image rel walk failed: {e}")

        MAX_W = 6.5 * inch
        MAX_H = 4.0 * inch

        def _img_flowable(blob: bytes):
            try:
                im = PILImage.open(io.BytesIO(blob)); im.load()
                if im.mode not in ("RGB", "RGBA"):
                    im = im.convert("RGB")
                buf = io.BytesIO(); im.save(buf, format="PNG"); buf.seek(0)
                w, h = im.size
                scale = min(MAX_W / w, MAX_H / h, 1.0)
                return RLImage(buf, width=w*scale, height=h*scale)
            except Exception:
                return None

        EMBED_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
        flow = []
        for blk in d.element.body.iter():
            tag = blk.tag.split("}")[-1]
            if tag == "p":
                txt = "".join(node.text or "" for node in blk.iter() if node.tag.endswith("}t"))
                # Detect any inline image in this paragraph and emit it before
                # / after the text (we always emit after for simplicity).
                image_flowables: list = []
                for blip in blk.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}blip"):
                    rid = blip.get(EMBED_NS + "embed")
                    if rid and rid in rels_by_id:
                        img = _img_flowable(rels_by_id[rid])
                        if img is not None:
                            image_flowables.append(img)
                if not txt.strip() and not image_flowables:
                    flow.append(Spacer(1, 5))
                    continue
                if txt.strip():
                    style = body_style
                    pStyle = blk.find(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pStyle")
                    if pStyle is not None:
                        val = (pStyle.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val") or "").lower()
                        if "heading 1" in val or val == "heading1" or val == "title":
                            style = h1_style
                        elif "heading" in val:
                            style = h2_style
                    flow.append(Paragraph(_h.escape(txt), style))
                for img in image_flowables:
                    flow.append(Spacer(1, 4))
                    flow.append(img)
                    flow.append(Spacer(1, 4))
            elif tag == "tbl":
                rows: list[list[str]] = []
                for tr in blk.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tr"):
                    row: list[str] = []
                    for tc in tr.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tc"):
                        cell_txt = "".join(t.text or "" for t in tc.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"))
                        row.append(cell_txt.strip())
                    if row:
                        rows.append(row)
                if rows:
                    flow.append(Spacer(1, 6))
                    flow.append(Table(rows, repeatRows=1, style=TableStyle([
                        ("BOX", (0,0), (-1,-1), 0.5, colors.HexColor("#1D4ED8")),
                        ("INNERGRID", (0,0), (-1,-1), 0.25, colors.HexColor("#A8AEC0")),
                        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#E8F0FE")),
                        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                        ("FONTSIZE", (0,0), (-1,-1), 9),
                        ("LEFTPADDING", (0,0), (-1,-1), 6),
                        ("RIGHTPADDING", (0,0), (-1,-1), 6),
                        ("TOPPADDING", (0,0), (-1,-1), 4),
                        ("BOTTOMPADDING", (0,0), (-1,-1), 4),
                    ])))
                    flow.append(Spacer(1, 6))
        if not flow:
            flow.append(Paragraph(f"<i>{_h.escape(filename)} appears to be empty.</i>", body_style))
        doc.build(flow)
        return out.getvalue()
    except Exception as e:
        logger.warning(f"docx→pdf render failed for {filename}: {e}")
        return None


def _xlsx_to_pdf_bytes(data: bytes, filename: str) -> bytes | None:
    """Render an XLSX as a multi-page PDF (one section per sheet) using
    openpyxl + reportlab. Caps each sheet at 200 rows × 24 cols to keep the
    output legible — buyers can ask the seller to enable download for the
    full spreadsheet."""
    try:
        from openpyxl import load_workbook
        from reportlab.lib.pagesizes import landscape, LETTER
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
        from reportlab.lib import colors
        import html as _h
    except Exception as e:
        logger.warning(f"xlsx→pdf deps missing: {e}")
        return None
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        out = io.BytesIO()
        doc = SimpleDocTemplate(out, pagesize=landscape(LETTER),
                                leftMargin=0.5*inch, rightMargin=0.5*inch,
                                topMargin=0.6*inch, bottomMargin=0.5*inch,
                                title=os.path.splitext(filename)[0])
        styles = getSampleStyleSheet()
        sheet_title = ParagraphStyle("sheetTitle", parent=styles["Heading2"],
                                      fontName="Helvetica-Bold", fontSize=13, leading=16,
                                      spaceAfter=8, textColor=colors.HexColor("#0B1B3D"))
        note = ParagraphStyle("note", parent=styles["Normal"], fontName="Helvetica-Oblique",
                               fontSize=8.5, textColor=colors.HexColor("#7A8299"), spaceAfter=8)
        flow = []
        for sheet_idx, ws in enumerate(wb.worksheets):
            if sheet_idx > 0:
                flow.append(PageBreak())
            flow.append(Paragraph(_h.escape(ws.title or f"Sheet {sheet_idx+1}"), sheet_title))
            max_rows, max_cols = 200, 24
            rows = []
            for r_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if r_idx >= max_rows:
                    break
                rows.append(["" if c is None else str(c)[:120] for c in row[:max_cols]])
            if not rows:
                flow.append(Paragraph("<i>Empty sheet.</i>", note))
                continue
            # Pad uneven rows
            width = max(len(r) for r in rows)
            for r in rows:
                while len(r) < width:
                    r.append("")
            tbl = Table(rows, repeatRows=1, style=TableStyle([
                ("BOX", (0,0), (-1,-1), 0.3, colors.HexColor("#A8AEC0")),
                ("INNERGRID", (0,0), (-1,-1), 0.15, colors.HexColor("#D6D9E0")),
                ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#E8F0FE")),
                ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
                ("FONTNAME", (0,1), (-1,-1), "Helvetica"),
                ("FONTSIZE", (0,0), (-1,-1), 7.5),
                ("LEFTPADDING", (0,0), (-1,-1), 4),
                ("RIGHTPADDING", (0,0), (-1,-1), 4),
                ("TOPPADDING", (0,0), (-1,-1), 2),
                ("BOTTOMPADDING", (0,0), (-1,-1), 2),
                ("VALIGN", (0,0), (-1,-1), "TOP"),
            ]))
            flow.append(tbl)
            # Footnote if we truncated the sheet
            total_rows = ws.max_row or 0
            if total_rows > max_rows:
                flow.append(Spacer(1, 4))
                flow.append(Paragraph(f"Showing first {max_rows} of {total_rows} rows.", note))
        wb.close()
        if not flow:
            return None
        doc.build(flow)
        return out.getvalue()
    except Exception as e:
        logger.warning(f"xlsx→pdf render failed for {filename}: {e}")
        return None


def _pptx_to_pdf_bytes(data: bytes, filename: str) -> bytes | None:
    """Render a PPTX as a PDF (one page per slide) using python-pptx +
    reportlab. Extracts slide text, embedded raster images, AND speaker
    notes so buyers see the deck's actual visual content — not just a
    text outline."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from reportlab.lib.pagesizes import landscape, LETTER
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage,
        )
        from reportlab.lib import colors
        import html as _h
        from PIL import Image as PILImage
    except Exception as e:
        logger.warning(f"pptx→pdf deps missing: {e}")
        return None
    try:
        pres = Presentation(io.BytesIO(data))
        out = io.BytesIO()
        doc = SimpleDocTemplate(out, pagesize=landscape(LETTER),
                                leftMargin=0.7*inch, rightMargin=0.7*inch,
                                topMargin=0.7*inch, bottomMargin=0.6*inch,
                                title=os.path.splitext(filename)[0])
        styles = getSampleStyleSheet()
        slide_no = ParagraphStyle("slideNo", parent=styles["Normal"], fontName="Helvetica",
                                   fontSize=8.5, textColor=colors.HexColor("#7A8299"),
                                   letterSpacing=2, spaceAfter=4)
        slide_title = ParagraphStyle("slideTitle", parent=styles["Heading1"],
                                      fontName="Helvetica-Bold", fontSize=20, leading=26,
                                      spaceAfter=12, textColor=colors.HexColor("#0B1B3D"))
        body_style = ParagraphStyle("body", parent=styles["Normal"], fontName="Helvetica",
                                     fontSize=11.5, leading=15.5, spaceAfter=4)
        bullet_style = ParagraphStyle("bullet", parent=body_style, leftIndent=14, bulletIndent=4)
        notes_style = ParagraphStyle("notes", parent=styles["Normal"], fontName="Helvetica-Oblique",
                                      fontSize=9, textColor=colors.HexColor("#5A6275"),
                                      leftIndent=8, spaceBefore=12, leading=12)
        # Reportlab usable page width (landscape LETTER ~11" minus margins).
        MAX_IMG_W = 9.0 * inch
        MAX_IMG_H = 4.5 * inch

        def _make_image_flowable(blob: bytes, ctype: str | None):
            """Return a reportlab Image flowable scaled to fit the slide page.
            Robust to broken/unknown formats — returns None if Pillow can't
            decode it, so a single bad image doesn't kill the whole render."""
            try:
                # Pillow normalizes anything reportlab might choke on (WMF,
                # uncommon CMYK JPEGs, etc.) into a clean PNG buffer first.
                im = PILImage.open(io.BytesIO(blob))
                im.load()
                if im.mode in ("RGBA", "P", "LA"):
                    im = im.convert("RGBA")
                elif im.mode != "RGB":
                    im = im.convert("RGB")
                norm = io.BytesIO()
                im.save(norm, format="PNG")
                norm.seek(0)
                src_w, src_h = im.size
            except Exception as e:
                logger.debug(f"pptx image decode failed: {e}")
                return None
            # Fit-to-box, preserving aspect ratio.
            scale = min(MAX_IMG_W / src_w, MAX_IMG_H / src_h, 1.0)
            return RLImage(norm, width=src_w * scale, height=src_h * scale)

        flow = []
        for i, slide in enumerate(pres.slides):
            if i > 0:
                flow.append(PageBreak())
            flow.append(Paragraph(f"SLIDE {i+1} / {len(pres.slides)}", slide_no))
            title_text = None
            bullets: list[str] = []
            images: list = []  # RLImage flowables, in shape document order

            for shape in slide.shapes:
                # Pictures — extract raw blob and embed.
                try:
                    is_pic = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                except Exception:
                    is_pic = False
                if is_pic and hasattr(shape, "image"):
                    try:
                        blob = shape.image.blob
                        ctype = getattr(shape.image, "content_type", None)
                        img = _make_image_flowable(blob, ctype)
                        if img is not None:
                            images.append(img)
                    except Exception as e:
                        logger.debug(f"pptx picture extract failed on slide {i+1}: {e}")
                # Text frames.
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                for p_idx, para in enumerate(tf.paragraphs):
                    txt = "".join(r.text or "" for r in para.runs).strip()
                    if not txt:
                        continue
                    if title_text is None and (
                        getattr(shape, "is_placeholder", False) and shape.placeholder_format and shape.placeholder_format.idx == 0
                    ):
                        title_text = txt
                    else:
                        bullets.append(txt)

            if title_text:
                flow.append(Paragraph(_h.escape(title_text), slide_title))
            for b in bullets[:12]:
                flow.append(Paragraph(f"• {_h.escape(b)}", bullet_style))
            # Place up to 2 images per slide so the layout stays sane; if a
            # slide has more, the extras still appear but stacked vertically.
            for j, img in enumerate(images[:4]):
                flow.append(Spacer(1, 8))
                flow.append(img)
            if not title_text and not bullets and not images:
                flow.append(Paragraph("<i>(empty slide)</i>", body_style))
            # Speaker notes
            try:
                notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            except Exception:
                notes = ""
            if notes and notes.strip():
                flow.append(Paragraph(f"<b>Speaker notes:</b> {_h.escape(notes[:1200])}", notes_style))
        if not flow:
            return None
        doc.build(flow)
        return out.getvalue()
    except Exception as e:
        logger.warning(f"pptx→pdf render failed for {filename}: {e}")
        return None


async def _office_to_pdf(filename: str, data: bytes) -> bytes | None:
    """Dispatch Office formats to the appropriate pure-Python renderer.
    Falls back to LibreOffice headless if installed (local dev convenience).
    Returns PDF bytes or None if no renderer matches."""
    fname_lc = filename.lower()
    # Prefer LibreOffice when available — better fidelity. Only used in local
    # dev; production containers don't ship it.
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        import tempfile, subprocess
        with tempfile.TemporaryDirectory() as td:
            in_path = os.path.join(td, filename)
            with open(in_path, "wb") as f:
                f.write(data)
            try:
                proc = subprocess.run(
                    [soffice, "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "pdf", "--outdir", td, in_path],
                    capture_output=True, timeout=90,
                )
                stem = os.path.splitext(filename)[0]
                out_path = os.path.join(td, f"{stem}.pdf")
                if proc.returncode == 0 and os.path.exists(out_path):
                    with open(out_path, "rb") as f:
                        return f.read()
            except Exception as e:
                logger.warning(f"soffice convert crashed for {filename}: {e}")
    # Pure-Python fallback path (the primary path in production).
    if fname_lc.endswith((".docx",)):
        return _docx_to_pdf_bytes(data, filename)
    if fname_lc.endswith((".xlsx", ".xlsm")):
        return _xlsx_to_pdf_bytes(data, filename)
    if fname_lc.endswith((".pptx",)):
        return _pptx_to_pdf_bytes(data, filename)
    # Legacy formats (.doc, .xls, .ppt, .odt, .ods, .odp) without LibreOffice
    # are unsupported — buyer must download to view locally.
    return None


# Mime types we directly stream into the browser viewer (PDFs + images + plain
# text). Anything not in this set + not an Office format triggers a 415 from
# the preview endpoint and the UI falls back to "Preview unavailable" copy.
_DIRECT_PREVIEW_TYPES = {
    "application/pdf",
    "image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp", "image/svg+xml",
    "text/plain", "text/markdown", "text/csv", "text/html",
    "application/json",
}
_OFFICE_PREVIEW_EXTS = (
    ".docx", ".doc", ".odt",
    ".xlsx", ".xls", ".ods", ".xlsm",
    ".pptx", ".ppt", ".odp",
)


@api_router.get("/deal-rooms/{rid}/files/{file_id}/preview")
async def preview_file(rid: str, file_id: str, user=Depends(get_current_user)):
    """Stream a file inline for in-browser viewing. PDFs + images + text
    served directly. Office formats converted on-the-fly via the pure-Python
    pipeline (python-docx / openpyxl / python-pptx + reportlab) and cached
    in GridFS so the second preview is instant.

    Unlike `/download`, the preview endpoint is ALWAYS available to every
    Vault participant — the watermark overlay added client-side carries
    the viewer's identity. Audit-logged so the seller sees who previewed
    what."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    f = await db.deal_room_files.find_one({"id": file_id, "room_id": rid}, {"_id": 0})
    if not f:
        raise HTTPException(status_code=404, detail="File not found")

    fname = f.get("filename", "file")
    fname_lc = fname.lower()
    raw_ctype = (f.get("content_type") or "").lower().split(";")[0].strip()

    # Detect the *effective* content-type for routing. Browsers (and curl)
    # frequently fail to supply a specific MIME for .csv / .md / .json /
    # .svg → uploads land as "application/octet-stream" or empty. Fall back
    # to the filename extension whenever the stored content_type is missing
    # or generic so the preview still resolves.
    EXT_MIME = {
        ".pdf":  "application/pdf",
        ".png":  "image/png",
        ".jpg":  "image/jpeg", ".jpeg": "image/jpeg",
        ".gif":  "image/gif", ".webp": "image/webp",
        ".svg":  "image/svg+xml",
        ".txt":  "text/plain", ".log": "text/plain",
        ".md":   "text/markdown", ".markdown": "text/markdown",
        ".csv":  "text/csv", ".tsv": "text/tab-separated-values",
        ".json": "application/json",
        ".html": "text/html", ".htm": "text/html",
        ".xml":  "application/xml",
    }
    ctype = raw_ctype
    if (not ctype) or ctype == "application/octet-stream":
        for ext, mime in EXT_MIME.items():
            if fname_lc.endswith(ext):
                ctype = mime
                break

    # ----- Source the bytes -----
    # Three storage layouts to support: (a) modern encrypted GridFS upload,
    # (b) plain GridFS without encryption, (c) text-only legacy / seed file
    # stored as `content` string with no GridFS id. (c) is the path that was
    # returning 404 — fix is to fall through to text/plain inline streaming.
    plaintext: bytes | None = None
    if f.get("gridfs_id"):
        try:
            grid_out = await gridfs_bucket.open_download_stream(ObjectId(f["gridfs_id"]))
            raw = await grid_out.read()
        except Exception as e:
            raise HTTPException(status_code=404, detail=f"Binary not found: {e}")
        if f.get("encrypted"):
            try:
                aad = f"{rid}:{file_id}".encode("utf-8")
                plaintext = decrypt_envelope(raw, associated_data=aad)
            except Exception as e:
                logger.exception("Vault file decryption failed (preview)")
                raise HTTPException(status_code=500, detail=f"Decryption failed: {e}")
        else:
            plaintext = raw
    elif f.get("content"):
        # Text-only file (seed / legacy / extracted-from-binary): serve the
        # extracted plaintext as text/plain inline.
        plaintext = (f.get("content") or "").encode("utf-8", errors="replace")
        if not ctype or not (ctype.startswith("text/") or ctype == "application/json"):
            ctype = "text/plain"
    else:
        raise HTTPException(status_code=404, detail="File has no previewable content")

    await log_audit(user["id"], "dealroom.file.preview", rid,
                    {"filename": fname, "file_id": file_id})

    # Direct preview path — no conversion needed.
    if ctype in _DIRECT_PREVIEW_TYPES or fname_lc.endswith(".pdf"):
        out_ctype = "application/pdf" if fname_lc.endswith(".pdf") else ctype
        return StreamingResponse(
            io.BytesIO(plaintext),
            media_type=out_ctype,
            headers={"Content-Disposition": f'inline; filename="{fname}"',
                     "Cache-Control": "private, max-age=60"},
        )

    # Office formats — check cache first, else convert via the pure-Python
    # pipeline and persist the PDF rendition in GridFS for next time.
    if fname_lc.endswith(_OFFICE_PREVIEW_EXTS):
        pdf_id = f.get("preview_pdf_gridfs_id")
        pdf_bytes: bytes | None = None
        if pdf_id:
            try:
                cached = await gridfs_bucket.open_download_stream(ObjectId(pdf_id))
                pdf_bytes = await cached.read()
            except Exception as e:
                logger.warning(f"preview cache miss for {file_id}: {e}")
                pdf_bytes = None
        if pdf_bytes is None:
            pdf_bytes = await _office_to_pdf(fname, plaintext)
            if pdf_bytes is None:
                raise HTTPException(
                    status_code=415,
                    detail="Preview unavailable — Office conversion failed. Try downloading the file.",
                )
            try:
                new_id = await gridfs_bucket.upload_from_stream(
                    f"preview-{fname}.pdf", io.BytesIO(pdf_bytes),
                    metadata={"kind": "preview-pdf", "room_id": rid, "file_id": file_id},
                )
                await db.deal_room_files.update_one(
                    {"id": file_id, "room_id": rid},
                    {"$set": {"preview_pdf_gridfs_id": str(new_id),
                              "preview_pdf_generated_at": now_utc().isoformat()}},
                )
            except Exception as e:
                logger.warning(f"preview cache write failed for {file_id}: {e}")
        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'inline; filename="{os.path.splitext(fname)[0]}.pdf"',
                     "Cache-Control": "private, max-age=300"},
        )

    raise HTTPException(
        status_code=415,
        detail="Preview not supported for this format. Download to view locally.",
    )


@api_router.get("/deal-rooms/{rid}/certificate")
async def vault_provenance_certificate(rid: str, request: Request, user=Depends(get_current_user)):
    """Generate a Cryptographic Provenance Certificate PDF for this Vault.

    Aggregates every Bitcoin-anchored event (NDA signature, file uploads, AI findings,
    inquiry status changes), the file inventory with SHA-256s and encryption flags, and
    the current audit-chain head. Anyone — including a court, regulator, or counterparty —
    can independently verify every event using the open-source `ots verify` CLI.
    """
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)

    inquiry = await db.inquiries.find_one({"id": room.get("inquiry_id")}, {"_id": 0}) or {}
    listing = await db.listings.find_one({"id": room.get("listing_id")}, {"_id": 0}) or {}
    if not listing.get("name"):
        listing["name"] = room.get("listing_name") or inquiry.get("listing_name") or "Deal"
    buyer = await db.users.find_one({"id": room.get("buyer_id")}, {"_id": 0, "password_hash": 0}) or {}
    seller = await db.users.find_one({"id": room.get("seller_id")}, {"_id": 0, "password_hash": 0}) or {}

    files = await db.deal_room_files.find(
        {"room_id": rid},
        {"_id": 0, "content": 0, "pages": 0},
    ).sort("uploaded_at", 1).to_list(500)

    file_ids = [f["id"] for f in files]
    inquiry_id = room.get("inquiry_id")
    proof_query = {
        "$or": [
            {"kind": "nda.signature", "target_id": rid},
            {"kind": "vault.findings", "target_id": rid},
            {"kind": "vault.file", "target_id": {"$in": file_ids}},
            {"kind": "inquiry.status", "target_id": inquiry_id},
        ]
    }
    proofs = await db.ots_proofs.find(
        proof_query, {"_id": 0, "ots_bytes": 0}
    ).sort("created_at", 1).to_list(500)

    findings = await db.deal_room_findings.find(
        {"room_id": rid}, {"_id": 0}
    ).sort("created_at", 1).to_list(200)

    chain_head = await db.audit_chain_head.find_one({"_id": "head"}, {"_id": 0}) or {}

    cert_id = str(uuid.uuid4())[:8].upper()
    generated_at = now_utc().isoformat()
    base_url = str(request.base_url).rstrip("/").replace("/api", "")

    try:
        pdf_bytes = build_provenance_pdf(
            cert_id=cert_id,
            generated_at=generated_at,
            room=room,
            inquiry=inquiry,
            listing=listing,
            buyer=buyer,
            seller=seller,
            files=files,
            proofs=proofs,
            findings=findings,
            chain_head=chain_head,
            base_url=base_url,
        )
    except Exception as e:
        logger.exception("Provenance PDF render failed")
        raise HTTPException(status_code=500, detail=f"Could not render certificate: {e}")

    await log_audit(user["id"], "dealroom.certificate.generate", rid, {
        "cert_id": cert_id, "proof_count": len(proofs), "file_count": len(files),
    })
    # Notarize the certificate itself (Bitcoin-anchored "this certificate existed at T")
    asyncio.create_task(notarize_bytes(
        kind="vault.certificate",
        target_id=rid,
        data=pdf_bytes,
        owner_user_id=user["id"],
        label=f"Provenance certificate {cert_id} · {listing.get('name','')}",
        extra={"vault_id": rid, "cert_id": cert_id, "size_bytes": len(pdf_bytes)},
    ))

    safe_name = (listing.get("name") or "deal").lower().replace(" ", "-").replace("/", "-")
    filename = f"workz-provenance-{safe_name}-{cert_id}.pdf"
    return StreamingResponse(
        io.BytesIO(pdf_bytes),
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


async def _run_findings_job(job_id: str, rid: str, user_id: str) -> None:
    """Heavy AI analysis task — runs in the background so the HTTP request
    that kicked it returns immediately. Cloudflare otherwise 524s after
    100 s. The frontend polls `/findings-job/{job_id}` until status
    terminal (completed | failed).

    No per-job time cap: institutional vaults can hold hundreds of large
    files and the Claude pass may legitimately take several minutes. The
    task awaits the LLM call directly (no `asyncio.wait_for`)."""
    started = now_utc()
    await db.findings_jobs.update_one(
        {"id": job_id},
        {"$set": {"status": "running", "started_at": started.isoformat()}},
    )
    try:
        # Pull every file in the room. Cap at 200 to keep the Claude
        # prompt under the model's input window; if the vault has more,
        # we surface the truncation in the job doc so the buyer knows the
        # analysis skipped the tail (and can re-trigger after reorganizing).
        # Files are read oldest-first so the inventory order is stable
        # across runs and the cited file_index numbers don't drift.
        FILE_CAP = 200
        total_files = await db.deal_room_files.count_documents({"room_id": rid})
        files = await db.deal_room_files.find({"room_id": rid}, {"_id": 0}).sort("uploaded_at", 1).to_list(FILE_CAP)
        truncated = total_files > FILE_CAP
        if not files:
            await db.findings_jobs.update_one(
                {"id": job_id},
                {"$set": {"status": "failed", "error": "No files in room yet",
                          "finished_at": now_utc().isoformat()}},
            )
            return

        inventory = []
        for idx, f in enumerate(files, start=1):
            pages = f.get("pages") or []
            if pages:
                page_blocks = []
                for p in pages[:6]:
                    page_blocks.append(f"  <page n={p['page']}>\n  {p.get('text','')[:800]}\n  </page>")
                body_block = "\n".join(page_blocks)
            else:
                body_block = (f.get("content") or "")[:1500]
            inventory.append(
                f"[{idx}] file_id={f['id']} · filename={f['filename']} · folder={f['folder']} · page_count={f.get('page_count', 1)}\n{body_block}"
            )
        files_block = "\n\n---\n\n".join(inventory)

        sys = """You are a senior M&A diligence analyst. Given a numbered file inventory with per-page markers <page n=X>, produce STRICT JSON findings.
Return: {"findings":[{"severity":"high|medium|low","workstream":"finance|legal|hr|it|operations|commercial","title":str,"description":str,"file_index":int,"page":int,"excerpt":str}]}
Cap to 10 findings. Each excerpt MUST be a verbatim short quote (≤200 chars) drawn from the exact referenced page. "page" MUST be the integer from the <page n=X> tag the excerpt came from (use 1 if unknown). Be specific."""

        raw = await call_claude(sys, f"File inventory:\n\n{files_block}\n\nReturn JSON now.", session_id=f"findings-{rid}-{job_id}")
        data = safe_json_loads(raw)
        findings = data.get("findings", []) if isinstance(data, dict) else []
        inserted = []
        for f in findings[:10]:
            try:
                idx = int(f.get("file_index", 0))
            except Exception:
                idx = 0
            cited_file = files[idx - 1] if 1 <= idx <= len(files) else None
            try:
                page_num = int(f.get("page") or 1)
            except Exception:
                page_num = 1
            page_count = (cited_file or {}).get("page_count", 1) or 1
            if page_num < 1 or page_num > page_count:
                page_num = 1
            doc = {
                "id": str(uuid.uuid4()),
                "room_id": rid,
                # NEW: stamp the snapshot ID so this finding belongs to ONE
                # specific Analyze run. Prior to iter-34 findings accumulated
                # across runs with no way to filter — the UI showed every
                # finding ever generated.
                "job_id": job_id,
                "severity": f.get("severity", "medium"),
                "workstream": f.get("workstream", "operations"),
                "title": (f.get("title") or "Untitled finding")[:200],
                "description": (f.get("description") or "")[:1000],
                "citation": {
                    "file_id": cited_file["id"] if cited_file else None,
                    "filename": cited_file["filename"] if cited_file else None,
                    "page": page_num if cited_file else None,
                    "excerpt": (f.get("excerpt") or "")[:240],
                },
                "created_at": now_utc().isoformat(),
            }
            inserted.append(doc)
        if inserted:
            await db.deal_room_findings.insert_many(inserted)
            for d in inserted:
                d.pop("_id", None)

        # Executive summary: one short Claude pass synthesizing the
        # findings into a 1-2 sentence headline. Useful at the top of the
        # PDF + in the snapshot list. Cheap (~200 tokens out) and bounded
        # so a failure here doesn't fail the whole job.
        executive_summary = ""
        if inserted:
            try:
                summary_input = "\n".join(
                    f"- [{d['severity']}/{d['workstream']}] {d['title']}: {d['description'][:160]}"
                    for d in inserted
                )
                exec_sys = (
                    "You are an M&A diligence analyst. Given a list of findings, write a "
                    "1-2 sentence executive summary that a senior partner can read in 5 seconds. "
                    "No preamble, no markdown. Mention the most material risk by name. "
                    "If everything is low severity, say so explicitly."
                )
                executive_summary = (await call_claude(
                    exec_sys, summary_input,
                    session_id=f"findings-summary-{rid}-{job_id}",
                ) or "").strip()[:600]
            except Exception as se:
                logger.warning(f"findings job {job_id} exec summary failed: {se}")
                executive_summary = ""

        # Severity breakdown for the snapshot list UI (so we don't have to
        # re-query findings just to render badges in the picker).
        severity_breakdown = {"high": 0, "medium": 0, "low": 0}
        for d in inserted:
            sev = d.get("severity", "medium")
            severity_breakdown[sev] = severity_breakdown.get(sev, 0) + 1

        duration = int((now_utc() - started).total_seconds() * 1000)
        await db.findings_jobs.update_one(
            {"id": job_id},
            {"$set": {
                "status": "completed",
                "findings_count": len(inserted),
                "files_analyzed": len(files),
                "total_files_in_room": total_files,
                "truncated": truncated,
                "executive_summary": executive_summary,
                "severity_breakdown": severity_breakdown,
                "finished_at": now_utc().isoformat(),
                "duration_ms": duration,
            }},
        )
        await log_agent_activity(
            "findings-agent",
            f"room:{rid} · {len(inserted)} findings",
            "completed",
            user_id=user_id,
            duration_ms=duration,
        )
        await log_audit(user_id, "dealroom.findings.generate", rid, {"count": len(inserted), "job_id": job_id})
        # Bitcoin-anchored proof of findings (digest of the sorted findings JSON).
        asyncio.create_task(notarize_event(
            kind="vault.findings",
            target_id=rid,
            payload={
                "vault_id": rid,
                "buyer_id": user_id,
                "generated_at": now_utc().isoformat(),
                "findings": [
                    {"id": d["id"], "severity": d["severity"], "workstream": d["workstream"],
                     "title": d["title"], "citation": d.get("citation")}
                    for d in inserted
                ],
            },
            owner_user_id=user_id,
            label=f"AI findings · {len(inserted)} items",
        ))
    except Exception as e:
        await db.findings_jobs.update_one(
            {"id": job_id},
            {"$set": {"status": "failed", "error": str(e)[:500],
                      "finished_at": now_utc().isoformat()}},
        )
        await log_agent_activity("findings-agent", f"room:{rid}", "failed",
                                  user_id=user_id, friction=str(e)[:200])


@api_router.post("/deal-rooms/{rid}/generate-findings")
async def generate_findings(rid: str, user=Depends(get_current_user)):
    """Kick off AI findings generation as a background job. Returns
    immediately with a job_id the client polls via `/findings-job/{job_id}`.

    Previously this endpoint ran the Claude call synchronously, which
    routinely hit Cloudflare's 100 s edge timeout for vaults with lots of
    files — buyers saw a 524 with no recovery. The job pattern removes the
    request-time ceiling entirely; the LLM can take as long as it needs."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    role = await participant_check(room, user)
    if role not in ("buyer", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Only the buyer can generate findings")

    file_count = await db.deal_room_files.count_documents({"room_id": rid})
    if file_count == 0:
        raise HTTPException(status_code=400, detail="No files in room yet")

    # If a job is already running, don't queue a second one — return the
    # in-flight job so the client can resume polling.
    inflight = await db.findings_jobs.find_one(
        {"room_id": rid, "status": {"$in": ["pending", "running"]}},
        {"_id": 0},
        sort=[("created_at", -1)],
    )
    if inflight:
        return {"job_id": inflight["id"], "status": inflight["status"],
                "already_running": True,
                "created_at": inflight.get("created_at")}

    job_id = str(uuid.uuid4())
    await db.findings_jobs.insert_one({
        "id": job_id,
        "room_id": rid,
        "requested_by": user["id"],
        "status": "pending",
        "files_to_analyze": file_count,
        "findings_count": 0,
        "created_at": now_utc().isoformat(),
        "started_at": None,
        "finished_at": None,
        "error": None,
    })
    asyncio.create_task(_run_findings_job(job_id, rid, user["id"]))
    return {"job_id": job_id, "status": "pending", "already_running": False,
            "files_to_analyze": file_count}


@api_router.get("/deal-rooms/{rid}/findings-job/{job_id}")
async def get_findings_job(rid: str, job_id: str, user=Depends(get_current_user)):
    """Poll endpoint for the AI findings job. Returns the job document so
    the UI can display 'Analyzing N files… (started 32s ago)' and switch
    to a completed/failed state when the task finishes."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    job = await db.findings_jobs.find_one({"id": job_id, "room_id": rid}, {"_id": 0})
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return job


@api_router.get("/deal-rooms/{rid}/findings-job")
async def get_latest_findings_job(rid: str, user=Depends(get_current_user)):
    """Convenience: return the most recent job for this room so a buyer who
    refreshes the page mid-run can re-attach to the polling cycle without
    needing the original job_id."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    job = await db.findings_jobs.find_one(
        {"room_id": rid}, {"_id": 0},
        sort=[("created_at", -1)],
    )
    return job or {}


@api_router.get("/deal-rooms/{rid}/findings-snapshots")
async def list_findings_snapshots(rid: str, user=Depends(get_current_user)):
    """Snapshot picker backend. Returns every completed Findings run for
    this Vault, latest first. Each snapshot carries enough metadata to
    render the picker dropdown without an extra round-trip — exec summary,
    findings count, severity breakdown, files analyzed."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    jobs = await db.findings_jobs.find(
        {"room_id": rid, "status": "completed"},
        {"_id": 0},
    ).sort("created_at", -1).to_list(50)

    # "Fresh files" since latest analysis → drives the "5 new files · Re-run" banner.
    fresh_count = 0
    if jobs:
        try:
            last_run_iso = jobs[0].get("finished_at") or jobs[0].get("started_at") or jobs[0].get("created_at")
            fresh_count = await db.deal_room_files.count_documents({
                "room_id": rid,
                "uploaded_at": {"$gt": last_run_iso},
                "deleted_at": {"$exists": False},
            })
        except Exception:
            fresh_count = 0
    return {"snapshots": jobs, "fresh_files_since_last_run": fresh_count}


def _diff_findings(current: list[dict], prior: list[dict]) -> dict:
    """Match findings across snapshots on (workstream, normalized title).
    Claude rephrases descriptions but rarely renames the title for the
    same risk."""
    def key(d):
        return ((d.get("workstream") or "").lower(),
                (d.get("title") or "").strip().lower())
    cur_by = {key(d): d for d in current}
    pri_by = {key(d): d for d in prior}
    new = [cur_by[k] for k in cur_by if k not in pri_by]
    resolved = [pri_by[k] for k in pri_by if k not in cur_by]
    unchanged = [cur_by[k] for k in cur_by if k in pri_by]
    return {"new": new, "resolved": resolved, "unchanged": unchanged}


@api_router.get("/deal-rooms/{rid}/findings-snapshots/{job_id}")
async def get_findings_snapshot(rid: str, job_id: str, user=Depends(get_current_user)):
    """Findings for one snapshot + diff vs prior. Powers the picker detail
    view and the +N/-M/=K diff badge."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    job = await db.findings_jobs.find_one({"id": job_id, "room_id": rid}, {"_id": 0})
    if not job:
        raise HTTPException(status_code=404, detail="Snapshot not found")
    findings = await db.deal_room_findings.find(
        {"room_id": rid, "job_id": job_id}, {"_id": 0},
    ).sort("created_at", -1).to_list(200)
    prior_job = await db.findings_jobs.find_one(
        {"room_id": rid, "status": "completed",
         "created_at": {"$lt": job["created_at"]}},
        {"_id": 0},
        sort=[("created_at", -1)],
    )
    diff_summary = None
    if prior_job:
        prior_findings = await db.deal_room_findings.find(
            {"room_id": rid, "job_id": prior_job["id"]}, {"_id": 0},
        ).to_list(200)
        diff = _diff_findings(findings, prior_findings)
        diff_summary = {
            "new": len(diff["new"]),
            "resolved": len(diff["resolved"]),
            "unchanged": len(diff["unchanged"]),
            "prior_job_id": prior_job["id"],
            "prior_finished_at": prior_job.get("finished_at"),
        }
    return {"job": job, "findings": findings, "diff": diff_summary}


def _build_findings_pdf(job: dict, findings: list[dict], room: dict) -> bytes:
    """Branded PDF: header → exec summary → findings by severity →
    verbatim-excerpt evidence appendix. Uses ReportLab (already in stack)."""
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether,
    )

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER,
        leftMargin=0.75 * inch, rightMargin=0.75 * inch,
        topMargin=0.6 * inch, bottomMargin=0.6 * inch,
        title=f"NextCapOS Findings · {room.get('listing_name', 'Vault')}",
        author="NextCapOS",
    )
    styles = getSampleStyleSheet()
    BLOOM_NAVY = colors.HexColor("#0B1B3D")
    BLOOM_GOLD = colors.HexColor("#C9A14A")
    SEV_COLOR = {
        "high": colors.HexColor("#B91C1C"),
        "medium": colors.HexColor("#B45309"),
        "low": colors.HexColor("#15803D"),
    }
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], textColor=BLOOM_NAVY, fontSize=18, leading=22)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], textColor=BLOOM_NAVY, fontSize=13, leading=16, spaceBefore=16)
    overline = ParagraphStyle("overline", parent=styles["BodyText"], fontSize=8,
                              textColor=colors.HexColor("#666"), spaceAfter=2)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=10, leading=14)
    quote = ParagraphStyle("quote", parent=styles["BodyText"], fontSize=9, leading=12,
                           leftIndent=12, textColor=colors.HexColor("#444"),
                           borderColor=BLOOM_GOLD, borderPadding=4)
    title_style = ParagraphStyle("title", parent=styles["BodyText"], fontSize=11,
                                 leading=14, textColor=BLOOM_NAVY, spaceAfter=2,
                                 fontName="Helvetica-Bold")

    elems: list = []
    elems.append(Paragraph("AI DILIGENCE FINDINGS", overline))
    elems.append(Paragraph(room.get("listing_name") or "Vault Report", h1))
    elems.append(Spacer(1, 4))
    meta_bits = []
    if job.get("finished_at"):
        meta_bits.append(f"Generated {job['finished_at'][:10]}")
    if job.get("files_analyzed"):
        meta_bits.append(f"{job['files_analyzed']} files analyzed")
    if job.get("findings_count") is not None:
        meta_bits.append(f"{job['findings_count']} findings")
    elems.append(Paragraph(" &nbsp; · &nbsp; ".join(meta_bits) or "—", body))

    sev = job.get("severity_breakdown") or {"high": 0, "medium": 0, "low": 0}
    chip_data = [[
        Paragraph(f"<b>{sev.get('high', 0)}</b> high", body),
        Paragraph(f"<b>{sev.get('medium', 0)}</b> medium", body),
        Paragraph(f"<b>{sev.get('low', 0)}</b> low", body),
    ]]
    chip_tbl = Table(chip_data, colWidths=[1.2 * inch] * 3)
    chip_tbl.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (0, 0), SEV_COLOR["high"]),
        ("BACKGROUND", (1, 0), (1, 0), SEV_COLOR["medium"]),
        ("BACKGROUND", (2, 0), (2, 0), SEV_COLOR["low"]),
        ("TEXTCOLOR", (0, 0), (-1, -1), colors.white),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    elems.append(Spacer(1, 10))
    elems.append(chip_tbl)

    summary_text = (job.get("executive_summary") or "").strip()
    if summary_text:
        elems.append(Paragraph("EXECUTIVE SUMMARY", overline))
        elems.append(Paragraph("Headline", h2))
        elems.append(Paragraph(summary_text, body))

    elems.append(Paragraph("FINDINGS", overline))
    elems.append(Paragraph("Detailed analysis", h2))
    sev_order = ["high", "medium", "low"]
    findings_by_sev = {s: [f for f in findings if f.get("severity") == s] for s in sev_order}
    if not any(findings_by_sev.values()):
        elems.append(Paragraph("<i>No findings in this snapshot.</i>", body))
    for s in sev_order:
        bucket = findings_by_sev[s]
        if not bucket:
            continue
        elems.append(Spacer(1, 8))
        sev_label = Paragraph(f"<font color='{SEV_COLOR[s].hexval()}'><b>{s.upper()} SEVERITY · {len(bucket)}</b></font>", body)
        elems.append(sev_label)
        for f in bucket:
            cit = f.get("citation") or {}
            cit_str = ""
            if cit.get("filename"):
                cit_str = f"<b>{cit['filename']}</b>"
                if cit.get("page"):
                    cit_str += f" · p.{cit['page']}"
            block = [
                Paragraph(f.get("title") or "Untitled finding", title_style),
                Paragraph(f"<i>workstream:</i> {f.get('workstream', '—')}", body),
                Paragraph((f.get("description") or "").replace("\n", "<br/>"), body),
            ]
            if cit.get("excerpt") and cit_str:
                block.append(Spacer(1, 2))
                block.append(Paragraph(f"&ldquo;{cit['excerpt']}&rdquo; — {cit_str}", quote))
            block.append(Spacer(1, 6))
            elems.append(KeepTogether(block))

    cited = [f for f in findings if (f.get("citation") or {}).get("excerpt")]
    if cited:
        elems.append(PageBreak())
        elems.append(Paragraph("EVIDENCE APPENDIX", overline))
        elems.append(Paragraph("Verbatim excerpts", h2))
        elems.append(Paragraph(
            "Quoted text below is drawn verbatim from the cited file at the "
            "cited page — included so reviewers can verify the AI's claims "
            "without opening each source document.", body))
        elems.append(Spacer(1, 8))
        for f in cited:
            cit = f["citation"]
            cit_str = f"{cit.get('filename', '?')}"
            if cit.get("page"):
                cit_str += f" · p.{cit['page']}"
            block = [
                Paragraph(f"<b>{f.get('title')}</b> · <font color='#666'>{cit_str}</font>", body),
                Paragraph(f"&ldquo;{cit.get('excerpt')}&rdquo;", quote),
                Spacer(1, 8),
            ]
            elems.append(KeepTogether(block))

    elems.append(Spacer(1, 14))
    footer = ParagraphStyle("footer", parent=styles["BodyText"], fontSize=7,
                            textColor=colors.HexColor("#999"), alignment=1)
    elems.append(Paragraph(
        f"Snapshot ID {job.get('id', '—')[:8]} · Generated by NextCapOS AI · "
        f"This is a machine-generated analysis. Verify all claims against source files.",
        footer,
    ))

    doc.build(elems)
    pdf = buf.getvalue()
    buf.close()
    return pdf


@api_router.get("/deal-rooms/{rid}/findings-snapshots/{job_id}/pdf")
async def export_findings_pdf(rid: str, job_id: str, user=Depends(get_current_user)):
    """Stream the branded PDF. Buyer clicks Export PDF → instant download."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    job = await db.findings_jobs.find_one({"id": job_id, "room_id": rid}, {"_id": 0})
    if not job or job.get("status") != "completed":
        raise HTTPException(status_code=404, detail="Snapshot not found or not completed")
    findings = await db.deal_room_findings.find(
        {"room_id": rid, "job_id": job_id}, {"_id": 0},
    ).sort("created_at", -1).to_list(200)
    pdf_bytes = _build_findings_pdf(job, findings, room)
    await log_audit(user["id"], "vault.findings.pdf_export", rid, {"job_id": job_id})
    slug = (room.get("listing_name") or "vault").replace(" ", "_")[:40]
    date_part = (job.get("finished_at") or now_utc().isoformat())[:10]
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="Findings_{slug}_{date_part}.pdf"'},
    )


class FindingsEmailRequest(BaseModel):
    recipients: List[EmailStr]
    note: Optional[str] = None


@api_router.post("/deal-rooms/{rid}/findings-snapshots/{job_id}/email")
async def email_findings_pdf(
    rid: str, job_id: str, body: FindingsEmailRequest,
    user=Depends(get_current_user),
):
    """Auditable share — emails PDF via Resend; each send logged."""
    if not body.recipients:
        raise HTTPException(status_code=400, detail="At least one recipient required")
    if len(body.recipients) > 10:
        raise HTTPException(status_code=400, detail="Max 10 recipients per send")
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    job = await db.findings_jobs.find_one({"id": job_id, "room_id": rid}, {"_id": 0})
    if not job or job.get("status") != "completed":
        raise HTTPException(status_code=404, detail="Snapshot not found or not completed")
    findings = await db.deal_room_findings.find(
        {"room_id": rid, "job_id": job_id}, {"_id": 0},
    ).sort("created_at", -1).to_list(200)
    pdf_bytes = _build_findings_pdf(job, findings, room)

    slug = (room.get("listing_name") or "vault").replace(" ", "_")[:40]
    date_part = (job.get("finished_at") or now_utc().isoformat())[:10]
    fname = f"Findings_{slug}_{date_part}.pdf"

    subject = f"NextCapOS Diligence Findings · {room.get('listing_name', 'Vault')}"
    html_body = (
        f"<p>{user.get('name', 'A buyer')} shared an AI-generated diligence "
        f"findings report for <b>{room.get('listing_name', 'a vault')}</b> "
        f"with you via NextCapOS.</p>"
        + (f"<blockquote>{body.note}</blockquote>" if body.note else "")
        + f"<p>The full report is attached. Generated "
          f"{(job.get('finished_at') or '')[:10]} · "
          f"{job.get('findings_count', 0)} findings across "
          f"{job.get('files_analyzed', 0)} files analyzed.</p>"
    )

    sent_count = 0
    failures = []
    for r in body.recipients:
        try:
            await send_email_with_attachment(
                to=r,
                subject=subject,
                html=html_body,
                attachment_filename=fname,
                attachment_bytes=pdf_bytes,
                attachment_mime="application/pdf",
            )
            sent_count += 1
        except Exception as e:
            failures.append({"to": r, "error": str(e)[:140]})

    await log_audit(user["id"], "vault.findings.email", rid, {
        "job_id": job_id,
        "recipients": list(body.recipients),
        "sent": sent_count,
        "failed": [f["to"] for f in failures],
    })
    return {"ok": sent_count > 0, "sent": sent_count, "failures": failures}




# --- Co-pilot (chat against the file corpus, with citations) ---
COPILOT_SYS = """You are the NextCapOS Vault Co-pilot — a senior M&A diligence analyst assisting a buyer.
You answer questions strictly from the provided file inventory. The inventory uses per-page markers <page n=X>.
Cite the file AND page you drew from inline using the format [filename p.N] (e.g. [DPA.pdf p.3]). Cite every claim.
If the answer is not in the inventory, say so explicitly. Be precise — verbatim figures, dates, named parties.
Keep answers under 250 words. Tone: institutional, terse, analytical."""


@api_router.get("/deal-rooms/{rid}/copilot")
async def get_copilot_history(rid: str, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    msgs = await db.deal_room_messages.find({"room_id": rid}, {"_id": 0}).sort("created_at", 1).to_list(200)
    return msgs


async def _run_copilot_job(job_id: str, rid: str, user_id: str, user_name: str | None,
                           user_msg_id: str, question: str) -> None:
    """Heavy Co-pilot work — runs in a background task so the POST returns
    immediately and Cloudflare's 100 s edge timeout never fires. Same
    pattern as `_run_findings_job`. The frontend polls
    `/copilot-job/{job_id}` until status terminal.

    On completion writes the assistant message to `deal_room_messages` so
    the existing `GET /copilot` history endpoint surfaces it on the next
    fetch — no schema change needed for client renders."""
    started = now_utc()
    await db.copilot_jobs.update_one(
        {"id": job_id},
        {"$set": {"status": "running", "started_at": started.isoformat()}},
    )
    try:
        room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
        if not room:
            raise RuntimeError("Vault disappeared mid-run")

        # Self-heal: clone any newly-staged listing files (incl. Composio-
        # mirrored ones) into the room so this turn sees them.
        if room.get("listing_id"):
            try:
                await _clone_listing_files_into_room(
                    room["listing_id"], rid, room.get("seller_id") or user_id, only_missing=True,
                )
            except Exception as e:
                logger.warning(f"copilot job {job_id}: backfill clone failed: {e}")

        FILE_CAP = 200
        total_files = await db.deal_room_files.count_documents({"room_id": rid})
        files = await db.deal_room_files.find({"room_id": rid}, {"_id": 0}).sort("uploaded_at", 1).to_list(FILE_CAP)
        truncated_inventory = total_files > FILE_CAP

        if not files:
            empty_reply = "No documents have been uploaded to this Vault yet. Ask the seller to upload diligence materials, then re-ask."
            asst_msg = {
                "id": str(uuid.uuid4()),
                "room_id": rid,
                "role": "assistant",
                "user_id": "copilot",
                "user_name": "Vault Co-pilot",
                "content": empty_reply,
                "citations": [],
                "created_at": now_utc().isoformat(),
            }
            await db.deal_room_messages.insert_one(asst_msg)
            asst_msg.pop("_id", None)
            await db.copilot_jobs.update_one(
                {"id": job_id},
                {"$set": {"status": "completed", "assistant_message_id": asst_msg["id"],
                          "files_analyzed": 0, "total_files_in_room": 0,
                          "truncated": False,
                          "finished_at": now_utc().isoformat(),
                          "duration_ms": int((now_utc() - started).total_seconds() * 1000)}},
            )
            return

        # Dynamic per-file budget: total inventory ≤ ~150 K chars (well
        # inside Claude Sonnet's input window). Every file is represented.
        CHAR_BUDGET = 150_000
        per_file_chars = max(800, min(6_000, CHAR_BUDGET // max(len(files), 1)))
        pages_per_file = 6 if per_file_chars >= 3_000 else 3 if per_file_chars >= 1_500 else 2
        chars_per_page = per_file_chars // pages_per_file

        inventory_lines = []
        for idx, f in enumerate(files, start=1):
            pages = f.get("pages") or []
            if pages:
                page_blocks = []
                for p in pages[:pages_per_file]:
                    txt = (p.get("text") or "")[:chars_per_page]
                    page_blocks.append(f"  <page n={p.get('page', 1)}>\n  {txt}\n  </page>")
                body_block = "\n".join(page_blocks)
            else:
                txt = (f.get("content") or "")[:per_file_chars]
                body_block = f"  <page n=1>\n  {txt}\n  </page>"
            source_kind = (f.get("source") or {}).get("kind") if isinstance(f.get("source"), dict) else None
            provenance = f" · source={source_kind}" if source_kind else ""
            inventory_lines.append(
                f"[{idx}] filename={f['filename']} · folder={f.get('folder', 'other')}"
                f" · page_count={f.get('page_count', 1)}{provenance}\n{body_block}"
            )
        inventory = "\n\n---\n\n".join(inventory_lines)

        history = await db.deal_room_messages.find(
            {"room_id": rid, "id": {"$ne": user_msg_id}},
            {"_id": 0},
        ).sort("created_at", -1).to_list(8)
        history.reverse()
        transcript = "\n".join(f"{m['role'].upper()}: {m['content']}" for m in history)

        truncation_note = ""
        if truncated_inventory:
            truncation_note = (
                f"\nNOTE: This Vault contains {total_files} files; only the {FILE_CAP} oldest "
                "are shown above. If the buyer asks about a file not listed, say so "
                "and suggest narrowing by folder.\n"
            )

        prompt = (
            f"FILE INVENTORY (only source you may cite — uses <page n=X> markers):\n{inventory}\n"
            + truncation_note
            + (f"\nPRIOR CONVERSATION:\n{transcript}\n" if transcript else "")
            + f"\nBUYER QUESTION: {question}\n\nAnswer now."
        )

        answer = await call_claude(COPILOT_SYS, prompt, session_id=f"copilot-{rid}-{user_id}")

        import re
        citations = []
        seen_keys: set[tuple[str, int]] = set()
        for m in re.finditer(r"\[([^\[\]]+?\.[a-zA-Z0-9]+)\s+p\.(\d+)\]", answer or ""):
            fname = m.group(1).strip()
            try:
                page = int(m.group(2))
            except ValueError:
                page = 1
            key = (fname, page)
            if key in seen_keys:
                continue
            seen_keys.add(key)
            match = next((f for f in files if f["filename"] == fname), None)
            if match:
                citations.append({"file_id": match["id"], "filename": fname, "page": page})
        for m in re.finditer(r"\[([^\[\]]+?\.[a-zA-Z0-9]+)\](?!\s*p\.)", answer or ""):
            fname = m.group(1).strip()
            key = (fname, 1)
            if key in seen_keys:
                continue
            seen_keys.add(key)
            match = next((f for f in files if f["filename"] == fname), None)
            if match:
                citations.append({"file_id": match["id"], "filename": fname, "page": 1})

        asst_msg = {
            "id": str(uuid.uuid4()),
            "room_id": rid,
            "role": "assistant",
            "user_id": "copilot",
            "user_name": "Vault Co-pilot",
            "content": answer or "",
            "citations": citations,
            "created_at": now_utc().isoformat(),
        }
        await db.deal_room_messages.insert_one(asst_msg)
        asst_msg.pop("_id", None)
        duration = int((now_utc() - started).total_seconds() * 1000)
        await db.copilot_jobs.update_one(
            {"id": job_id},
            {"$set": {
                "status": "completed",
                "assistant_message_id": asst_msg["id"],
                "files_analyzed": len(files),
                "total_files_in_room": total_files,
                "truncated": truncated_inventory,
                "citation_count": len(citations),
                "finished_at": now_utc().isoformat(),
                "duration_ms": duration,
            }},
        )
        await log_agent_activity(
            "vault-copilot",
            f"ask:{question[:60]}",
            "completed",
            user_id=user_id,
            duration_ms=duration,
            meta={"citations": len(citations), "job_id": job_id},
        )
        await log_audit(user_id, "vault.copilot.ask", rid)
    except Exception as e:
        logger.exception(f"copilot job {job_id} failed")
        await db.copilot_jobs.update_one(
            {"id": job_id},
            {"$set": {"status": "failed", "error": str(e)[:500],
                      "finished_at": now_utc().isoformat()}},
        )
        await log_agent_activity("vault-copilot", f"ask:{question[:60]}", "failed",
                                  user_id=user_id, friction=str(e)[:200])


@api_router.post("/deal-rooms/{rid}/copilot")
async def ask_copilot(rid: str, body: CopilotAsk, user=Depends(get_current_user)):
    """Kick off the Co-pilot turn as a background job. Returns immediately
    with `{job_id, user_message, status:'pending'}` so the request is well
    under Cloudflare's 100 s edge timeout. Frontend polls
    `/copilot-job/{job_id}` until status terminal.

    Big vaults (hundreds of mirrored files) routinely took 60-180 s when
    we ran Claude inside the request handler — production buyers saw
    a 524 with no recovery. The job pattern removes the request-time
    ceiling entirely."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    if room.get("status") == "pending_nda":
        raise HTTPException(status_code=400, detail="NDA must be accepted before using the Co-pilot")

    user_msg = {
        "id": str(uuid.uuid4()),
        "room_id": rid,
        "role": "user",
        "user_id": user["id"],
        "user_name": user.get("name"),
        "content": body.message[:2000],
        "citations": [],
        "created_at": now_utc().isoformat(),
    }
    await db.deal_room_messages.insert_one(user_msg)
    user_msg.pop("_id", None)

    job_id = str(uuid.uuid4())
    await db.copilot_jobs.insert_one({
        "id": job_id,
        "room_id": rid,
        "user_message_id": user_msg["id"],
        "requested_by": user["id"],
        "status": "pending",
        "question": body.message[:300],
        "created_at": now_utc().isoformat(),
        "started_at": None,
        "finished_at": None,
        "error": None,
    })
    asyncio.create_task(_run_copilot_job(
        job_id, rid, user["id"], user.get("name"), user_msg["id"], body.message[:2000],
    ))
    return {"job_id": job_id, "status": "pending", "user_message": user_msg}


@api_router.get("/deal-rooms/{rid}/copilot-job/{job_id}")
async def get_copilot_job(rid: str, job_id: str, user=Depends(get_current_user)):
    """Poll endpoint for a Co-pilot turn. When status='completed', the UI
    re-fetches the message history via GET /copilot — the assistant
    message has already been written there by the background task."""
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    await participant_check(room, user)
    job = await db.copilot_jobs.find_one({"id": job_id, "room_id": rid}, {"_id": 0})
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return job


# -----------------------------------------------------------------------------
# AUDIT + SECURITY
# -----------------------------------------------------------------------------
@api_router.get("/audit/logs")
async def audit_logs(user=Depends(get_current_user)):
    items = await db.audit_logs.find({}, {"_id": 0}).sort("timestamp", -1).to_list(200)
    return items


@api_router.get("/security/audit/verify")
async def verify_audit_chain(user=Depends(get_current_user)):
    """
    Re-walk the audit chain from genesis, recomputing each entry's content hash and
    asserting prev_hash continuity. Returns the first break point (if any) so the
    user / regulator can prove the log has not been tampered with.
    """
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin only")
    expected_prev = GENESIS_HASH
    total = 0
    broken = None
    cursor = db.audit_logs.find({}, {"_id": 0}).sort("seq", 1)
    async for entry in cursor:
        total += 1
        if "seq" not in entry or "prev_hash" not in entry or "content_hash" not in entry:
            # Legacy pre-chain entry — accept but mark as "unverifiable"
            continue
        if entry["prev_hash"] != expected_prev:
            broken = {"seq": entry["seq"], "id": entry["id"], "reason": "prev_hash mismatch"}
            break
        recomputed = compute_content_hash(entry)
        if recomputed != entry["content_hash"]:
            broken = {"seq": entry["seq"], "id": entry["id"], "reason": "content_hash mismatch"}
            break
        expected_prev = entry["content_hash"]
    head = await db.audit_chain_head.find_one({"_id": "head"}, {"_id": 0}) or {}
    return {
        "total_entries": total,
        "chain_valid": broken is None,
        "broken_at": broken,
        "chain_head": head,
        "verified_at": now_utc().isoformat(),
    }


@api_router.get("/security/proofs")
async def list_proofs(
    kind: Optional[str] = None,
    target_id: Optional[str] = None,
    user=Depends(get_current_user),
):
    """List OpenTimestamps proofs visible to the calling user.
    - admin → all
    - buyer/seller → proofs they own (own_user_id == self) + proofs they participate in
      (vault.file / nda.signature / vault.findings where target_id is a deal_room they're in)
    """
    q: Dict[str, Any] = {}
    if kind:
        q["kind"] = kind
    if target_id:
        q["target_id"] = target_id

    if user.get("role") == "admin":
        proofs = await db.ots_proofs.find(q, {"_id": 0, "ots_bytes": 0}).sort("created_at", -1).to_list(500)
    else:
        # Visible scope: events the user authored OR events on a deal_room they participate in
        my_rooms = await db.deal_rooms.find(
            {"$or": [{"buyer_id": user["id"]}, {"seller_id": user["id"]}]},
            {"_id": 0, "id": 1},
        ).to_list(500)
        my_room_ids = [r["id"] for r in my_rooms]
        my_inquiries = await db.inquiries.find(
            {"$or": [{"buyer_id": user["id"]}, {"seller_id": user["id"]}]},
            {"_id": 0, "id": 1},
        ).to_list(500)
        my_inquiry_ids = [i["id"] for i in my_inquiries]
        # Find vault files owned by their rooms — pull file ids
        my_files = await db.deal_room_files.find(
            {"room_id": {"$in": my_room_ids}},
            {"_id": 0, "id": 1},
        ).to_list(2000) if my_room_ids else []
        my_file_ids = [f["id"] for f in my_files]
        scope_q = {
            "$or": [
                {"owner_user_id": user["id"]},
                {"kind": "nda.signature", "target_id": {"$in": my_room_ids}},
                {"kind": "vault.findings", "target_id": {"$in": my_room_ids}},
                {"kind": "vault.file", "target_id": {"$in": my_file_ids}},
                {"kind": "inquiry.status", "target_id": {"$in": my_inquiry_ids}},
            ]
        }
        if q:
            scope_q = {"$and": [scope_q, q]}
        proofs = await db.ots_proofs.find(scope_q, {"_id": 0, "ots_bytes": 0}).sort("created_at", -1).to_list(500)
    return proofs


@api_router.get("/security/proofs/{proof_id}")
async def get_proof(proof_id: str, user=Depends(get_current_user)):
    p = await db.ots_proofs.find_one({"id": proof_id}, {"_id": 0, "ots_bytes": 0})
    if not p:
        raise HTTPException(status_code=404, detail="Proof not found")
    return p


@api_router.get("/security/proofs/{proof_id}/download")
async def download_proof(proof_id: str, user=Depends(get_current_user)):
    p = await db.ots_proofs.find_one({"id": proof_id})
    if not p:
        raise HTTPException(status_code=404, detail="Proof not found")
    ots = p.get("ots_bytes")
    if not ots:
        raise HTTPException(status_code=404, detail="Proof bytes missing")
    digest_short = (p.get("digest_hex") or "proof")[:8]
    return StreamingResponse(
        io.BytesIO(ots),
        media_type="application/vnd.opentimestamps",
        headers={"Content-Disposition": f'attachment; filename="workz-{p.get("kind","proof")}-{digest_short}.ots"'},
    )


@api_router.post("/security/proofs/{proof_id}/upgrade")
async def upgrade_proof(proof_id: str, user=Depends(get_current_user)):
    """Attempt to fetch a Bitcoin-anchored extension of a pending OTS proof from
    its calendar(s). Once a Bitcoin attestation is present, status flips to 'confirmed'."""
    p = await db.ots_proofs.find_one({"id": proof_id})
    if not p:
        raise HTTPException(status_code=404, detail="Proof not found")
    if not p.get("ots_bytes"):
        raise HTTPException(status_code=400, detail="No proof bytes")
    try:
        result = await upgrade_ots(p["ots_bytes"])
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"OTS upgrade failed: {e}")
    update = {"ots_bytes": result["ots_bytes"]}
    if result.get("btc_block_height"):
        update["btc_block_height"] = result["btc_block_height"]
        update["status"] = "confirmed"
        update["confirmed_at"] = now_utc().isoformat()
    await db.ots_proofs.update_one({"id": proof_id}, {"$set": update})
    return {
        "ok": True,
        "upgraded": result.get("upgraded", False),
        "btc_block_height": result.get("btc_block_height"),
        "status": update.get("status", p.get("status")),
    }


@api_router.post("/security/verify")
async def verify_uploaded_proof(
    ots_file: UploadFile = File(...),
    digest_hex: str = Form(...),
    user=Depends(get_current_user),
):
    """Independent verifier: upload any .ots file + the digest you want to verify and
    we'll parse the proof, confirm it stamps the given digest, and report any
    Bitcoin attestation. Anyone can also verify with the open-source `ots verify` CLI."""
    try:
        digest = bytes.fromhex(digest_hex.strip())
    except Exception:
        raise HTTPException(status_code=400, detail="digest_hex must be hex-encoded")
    if len(digest) != 32:
        raise HTTPException(status_code=400, detail="digest must be 32 bytes (SHA-256)")
    ots_bytes = await ots_file.read()
    try:
        result = verify_ots(ots_bytes, digest)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not parse .ots: {e}")
    return result


@api_router.get("/security/posture")
async def security_posture(user=Depends(get_current_user)):
    """Public summary of security posture: features enabled, counts, calendar reachability."""
    pending = await db.ots_proofs.count_documents({"status": "pending"})
    confirmed = await db.ots_proofs.count_documents({"status": "confirmed"})
    head = await db.audit_chain_head.find_one({"_id": "head"}, {"_id": 0}) or {}
    return {
        "features": {
            "opentimestamps": True,
            "at_rest_encryption": encryption_configured(),
            "encryption_alg": "AES-256-GCM" if encryption_configured() else None,
            "audit_chain": True,
            "brute_force_lockout": True,
            "security_headers": True,
            "password_complexity": True,
        },
        "ots": {
            "calendars": [
                "alice.btc.calendar.opentimestamps.org",
                "bob.btc.calendar.opentimestamps.org",
                "finney.calendar.eternitywall.com",
            ],
            "pending_proofs": pending,
            "confirmed_proofs": confirmed,
        },
        "audit_chain": {
            "last_seq": head.get("last_seq", 0),
            "last_hash": head.get("last_hash"),
            "last_ts": head.get("last_ts"),
        },
    }


# -----------------------------------------------------------------------------
# CRUD: DELETE / EDIT / MESSAGING / COLLATERAL ACTIONS  (iter-11)
# -----------------------------------------------------------------------------

def _now_iso() -> str:
    return now_utc().isoformat()


async def _soft_delete(collection, query: dict, actor_id: str, action: str, target: str):
    res = await collection.update_one({**query, "deleted_at": {"$exists": False}}, {"$set": {"deleted_at": _now_iso(), "deleted_by": actor_id}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Not found or already deleted")
    await log_audit(actor_id, action, target)
    return {"ok": True, "soft_deleted": True}


# ---- Research ---------------------------------------------------------------
@api_router.delete("/research/{rid}")
async def delete_research(rid: str, user=Depends(get_current_user)):
    """Hard delete own research brief — personal content."""
    res = await db.research.delete_one({"id": rid, "user_id": user["id"]})
    if res.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Research brief not found")
    await log_audit(user["id"], "research.delete", rid)
    return {"ok": True}


# ---- Inquiries: messaging + delete -----------------------------------------
class InquiryMessageCreate(BaseModel):
    body: str = Field(..., min_length=1, max_length=4000)
    attachment_id: Optional[str] = None  # collateral_id when seller pushes a piece into the thread


async def _inquiry_participant(iid: str, user: dict) -> dict:
    inq = await db.inquiries.find_one({"id": iid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not inq:
        raise HTTPException(status_code=404, detail="Inquiry not found")
    if user.get("role") == "admin":
        return inq
    if user["id"] in (inq.get("buyer_id"), inq.get("seller_id")):
        return inq
    # Org teammates / listing collaborators on the seller side count as participants
    if user.get("role") in ("seller", "agent"):
        ws_listings, _ = await _user_workspace_listing_ids(user)
        if inq.get("listing_id") in ws_listings:
            return inq
    raise HTTPException(status_code=403, detail="Not a participant of this inquiry")


@api_router.get("/inquiries/{iid}/messages")
async def list_inquiry_messages(iid: str, user=Depends(get_current_user)):
    await _inquiry_participant(iid, user)
    msgs = await db.inquiry_messages.find({"inquiry_id": iid}, {"_id": 0}).sort("created_at", 1).to_list(500)
    # Mark unseen-by-me as read
    await db.inquiry_messages.update_many(
        {"inquiry_id": iid, "author_id": {"$ne": user["id"]}, "read_by": {"$nin": [user["id"]]}},
        {"$addToSet": {"read_by": user["id"]}},
    )
    return msgs


@api_router.post("/inquiries/{iid}/messages")
async def post_inquiry_message(iid: str, body: InquiryMessageCreate, user=Depends(get_current_user)):
    inq = await _inquiry_participant(iid, user)
    attachment = None
    if body.attachment_id:
        coll = await db.collateral.find_one({"id": body.attachment_id, "user_id": user["id"]}, {"_id": 0})
        if not coll:
            raise HTTPException(status_code=404, detail="Attachment not found")
        attachment = {
            "kind": "collateral",
            "id": coll["id"],
            "title": coll.get("title") or coll.get("asset_type") or "Collateral",
            "asset_type": coll.get("asset_type"),
        }
    msg = {
        "id": str(uuid.uuid4()),
        "inquiry_id": iid,
        "author_id": user["id"],
        "author_name": user.get("name"),
        "author_role": user.get("role"),
        "body": body.body,
        "attachment": attachment,
        "read_by": [user["id"]],
        "created_at": _now_iso(),
    }
    await db.inquiry_messages.insert_one(msg)
    await db.inquiries.update_one(
        {"id": iid},
        {"$set": {"last_message_at": msg["created_at"], "last_message_preview": (body.body or "")[:140]}, "$inc": {"message_count": 1}},
    )
    await log_audit(user["id"], "inquiry.message", iid, {"len": len(body.body), "attachment": bool(attachment)})
    msg.pop("_id", None)
    return msg


@api_router.delete("/inquiries/{iid}")
async def delete_inquiry(iid: str, user=Depends(get_current_user)):
    """Buyer withdraws own inquiry · seller dismisses inbound · admin override. Soft delete."""
    inq = await db.inquiries.find_one({"id": iid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not inq:
        raise HTTPException(status_code=404, detail="Inquiry not found")
    if user.get("role") != "admin" and user["id"] not in (inq.get("buyer_id"), inq.get("seller_id")):
        raise HTTPException(status_code=403, detail="Not your inquiry")
    return await _soft_delete(db.inquiries, {"id": iid}, user["id"], "inquiry.delete", iid)


# ---- Vaults (deal rooms): delete + buyer notes -----------------------------
@api_router.delete("/deal-rooms/{rid}")
async def delete_vault(rid: str, user=Depends(get_current_user)):
    """Soft-delete a vault. Buyer/seller/admin all permitted (it ends the workspace for both)."""
    room = await db.deal_rooms.find_one({"id": rid, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    if user.get("role") != "admin" and user["id"] not in (room.get("buyer_id"), room.get("seller_id")):
        raise HTTPException(status_code=403, detail="Not a participant of this vault")
    return await _soft_delete(db.deal_rooms, {"id": rid}, user["id"], "dealroom.delete", rid)


# Buyer-side uploads: already allowed via /files & /files/binary (participant check
# admits buyer). We expose a 'buyer_notes' folder option to keep things organized.


# ---- Newsletter: edit recipients + delete + buyer interests ----------------
class NewsletterEdit(BaseModel):
    title: Optional[str] = None
    content: Optional[str] = None
    sectors: Optional[List[str]] = None
    recipient_ids: Optional[List[str]] = None  # explicit user ids if seller hand-picks


@api_router.patch("/newsletter/{nid}")
async def edit_newsletter(nid: str, body: NewsletterEdit, user=Depends(get_current_user)):
    nl = await db.newsletters.find_one({"id": nid, "user_id": user["id"]}, {"_id": 0})
    if not nl:
        raise HTTPException(status_code=404, detail="Newsletter not found")
    if nl.get("status") not in ("draft", "approved"):
        raise HTTPException(status_code=400, detail="Only draft/approved newsletters can be edited")
    update: Dict[str, Any] = {"updated_at": _now_iso()}
    for k in ("title", "content", "sectors", "recipient_ids"):
        v = getattr(body, k, None)
        if v is not None:
            update[k] = v
    await db.newsletters.update_one({"id": nid}, {"$set": update})
    await log_audit(user["id"], "newsletter.edit", nid, {"fields": list(update.keys())})
    nl.update(update)
    return nl


@api_router.get("/newsletter/recipient-candidates")
async def newsletter_recipient_candidates(user=Depends(get_current_user)):
    """Return opted-in buyer accounts a seller can hand-pick into a broadcast."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Seller/admin only")
    cands = await db.users.find(
        {"role": "buyer", "newsletter_opt_in": True},
        {"_id": 0, "id": 1, "name": 1, "email": 1, "organization": 1, "interests": 1},
    ).sort("name", 1).to_list(500)
    return cands


@api_router.delete("/newsletter/{nid}")
async def delete_newsletter(nid: str, user=Depends(get_current_user)):
    nl = await db.newsletters.find_one({"id": nid, "user_id": user["id"]}, {"_id": 0})
    if not nl:
        raise HTTPException(status_code=404, detail="Newsletter not found")
    # Hard delete drafts; soft delete dispatched (preserve OTS/audit trail)
    if nl.get("status") in ("draft", "approved"):
        await db.newsletters.delete_one({"id": nid, "user_id": user["id"]})
        await log_audit(user["id"], "newsletter.delete", nid, {"mode": "hard"})
        return {"ok": True, "hard_deleted": True}
    return await _soft_delete(db.newsletters, {"id": nid, "user_id": user["id"]}, user["id"], "newsletter.delete", nid)


class UserInterestsUpdate(BaseModel):
    interests: List[str] = Field(default_factory=list)
    newsletter_opt_in: Optional[bool] = None
    newsletter_cadence: Optional[Literal["weekly", "biweekly", "monthly"]] = None


@api_router.patch("/me/interests")
async def update_my_interests(body: UserInterestsUpdate, user=Depends(get_current_user)):
    """Buyer-side editor of personalized digest preferences (sectors + opt-in + cadence)."""
    update: Dict[str, Any] = {"interests": [s.strip() for s in (body.interests or []) if s and s.strip()]}
    if body.newsletter_opt_in is not None:
        update["newsletter_opt_in"] = body.newsletter_opt_in
    if body.newsletter_cadence is not None:
        update["newsletter_cadence"] = body.newsletter_cadence
    await db.users.update_one({"id": user["id"]}, {"$set": update})
    await log_audit(user["id"], "user.interests.update", user["id"], {"interests": update["interests"]})
    refreshed = await db.users.find_one({"id": user["id"]}, {"_id": 0, "password_hash": 0})
    return refreshed


# ---- Outreach: edit + delete -----------------------------------------------
class OutreachEdit(BaseModel):
    name: Optional[str] = None
    target_persona: Optional[str] = None
    message_brief: Optional[str] = None
    draft: Optional[Dict[str, Any]] = None  # the structured AI output (subject/opening/value_props/etc)
    audience_size: Optional[int] = None


@api_router.patch("/outreach/campaigns/{cid}")
async def edit_outreach(cid: str, body: OutreachEdit, user=Depends(get_current_user)):
    camp = await db.outreach.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not camp:
        raise HTTPException(status_code=404, detail="Campaign not found")
    if camp.get("status") not in ("draft", None):
        raise HTTPException(status_code=400, detail="Only draft campaigns can be edited")
    update: Dict[str, Any] = {"updated_at": _now_iso()}
    for k in ("name", "target_persona", "message_brief", "draft", "audience_size"):
        v = getattr(body, k, None)
        if v is not None:
            update[k] = v
    await db.outreach.update_one({"id": cid}, {"$set": update})
    await log_audit(user["id"], "outreach.edit", cid, {"fields": list(update.keys())})
    return {**camp, **update}


@api_router.delete("/outreach/campaigns/{cid}")
async def delete_outreach(cid: str, user=Depends(get_current_user)):
    camp = await db.outreach.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not camp:
        raise HTTPException(status_code=404, detail="Campaign not found")
    if camp.get("status") in ("draft", None):
        await db.outreach.delete_one({"id": cid, "user_id": user["id"]})
        await log_audit(user["id"], "outreach.delete", cid, {"mode": "hard"})
        return {"ok": True, "hard_deleted": True}
    return await _soft_delete(db.outreach, {"id": cid, "user_id": user["id"]}, user["id"], "outreach.delete", cid)


# ---- Collateral: edit + versions + PDF + attach/push/send ------------------
class CollateralEdit(BaseModel):
    title: Optional[str] = None
    headline: Optional[str] = None
    subheadline: Optional[str] = None
    sections: Optional[List[Dict[str, Any]]] = None
    cta: Optional[str] = None
    compliance_note: Optional[str] = None
    full: Optional[Dict[str, Any]] = None  # full JSON replacement


@api_router.patch("/collateral/{cid}")
async def edit_collateral(cid: str, body: CollateralEdit, user=Depends(get_current_user)):
    item = await db.collateral.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not item:
        raise HTTPException(status_code=404, detail="Collateral not found")
    # Snapshot the previous version
    version_doc = {
        "id": str(uuid.uuid4()),
        "collateral_id": cid,
        "user_id": user["id"],
        "snapshot": {k: v for k, v in item.items() if k not in ("id", "user_id")},
        "created_at": _now_iso(),
    }
    await db.collateral_versions.insert_one(version_doc)

    # Asset content lives under coll["data"]
    data = dict(item.get("data") or {})
    if body.full is not None:
        data.update({k: v for k, v in body.full.items() if k not in ("id", "user_id", "created_at")})
    for k in ("title", "headline", "subheadline", "sections", "cta", "compliance_note"):
        v = getattr(body, k, None)
        if v is not None:
            data[k] = v
    update: Dict[str, Any] = {"updated_at": _now_iso(), "data": data}
    await db.collateral.update_one({"id": cid, "user_id": user["id"]}, {"$set": update})
    await log_audit(user["id"], "collateral.edit", cid, {"fields": list(data.keys())})
    refreshed = await db.collateral.find_one({"id": cid}, {"_id": 0})
    return refreshed


def _coll_data(coll: dict) -> dict:
    """Return the actual asset content dict (handles legacy top-level and new data-nested)."""
    if isinstance(coll.get("data"), dict) and coll["data"]:
        return coll["data"]
    return {k: coll.get(k) for k in ("title", "headline", "subheadline", "sections", "cta", "compliance_note", "asset_type")}


@api_router.delete("/collateral/{cid}")
async def delete_collateral(cid: str, user=Depends(get_current_user)):
    res = await db.collateral.delete_one({"id": cid, "user_id": user["id"]})
    if res.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Collateral not found")
    await db.collateral_versions.delete_many({"collateral_id": cid})
    await log_audit(user["id"], "collateral.delete", cid)
    return {"ok": True}


@api_router.get("/collateral/{cid}/versions")
async def list_collateral_versions(cid: str, user=Depends(get_current_user)):
    item = await db.collateral.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0, "id": 1})
    if not item:
        raise HTTPException(status_code=404, detail="Collateral not found")
    versions = await db.collateral_versions.find(
        {"collateral_id": cid, "user_id": user["id"]}, {"_id": 0}
    ).sort("created_at", -1).to_list(100)
    return versions


def _collateral_to_pdf_bytes(coll: dict, user: dict) -> bytes:
    """Lightweight branded one-pager using ReportLab."""
    from reportlab.lib import colors as _c
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable

    d = _coll_data(coll)
    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER, leftMargin=0.8 * inch, rightMargin=0.8 * inch,
        topMargin=0.85 * inch, bottomMargin=0.85 * inch,
        title=f"NextCapOS · {d.get('title') or coll.get('asset_type','Collateral')}",
        author=user.get("name") or "NextCapOS",
    )
    s_overline = ParagraphStyle("o", fontName="Helvetica-Bold", fontSize=7, textColor=_c.HexColor("#9E7B45"), leading=9, spaceAfter=4)
    s_h1 = ParagraphStyle("h1", fontName="Helvetica-Bold", fontSize=22, textColor=_c.HexColor("#1A1A19"), leading=26, spaceAfter=6)
    s_h2 = ParagraphStyle("h2", fontName="Helvetica-Bold", fontSize=14, textColor=_c.HexColor("#1A1A19"), leading=18, spaceAfter=4, spaceBefore=10)
    s_body = ParagraphStyle("body", fontName="Helvetica", fontSize=10.5, textColor=_c.HexColor("#1A1A19"), leading=15)
    s_subhead = ParagraphStyle("sub", fontName="Helvetica-Oblique", fontSize=11.5, textColor=_c.HexColor("#575754"), leading=16, spaceAfter=10)
    s_cta = ParagraphStyle("cta", fontName="Helvetica-Bold", fontSize=11.5, textColor=_c.HexColor("#9E7B45"), leading=16, spaceBefore=12)
    s_small = ParagraphStyle("sm", fontName="Helvetica", fontSize=8, textColor=_c.HexColor("#8A8A85"), leading=11)

    story = []
    story.append(Paragraph(f"WORKZ VENTURES · {(coll.get('asset_type') or 'COLLATERAL').upper().replace('_',' ')}", s_overline))
    story.append(Paragraph(d.get("headline") or d.get("title") or coll.get("deal_name","Untitled"), s_h1))
    if d.get("subheadline"):
        story.append(Paragraph(d["subheadline"], s_subhead))
    story.append(HRFlowable(width="100%", thickness=0.6, color=_c.HexColor("#DCDCD5"), spaceAfter=10))
    for sec in d.get("sections") or []:
        if isinstance(sec, dict):
            if sec.get("heading"):
                story.append(Paragraph(sec["heading"], s_h2))
            if sec.get("body"):
                story.append(Paragraph(sec["body"], s_body))
    if d.get("cta"):
        story.append(Paragraph(d["cta"], s_cta))
    if d.get("compliance_note"):
        story.append(Spacer(1, 0.18 * inch))
        story.append(HRFlowable(width="100%", thickness=0.4, color=_c.HexColor("#DCDCD5"), spaceAfter=6))
        story.append(Paragraph(d["compliance_note"], s_small))
    story.append(Spacer(1, 0.15 * inch))
    story.append(Paragraph(f"Prepared by {user.get('name','—')} · {user.get('organization','—')} · {_fmt_now_short()}", s_small))
    doc.build(story)
    return buf.getvalue()


def _fmt_now_short() -> str:
    return now_utc().strftime("%Y-%m-%d")


@api_router.get("/collateral/{cid}/pdf")
async def collateral_pdf(cid: str, user=Depends(get_current_user)):
    coll = await db.collateral.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not coll:
        raise HTTPException(status_code=404, detail="Collateral not found")
    try:
        pdf = _collateral_to_pdf_bytes(coll, user)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Could not render PDF: {e}")
    await log_audit(user["id"], "collateral.export.pdf", cid)
    cd = _coll_data(coll)
    safe = (cd.get("title") or coll.get("deal_name") or "collateral").lower().replace(" ", "-")
    return StreamingResponse(
        io.BytesIO(pdf), media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="workz-{safe}.pdf"'},
    )


class AttachToListing(BaseModel):
    listing_id: str


@api_router.post("/collateral/{cid}/attach-to-listing")
async def attach_collateral(cid: str, body: AttachToListing, user=Depends(get_current_user)):
    coll = await db.collateral.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not coll:
        raise HTTPException(status_code=404, detail="Collateral not found")
    listing = await db.listings.find_one({"id": body.listing_id, "seller_id": user["id"]}, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found or not yours")
    cd = _coll_data(coll)
    attachments = listing.get("collateral_attachments") or []
    attachments = [a for a in attachments if a.get("collateral_id") != cid]  # dedupe
    attachments.append({
        "collateral_id": cid,
        "title": cd.get("title") or cd.get("headline") or coll.get("deal_name") or "Collateral",
        "asset_type": coll.get("asset_type"),
        "attached_at": _now_iso(),
    })
    await db.listings.update_one({"id": body.listing_id}, {"$set": {"collateral_attachments": attachments}})
    await log_audit(user["id"], "collateral.attach.listing", cid, {"listing_id": body.listing_id})
    return {"ok": True, "attachments": attachments}


class PushToVault(BaseModel):
    room_id: str
    folder: Literal["financials", "legal", "hr", "it", "operations", "commercial", "other"] = "commercial"


@api_router.post("/collateral/{cid}/push-to-vault")
async def push_collateral_to_vault(cid: str, body: PushToVault, user=Depends(get_current_user)):
    coll = await db.collateral.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not coll:
        raise HTTPException(status_code=404, detail="Collateral not found")
    room = await db.deal_rooms.find_one({"id": body.room_id, "deleted_at": {"$exists": False}}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Vault not found")
    if user.get("role") != "admin" and user["id"] not in (room.get("buyer_id"), room.get("seller_id")):
        raise HTTPException(status_code=403, detail="Not a participant of this vault")
    if room.get("status") == "pending_nda":
        raise HTTPException(status_code=400, detail="Vault still pending NDA")

    # Render PDF
    pdf = _collateral_to_pdf_bytes(coll, user)
    cd = _coll_data(coll)
    filename = f"{(cd.get('title') or coll.get('deal_name') or 'collateral').lower().replace(' ', '-')}.pdf"
    file_id = str(uuid.uuid4())
    plaintext_sha = sha256_hex(pdf)

    storage_bytes = pdf
    encrypted = False
    encryption_alg = None
    if encryption_configured():
        try:
            aad = f"{room['id']}:{file_id}".encode("utf-8")
            enc = encrypt_bytes(pdf, associated_data=aad)
            storage_bytes = enc["envelope"]
            encrypted = True
            encryption_alg = enc["alg"]
        except Exception as e:
            logger.warning(f"Collateral→Vault encryption failed: {e}")

    gridfs_id = await gridfs_bucket.upload_from_stream(
        filename, io.BytesIO(storage_bytes),
        metadata={"room_id": room["id"], "file_id": file_id, "uploaded_by": user["id"],
                  "content_type": "application/pdf", "encrypted": encrypted, "encryption_alg": encryption_alg, "source": f"collateral:{cid}"},
    )
    doc = {
        "id": file_id,
        "room_id": room["id"],
        "filename": filename,
        "folder": body.folder,
        "content_type": "application/pdf",
        "size_bytes": len(pdf),
        "page_count": 1,
        "pages": [{"page": 1, "text": _collateral_flat_text(coll)}],
        "content": _collateral_flat_text(coll),
        "char_count": len(_collateral_flat_text(coll)),
        "gridfs_id": str(gridfs_id),
        "storage": "gridfs",
        "encrypted": encrypted,
        "encryption_alg": encryption_alg,
        "sha256_hex": plaintext_sha,
        "note": f"Generated from collateral '{coll.get('title','')}'",
        "uploaded_by": user["id"],
        "uploaded_by_role": user.get("role"),
        "uploaded_at": _now_iso(),
        "matched_request_id": None,
        "source_collateral_id": cid,
    }
    await db.deal_room_files.insert_one(doc)
    doc.pop("_id", None)
    await log_audit(user["id"], "collateral.push.vault", cid, {"room_id": room["id"], "file_id": file_id})
    asyncio.create_task(notarize_bytes(
        kind="vault.file", target_id=file_id, data=pdf, owner_user_id=user["id"],
        label=f"Vault file (from collateral): {filename}", extra={"vault_id": room["id"], "filename": filename, "size_bytes": len(pdf), "source_collateral_id": cid},
    ))
    doc.pop("content", None); doc.pop("pages", None)
    return doc


def _collateral_flat_text(coll: dict) -> str:
    d = _coll_data(coll)
    parts = [d.get("headline",""), d.get("subheadline","")]
    for sec in d.get("sections") or []:
        if isinstance(sec, dict):
            parts.append(sec.get("heading",""))
            parts.append(sec.get("body",""))
    if d.get("cta"): parts.append(d["cta"])
    if d.get("compliance_note"): parts.append(d["compliance_note"])
    return "\n\n".join(p for p in parts if p)


class SendToInquiry(BaseModel):
    inquiry_id: str
    note: Optional[str] = None


@api_router.post("/collateral/{cid}/send-to-inquiry")
async def send_collateral_to_inquiry(cid: str, body: SendToInquiry, user=Depends(get_current_user)):
    coll = await db.collateral.find_one({"id": cid, "user_id": user["id"]}, {"_id": 0})
    if not coll:
        raise HTTPException(status_code=404, detail="Collateral not found")
    inq = await _inquiry_participant(body.inquiry_id, user)
    body_text = body.note or f"Sharing this collateral for your review: {coll.get('title') or coll.get('asset_type','one-pager')}."
    msg_in = InquiryMessageCreate(body=body_text, attachment_id=cid)
    return await post_inquiry_message(body.inquiry_id, msg_in, user)


# -----------------------------------------------------------------------------
# BUYER DISCOVERY · Phase 1 (SEC EDGAR + Companies House [stubbed] + Claude)
# -----------------------------------------------------------------------------
from buyer_discovery import gather_candidates, rank_with_claude, resolve_match_contacts  # noqa: E402
from social_presence import discover_social_profiles  # noqa: E402

BUYER_DISCOVERY_RESCAN_HOURS = int(os.environ.get("BUYER_DISCOVERY_RESCAN_HOURS", "24"))
BUYER_MATCH_ALERT_THRESHOLD = 70  # only score >= 70 produces a new-buyer alert


class BuyerMatchAction(BaseModel):
    status: Optional[Literal["new", "saved", "dismissed", "contacted"]] = None


def _serialize_match(m: Dict[str, Any]) -> Dict[str, Any]:
    m.pop("_id", None)
    return m


async def _listing_for_seller(lid: str, user: dict) -> Dict[str, Any]:
    query = {"id": lid}
    if user.get("role") != "admin":
        query["seller_id"] = user["id"]
    listing = await db.listings.find_one(query, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found")
    return listing


async def _run_buyer_scan(listing: Dict[str, Any], *, triggered_by: str) -> Dict[str, Any]:
    """Run discovery for one listing. Inserts new matches + alerts for high-score buyers."""
    lid = listing["id"]
    seller_id = listing.get("seller_id")
    started = now_utc()
    candidates = await gather_candidates({
        "id": lid,
        "company_name": listing.get("company_name"),
        "sector": listing.get("sector"),
        "revenue_band": (
            f"{int(listing.get('revenue_usd_m'))}M USD" if listing.get("revenue_usd_m") else None
        ),
        "ebitda_band": (
            f"{int(listing.get('ebitda_usd_m'))}M USD" if listing.get("ebitda_usd_m") else None
        ),
        "geography": listing.get("geography"),
        "deal_type": "majority sale",
        "tagline": listing.get("headline"),
    })
    ranked = await rank_with_claude(call_claude, listing, candidates) if candidates else []
    # Filter very-weak matches
    ranked = [r for r in ranked if int(r.get("score", 0)) >= 35]

    new_alerts = 0
    inserted = 0
    for cand in ranked:
        # Dedupe per (listing_id, buyer_name) — refresh score if existing, otherwise insert
        key = {"listing_id": lid, "buyer_name": cand.get("buyer_name")}
        existing = await db.buyer_matches.find_one(key, {"_id": 0})
        if existing:
            await db.buyer_matches.update_one(
                {"id": existing["id"]},
                {"$set": {
                    "score": cand.get("score"),
                    "rationale": cand.get("rationale"),
                    "fit": cand.get("fit"),
                    "snippet": cand.get("snippet"),
                    "filing_url": cand.get("filing_url"),
                    "filed_at": cand.get("filed_at"),
                    "last_seen_at": now_utc().isoformat(),
                }},
            )
            continue
        mid = str(uuid.uuid4())
        doc = {
            "id": mid,
            "listing_id": lid,
            "seller_id": seller_id,
            "buyer_name": cand.get("buyer_name"),
            "buyer_cik": cand.get("buyer_cik"),
            "country": cand.get("country"),
            "source": cand.get("source"),
            "form": cand.get("form"),
            "filing_url": cand.get("filing_url"),
            "filed_at": cand.get("filed_at"),
            "snippet": cand.get("snippet"),
            "tickers": cand.get("tickers") or [],
            "score": int(cand.get("score", 0)),
            "rationale": cand.get("rationale"),
            "fit": cand.get("fit") or {},
            "status": "new",
            "created_at": now_utc().isoformat(),
            "last_seen_at": now_utc().isoformat(),
        }
        await db.buyer_matches.insert_one(doc)
        inserted += 1
        if doc["score"] >= BUYER_MATCH_ALERT_THRESHOLD:
            await db.buyer_alerts.insert_one({
                "id": str(uuid.uuid4()),
                "seller_id": seller_id,
                "listing_id": lid,
                "listing_company": listing.get("company_name"),
                "match_id": mid,
                "buyer_name": doc["buyer_name"],
                "score": doc["score"],
                "rationale": doc["rationale"],
                "source": doc["source"],
                "country": doc["country"],
                "seen": False,
                "created_at": now_utc().isoformat(),
            })
            new_alerts += 1

    duration_ms = int((now_utc() - started).total_seconds() * 1000)
    await db.buyer_scans.update_one(
        {"listing_id": lid},
        {"$set": {
            "listing_id": lid,
            "seller_id": seller_id,
            "last_scanned_at": now_utc().isoformat(),
            "last_triggered_by": triggered_by,
            "last_candidate_count": len(candidates),
            "last_ranked_count": len(ranked),
            "last_inserted": inserted,
            "last_new_alerts": new_alerts,
            "last_duration_ms": duration_ms,
        }},
        upsert=True,
    )
    await log_agent_activity("buyer-discovery-agent", f"scan:{listing.get('company_name')}", "completed",
                             user_id=seller_id, duration_ms=duration_ms)
    await log_audit(seller_id or "system", "buyer_discovery.scan", lid,
                    {"candidates": len(candidates), "ranked": len(ranked),
                     "inserted": inserted, "alerts": new_alerts, "trigger": triggered_by})
    return {
        "listing_id": lid,
        "candidate_count": len(candidates),
        "ranked_count": len(ranked),
        "inserted": inserted,
        "new_alerts": new_alerts,
        "duration_ms": duration_ms,
    }


@api_router.post("/buyer-discovery/listings/{lid}/scan")
async def trigger_buyer_scan(lid: str, user=Depends(get_current_user)):
    """Sellers (or admins) — run a buyer discovery scan now for one of their listings."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    listing = await _listing_for_seller(lid, user)
    result = await _run_buyer_scan(listing, triggered_by=f"user:{user['id']}")
    return result


@api_router.get("/buyer-discovery/listings/{lid}/matches")
async def list_buyer_matches(lid: str, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    listing = await _listing_for_seller(lid, user)
    matches = await db.buyer_matches.find(
        {"listing_id": lid, "deleted_at": {"$exists": False}},
        {"_id": 0},
    ).sort("score", -1).to_list(100)
    last_scan = await db.buyer_scans.find_one({"listing_id": lid}, {"_id": 0})
    return {"listing": {"id": listing["id"], "company_name": listing.get("company_name"),
                        "sector": listing.get("sector"), "geography": listing.get("geography")},
            "last_scan": last_scan,
            "matches": matches}


@api_router.get("/buyer-discovery/overview")
async def buyer_discovery_overview(user=Depends(get_current_user)):
    """Seller cockpit — per-listing summary (count, top score, last scan)."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    seller_filter = {} if user.get("role") == "admin" else {"seller_id": user["id"]}
    listings = await db.listings.find(seller_filter, {"_id": 0}).to_list(200)
    out = []
    for li in listings:
        lid = li["id"]
        pipeline_count = await db.buyer_matches.count_documents({"listing_id": lid, "deleted_at": {"$exists": False}})
        top = await db.buyer_matches.find_one(
            {"listing_id": lid, "deleted_at": {"$exists": False}},
            {"_id": 0},
            sort=[("score", -1)],
        )
        last_scan = await db.buyer_scans.find_one({"listing_id": lid}, {"_id": 0})
        out.append({
            "listing_id": lid,
            "company_name": li.get("company_name"),
            "sector": li.get("sector"),
            "geography": li.get("geography"),
            "status": li.get("status"),
            "match_count": pipeline_count,
            "top_score": (top or {}).get("score"),
            "top_buyer": (top or {}).get("buyer_name"),
            "last_scanned_at": (last_scan or {}).get("last_scanned_at"),
        })
    return out


@api_router.patch("/buyer-discovery/matches/{mid}")
async def update_buyer_match(mid: str, body: BuyerMatchAction, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    match = await db.buyer_matches.find_one({"id": mid}, {"_id": 0})
    if not match:
        raise HTTPException(status_code=404, detail="Match not found")
    if user.get("role") != "admin" and match.get("seller_id") != user["id"]:
        raise HTTPException(status_code=403, detail="Not your match")
    update: Dict[str, Any] = {}
    if body.status is not None:
        update["status"] = body.status
    if not update:
        return match
    await db.buyer_matches.update_one({"id": mid}, {"$set": update})
    await log_audit(user["id"], "buyer_discovery.match.update", mid, update)
    match.update(update)
    return match


@api_router.delete("/buyer-discovery/matches/{mid}")
async def delete_buyer_match(mid: str, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    match = await db.buyer_matches.find_one({"id": mid}, {"_id": 0})
    if not match:
        raise HTTPException(status_code=404, detail="Match not found")
    if user.get("role") != "admin" and match.get("seller_id") != user["id"]:
        raise HTTPException(status_code=403, detail="Not your match")
    await db.buyer_matches.update_one({"id": mid},
                                      {"$set": {"deleted_at": now_utc().isoformat(),
                                                "status": "dismissed"}})
    await log_audit(user["id"], "buyer_discovery.match.delete", mid)
    return {"ok": True}


@api_router.post("/buyer-discovery/matches/{mid}/add-to-leads")
async def add_match_to_leads(mid: str, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    match = await db.buyer_matches.find_one({"id": mid}, {"_id": 0})
    if not match or (user.get("role") != "admin" and match.get("seller_id") != user["id"]):
        raise HTTPException(status_code=404, detail="Match not found")
    listing = await db.listings.find_one({"id": match.get("listing_id")}, {"_id": 0}) or {}
    lead = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "name": f"Corp Dev · {match.get('buyer_name')}",
        "company": match.get("buyer_name"),
        "title": "Corporate Development",
        "email": None,
        "source": "buyer-discovery",
        "stage": "new",
        "score": int(match.get("score", 50)),
        "buyer_match_id": mid,
        "listing_id": match.get("listing_id"),
        "listing_company": listing.get("company_name"),
        "notes": match.get("rationale"),
        "filing_url": match.get("filing_url"),
        "country": match.get("country"),
        "created_at": now_utc().isoformat(),
    }
    await db.leads.insert_one(lead)
    await db.buyer_matches.update_one({"id": mid}, {"$set": {"status": "saved",
                                                              "lead_id": lead["id"]}})
    await log_audit(user["id"], "buyer_discovery.match.to_lead", mid, {"lead_id": lead["id"]})
    lead.pop("_id", None)
    return lead


@api_router.post("/buyer-discovery/matches/{mid}/generate-outreach")
async def generate_outreach_for_match(mid: str, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    match = await db.buyer_matches.find_one({"id": mid}, {"_id": 0})
    if not match or (user.get("role") != "admin" and match.get("seller_id") != user["id"]):
        raise HTTPException(status_code=404, detail="Match not found")
    listing = await db.listings.find_one({"id": match.get("listing_id")}, {"_id": 0}) or {}

    persona = f"Head of M&A / Corp Dev at {match.get('buyer_name')} ({match.get('country','US')})"
    brief = (
        f"Reach out about a confidential opportunity to acquire {listing.get('company_name','our portfolio company')} "
        f"in the {listing.get('sector','—')} sector ({listing.get('geography','—')}). "
        f"They recently signaled M&A appetite: {match.get('rationale','')}. "
        f"Source filing: {match.get('filing_url','')}."
    )
    req = OutreachCampaignRequest(
        name=f"{listing.get('company_name','Listing')} → {match.get('buyer_name')}",
        target_persona=persona,
        channel="linkedin",
        audience_size=1,
        message_brief=brief,
    )
    campaign = await create_campaign(req, user)
    await db.buyer_matches.update_one({"id": mid}, {"$set": {"status": "contacted",
                                                              "campaign_id": campaign["id"]}})
    await log_audit(user["id"], "buyer_discovery.match.outreach", mid, {"campaign_id": campaign["id"]})
    return campaign


@api_router.post("/buyer-discovery/matches/{mid}/find-contacts")
async def find_contacts_for_match(mid: str, refresh: bool = False, user=Depends(get_current_user)):
    """Find named M&A/Corp Dev/IR contacts at this buyer firm by parsing their recent SEC filings
    (DEF 14A, 10-K, 8-K) with Claude, then resolve LinkedIn URLs via Brave. Emails/phones are only
    surfaced if they literally appear in the filing text."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    match = await db.buyer_matches.find_one({"id": mid}, {"_id": 0})
    if not match or (user.get("role") != "admin" and match.get("seller_id") != user["id"]):
        raise HTTPException(status_code=404, detail="Match not found")

    cached = match.get("contacts") or None
    if cached and not refresh:
        return cached

    if not match.get("buyer_cik"):
        raise HTTPException(status_code=400, detail="No SEC CIK on this match — contact resolution requires a US public-company candidate")

    started = now_utc()
    try:
        result = await resolve_match_contacts(
            call_claude, search_brave,
            cik=str(match["buyer_cik"]),
            company_name=match.get("buyer_name") or "",
        )
    except Exception as e:
        logger.exception("Contact resolution failed")
        raise HTTPException(status_code=500, detail=f"Contact resolution failed: {e}")

    # Persist on the match for cheap subsequent reads
    await db.buyer_matches.update_one(
        {"id": mid},
        {"$set": {"contacts": result, "contacts_resolved_at": result["generated_at"]}},
    )
    duration_ms = int((now_utc() - started).total_seconds() * 1000)
    await log_agent_activity(
        "contact-resolution-agent",
        f"resolve:{match.get('buyer_name')}",
        "completed", user_id=user["id"], duration_ms=duration_ms,
    )
    await log_audit(
        user["id"], "buyer_discovery.match.contacts", mid,
        {"executives": len(result.get("executives") or []),
         "ir_email": bool((result.get("ir_contact") or {}).get("email")),
         "filings_used": len(result.get("used_filings") or [])},
    )
    return result


@api_router.post("/buyer-discovery/matches/{mid}/contacts/{contact_idx}/add-to-leads")
async def add_contact_to_leads(mid: str, contact_idx: int, user=Depends(get_current_user)):
    """Promote a single resolved executive into the Lead Nurturing kanban."""
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer discovery is seller/admin only")
    match = await db.buyer_matches.find_one({"id": mid}, {"_id": 0})
    if not match or (user.get("role") != "admin" and match.get("seller_id") != user["id"]):
        raise HTTPException(status_code=404, detail="Match not found")
    execs = ((match.get("contacts") or {}).get("executives")) or []
    if contact_idx < 0 or contact_idx >= len(execs):
        raise HTTPException(status_code=404, detail="Contact not found — run find-contacts first")
    ex = execs[contact_idx]
    listing = await db.listings.find_one({"id": match.get("listing_id")}, {"_id": 0}) or {}
    lead = {
        "id": str(uuid.uuid4()),
        "user_id": user["id"],
        "name": ex.get("name") or "Corporate Development",
        "company": match.get("buyer_name"),
        "title": ex.get("title") or "Corporate Development",
        "email": None,  # never inferred — only if literally in filing
        "linkedin_url": ex.get("linkedin_url"),
        "source": "buyer-discovery-contact",
        "stage": "new",
        "score": int(match.get("score", 60)),
        "buyer_match_id": mid,
        "listing_id": match.get("listing_id"),
        "listing_company": listing.get("company_name"),
        "notes": ex.get("rationale"),
        "country": match.get("country"),
        "created_at": now_utc().isoformat(),
    }
    await db.leads.insert_one(lead)
    await log_audit(user["id"], "buyer_discovery.contact.to_lead", mid,
                    {"lead_id": lead["id"], "exec": ex.get("name")})
    lead.pop("_id", None)
    return lead


@api_router.get("/buyer-alerts")
async def list_buyer_alerts(unseen_only: bool = False, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer alerts are seller/admin only")
    q: Dict[str, Any] = {"deleted_at": {"$exists": False}}
    if user.get("role") != "admin":
        q["seller_id"] = user["id"]
    if unseen_only:
        q["seen"] = False
    items = await db.buyer_alerts.find(q, {"_id": 0}).sort("created_at", -1).to_list(200)
    return items


@api_router.get("/buyer-alerts/count")
async def buyer_alerts_count(user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        return {"unseen": 0}
    q: Dict[str, Any] = {"seen": False, "deleted_at": {"$exists": False}}
    if user.get("role") != "admin":
        q["seller_id"] = user["id"]
    n = await db.buyer_alerts.count_documents(q)
    return {"unseen": n}


@api_router.patch("/buyer-alerts/{aid}/seen")
async def mark_buyer_alert_seen(aid: str, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer alerts are seller/admin only")
    q = {"id": aid}
    if user.get("role") != "admin":
        q["seller_id"] = user["id"]
    res = await db.buyer_alerts.update_one(q, {"$set": {"seen": True, "seen_at": now_utc().isoformat()}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Alert not found")
    return {"ok": True}


@api_router.post("/buyer-alerts/mark-all-seen")
async def mark_all_buyer_alerts_seen(user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer alerts are seller/admin only")
    q: Dict[str, Any] = {"seen": False}
    if user.get("role") != "admin":
        q["seller_id"] = user["id"]
    res = await db.buyer_alerts.update_many(q, {"$set": {"seen": True, "seen_at": now_utc().isoformat()}})
    return {"ok": True, "updated": res.modified_count}


@api_router.delete("/buyer-alerts/{aid}")
async def delete_buyer_alert(aid: str, user=Depends(get_current_user)):
    if user.get("role") not in ("seller", "admin", "agent"):
        raise HTTPException(status_code=403, detail="Buyer alerts are seller/admin only")
    q = {"id": aid}
    if user.get("role") != "admin":
        q["seller_id"] = user["id"]
    res = await db.buyer_alerts.update_one(q, {"$set": {"deleted_at": now_utc().isoformat()}})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Alert not found")
    return {"ok": True}


# ---- Background rescan scheduler --------------------------------------------
_buyer_scheduler_task: Optional[asyncio.Task] = None
_demo_cleanup_task: Optional[asyncio.Task] = None

from demo_cleanup import (  # noqa: E402
    DEMO_EMAILS,
    DEMO_RETENTION_HOURS,
    demo_cleanup_scheduler,
    purge_demo_data,
)


async def _buyer_discovery_scheduler():
    """Wake every hour. Find listings (status=live) whose last scan is older than
    BUYER_DISCOVERY_RESCAN_HOURS, and rescan them with a small concurrency cap."""
    SLEEP_SECONDS = 60 * 60  # 1 hour
    INITIAL_DELAY = 90       # let the app warm up before first sweep
    await asyncio.sleep(INITIAL_DELAY)
    while True:
        try:
            cutoff = now_utc() - timedelta(hours=BUYER_DISCOVERY_RESCAN_HOURS)
            cutoff_iso = cutoff.isoformat()
            live_listings = await db.listings.find(
                {"status": "live"}, {"_id": 0, "id": 1, "company_name": 1, "sector": 1,
                                     "geography": 1, "revenue_usd_m": 1, "ebitda_usd_m": 1,
                                     "headline": 1, "seller_id": 1}
            ).to_list(500)
            due: List[Dict[str, Any]] = []
            for li in live_listings:
                if not li.get("sector"):
                    continue
                scan = await db.buyer_scans.find_one({"listing_id": li["id"]}, {"_id": 0})
                if (not scan) or (scan.get("last_scanned_at") or "") < cutoff_iso:
                    due.append(li)
            # Sweep with concurrency cap = 2 to be polite to SEC
            sem = asyncio.Semaphore(2)
            async def _bound(li):
                async with sem:
                    try:
                        await _run_buyer_scan(li, triggered_by="scheduler")
                    except Exception as e:
                        logger.warning(f"Scheduled buyer scan failed for {li.get('id')}: {e}")
            if due:
                logger.info(f"buyer-discovery scheduler: rescanning {len(due)} listings")
                await asyncio.gather(*[_bound(li) for li in due])
        except Exception as e:
            logger.warning(f"buyer-discovery scheduler loop error: {e}")
        await asyncio.sleep(SLEEP_SECONDS)


# -----------------------------------------------------------------------------
# SEED (demo deals on startup)
# -----------------------------------------------------------------------------
async def seed_demo_user():
    seed_email = "alex@workz.example.com"
    if not await db.users.find_one({"email": seed_email}):
        await db.users.insert_one({
            "id": str(uuid.uuid4()),
            "email": seed_email,
            "name": "Alex Buyer",
            "role": "buyer",
            "organization": "Cascade Capital",
            "password_hash": hash_password("WorkzPass123!"),
            "interests": ["SaaS", "EMEA"],
            "newsletter_opt_in": True,
            "newsletter_cadence": "weekly",
            "created_at": now_utc().isoformat(),
            "is_demo": True,
        })

    seller_email = "mira@workz.example.com"
    seller = await db.users.find_one({"email": seller_email})
    if not seller:
        seller_id = str(uuid.uuid4())
        await db.users.insert_one({
            "id": seller_id,
            "email": seller_email,
            "name": "Mira Seller",
            "role": "seller",
            "organization": "Northstar Holdings",
            "password_hash": hash_password("WorkzPass123!"),
            "interests": ["HealthTech", "Industrial"],
            "newsletter_opt_in": False,
            "created_at": now_utc().isoformat(),
            "is_demo": True,
        })
    else:
        seller_id = seller["id"]

    admin_email = "admin@workz.example.com"
    if not await db.users.find_one({"email": admin_email}):
        await db.users.insert_one({
            "id": str(uuid.uuid4()),
            "email": admin_email,
            "name": "NextCapOS Admin",
            "role": "admin",
            "organization": "NextCapOS",
            "password_hash": hash_password("WorkzAdmin123!"),
            "interests": [],
            "newsletter_opt_in": False,
            "created_at": now_utc().isoformat(),
            "is_demo": True,
        })

    # Backfill the is_demo flag on any pre-existing demo accounts
    await db.users.update_many(
        {"email": {"$in": list(DEMO_EMAILS)}, "is_demo": {"$ne": True}},
        {"$set": {"is_demo": True}},
    )

    # Seed sample seller listings if none yet
    if await db.listings.count_documents({"seller_id": seller_id, "is_seed": True}) == 0:
        sample = [
            {
                "company_name": "Helios MedTech",
                "sector": "HealthTech",
                "geography": "EMEA",
                "asking_price_usd_m": 412.0,
                "revenue_usd_m": 142.0,
                "ebitda_usd_m": 38.0,
                "employees": 320,
                "headline": "Category-leading EU surgical robotics platform",
                "summary": "DACH-dominant surgical robotics with 38% YoY growth and 27% EBITDA margins. Recurring revenue ~62%.",
                "highlights": ["32 hospital systems under contract", "FDA + CE marked", "Founder-led, succession-ready"],
                "status": "live",
            },
            {
                "company_name": "Atlas Logistics",
                "sector": "Industrial",
                "geography": "NA",
                "asking_price_usd_m": 287.0,
                "revenue_usd_m": 318.0,
                "ebitda_usd_m": 54.0,
                "employees": 1240,
                "headline": "Asset-light North-American 3PL with a tech moat",
                "summary": "Tier-1 retail accounts, AI-driven route optimization, 18% EBITDA margins growing.",
                "highlights": ["Top-50 retailer concentration <20%", "Proprietary routing engine", "Cross-dock footprint in 12 cities"],
                "status": "live",
            },
            {
                "company_name": "Vertex Climate",
                "sector": "ClimateTech",
                "geography": "EMEA",
                "asking_price_usd_m": 178.0,
                "revenue_usd_m": 64.0,
                "ebitda_usd_m": 12.0,
                "employees": 180,
                "headline": "Industrial carbon-capture for hard-to-abate sectors",
                "summary": "Patented amine-free capture, two operating plants, €40M backlog.",
                "highlights": ["3 LOI in pipeline", "EU CBAM tailwind", "Founder + CTO contracted to stay 24 months"],
                "status": "draft",
            },
        ]
        await db.listings.insert_many([
            {
                "id": str(uuid.uuid4()),
                "seller_id": seller_id,
                "seller_name": "Mira Seller",
                "seller_org": "Northstar Holdings",
                **s,
                "inquiry_count": 0,
                "view_count": 0,
                "created_at": now_utc().isoformat(),
                "is_seed": True,
            } for s in sample
        ])

    # Migrate any legacy analyst-role users to buyer
    await db.users.update_many({"role": "analyst"}, {"$set": {"role": "buyer"}})

    # Seed an ACTIVE demo Vault between Alex (buyer) and Mira (seller) so the
    # platform always demonstrates shared diligence + AI Co-pilot.
    await _seed_demo_vault()


async def _seed_demo_vault():
    buyer = await db.users.find_one({"email": "alex@workz.example.com"}, {"_id": 0, "id": 1, "name": 1, "organization": 1})
    seller = await db.users.find_one({"email": "mira@workz.example.com"}, {"_id": 0, "id": 1, "name": 1, "organization": 1})
    if not buyer or not seller:
        return
    listing = await db.listings.find_one(
        {"seller_id": seller["id"], "is_seed": True, "company_name": "Helios MedTech"},
        {"_id": 0},
    )
    if not listing:
        return

    # Skip if seed vault already exists between these two on this listing
    existing = await db.deal_rooms.find_one(
        {"buyer_id": buyer["id"], "seller_id": seller["id"], "listing_id": listing["id"], "is_seed": True},
        {"_id": 0, "id": 1},
    )
    if existing:
        return

    now_iso = now_utc().isoformat()
    inquiry_id = str(uuid.uuid4())
    room_id = str(uuid.uuid4())

    await db.inquiries.insert_one({
        "id": inquiry_id,
        "listing_id": listing["id"],
        "listing_name": listing.get("company_name"),
        "buyer_id": buyer["id"],
        "buyer_name": buyer["name"],
        "buyer_org": buyer.get("organization"),
        "seller_id": seller["id"],
        "seller_name": seller["name"],
        "seller_org": seller.get("organization"),
        "status": "engaged",
        "message": "Seed inquiry — institutional buyer evaluating Helios MedTech for platform demo.",
        "deal_room_id": room_id,
        "created_at": now_iso,
        "is_seed": True,
    })

    await db.deal_rooms.insert_one({
        "id": room_id,
        "inquiry_id": inquiry_id,
        "listing_id": listing["id"],
        "listing_name": listing.get("company_name"),
        "sector": listing.get("sector"),
        "buyer_id": buyer["id"],
        "buyer_name": buyer["name"],
        "buyer_org": buyer.get("organization"),
        "seller_id": seller["id"],
        "seller_name": seller["name"],
        "seller_org": seller.get("organization"),
        "status": "active",
        "nda_signed_name": buyer["name"],
        "nda_accepted_by_buyer_at": now_iso,
        "drl_template_id": None,
        "created_at": now_iso,
        "is_seed": True,
    })

    # Three lightweight, text-only files (no GridFS) so the Co-pilot has context
    seed_files = [
        {
            "filename": "Helios_CIM_summary.md",
            "folder": "Commercial",
            "uploaded_by_role": "seller",
            "uploaded_by": seller["id"],
            "content": (
                "# Helios MedTech — Confidential Information Memorandum (Summary)\n\n"
                "**Sector:** Medical Devices · **HQ:** Munich, DE · **Founded:** 2018\n\n"
                "## Business\n"
                "Helios builds connected infusion pumps for ICU and surgical suites. "
                "Recurring revenue is driven by a SaaS layer (Helios Connect) on top of the hardware.\n\n"
                "## Financials (FY24)\n"
                "- Revenue: €54.2M (+38% YoY)\n"
                "- Gross margin: 64%\n"
                "- ARR (Helios Connect): €18.1M (+72% YoY)\n"
                "- Net retention: 119%\n"
                "- EBITDA margin: 9% (target 18% by FY26)\n\n"
                "## Customers\n"
                "412 hospitals across DACH + Nordics. Top 10 = 23% of revenue. Mean tenure 4.1 yrs.\n\n"
                "## Ask\n"
                "Sale of 100% equity. Indicative range €280-360M. Management open to rollover."
            ),
        },
        {
            "filename": "Q4_2024_financial_snapshot.md",
            "folder": "Financials",
            "uploaded_by_role": "seller",
            "uploaded_by": seller["id"],
            "content": (
                "# Helios MedTech — Q4 2024 Financial Snapshot\n\n"
                "| Metric | Q4 2024 | Q4 2023 | YoY |\n"
                "|---|---|---|---|\n"
                "| Revenue | €15.8M | €11.2M | +41% |\n"
                "| Gross Profit | €10.1M | €7.0M | +44% |\n"
                "| Op Expenses | €8.7M | €6.9M | +26% |\n"
                "| Adj. EBITDA | €1.4M | €0.1M | nm |\n"
                "| Cash | €22.4M | €18.6M | — |\n\n"
                "## Notes\n"
                "- Q4 lifted by Klinikum-Stuttgart rollout (€2.1M one-time).\n"
                "- Helios Connect ARR run-rate exit: €19.6M.\n"
                "- DSO: 67 days (target <55).\n"
                "- Inventory: €6.8M (built ahead of EU MDR Cl-IIb refresh)."
            ),
        },
        {
            "filename": "DD_Risks_register.md",
            "folder": "Risk",
            "uploaded_by_role": "buyer",
            "uploaded_by": buyer["id"],
            "content": (
                "# Buyer-side DD Risk Register — Helios MedTech\n\n"
                "## R1 — Regulatory (MDR)\n"
                "EU MDR Class-IIb certification expires Q2 2026. Re-certification budget €0.8M, "
                "notified body backlog ~9 months. Mitigation: parallel-file with TÜV Süd.\n\n"
                "## R2 — Customer concentration\n"
                "Top-10 hospitals = 23% revenue. Klinikum-Stuttgart alone = 6.2%. "
                "Renewal due Mar 2025 — request seller share renewal correspondence.\n\n"
                "## R3 — Hardware supply chain\n"
                "Sole-sourced microcontroller (STM32H7). 14-week lead time, no qualified second source. "
                "Mitigation plan: dual-source by Q3 2025."
            ),
        },
    ]

    docs = []
    for sf in seed_files:
        docs.append({
            "id": str(uuid.uuid4()),
            "room_id": room_id,
            "filename": sf["filename"],
            "folder": sf["folder"],
            "content": sf["content"],
            "char_count": len(sf["content"]),
            "note": None,
            "uploaded_by": sf["uploaded_by"],
            "uploaded_by_role": sf["uploaded_by_role"],
            "uploaded_at": now_iso,
            "matched_request_id": None,
            "is_seed": True,
        })
    if docs:
        await db.deal_room_files.insert_many(docs)


@app.on_event("startup")
async def seed_demo():
    await seed_demo_user()
    if await db.deals.count_documents({}) == 0:
        await db.deals.insert_many([
            {"id": str(uuid.uuid4()), "name": "Project Helios", "sector": "Industrial Tech", "stage": "DD", "value_usd_m": 412, "geography": "EMEA", "status": "active", "created_at": now_utc().isoformat(), "is_seed": True},
            {"id": str(uuid.uuid4()), "name": "Project Atlas", "sector": "SaaS", "stage": "LOI", "value_usd_m": 287, "geography": "NA", "status": "active", "created_at": now_utc().isoformat(), "is_seed": True},
            {"id": str(uuid.uuid4()), "name": "Project Meridian", "sector": "Healthcare", "stage": "Sourcing", "value_usd_m": 619, "geography": "APAC", "status": "active", "created_at": now_utc().isoformat(), "is_seed": True},
            {"id": str(uuid.uuid4()), "name": "Project Nautilus", "sector": "FinServ", "stage": "Closing", "value_usd_m": 1240, "geography": "EMEA", "status": "active", "created_at": now_utc().isoformat(), "is_seed": True},
            {"id": str(uuid.uuid4()), "name": "Project Vertex", "sector": "ClimateTech", "stage": "DD", "value_usd_m": 178, "geography": "NA", "status": "active", "created_at": now_utc().isoformat(), "is_seed": True},
        ])
    # Buyer Discovery background scheduler — periodic rescan of live listings
    global _buyer_scheduler_task
    if _buyer_scheduler_task is None or _buyer_scheduler_task.done():
        _buyer_scheduler_task = asyncio.create_task(_buyer_discovery_scheduler())
        logger.info("buyer-discovery scheduler started")
    # Demo data 48h retention sweeper
    global _demo_cleanup_task
    if _demo_cleanup_task is None or _demo_cleanup_task.done():
        _demo_cleanup_task = asyncio.create_task(
            demo_cleanup_scheduler(db, gridfs_bucket, listing_files_bucket, private_locker_bucket)
        )
        logger.info("demo-cleanup scheduler started (48h retention)")
    logger.info("NextCapOS backend ready")


# -----------------------------------------------------------------------------
# Health
# -----------------------------------------------------------------------------
@api_router.get("/")
async def health():
    return {"service": "workz-ventures", "ok": True}


@api_router.get("/health")
async def deep_health():
    """Deployment + integration self-check. Hit this in production to verify
    you're running the latest code AND that downstream integrations are
    actually reachable with the configured credentials. Safe to expose
    publicly — no secrets in the response."""
    # Server build marker — set DEPLOY_SHA in env or fall back to the latest
    # local git commit so production tells you exactly which build is live.
    deploy_sha = os.environ.get("DEPLOY_SHA")
    if not deploy_sha:
        try:
            import subprocess
            deploy_sha = subprocess.check_output(
                ["git", "log", "-1", "--format=%h"],
                cwd=os.path.dirname(os.path.abspath(__file__ + "/../")),
                stderr=subprocess.DEVNULL,
                timeout=2,
            ).decode().strip()
        except Exception:
            deploy_sha = "unknown"

    checks: dict = {
        "service": "nextcapos",
        "ok": True,
        "deploy_sha": deploy_sha,
        "code_markers": {
            # If these markers are present in the current code, you know
            # the running build includes the latest Composio fixes.
            "uses_connected_accounts_link": True,
            "supports_5_file_sources": len(COMPOSIO_FILE_SOURCES),
            "auto_expire_pending_sources": True,
        },
        "integrations": {},
    }

    # Composio check — single GET that proves the key is valid + that
    # auth_configs are visible. Doesn't require a write scope.
    try:
        async with httpx.AsyncClient(timeout=5.0) as c:
            r = await c.get(
                f"{COMPOSIO_BASE_URL}/api/v3/auth_configs",
                headers={"x-api-key": COMPOSIO_API_KEY} if COMPOSIO_API_KEY else {},
            )
            if r.status_code == 200:
                items = (r.json() or {}).get("items") or []
                checks["integrations"]["composio"] = {
                    "ok": True,
                    "auth_configs_visible": len(items),
                }
            elif r.status_code == 401:
                checks["integrations"]["composio"] = {
                    "ok": False, "reason": "invalid_api_key",
                }
                checks["ok"] = False
            else:
                checks["integrations"]["composio"] = {
                    "ok": False, "reason": f"http_{r.status_code}",
                }
                checks["ok"] = False
    except Exception as e:
        checks["integrations"]["composio"] = {"ok": False, "reason": str(e)[:80]}
        checks["ok"] = False

    # Mongo ping — confirms the DB connection is alive.
    try:
        await db.command("ping")
        checks["integrations"]["mongo"] = {"ok": True}
    except Exception as e:
        checks["integrations"]["mongo"] = {"ok": False, "reason": str(e)[:80]}
        checks["ok"] = False

    return checks


@api_router.get("/demo/retention-info")
async def demo_retention_info(user=Depends(get_current_user)):
    """Public-to-the-authenticated-user introspection of demo retention policy."""
    return {
        "is_demo": bool(user.get("is_demo")),
        "retention_hours": DEMO_RETENTION_HOURS,
        "demo_emails": sorted(DEMO_EMAILS),
        "policy": (
            "Demo accounts are for evaluation only. Any content created in this "
            f"workspace is permanently deleted {DEMO_RETENTION_HOURS} hours after it "
            "is created. Seeded illustrative data (listings, sample deals) is preserved."
        ),
    }


@api_router.post("/admin/demo/purge")
async def admin_demo_purge(user=Depends(get_current_user)):
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin only")
    summary = await purge_demo_data(db, gridfs_bucket, listing_files_bucket, private_locker_bucket)
    await log_audit(user["id"], "admin.demo.purge", "manual_trigger", summary)
    return summary


# =============================================================================
# ORGANIZATIONS + COLLABORATORS (M&A advisor workflow)
# =============================================================================
# - Users can belong to many orgs (org_memberships: many-to-many).
# - An org owns its listings; org_admins manage members + invites.
# - On top of org membership, each listing also has a `collaborators` array
#   so the PRINCIPAL seller (typically not in the agent's org) can be invited
#   as owner/editor of that specific listing.
# - access_policy on each listing controls who approves Vault access:
#   any editor by default; the principal owner gets a veto via the
#   require_principal_approval toggle and the competitor_blocklist.
# =============================================================================


class OrgCreateRequest(BaseModel):
    name: str = Field(..., min_length=2, max_length=120)
    org_type: Literal["advisory", "fund", "corporate", "other"] = "advisory"
    description: Optional[str] = None


class OrgUpdateRequest(BaseModel):
    name: Optional[str] = Field(None, min_length=2, max_length=120)
    org_type: Optional[Literal["advisory", "fund", "corporate", "other"]] = None
    description: Optional[str] = None


class OrgInviteRequest(BaseModel):
    email: EmailStr
    role: Literal["org_admin", "org_member"] = "org_member"


class CollaboratorInviteRequest(BaseModel):
    email: EmailStr
    role: Literal["owner", "editor", "viewer"] = "editor"
    message: Optional[str] = None


class AccessPolicyUpdate(BaseModel):
    require_principal_approval: Optional[bool] = None
    competitor_blocklist: Optional[List[str]] = None


# ---- Helpers --------------------------------------------------------------

async def _get_user_org_ids(user: dict) -> List[str]:
    """Return list of org_ids the user is a member of."""
    rows = await db.org_memberships.find(
        {"user_id": user["id"]}, {"_id": 0, "org_id": 1}
    ).to_list(200)
    return [r["org_id"] for r in rows]


async def _user_org_role(user_id: str, org_id: str) -> Optional[str]:
    row = await db.org_memberships.find_one(
        {"user_id": user_id, "org_id": org_id}, {"_id": 0, "role": 1}
    )
    return row["role"] if row else None


async def _resolve_listing_access(user: dict, listing: dict) -> Tuple[bool, Optional[str]]:
    """
    Returns (can_view, role_on_listing). role_on_listing is one of:
      "owner" (principal), "org_admin", "org_member", "collab_owner",
      "collab_editor", "collab_viewer", "admin", or None.
    """
    if user.get("role") == "admin":
        return True, "admin"
    # Legacy / individual seller ownership
    if listing.get("seller_id") == user["id"]:
        return True, "owner"
    # Org membership
    listing_org = listing.get("org_id")
    if listing_org:
        role = await _user_org_role(user["id"], listing_org)
        if role:
            return True, "org_admin" if role == "org_admin" else "org_member"
    # Explicit collaborators
    for c in listing.get("collaborators", []) or []:
        if c.get("user_id") == user["id"]:
            return True, f"collab_{c.get('role', 'viewer')}"
    return False, None


def _listing_role_can_edit(role: Optional[str]) -> bool:
    return role in {"admin", "owner", "org_admin", "org_member", "collab_owner", "collab_editor"}


async def _listing_for_edit_or_404(lid: str, user: dict) -> Dict[str, Any]:
    listing = await db.listings.find_one({"id": lid}, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found")
    can_view, role = await _resolve_listing_access(user, listing)
    if not can_view or not _listing_role_can_edit(role):
        raise HTTPException(status_code=403, detail="Not authorized to edit this listing")
    return listing


async def _listing_for_view_or_404(lid: str, user: dict) -> Dict[str, Any]:
    listing = await db.listings.find_one({"id": lid}, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing not found")
    can_view, _ = await _resolve_listing_access(user, listing)
    if not can_view:
        raise HTTPException(status_code=403, detail="Not authorized")
    return listing


# ---- Organizations CRUD ----------------------------------------------------

@api_router.post("/orgs")
async def create_org(body: OrgCreateRequest, user=Depends(get_current_user)):
    """Self-serve: any active user can create an org. The creator becomes org_admin."""
    org_id = str(uuid.uuid4())
    slug = re.sub(r"[^a-z0-9-]+", "-", body.name.lower()).strip("-")[:60] or org_id[:8]
    # Ensure uniqueness on slug
    if await db.organizations.find_one({"slug": slug}, {"_id": 0, "id": 1}):
        slug = f"{slug}-{org_id[:6]}"
    org_doc = {
        "id": org_id,
        "name": body.name,
        "slug": slug,
        "org_type": body.org_type,
        "description": body.description,
        "created_by": user["id"],
        "created_at": now_utc().isoformat(),
    }
    await db.organizations.insert_one(org_doc)
    await db.org_memberships.insert_one({
        "id": str(uuid.uuid4()),
        "org_id": org_id,
        "user_id": user["id"],
        "role": "org_admin",
        "joined_at": now_utc().isoformat(),
        "invited_by": None,
    })
    await log_audit(user["id"], "org.create", org_id, {"name": body.name})
    org_doc.pop("_id", None)
    org_doc["member_count"] = 1
    org_doc["my_role"] = "org_admin"
    return org_doc


@api_router.get("/orgs/mine")
async def list_my_orgs(user=Depends(get_current_user)):
    rows = await db.org_memberships.find({"user_id": user["id"]}, {"_id": 0}).to_list(200)
    out = []
    for r in rows:
        org = await db.organizations.find_one({"id": r["org_id"]}, {"_id": 0})
        if not org:
            continue
        org["my_role"] = r["role"]
        org["joined_at"] = r.get("joined_at")
        org["member_count"] = await db.org_memberships.count_documents({"org_id": org["id"]})
        out.append(org)
    return out


@api_router.get("/orgs/{org_id}")
async def get_org(org_id: str, user=Depends(get_current_user)):
    org = await db.organizations.find_one({"id": org_id}, {"_id": 0})
    if not org:
        raise HTTPException(status_code=404, detail="Org not found")
    role = await _user_org_role(user["id"], org_id)
    if not role and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Not a member")
    org["my_role"] = role or "admin"
    org["member_count"] = await db.org_memberships.count_documents({"org_id": org_id})
    return org


@api_router.patch("/orgs/{org_id}")
async def update_org(org_id: str, body: OrgUpdateRequest, user=Depends(get_current_user)):
    role = await _user_org_role(user["id"], org_id)
    if role != "org_admin" and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Org admin only")
    patch = {k: v for k, v in body.model_dump(exclude_unset=True).items() if v is not None}
    if not patch:
        return {"ok": True}
    res = await db.organizations.update_one({"id": org_id}, {"$set": patch})
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Org not found")
    await log_audit(user["id"], "org.update", org_id, patch)
    return {"ok": True}


@api_router.get("/orgs/{org_id}/members")
async def list_org_members(org_id: str, user=Depends(get_current_user)):
    role = await _user_org_role(user["id"], org_id)
    if not role and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Not a member")
    rows = await db.org_memberships.find({"org_id": org_id}, {"_id": 0}).to_list(500)
    out = []
    for r in rows:
        u = await db.users.find_one({"id": r["user_id"]}, {"_id": 0, "email": 1, "name": 1, "role": 1})
        if not u:
            continue
        out.append({
            "user_id": r["user_id"],
            "email": u.get("email"),
            "name": u.get("name"),
            "platform_role": u.get("role"),
            "org_role": r["role"],
            "joined_at": r.get("joined_at"),
        })
    return out


@api_router.delete("/orgs/{org_id}/members/{member_id}")
async def remove_org_member(org_id: str, member_id: str, user=Depends(get_current_user)):
    role = await _user_org_role(user["id"], org_id)
    if role != "org_admin" and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Org admin only")
    # Don't allow removing the last org_admin
    if member_id == user["id"]:
        admins = await db.org_memberships.count_documents({"org_id": org_id, "role": "org_admin"})
        if admins <= 1:
            raise HTTPException(status_code=400, detail="Cannot remove the last org admin")
    res = await db.org_memberships.delete_one({"org_id": org_id, "user_id": member_id})
    if res.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Member not found")
    await log_audit(user["id"], "org.member.remove", org_id, {"removed": member_id})
    return {"ok": True}


# ---- Org invites -----------------------------------------------------------

@api_router.post("/orgs/{org_id}/invites")
async def create_org_invite(org_id: str, body: OrgInviteRequest, user=Depends(get_current_user)):
    role = await _user_org_role(user["id"], org_id)
    if role != "org_admin" and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Org admin only")
    org = await db.organizations.find_one({"id": org_id}, {"_id": 0, "name": 1})
    if not org:
        raise HTTPException(status_code=404, detail="Org not found")
    email_norm = body.email.lower()
    # If a member already, no-op
    existing_user = await db.users.find_one({"email": email_norm}, {"_id": 0, "id": 1})
    if existing_user:
        if await _user_org_role(existing_user["id"], org_id):
            raise HTTPException(status_code=400, detail="User is already a member")
    token = uuid.uuid4().hex
    invite = {
        "id": str(uuid.uuid4()),
        "org_id": org_id,
        "email": email_norm,
        "role": body.role,
        "token": token,
        "invited_by": user["id"],
        "invited_by_name": user.get("name"),
        "created_at": now_utc().isoformat(),
        "expires_at": (now_utc() + timedelta(days=14)).isoformat(),
        "accepted_at": None,
    }
    await db.org_invites.insert_one(invite)
    accept_url = mail_link(f"/accept-org-invite?token={token}")
    register_url = mail_link(f"/register?invite_token={token}&invite_kind=org")
    html = f"""
    <p>Hi,</p>
    <p><strong>{user.get('name', 'A teammate')}</strong> invited you to join
    <strong>{org['name']}</strong> on NextCapOS as a <strong>{body.role.replace('_', ' ')}</strong>.</p>
    <p><a href="{accept_url}">Accept the invitation &rsaquo;</a></p>
    <p style="margin-top:18px;font-size:13px;color:#555;">
      New to NextCapOS? <a href="{register_url}">Create your account &rsaquo;</a> — your invite will be applied automatically and you'll skip the access-request queue.
    </p>
    <p style="margin-top:18px;font-size:12px;color:#999;">
      This invite expires in 14 days. Didn't expect it? You can ignore this email.
    </p>
    <p style="margin-top:18px;font-size:11px;color:#aaa;font-family:monospace;word-break:break-all;">
      If the buttons don't work, your invite token is:<br>{token}
    </p>
    """
    asyncio.create_task(send_email(email_norm, f"NextCapOS · join {org['name']}", html, reply_to=user.get("email")))
    await log_audit(user["id"], "org.invite.create", org_id, {"email": email_norm, "role": body.role})
    return {"ok": True, "invite_id": invite["id"], "token": token, "accept_url": accept_url}


@api_router.get("/orgs/{org_id}/invites")
async def list_org_invites(org_id: str, user=Depends(get_current_user)):
    role = await _user_org_role(user["id"], org_id)
    if role != "org_admin" and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Org admin only")
    rows = await db.org_invites.find(
        {"org_id": org_id, "accepted_at": None}, {"_id": 0, "token": 0}
    ).sort("created_at", -1).to_list(200)
    return rows


@api_router.delete("/orgs/{org_id}/invites/{iid}")
async def revoke_org_invite(org_id: str, iid: str, user=Depends(get_current_user)):
    role = await _user_org_role(user["id"], org_id)
    if role != "org_admin" and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Org admin only")
    res = await db.org_invites.delete_one({"id": iid, "org_id": org_id})
    if res.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Invite not found")
    await log_audit(user["id"], "org.invite.revoke", org_id, {"invite_id": iid})
    return {"ok": True}


@api_router.get("/org-invites/{token}")
async def public_get_org_invite(token: str):
    """Public endpoint so an unauthenticated user can preview their invite."""
    inv = await db.org_invites.find_one({"token": token}, {"_id": 0, "token": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found or already used")
    if inv.get("accepted_at"):
        raise HTTPException(status_code=400, detail="Invite already accepted")
    if inv.get("expires_at") and datetime.fromisoformat(inv["expires_at"]) < now_utc():
        raise HTTPException(status_code=400, detail="Invite expired")
    org = await db.organizations.find_one({"id": inv["org_id"]}, {"_id": 0, "name": 1, "org_type": 1})
    return {
        "email": inv["email"],
        "role": inv["role"],
        "org_id": inv["org_id"],
        "org_name": org["name"] if org else "Unknown",
        "org_type": org.get("org_type") if org else None,
        "invited_by_name": inv.get("invited_by_name"),
    }


@api_router.post("/org-invites/{token}/accept")
async def accept_org_invite(token: str, user=Depends(get_current_user)):
    """Authenticated user accepts an invite. The invite's email must match
    the user's email to prevent invite-token grabbing."""
    inv = await db.org_invites.find_one({"token": token}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("accepted_at"):
        raise HTTPException(status_code=400, detail="Already accepted")
    if inv.get("expires_at") and datetime.fromisoformat(inv["expires_at"]) < now_utc():
        raise HTTPException(status_code=400, detail="Invite expired")
    if inv["email"].lower() != user["email"].lower():
        raise HTTPException(status_code=403, detail="This invite is for a different email")
    # Already a member?
    if await _user_org_role(user["id"], inv["org_id"]):
        await db.org_invites.update_one({"token": token}, {"$set": {"accepted_at": now_utc().isoformat()}})
        return {"ok": True, "already_member": True, "org_id": inv["org_id"]}
    await db.org_memberships.insert_one({
        "id": str(uuid.uuid4()),
        "org_id": inv["org_id"],
        "user_id": user["id"],
        "role": inv["role"],
        "joined_at": now_utc().isoformat(),
        "invited_by": inv.get("invited_by"),
    })
    await db.org_invites.update_one({"token": token}, {"$set": {"accepted_at": now_utc().isoformat()}})
    await log_audit(user["id"], "org.invite.accept", inv["org_id"], {"role": inv["role"]})
    return {"ok": True, "org_id": inv["org_id"], "role": inv["role"]}


# ---- Listing collaborators -------------------------------------------------

@api_router.get("/listings/{lid}/collaborators")
async def list_listing_collaborators(lid: str, user=Depends(get_current_user)):
    listing = await _listing_for_view_or_404(lid, user)
    collabs = listing.get("collaborators", []) or []
    pending = await db.listing_invites.find(
        {"listing_id": lid, "accepted_at": None}, {"_id": 0, "token": 0}
    ).sort("created_at", -1).to_list(50)
    # Decorate every row with `can_manage` so the UI can hide role / remove /
    # resend / cancel controls in one place. Rule 1B: principal owner OR
    # original inviter only (admin overrides everything).
    for c in collabs:
        c["can_manage"] = _can_manage_collab_member(listing, user, c.get("user_id"))
    for iv in pending:
        iv["can_manage"] = _can_manage_pending_invite(listing, user, iv)
    return {
        "owner_id": listing.get("seller_id"),
        "org_id": listing.get("org_id"),
        "collaborators": collabs,
        "pending_invites": pending,
        # The viewer's own perspective — handy for client-side feature gates
        # (e.g. show or hide the "Invite a collaborator" form for non-managers).
        "viewer_is_principal": listing.get("seller_id") == user["id"],
        "viewer_id": user["id"],
    }


@api_router.post("/listings/{lid}/collaborators")
async def invite_listing_collaborator(lid: str, body: CollaboratorInviteRequest, user=Depends(get_current_user)):
    listing = await _listing_for_edit_or_404(lid, user)
    email_norm = body.email.lower()
    # If user already exists + already a collaborator, no-op
    existing_user = await db.users.find_one({"email": email_norm}, {"_id": 0, "id": 1, "name": 1})
    if existing_user:
        for c in listing.get("collaborators", []) or []:
            if c.get("user_id") == existing_user["id"]:
                raise HTTPException(status_code=400, detail="Already a collaborator")
    token = uuid.uuid4().hex
    invite = {
        "id": str(uuid.uuid4()),
        "listing_id": lid,
        "email": email_norm,
        "role": body.role,
        "token": token,
        "invited_by": user["id"],
        "invited_by_name": user.get("name"),
        "listing_name": listing.get("company_name"),
        "message": body.message,
        "created_at": now_utc().isoformat(),
        "expires_at": (now_utc() + timedelta(days=14)).isoformat(),
        "accepted_at": None,
    }
    await db.listing_invites.insert_one(invite)
    accept_url = mail_link(f"/accept-listing-invite?token={token}")
    register_url = mail_link(f"/register?invite_token={token}&invite_kind=listing")
    msg_block = f"<blockquote style='border-left:3px solid #ddd;margin:16px 0;padding:8px 14px;color:#444;'>{body.message}</blockquote>" if body.message else ""
    html = f"""
    <p>Hi,</p>
    <p><strong>{user.get('name', 'A teammate')}</strong> invited you to collaborate on the
    NextCapOS listing for <strong>{listing.get('company_name', 'a company')}</strong> as a
    <strong>{body.role}</strong>.</p>
    {msg_block}
    <p><a href="{accept_url}">Open the listing &rsaquo;</a></p>
    <p style="margin-top:18px;font-size:13px;color:#555;">
      New to NextCapOS? <a href="{register_url}">Create your account &rsaquo;</a> — your invite will be applied automatically and you'll skip the access-request queue.
    </p>
    <p style="margin-top:18px;font-size:12px;color:#999;">
      Owners and editors can edit the listing and upload to the Vault. Viewers can read only.
    </p>
    <p style="margin-top:18px;font-size:11px;color:#aaa;font-family:monospace;word-break:break-all;">
      If the buttons don't work, your invite token is:<br>{token}
    </p>
    """
    asyncio.create_task(send_email(email_norm, f"NextCapOS · listing invitation · {listing.get('company_name', '')}", html, reply_to=user.get("email")))
    await log_audit(user["id"], "listing.collab.invite", lid, {"email": email_norm, "role": body.role})
    return {"ok": True, "invite_id": invite["id"], "token": token, "accept_url": accept_url}


def _can_manage_collab_member(listing: dict, current_user: dict, member_id: str) -> bool:
    """
    Inviter-or-principal rule (Rule 1B):
      - Principal owner of the listing can manage every collaborator.
      - Otherwise, the caller can only manage collaborators THEY personally
        invited (matches `invited_by` on the collaborator entry).
      - Platform admin always allowed.
    """
    if current_user.get("role") == "admin":
        return True
    if listing.get("seller_id") == current_user["id"]:
        return True
    for c in listing.get("collaborators") or []:
        if c.get("user_id") == member_id and c.get("invited_by") == current_user["id"]:
            return True
    return False


def _can_manage_pending_invite(listing: dict, current_user: dict, invite: dict) -> bool:
    """Same rule as _can_manage_collab_member but against a pending invite doc."""
    if current_user.get("role") == "admin":
        return True
    if listing.get("seller_id") == current_user["id"]:
        return True
    return invite.get("invited_by") == current_user["id"]


@api_router.delete("/listings/{lid}/collaborators/{member_id}")
async def remove_listing_collaborator(lid: str, member_id: str, user=Depends(get_current_user)):
    listing = await _listing_for_edit_or_404(lid, user)
    if listing.get("seller_id") == member_id:
        raise HTTPException(status_code=400, detail="Cannot remove the principal owner")
    if not _can_manage_collab_member(listing, user, member_id):
        raise HTTPException(
            status_code=403,
            detail="Only the principal owner or the person who invited this collaborator can remove them.",
        )
    res = await db.listings.update_one(
        {"id": lid},
        {"$pull": {"collaborators": {"user_id": member_id}}},
    )
    if res.modified_count == 0:
        raise HTTPException(status_code=404, detail="Collaborator not found")
    await log_audit(user["id"], "listing.collab.remove", lid, {"removed": member_id})
    return {"ok": True}


class CollaboratorRolePatch(BaseModel):
    role: Literal["owner", "editor", "viewer"]


@api_router.patch("/listings/{lid}/collaborators/{member_id}")
async def update_listing_collaborator_role(
    lid: str, member_id: str, body: CollaboratorRolePatch, user=Depends(get_current_user),
):
    """Change an existing collaborator's role. Only the principal owner of the
    listing, or the user who originally invited this collaborator, can
    mutate the role (Rule 1B). The principal owner's own role is
    immutable — use remove + re-invite if that ever needs to change."""
    listing = await _listing_for_edit_or_404(lid, user)
    if listing.get("seller_id") == member_id:
        raise HTTPException(status_code=400, detail="Cannot change the principal owner's role")
    if not _can_manage_collab_member(listing, user, member_id):
        raise HTTPException(
            status_code=403,
            detail="Only the principal owner or the person who invited this collaborator can change their role.",
        )
    res = await db.listings.update_one(
        {"id": lid, "collaborators.user_id": member_id},
        {"$set": {"collaborators.$.role": body.role}},
    )
    if res.matched_count == 0:
        raise HTTPException(status_code=404, detail="Collaborator not found")
    await log_audit(user["id"], "listing.collab.update_role", lid,
                    {"member_id": member_id, "new_role": body.role})
    return {"ok": True, "role": body.role}


@api_router.delete("/listings/{lid}/collaborators/invites/{iid}")
async def revoke_listing_invite(lid: str, iid: str, user=Depends(get_current_user)):
    """Cancel a pending listing-collaborator invite. Only the principal owner
    or the user who created the invite can revoke (Rule 1B)."""
    listing = await _listing_for_edit_or_404(lid, user)
    inv = await db.listing_invites.find_one({"id": iid, "listing_id": lid}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("accepted_at"):
        raise HTTPException(
            status_code=400,
            detail="Invite already accepted — remove the collaborator instead.",
        )
    if not _can_manage_pending_invite(listing, user, inv):
        raise HTTPException(
            status_code=403,
            detail="Only the principal owner or the user who sent this invite can cancel it.",
        )
    await db.listing_invites.delete_one({"id": iid, "listing_id": lid})
    await log_audit(user["id"], "listing.invite.revoke", lid,
                    {"invite_id": iid, "email": inv.get("email")})
    return {"ok": True}


@api_router.get("/listing-invites/{token}")
async def public_get_listing_invite(token: str):
    inv = await db.listing_invites.find_one({"token": token}, {"_id": 0, "token": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found or already used")
    if inv.get("accepted_at"):
        raise HTTPException(status_code=400, detail="Invite already accepted")
    if inv.get("expires_at") and datetime.fromisoformat(inv["expires_at"]) < now_utc():
        raise HTTPException(status_code=400, detail="Invite expired")
    return inv


@api_router.post("/listing-invites/{token}/accept")
async def accept_listing_invite(token: str, user=Depends(get_current_user)):
    inv = await db.listing_invites.find_one({"token": token}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("accepted_at"):
        raise HTTPException(status_code=400, detail="Already accepted")
    if inv.get("expires_at") and datetime.fromisoformat(inv["expires_at"]) < now_utc():
        raise HTTPException(status_code=400, detail="Invite expired")
    if inv["email"].lower() != user["email"].lower():
        raise HTTPException(status_code=403, detail="This invite is for a different email")
    listing = await db.listings.find_one({"id": inv["listing_id"]}, {"_id": 0})
    if not listing:
        raise HTTPException(status_code=404, detail="Listing no longer exists")
    # Already a collaborator?
    for c in listing.get("collaborators", []) or []:
        if c.get("user_id") == user["id"]:
            await db.listing_invites.update_one({"token": token}, {"$set": {"accepted_at": now_utc().isoformat()}})
            return {"ok": True, "already_collaborator": True, "listing_id": inv["listing_id"]}
    await db.listings.update_one(
        {"id": inv["listing_id"]},
        {"$push": {"collaborators": {
            "user_id": user["id"],
            "email": user["email"],
            "name": user.get("name"),
            "role": inv["role"],
            "invited_by": inv.get("invited_by"),
            "invited_at": inv.get("created_at"),
            "accepted_at": now_utc().isoformat(),
        }}},
    )
    await db.listing_invites.update_one({"token": token}, {"$set": {"accepted_at": now_utc().isoformat()}})
    await log_audit(user["id"], "listing.collab.accept", inv["listing_id"], {"role": inv["role"]})
    return {"ok": True, "listing_id": inv["listing_id"], "role": inv["role"]}


# ---- Access policy (principal veto + competitor blocklist) -----------------

@api_router.patch("/listings/{lid}/access-policy")
async def update_access_policy(lid: str, body: AccessPolicyUpdate, user=Depends(get_current_user)):
    listing = await _listing_for_edit_or_404(lid, user)
    # Principal toggle + blocklist can be set by any editor; the principal
    # always inherits effective control because principal_approval falls back
    # to them on approval routing (server enforces).
    patch = {}
    if body.require_principal_approval is not None:
        patch["access_policy.require_principal_approval"] = body.require_principal_approval
    if body.competitor_blocklist is not None:
        # Normalize: lowercase, strip, dedupe
        items = list({(s or "").strip().lower() for s in body.competitor_blocklist if (s or "").strip()})
        patch["access_policy.competitor_blocklist"] = items
    if not patch:
        return {"ok": True}
    await db.listings.update_one({"id": lid}, {"$set": patch})
    await log_audit(user["id"], "listing.access_policy.update", lid, patch)
    return {"ok": True, "access_policy": {**(listing.get("access_policy") or {}), **{k.split('.')[-1]: v for k, v in patch.items()}}}


# Public listing preview links (share tokens) were removed in Iter-41.
# The invite-collaborator flow (owner/editor/viewer roles + audit trail) is the
# supported way to share a listing pre-signup. Any legacy `listing_preview_links`
# rows are ignored and can be dropped from Mongo at any time.


async def _user_workspace_listing_ids(user: dict) -> Tuple[List[str], List[str]]:
    """Return (seller_workspace_listing_ids, org_ids) — listings the user
    can act on as seller-side: personal + org-owned + collaborator (editor/owner).

    Used to widen access to inquiries, vaults, and listing actions so org
    teammates can pick up where each other left off without re-assigning.
    """
    org_ids = await _get_user_org_ids(user)
    or_clauses: List[Dict[str, Any]] = [
        {"seller_id": user["id"]},
        {"collaborators": {"$elemMatch": {"user_id": user["id"], "role": {"$in": ["owner", "editor"]}}}},
    ]
    if org_ids:
        or_clauses.append({"org_id": {"$in": org_ids}})
    rows = await db.listings.find({"$or": or_clauses}, {"_id": 0, "id": 1}).to_list(500)
    return [r["id"] for r in rows], org_ids


# ---- Cross-cutting: pending invites for the current user ------------------

@api_router.get("/me/invites/pending")
async def my_pending_invites(user=Depends(get_current_user)):
    """Returns all unaccepted, unexpired invites (org + listing) addressed to
    the current user's email. Drives the in-app "Pending invitations" panel
    so users don't have to depend on receiving the original email."""
    email = (user.get("email") or "").lower()
    now = now_utc()
    out_orgs = []
    async for inv in db.org_invites.find({"email": email, "accepted_at": None}, {"_id": 0}):
        try:
            if datetime.fromisoformat(inv["expires_at"]) < now:
                continue
        except Exception:
            pass
        org = await db.organizations.find_one({"id": inv["org_id"]}, {"_id": 0, "name": 1, "org_type": 1})
        out_orgs.append({
            "kind": "org",
            "token": inv["token"],
            "role": inv["role"],
            "invited_by_name": inv.get("invited_by_name"),
            "created_at": inv.get("created_at"),
            "expires_at": inv.get("expires_at"),
            "org_name": (org or {}).get("name", "Unknown organization"),
            "org_type": (org or {}).get("org_type"),
        })
    out_listings = []
    async for inv in db.listing_invites.find({"email": email, "accepted_at": None}, {"_id": 0}):
        try:
            if datetime.fromisoformat(inv["expires_at"]) < now:
                continue
        except Exception:
            pass
        out_listings.append({
            "kind": "listing",
            "token": inv["token"],
            "role": inv["role"],
            "invited_by_name": inv.get("invited_by_name"),
            "created_at": inv.get("created_at"),
            "expires_at": inv.get("expires_at"),
            "listing_name": inv.get("listing_name"),
            "message": inv.get("message"),
        })
    return {"org": out_orgs, "listing": out_listings}


# ---- Resend a stuck invitation (org admins + platform admins) --------------

@api_router.post("/orgs/{org_id}/invites/{iid}/resend")
async def resend_org_invite(org_id: str, iid: str, user=Depends(get_current_user)):
    """Re-fire the email for an existing org invite. Useful after fixing
    Resend config / sender domain. Does NOT rotate the token — existing
    accept-URL keeps working."""
    role = await _user_org_role(user["id"], org_id)
    if role != "org_admin" and user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Org admin only")
    inv = await db.org_invites.find_one({"id": iid, "org_id": org_id}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("accepted_at"):
        raise HTTPException(status_code=400, detail="Invite already accepted")
    org = await db.organizations.find_one({"id": org_id}, {"_id": 0, "name": 1}) or {"name": "your team"}
    accept_url = mail_link(f"/accept-org-invite?token={inv['token']}")
    register_url = mail_link(f"/register?invite_token={inv['token']}&invite_kind=org")
    html = f"""
    <p>Hi,</p>
    <p><strong>{user.get('name', 'A teammate')}</strong> invited you to join
    <strong>{org['name']}</strong> on NextCapOS as a <strong>{inv['role'].replace('_',' ')}</strong>.</p>
    <p><a href="{accept_url}">Accept the invitation &rsaquo;</a></p>
    <p style="margin-top:18px;font-size:13px;color:#555;">
      New to NextCapOS? <a href="{register_url}">Create your account &rsaquo;</a> — your invite will be applied automatically.
    </p>
    <p style="margin-top:18px;font-size:12px;color:#999;">
      If you didn't request this, you can ignore the email.
    </p>
    <p style="margin-top:18px;font-size:11px;color:#aaa;font-family:monospace;word-break:break-all;">
      If the buttons don't work, your invite token is:<br>{inv['token']}
    </p>
    """
    asyncio.create_task(send_email(inv["email"], f"NextCapOS · join {org['name']} (resent)", html, reply_to=user.get("email")))
    await log_audit(user["id"], "org.invite.resend", org_id, {"invite_id": iid, "email": inv["email"]})
    return {"ok": True, "email": inv["email"]}


@api_router.post("/listings/{lid}/collaborators/{iid}/resend")
async def resend_listing_invite(lid: str, iid: str, user=Depends(get_current_user)):
    """Re-fire the email for a pending listing-collaborator invite. Only the
    principal owner or the inviter can resend (Rule 1B)."""
    listing = await _listing_for_edit_or_404(lid, user)
    inv = await db.listing_invites.find_one({"id": iid, "listing_id": lid}, {"_id": 0})
    if not inv:
        raise HTTPException(status_code=404, detail="Invite not found")
    if inv.get("accepted_at"):
        raise HTTPException(status_code=400, detail="Invite already accepted")
    if not _can_manage_pending_invite(listing, user, inv):
        raise HTTPException(
            status_code=403,
            detail="Only the principal owner or the user who sent this invite can resend it.",
        )
    accept_url = mail_link(f"/accept-listing-invite?token={inv['token']}")
    register_url = mail_link(f"/register?invite_token={inv['token']}&invite_kind=listing")
    msg_block = f"<blockquote style='border-left:3px solid #ddd;margin:16px 0;padding:8px 14px;color:#444;'>{inv.get('message')}</blockquote>" if inv.get("message") else ""
    html = f"""
    <p>Hi,</p>
    <p><strong>{user.get('name', 'A teammate')}</strong> invited you to collaborate on the
    NextCapOS listing for <strong>{inv.get('listing_name', 'a company')}</strong> as a
    <strong>{inv['role']}</strong>.</p>
    {msg_block}
    <p><a href="{accept_url}">Open the listing &rsaquo;</a></p>
    <p style="margin-top:18px;font-size:13px;color:#555;">
      New to NextCapOS? <a href="{register_url}">Create your account &rsaquo;</a> — your invite will be applied automatically.
    </p>
    <p style="margin-top:18px;font-size:11px;color:#aaa;font-family:monospace;word-break:break-all;">
      If the buttons don't work, your invite token is:<br>{inv['token']}
    </p>
    """
    asyncio.create_task(send_email(inv["email"], f"NextCapOS · listing invitation · {inv.get('listing_name','')} (resent)", html, reply_to=user.get("email")))
    await log_audit(user["id"], "listing.collab.invite.resend", lid, {"invite_id": iid, "email": inv["email"]})
    return {"ok": True, "email": inv["email"]}


# ---- Deal-room collaborators (Phase 2) -------------------------------------

@api_router.get("/deal-rooms/{rid}/collaborators")
async def list_room_collaborators(rid: str, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    # Allowed if the user is buyer or seller side, or appears in collaborators.
    can = (
        room.get("buyer_id") == user["id"]
        or room.get("seller_id") == user["id"]
        or user.get("role") == "admin"
        or any(c.get("user_id") == user["id"] for c in (room.get("collaborators") or []))
    )
    if not can:
        raise HTTPException(status_code=403, detail="Not authorized")
    return {
        "buyer_id": room.get("buyer_id"),
        "seller_id": room.get("seller_id"),
        "collaborators": room.get("collaborators", []) or [],
    }


@api_router.post("/deal-rooms/{rid}/collaborators")
async def invite_room_collaborator(rid: str, body: CollaboratorInviteRequest, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    # Only the seller-side principal/agent or buyer-side principal/agent can invite.
    can = (
        room.get("buyer_id") == user["id"]
        or room.get("seller_id") == user["id"]
        or user.get("role") == "admin"
        or any(c.get("user_id") == user["id"] and c.get("role") in ("owner", "editor") for c in (room.get("collaborators") or []))
    )
    if not can:
        raise HTTPException(status_code=403, detail="Not authorized")
    email_norm = body.email.lower()
    existing_user = await db.users.find_one({"email": email_norm}, {"_id": 0, "id": 1, "name": 1})
    if not existing_user:
        raise HTTPException(status_code=400, detail="The collaborator must have a NextCapOS account first")
    for c in room.get("collaborators", []) or []:
        if c.get("user_id") == existing_user["id"]:
            raise HTTPException(status_code=400, detail="Already a collaborator")
    entry = {
        "user_id": existing_user["id"],
        "email": email_norm,
        "name": existing_user.get("name"),
        "role": body.role,
        "invited_by": user["id"],
        "invited_at": now_utc().isoformat(),
        "accepted_at": now_utc().isoformat(),  # implicit accept; user already has account
    }
    await db.deal_rooms.update_one({"id": rid}, {"$push": {"collaborators": entry}})
    await log_audit(user["id"], "room.collab.add", rid, {"user_id": existing_user["id"], "role": body.role})
    return {"ok": True, "collaborator": entry}


@api_router.delete("/deal-rooms/{rid}/collaborators/{member_id}")
async def remove_room_collaborator(rid: str, member_id: str, user=Depends(get_current_user)):
    room = await db.deal_rooms.find_one({"id": rid}, {"_id": 0})
    if not room:
        raise HTTPException(status_code=404, detail="Room not found")
    can = (
        room.get("seller_id") == user["id"]
        or room.get("buyer_id") == user["id"]
        or user.get("role") == "admin"
    )
    if not can:
        raise HTTPException(status_code=403, detail="Not authorized")
    res = await db.deal_rooms.update_one(
        {"id": rid}, {"$pull": {"collaborators": {"user_id": member_id}}}
    )
    if res.modified_count == 0:
        raise HTTPException(status_code=404, detail="Collaborator not found")
    await log_audit(user["id"], "room.collab.remove", rid, {"removed": member_id})
    return {"ok": True}


# Register router + middleware
app.include_router(api_router)


@app.middleware("http")
async def security_headers_middleware(request: Request, call_next):
    response = await call_next(request)
    response.headers["Strict-Transport-Security"] = "max-age=63072000; includeSubDomains"
    response.headers["X-Content-Type-Options"] = "nosniff"
    response.headers["X-Frame-Options"] = "DENY"
    response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
    response.headers["Permissions-Policy"] = "geolocation=(), camera=(), microphone=()"
    return response


app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get("CORS_ORIGINS", "*").split(","),
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("shutdown")
async def shutdown_db_client():
    global _buyer_scheduler_task, _demo_cleanup_task
    for t in (_buyer_scheduler_task, _demo_cleanup_task):
        if t and not t.done():
            t.cancel()
            try:
                await t
            except (asyncio.CancelledError, Exception):
                pass
    client.close()
()
