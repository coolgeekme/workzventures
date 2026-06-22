"""Unit tests for the pure-Python PPTX/DOCX → PDF converters with embedded
images. Verifies that when a deck/document contains raster images, the
resulting PDF actually carries them through (instead of dropping them as the
earlier text-only renderer did).

We don't need the API server for these — we import the converters directly.
"""

import io
import os
import sys

import pytest

# Make `backend/server.py` importable by tests.
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))


def _make_dummy_png_bytes(width=80, height=60, color=(220, 38, 38)) -> bytes:
    from PIL import Image as PILImage
    im = PILImage.new("RGB", (width, height), color)
    out = io.BytesIO()
    im.save(out, format="PNG")
    return out.getvalue()


def _pdf_has_image_xobject(pdf_bytes: bytes) -> bool:
    """Heuristic: every embedded raster image in a reportlab-generated PDF
    is a `/Subtype /Image` XObject. Presence of this marker → at least one
    image made it into the PDF."""
    if not pdf_bytes or not pdf_bytes.startswith(b"%PDF-"):
        return False
    return b"/Subtype /Image" in pdf_bytes or b"/Subtype/Image" in pdf_bytes


@pytest.fixture()
def pptx_with_image_bytes():
    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    slide_layout = pres.slide_layouts[5]  # title-only
    slide = pres.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Embedded Image Test Deck"

    img_blob = _make_dummy_png_bytes(color=(11, 27, 61))  # nextcapos navy
    img_stream = io.BytesIO(img_blob)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(2), width=Inches(4))

    # A second slide with NO image so the iteration logic also covers the
    # empty-image branch.
    pres.slides.add_slide(pres.slide_layouts[5]).shapes.title.text = "Text Only Slide"

    out = io.BytesIO()
    pres.save(out)
    out.seek(0)
    return out.read()


@pytest.fixture()
def docx_with_image_bytes():
    from docx import Document
    from docx.shared import Inches

    d = Document()
    d.add_heading("Embedded Image Test Doc", level=0)
    d.add_paragraph("This document should preview with the image below intact.")

    img_blob = _make_dummy_png_bytes(color=(29, 78, 216))  # bloomberg blue
    d.add_picture(io.BytesIO(img_blob), width=Inches(3))

    d.add_paragraph("Trailing paragraph after the embedded image.")

    out = io.BytesIO()
    d.save(out)
    out.seek(0)
    return out.read()


class TestPptxImageRendering:
    def test_pptx_with_image_produces_pdf(self, pptx_with_image_bytes):
        from server import _pptx_to_pdf_bytes
        pdf = _pptx_to_pdf_bytes(pptx_with_image_bytes, "deck.pptx")
        assert pdf is not None, "pptx→pdf returned None"
        assert pdf.startswith(b"%PDF-"), "pptx→pdf output is not a PDF"

    def test_pptx_pdf_carries_embedded_image(self, pptx_with_image_bytes):
        from server import _pptx_to_pdf_bytes
        pdf = _pptx_to_pdf_bytes(pptx_with_image_bytes, "deck.pptx")
        assert _pdf_has_image_xobject(pdf), \
            "embedded image was dropped during pptx→pdf conversion"

    def test_pptx_without_image_still_renders(self):
        """Regression guard: a deck with zero images must not throw."""
        from pptx import Presentation
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[5])
        slide.shapes.title.text = "No Pictures Here"
        buf = io.BytesIO()
        pres.save(buf)
        buf.seek(0)

        from server import _pptx_to_pdf_bytes
        pdf = _pptx_to_pdf_bytes(buf.read(), "plain.pptx")
        assert pdf is not None
        assert pdf.startswith(b"%PDF-")


class TestDocxImageRendering:
    def test_docx_with_image_produces_pdf(self, docx_with_image_bytes):
        from server import _docx_to_pdf_bytes
        pdf = _docx_to_pdf_bytes(docx_with_image_bytes, "doc.docx")
        assert pdf is not None, "docx→pdf returned None"
        assert pdf.startswith(b"%PDF-"), "docx→pdf output is not a PDF"

    def test_docx_pdf_carries_embedded_image(self, docx_with_image_bytes):
        from server import _docx_to_pdf_bytes
        pdf = _docx_to_pdf_bytes(docx_with_image_bytes, "doc.docx")
        assert _pdf_has_image_xobject(pdf), \
            "embedded image was dropped during docx→pdf conversion"

    def test_docx_without_image_still_renders(self):
        from docx import Document
        d = Document()
        d.add_heading("No Image Doc", level=0)
        d.add_paragraph("Plain text only.")
        buf = io.BytesIO()
        d.save(buf)
        buf.seek(0)
        from server import _docx_to_pdf_bytes
        pdf = _docx_to_pdf_bytes(buf.read(), "plain.docx")
        assert pdf is not None
        assert pdf.startswith(b"%PDF-")


class TestCorruptImageResilience:
    """A single broken image must not kill the whole render — buyers should
    still see the rest of the deck."""

    def test_pptx_with_text_around_corrupt_image_falls_through(self):
        # Build a deck where we manually corrupt the image blob in the package
        # zip after python-pptx writes it. python-pptx itself won't let us
        # add invalid PNGs, so we tamper with the produced zip directly.
        import zipfile
        from pptx import Presentation
        from pptx.util import Inches

        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[5])
        slide.shapes.title.text = "Corrupt Image Resilience"
        slide.shapes.add_picture(
            io.BytesIO(_make_dummy_png_bytes()), Inches(1), Inches(2), width=Inches(3),
        )
        clean = io.BytesIO()
        pres.save(clean)
        clean.seek(0)

        # Rewrite the media file inside the .pptx zip with garbage bytes.
        in_zip = zipfile.ZipFile(clean, "r")
        tampered = io.BytesIO()
        with zipfile.ZipFile(tampered, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in in_zip.infolist():
                blob = in_zip.read(item.filename)
                if item.filename.startswith("ppt/media/"):
                    blob = b"\x00\x01\x02NOT-A-REAL-IMAGE\xff\xff"
                zout.writestr(item, blob)
        tampered.seek(0)

        from server import _pptx_to_pdf_bytes
        pdf = _pptx_to_pdf_bytes(tampered.read(), "corrupt.pptx")
        # Must still produce a valid PDF (text is preserved even if image
        # decode fails).
        assert pdf is not None
        assert pdf.startswith(b"%PDF-")
