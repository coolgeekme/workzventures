"""
Smoke tests for multi-format file uploads (XLSX, PPTX, images) across the
three upload surfaces: Listing Data Room, Vault (deal room), Private Locker.
"""
import io
import os
import uuid

import requests

BASE_URL = os.environ.get(
    "REACT_APP_BACKEND_URL",
    "https://buyer-intel-lab.preview.emergentagent.com",
).rstrip("/")
API = f"{BASE_URL}/api"

BUYER = {"email": "alex@workz.example.com", "password": "WorkzPass123!"}
SELLER = {"email": "mira@workz.example.com", "password": "WorkzPass123!"}


def _login(creds):
    r = requests.post(f"{API}/auth/login", json=creds, timeout=20)
    assert r.status_code == 200, r.text
    return r.json()["token"]


def _build_xlsx():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "P&L"
    ws.append(["Metric", "FY23", "FY24"])
    ws.append(["Revenue", 38.5, 54.2])
    ws.append(["Gross Margin %", 61, 64])
    ws.append(["EBITDA", 2.1, 4.9])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _build_pptx():
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from pptx.util import Inches, Pt
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = tb.text_frame
    tf.text = "Project Helios — Investor Deck"
    p = tf.add_paragraph()
    p.text = "FY24 revenue €54.2M, gross margin 64%, ARR €18.1M"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _png_bytes():
    # Minimal 1x1 PNG
    import base64
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="
    )


def test_locker_accepts_xlsx_with_text_extraction():
    token = _login(BUYER)
    payload = _build_xlsx()
    fname = f"financials-{uuid.uuid4().hex[:6]}.xlsx"
    r = requests.post(
        f"{API}/private-locker/files",
        headers={"Authorization": f"Bearer {token}"},
        files={"file": (fname, io.BytesIO(payload),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        data={"folder": "modeling", "note": "P&L extract"},
        timeout=30,
    )
    assert r.status_code == 200, r.text
    doc = r.json()
    assert doc["filename"] == fname
    assert doc["size_bytes"] == len(payload)
    assert doc["page_count"] >= 1, "xlsx must produce >=1 page during extraction"
    fid = doc["id"]

    # Round-trip download
    dl = requests.get(
        f"{API}/private-locker/files/{fid}/download",
        headers={"Authorization": f"Bearer {token}"},
        timeout=15,
    )
    assert dl.status_code == 200
    assert dl.content == payload

    requests.delete(f"{API}/private-locker/files/{fid}",
                    headers={"Authorization": f"Bearer {token}"}, timeout=15)


def test_locker_accepts_pptx_with_slide_extraction():
    token = _login(BUYER)
    payload = _build_pptx()
    fname = f"deck-{uuid.uuid4().hex[:6]}.pptx"
    r = requests.post(
        f"{API}/private-locker/files",
        headers={"Authorization": f"Bearer {token}"},
        files={"file": (fname, io.BytesIO(payload),
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        data={"folder": "memos"},
        timeout=30,
    )
    assert r.status_code == 200, r.text
    doc = r.json()
    assert doc["page_count"] >= 1
    fid = doc["id"]
    requests.delete(f"{API}/private-locker/files/{fid}",
                    headers={"Authorization": f"Bearer {token}"}, timeout=15)


def test_locker_accepts_image_files_without_extraction_crash():
    token = _login(BUYER)
    payload = _png_bytes()
    fname = f"shot-{uuid.uuid4().hex[:6]}.png"
    r = requests.post(
        f"{API}/private-locker/files",
        headers={"Authorization": f"Bearer {token}"},
        files={"file": (fname, io.BytesIO(payload), "image/png")},
        data={"folder": "other"},
        timeout=30,
    )
    assert r.status_code == 200, r.text
    doc = r.json()
    assert doc["filename"] == fname
    assert doc["page_count"] >= 1
    # Round-trip
    dl = requests.get(
        f"{API}/private-locker/files/{doc['id']}/download",
        headers={"Authorization": f"Bearer {token}"},
        timeout=15,
    )
    assert dl.status_code == 200
    assert dl.content == payload
    requests.delete(f"{API}/private-locker/files/{doc['id']}",
                    headers={"Authorization": f"Bearer {token}"}, timeout=15)


def test_listing_data_room_accepts_xlsx():
    stok = _login(SELLER)
    listings = requests.get(
        f"{API}/listings", headers={"Authorization": f"Bearer {stok}"}, timeout=15
    ).json()
    lid = listings[0]["id"]

    payload = _build_xlsx()
    fname = f"cohorts-{uuid.uuid4().hex[:6]}.xlsx"
    r = requests.post(
        f"{API}/listings/{lid}/staged-files/binary",
        headers={"Authorization": f"Bearer {stok}"},
        files={"file": (fname, io.BytesIO(payload),
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")},
        data={"folder": "financials", "note": "cohort drilldown"},
        timeout=30,
    )
    assert r.status_code == 200, r.text
    fid = r.json()["id"]

    # Cleanup
    requests.delete(
        f"{API}/listings/{lid}/staged-files/{fid}",
        headers={"Authorization": f"Bearer {stok}"},
        timeout=15,
    )
